from __future__ import annotations

from dataclasses import dataclass
from html import unescape
from io import BytesIO
import math
import re
import zipfile
from collections import Counter
from time import perf_counter
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.models import PptxDownloadRequest
from app.services.pptx_service import build_pptx_context


RENDERER_MVP_VERSION = "presentation_master_v3_renderer_mvp"
RENDERER_MVP_FEATURE_FLAG = "PRESENTATION_MASTER_V3_RENDERER_MVP_ENABLED"
RENDERER_MVP_CONTRACT_VERSION = "presentation_master_v3_renderer_mvp_adapter_v1"
SLIDE_W = 13.333
SLIDE_H = 7.5
EMU_PER_INCH = 914400

ALLOWED_PRIMITIVES = {
    "text",
    "rule",
    "semantic_container",
    "connector",
    "evidence_object",
    "path_approximation",
    "boundary",
}
ALLOWED_RELATIONSHIPS = {
    "contrast",
    "cause",
    "sequence",
    "dependency",
    "boundary",
    "tension",
    "convergence",
    "evidence_supports_decision",
}
PROHIBITED_VISIBLE_TOKENS = (
    "placeholder",
    "todo",
    "dummy",
    "lorem",
    "internal",
    "debug",
    "layout_id",
    "archetype_",
    "composition_",
    "仮ラベル",
    "内部ラベル",
)


@dataclass(frozen=True)
class RendererMvpBuildOutput:
    pptx_bytes: bytes
    quality_report: dict[str, Any]


class RendererMvpIntegrationError(RuntimeError):
    def __init__(self, reason_code: str, *, failure_stage: str, details: dict[str, Any] | None = None):
        self.reason_code = reason_code
        self.failure_stage = failure_stage
        self.details = details or {}
        super().__init__(reason_code)


def build_renderer_mvp_pptx(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
) -> RendererMvpBuildOutput:
    started = perf_counter()
    if payload.summary:
        raise RendererMvpIntegrationError(
            "summary_deck_uses_existing_renderer",
            failure_stage="routing",
        )

    contract = ProposalToRendererMvpAdapter().build_contract(payload)
    contract_issues = RendererMvpContractLinter().lint(contract)
    architecture_deviation_count = sum(1 for issue in contract_issues if issue["severity"] == "error")
    if architecture_deviation_count:
        raise RendererMvpIntegrationError(
            "renderer_mvp_contract_deviation",
            failure_stage="contract_validation",
            details={"contract_issues": contract_issues},
        )

    renderer = RendererMvpNativeRenderer()
    pptx_bytes, render_report = renderer.render_deck_to_bytes(contract)
    pptx_audit = inspect_pptx_bytes(pptx_bytes, source_payload=payload)
    validation = _validate_runtime_output(pptx_audit, render_report)
    if validation["blocking"]:
        raise RendererMvpIntegrationError(
            "renderer_mvp_runtime_validation_failed",
            failure_stage="runtime_validation",
            details=validation,
        )

    page_count = int(pptx_audit["page_count"])
    quality_report = {
        "requested_version": RENDERER_MVP_VERSION,
        "actual_version": RENDERER_MVP_VERSION,
        "fallback_used": False,
        "fallback_reason": "",
        "feature_flag": RENDERER_MVP_FEATURE_FLAG,
        "request_id": request_id or "",
        "project_id": project_id or "",
        "generation_time_ms": round((perf_counter() - started) * 1000),
        "production_native_module": "app.services.presentation_master.renderer_mvp",
        "artifact_runtime_dependency": False,
        "adapter_contract_version": RENDERER_MVP_CONTRACT_VERSION,
        "renderer_architecture": "native_editable_powerpoint_shapes",
        "source_of_truth": "presentation_master_v3_renderer_mvp_powerpoint_desktop_approved",
        "engine_requested": RENDERER_MVP_VERSION,
        "engine_used": RENDERER_MVP_VERSION,
        "mode": "enabled",
        "summary_or_normal": "normal",
        "customer": contract["case_summary"]["customer"],
        "category": contract["case_summary"]["category"],
        "page_count": page_count,
        "director_slide_count": page_count,
        "render_report": render_report,
        "pptx_audit": pptx_audit,
        "architecture_deviation_count": architecture_deviation_count,
        "contract_issues": contract_issues,
        "placeholder_internal_label_count": pptx_audit["placeholder_count"] + pptx_audit["internal_label_count"],
        "fake_evidence_count": pptx_audit["fake_evidence_count"],
        "tier1_editability": pptx_audit["tier1_editable_coverage"],
        "rasterization_ratio": pptx_audit["rasterization_ratio"],
        "template_collapse": _template_collapse_from_render_report(render_report),
        "architecture_guardrails": {
            "strategy_engine_bypassed": False,
            "evidence_controls_bypassed": False,
            "term_guard_bypassed": False,
            "human_review_gate_bypassed": False,
            "api_contract_changed": False,
        },
    }
    return RendererMvpBuildOutput(pptx_bytes=pptx_bytes, quality_report=quality_report)


class ProposalToRendererMvpAdapter:
    def build_contract(self, payload: PptxDownloadRequest) -> dict[str, Any]:
        context = build_pptx_context(payload)
        data = payload.powerpoint_generation_data
        text = "\n".join(
            [
                data.deck_title,
                data.client_name,
                payload.client_company_info,
                payload.project_brief,
                payload.hearing_result,
                payload.own_service_info,
                payload.special_function_required,
            ]
        )
        profile = self._profile_for_text(text)
        customer = _safe_text(data.client_name or context.client_name or payload.client_company_info or "顧客企業", 48)
        category = _safe_text(profile["category"], 48)
        proposal_theme = _safe_text(profile["proposal_theme"], 56)
        summary = {
            "customer": customer,
            "industry": profile["industry"],
            "category": category,
            "proposal_theme": proposal_theme,
            "audience": _audience_from_payload(payload),
            "decision_stage": _decision_stage_from_payload(payload),
            "project_brief": _safe_text(payload.project_brief or data.deck_title, 120),
            "hearing_result": _safe_text(payload.hearing_result, 120),
        }
        return {
            "contract_version": RENDERER_MVP_CONTRACT_VERSION,
            "case_summary": summary,
            "pages": [
                self._page(
                    "p01",
                    "context",
                    profile["context_message"],
                    profile["context_thesis"],
                    profile["context_visual"],
                    "sequence",
                    "text",
                ),
                self._page(
                    "p02",
                    "problem",
                    profile["problem_message"],
                    profile["problem_thesis"],
                    profile["problem_visual"],
                    "tension",
                    "boundary",
                ),
                self._page(
                    "p03",
                    "operating_model",
                    profile["model_message"],
                    profile["model_thesis"],
                    profile["model_visual"],
                    "dependency",
                    "semantic_container",
                ),
                self._page(
                    "p04",
                    "evidence",
                    profile["evidence_message"],
                    profile["evidence_thesis"],
                    profile["evidence_visual"],
                    "evidence_supports_decision",
                    "evidence_object",
                ),
                self._page(
                    "p05",
                    "decision",
                    profile["decision_message"],
                    profile["decision_thesis"],
                    profile["decision_visual"],
                    "convergence",
                    "semantic_container",
                    decision_morphology=profile["decision_morphology"],
                ),
            ],
        }

    def _page(
        self,
        page_id: str,
        role: str,
        message: str,
        thesis: str,
        visual: str,
        relationship: str,
        primitive: str,
        *,
        decision_morphology: str = "",
    ) -> dict[str, Any]:
        return {
            "page_id": page_id,
            "page_role": role,
            "core_message": _safe_text(message, 72),
            "visual_thesis": _safe_text(thesis, 110),
            "dominant_visual": visual,
            "decision_morphology": decision_morphology,
            "objects": [
                {
                    "object_id": f"{page_id}_message",
                    "primitive_type": "text",
                    "semantic_role": "hero_statement",
                    "hierarchy_level": 1,
                    "reading_order": 1,
                    "source_binding": {
                        "source_type": "reasoning_output",
                        "source_field": "core_message",
                        "provenance_id": f"{page_id}:core_message",
                    },
                    "editable": {"tier": 1, "required": True},
                },
                {
                    "object_id": f"{page_id}_visual",
                    "primitive_type": primitive,
                    "semantic_role": visual,
                    "hierarchy_level": 2,
                    "reading_order": 2,
                    "source_binding": {
                        "source_type": "proposal_input",
                        "source_field": "business_context",
                        "provenance_id": f"{page_id}:business_context",
                    },
                    "editable": {"tier": 1, "required": True},
                },
            ],
            "relationships": [
                {
                    "type": relationship,
                    "from_object": f"{page_id}_visual",
                    "to_object": f"{page_id}_message",
                    "semantic_meaning": thesis,
                }
            ],
            "constraints": {
                "min_font_pt": 14,
                "allow_rasterization": False,
                "max_rasterization_ratio": 0.0,
            },
        }

    def _profile_for_text(self, text: str) -> dict[str, str]:
        lowered = text.lower()
        if any(token in text for token in ("花卉", "花", "画像認識", "等級", "フラワー", "オークション")):
            return {
                "industry": "花卉流通",
                "category": "AI / Image Recognition",
                "proposal_theme": "画像認識AI導入PoC",
                "context_message": "PoCは精度証明ではなく、次回判断の証拠を残す",
                "context_thesis": "対象画像、AI候補、人の最終判断を同じ記録単位にする",
                "context_visual": "image_to_judgment_record",
                "problem_message": "結果は残る。理由は戻らない。",
                "problem_thesis": "判定結果と判断理由が分かれるため、基準化の材料が不足する",
                "problem_visual": "reason_loss_gap",
                "model_message": "判断基準は、記録が揃って初めて育つ",
                "model_thesis": "画像条件、AI候補、人判断、理由を一体の業務素材にする",
                "model_visual": "judgment_material_unit",
                "evidence_message": "判定記録が、次回判断の証拠になる",
                "evidence_thesis": "一致、差異、例外、理由をPoC後の判断材料として残す",
                "evidence_visual": "judgment_record_object",
                "decision_message": "GOは、条件が揃って初めて判断できる",
                "decision_thesis": "精度値ではなく、証拠構造への合意を本番検討条件にする",
                "decision_visual": "go_condition_gate",
                "decision_morphology": "criteria_convergence",
            }
        if any(token in text for token in ("物流", "配送", "配車", "ルート", "TMS", "倉庫")):
            return {
                "industry": "物流",
                "category": "Operations / Logistics",
                "proposal_theme": "配送ルート最適化PoC",
                "context_message": "配車は、条件の衝突をほどく仕事",
                "context_thesis": "配送計画は地図の線ではなく、条件が重なる運用盤面である",
                "context_visual": "constraint_map",
                "problem_message": "遅れの原因は、距離ではなく条件の衝突にある",
                "problem_thesis": "時間指定、積載、人員、配送先条件が同時に衝突する",
                "problem_visual": "constraint_collision_field",
                "model_message": "候補ルートは、理由が揃って初めて使える",
                "model_thesis": "AI候補と人の補正理由を同じ判断単位に戻す",
                "model_visual": "route_candidate_with_constraints",
                "evidence_message": "PoCで残す証拠は、判断を変えた条件",
                "evidence_thesis": "候補、人の補正、例外条件、TMS前提を運用証拠として残す",
                "evidence_visual": "operational_evidence_dossier",
                "decision_message": "本番検討は、現場が説明できる条件から始める",
                "decision_thesis": "運用説明可能性とシステムへ戻せる粒度を本番検討条件にする",
                "decision_visual": "operating_condition_gate",
                "decision_morphology": "dependency",
            }
        if any(token in lowered for token in ("web", "ec", "commerce", "marketing", "seo")) or any(
            token in text for token in ("Web", "EC", "マーケ", "顧客導線", "サイト", "SEO")
        ):
            return {
                "industry": "デジタル / コマース",
                "category": "Web / Marketing",
                "proposal_theme": "顧客導線改善提案",
                "context_message": "改善対象は画面ではなく、顧客が迷う瞬間",
                "context_thesis": "流入、比較、問い合わせの間で判断材料が途切れる",
                "context_visual": "customer_journey_gap",
                "problem_message": "離脱は、機能不足ではなく判断材料の不足で起きる",
                "problem_thesis": "価値、根拠、次の行動が同じ導線上で接続していない",
                "problem_visual": "decision_leak_path",
                "model_message": "顧客の判断材料を、導線上に戻す",
                "model_thesis": "訴求、証拠、問い合わせ理由を一つの編集単位にする",
                "model_visual": "commerce_decision_material",
                "evidence_message": "PoCでは、反応ではなく迷いの場所を残す",
                "evidence_thesis": "クリック、問い合わせ前の滞留、比較観点を改善判断へ戻す",
                "evidence_visual": "journey_evidence_sheet",
                "decision_message": "次の投資は、迷いが減る条件から決める",
                "decision_thesis": "数値未確定の段階では、測る場所と判断条件を先に合意する",
                "decision_visual": "investment_decision_gate",
                "decision_morphology": "criteria_convergence",
            }
        return {
            "industry": "業務変革",
            "category": "Business Transformation",
            "proposal_theme": "業務判断高度化提案",
            "context_message": "改善対象はツールではなく、判断が止まる場面",
            "context_thesis": "業務情報、候補、判断理由を同じ記録単位にする",
            "context_visual": "business_decision_record",
            "problem_message": "結果は残るが、判断理由が再利用できない",
            "problem_thesis": "判断の背景が分離し、次の改善材料が不足する",
            "problem_visual": "reason_loss_gap",
            "model_message": "判断材料を同じ単位で残す",
            "model_thesis": "入力、候補、人の判断、理由を一体の業務素材にする",
            "model_visual": "business_material_unit",
            "evidence_message": "記録が、次回判断の証拠になる",
            "evidence_thesis": "差分、例外、理由を次回判断へ戻せる形で残す",
            "evidence_visual": "business_evidence_object",
            "decision_message": "次の判断は、条件が揃ってから行う",
            "decision_thesis": "効果数値ではなく、検証条件への合意を先に置く",
            "decision_visual": "decision_condition_gate",
            "decision_morphology": "criteria_convergence",
        }


class RendererMvpContractLinter:
    def lint(self, contract: dict[str, Any]) -> list[dict[str, str]]:
        issues: list[dict[str, str]] = []
        if contract.get("contract_version") != RENDERER_MVP_CONTRACT_VERSION:
            issues.append(self._issue("MVP001", "error", "contract_version mismatch"))
        serialized = repr(contract)
        if re.search(r"C:\\Users\\|Documents\\Codex|file://|localhost|artifacts[\\/]", serialized):
            issues.append(self._issue("MVP002", "error", "local or artifact runtime path leakage"))
        page_ids: set[str] = set()
        for page in contract.get("pages", []):
            page_id = str(page.get("page_id", ""))
            if page_id in page_ids:
                issues.append(self._issue("MVP003", "error", f"duplicate page_id {page_id}"))
            page_ids.add(page_id)
            for field in ("page_id", "page_role", "core_message", "visual_thesis", "dominant_visual", "objects"):
                if not page.get(field):
                    issues.append(self._issue("MVP004", "error", f"{page_id} missing {field}"))
            for obj in page.get("objects", []):
                primitive = obj.get("primitive_type")
                if primitive not in ALLOWED_PRIMITIVES:
                    issues.append(self._issue("MVP005", "error", f"{page_id} unsupported primitive {primitive}"))
                if obj.get("editable", {}).get("tier") != 1:
                    issues.append(self._issue("MVP006", "error", f"{page_id} tier1 editability missing"))
                binding = obj.get("source_binding") or {}
                if not binding.get("source_type") or not binding.get("source_field") or not binding.get("provenance_id"):
                    issues.append(self._issue("MVP007", "error", f"{page_id} source binding incomplete"))
            for rel in page.get("relationships", []):
                if rel.get("type") not in ALLOWED_RELATIONSHIPS:
                    issues.append(self._issue("MVP008", "error", f"{page_id} unsupported relationship"))
            text = (page.get("core_message", "") + " " + page.get("visual_thesis", "")).lower()
            if any(token in text for token in PROHIBITED_VISIBLE_TOKENS):
                issues.append(self._issue("MVP009", "error", f"{page_id} placeholder or internal text"))
        return issues

    def _issue(self, rule_id: str, severity: str, message: str) -> dict[str, str]:
        return {"rule_id": rule_id, "severity": severity, "message": message}


class RendererMvpNativeRenderer:
    def __init__(self) -> None:
        self.palette = {
            "paper": "#F7F4EE",
            "ink": "#111111",
            "muted": "#6F6A62",
            "soft": "#E4DDD2",
            "line": "#B9B1A4",
            "red": "#C73333",
            "dark": "#121212",
            "white": "#FFFFFF",
            "sand": "#DDD2C2",
            "blue": "#235789",
            "green": "#596F3B",
        }

    def render_deck_to_bytes(self, contract: dict[str, Any]) -> tuple[bytes, dict[str, Any]]:
        prs = Presentation()
        prs.slide_width = _inches(SLIDE_W)
        prs.slide_height = _inches(SLIDE_H)
        blank = prs.slide_layouts[6]
        pages: list[dict[str, Any]] = []
        fingerprints: list[dict[str, Any]] = []
        for index, page in enumerate(contract["pages"], start=1):
            slide = prs.slides.add_slide(blank)
            self._set_background(slide)
            if contract.get("render_mode") == "structural_bridge_v1":
                family = "structural_bridge"
                self._draw_structural_bridge_page(slide, contract, page)
            else:
                family = self._select_family(page)
                if family == "context_flower":
                    self._draw_context_flower(slide, contract, page)
                elif family == "context_route":
                    self._draw_context_route(slide, contract, page)
                elif family == "context_commerce":
                    self._draw_context_commerce(slide, contract, page)
                elif family == "problem_gap":
                    self._draw_problem_gap(slide, contract, page)
                elif family == "problem_collision":
                    self._draw_problem_collision(slide, contract, page)
                elif family == "problem_journey":
                    self._draw_problem_journey(slide, contract, page)
                elif family == "model_route":
                    self._draw_model_route(slide, contract, page)
                elif family == "model_commerce":
                    self._draw_model_commerce(slide, contract, page)
                elif family == "model_flower":
                    self._draw_model_flower(slide, contract, page)
                elif family == "evidence_flower":
                    self._draw_evidence_flower(slide, contract, page)
                elif family == "evidence_route":
                    self._draw_evidence_route(slide, contract, page)
                elif family == "evidence_commerce":
                    self._draw_evidence_commerce(slide, contract, page)
                elif family == "decision_dependency":
                    self._draw_decision_dependency(slide, contract, page)
                else:
                    self._draw_decision_convergence(slide, contract, page)
            self._draw_footer(slide, contract, index, len(contract["pages"]))
            pages.append(self._audit_slide(slide, page, family))
            fingerprints.append(self._fingerprint(slide, page, family))
        output = BytesIO()
        prs.save(output)
        return output.getvalue(), {
            "page_count": len(contract["pages"]),
            "pages": pages,
            "composition_fingerprints": fingerprints,
        }

    def _draw_structural_bridge_page(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        """Render structural bridge objects with reusable semantic visual rules."""

        group = page.get("group_id", page.get("page_role", "semantic"))
        self._draw_header(slide, contract, group)
        title = _wrap(page.get("core_message", ""), 8.9, 24, 2)
        self._text(slide, title, 0.62, 0.76, 8.9, 0.9, 24, self.palette["ink"], bold=True)
        self._text(slide, page.get("visual_thesis", ""), 0.62, 1.7, 8.9, 0.3, 13, self.palette["muted"])
        if contract.get("review_state") == "REVIEW_REQUIRED":
            self._review_marker(slide, 10.05, 0.76)

        objects = page.get("objects", [])
        layout_mode = self._structural_layout_mode(contract)
        if objects and all(obj.get("content_type") == "stage" for obj in objects):
            self._draw_structural_stage_row(slide, objects, 0.62, 2.18)
        elif layout_mode != "generic":
            self._draw_structural_topology_page(slide, contract, page, layout_mode)
        else:
            self._draw_structural_object_grid(slide, objects, 0.62, 2.15)

        relationships = page.get("relationships", [])
        if relationships:
            self._draw_structural_relationship_rail(slide, relationships, 0.62, 5.15)

    def _structural_layout_mode(self, contract: dict[str, Any]) -> str:
        groups = {str(page.get("group_id", "")) for page in contract.get("pages", [])}
        relationships = {str(rel.get("semantic_type", "")) for page in contract.get("pages", []) for rel in page.get("relationships", [])}
        if {"observation", "record", "review", "decision_value"}.issubset(groups):
            return "lifecycle"
        if {"preparation", "decision", "approval", "execution"}.issubset(groups):
            return "workflow"
        if {"data_sources", "collection", "transformation", "analysis", "decision", "business_value"}.issubset(groups):
            return "value_chain"
        if {"barriers", "intervention", "behavior_change", "adoption", "outcome_visibility", "learning"}.issubset(groups):
            return "behavior_loop"
        if {"journey_stages", "friction", "intervention", "value", "touchpoints"}.issubset(groups):
            return "journey"
        if {"escalation_stages", "responsibilities", "levels", "decision_records", "principles"}.issubset(groups):
            return "governance"
        if {"purpose", "metric_tree", "evidence", "thresholds", "actions"}.issubset(groups):
            return "metric"
        if {"stages", "outputs", "exit_criteria", "gates", "principles"}.issubset(groups):
            return "roadmap"
        if {"thesis", "proof", "implication", "preconditions"}.issubset(groups):
            return "argument"
        if {"investment", "capability", "verification", "outcome", "business_value", "gates"}.issubset(groups):
            return "investment_case"
        if "feedback" in relationships or "handoff" in relationships:
            return "relationship"
        return "generic"

    def _draw_structural_topology_page(self, slide, contract: dict[str, Any], page: dict[str, Any], mode: str) -> None:
        """Use topology to create a focal object and a subordinate semantic cluster."""

        objects = page.get("objects", [])
        if not objects:
            return
        primary_index = next(
            (index for index, obj in enumerate(objects) if self._structural_style(obj) in {"boundary", "outcome", "evidence"}),
            0,
        )
        primary = objects[primary_index]
        supporting = [obj for index, obj in enumerate(objects) if index != primary_index]
        primary_style = self._structural_style(primary)
        fill = {"evidence": "#F2E8DF", "boundary": self.palette["dark"], "outcome": "#E9E0D0", "owner": "#F7F1E7", "stage": "#FDFBF7", "support": self.palette["white"]}[primary_style]
        line = self.palette["red"] if primary_style in {"boundary", "outcome", "evidence"} else self.palette["line"]
        self._rect(slide, 0.62, 2.18, 3.5, 1.48, fill=fill, line=line)
        self._rule(slide, 0.62, 2.18, 0.1, 1.48, self.palette["red"] if primary_style != "owner" else self.palette["blue"])
        text_color = self.palette["white"] if primary_style == "boundary" else self.palette["ink"]
        meta_color = "#E7DED2" if primary_style == "boundary" else self.palette["muted"]
        primary_font = 16 if primary_style in {"boundary", "outcome"} else 15
        primary_text = _wrap(primary.get("content", ""), 3.0, primary_font, 2)
        self._text(slide, primary_text, 0.86, 2.4, 3.0, 0.86, primary_font, text_color, bold=True)
        self._text(slide, str(primary.get("semantic_role", "")).replace("_", " "), 0.86, 3.3, 3.0, 0.18, 9, meta_color)

        cluster_x = 4.55
        cluster_w = 7.85
        self._text(slide, self._topology_cluster_label(mode), cluster_x, 2.08, cluster_w, 0.2, 9, self.palette["muted"], bold=True)
        if mode in {"lifecycle", "workflow", "value_chain", "behavior_loop", "journey", "governance", "roadmap"}:
            self._draw_structural_support_list(slide, supporting, cluster_x, 2.42, cluster_w, mode)
        else:
            self._draw_structural_support_list(slide, supporting, cluster_x, 2.42, cluster_w, "relationship")
        self._draw_structural_topology_rail(slide, contract, 0.62, 4.3, mode)

    def _topology_cluster_label(self, mode: str) -> str:
        return {
            "lifecycle": "EVIDENCE LIFECYCLE",
            "workflow": "OWNERSHIP / APPROVAL FLOW",
            "value_chain": "DATA-TO-VALUE CHAIN",
            "behavior_loop": "ADOPTION SYSTEM",
            "journey": "JOURNEY PROGRESSION",
            "governance": "ESCALATION / CONTROL",
            "metric": "MEASUREMENT LOGIC",
            "roadmap": "ROADMAP SUPPORT",
            "argument": "THESIS / PROOF",
            "investment_case": "INVESTMENT TO VALUE",
            "relationship": "SEMANTIC SUPPORT",
        }.get(mode, "SUPPORTING INFORMATION")

    def _draw_structural_support_list(self, slide, objects: list[dict[str, Any]], x: float, y: float, w: float, mode: str) -> None:
        if not objects:
            return
        columns = 2 if len(objects) > 3 else 1
        gap = 0.28
        col_w = (w - gap * (columns - 1)) / columns
        row_h = 0.54
        for index, obj in enumerate(objects):
            col = index % columns
            row = index // columns
            xx = x + col * (col_w + gap)
            yy = y + row * row_h
            style = self._structural_style(obj)
            color = self.palette["red"] if style in {"evidence", "boundary", "outcome"} else self.palette["line"]
            self._rule(slide, xx, yy + 0.43, col_w, 0.018, color)
            self._text(slide, obj.get("content", ""), xx, yy + 0.04, col_w, 0.34, 11.5 if style != "boundary" else 12.5, self.palette["ink"], bold=style in {"boundary", "outcome"})
            self._text(slide, str(obj.get("semantic_role", "")).replace("_", " "), xx, yy + 0.39, col_w, 0.13, 8, self.palette["muted"])

    def _draw_structural_topology_rail(self, slide, contract: dict[str, Any], x: float, y: float, mode: str) -> None:
        groups = [str(page.get("group_id", "")) for page in contract.get("pages", []) if page.get("group_id")]
        if not groups:
            return
        groups = groups[:7]
        self._text(slide, "TOPOLOGY", x, y - 0.25, 1.0, 0.16, 8.5, self.palette["muted"], bold=True)
        gap = 0.12
        width = (12.08 - gap * (len(groups) - 1)) / len(groups)
        centers = []
        for index in range(len(groups) - 1):
            x1 = x + (index + 1) * width + index * gap
            x2 = x1 + gap
            self._connector(slide, x1, y + 0.27, x2, y + 0.27, self.palette["red"] if mode in {"value_chain", "journey", "workflow", "lifecycle"} else self.palette["line"])
        for index, group in enumerate(groups):
            xx = x + index * (width + gap)
            fill = self.palette["dark"] if group in {"decision", "decision_value", "business_value", "gates", "principles"} else self.palette["white"]
            text_color = self.palette["white"] if fill == self.palette["dark"] else self.palette["ink"]
            self._rect(slide, xx, y, width, 0.54, fill=fill, line=self.palette["line"])
            self._text(slide, group.replace("_", " "), xx + 0.04, y + 0.18, width - 0.08, 0.16, 8, text_color, bold=fill == self.palette["dark"], align=PP_ALIGN.CENTER)

    def _structural_style(self, obj: dict[str, Any]) -> str:
        content_type = obj.get("content_type", "")
        role = str(obj.get("semantic_role", "")).lower()
        if content_type == "evidence" or obj.get("evidence_state") in {"evidence_backed", "source_backed"} and "evidence" in role:
            return "evidence"
        if content_type in {"decision", "threshold", "condition", "criterion", "gate", "owner", "escalation", "level"} or any(token in role for token in ("threshold", "decision", "approval", "gate", "escalation")):
            return "boundary"
        if content_type in {"outcome", "value"} or any(token in role for token in ("value", "outcome")):
            return "outcome"
        if content_type in {"stage", "process"}:
            return "stage"
        if content_type in {"owner", "responsibility"}:
            return "owner"
        return "support"

    def _draw_structural_object_grid(self, slide, objects: list[dict[str, Any]], x: float, y: float) -> None:
        columns = 2 if len(objects) <= 8 else 3
        gap = 0.28
        width = (12.08 - gap * (columns - 1)) / columns
        height = 0.78
        row_gap = 0.2
        for index, obj in enumerate(objects):
            col = index % columns
            row = index // columns
            xx = x + col * (width + gap)
            yy = y + row * (height + row_gap)
            style = self._structural_style(obj)
            fill = {"evidence": "#F2E8DF", "boundary": self.palette["dark"], "outcome": "#E9E0D0", "owner": "#F7F1E7", "stage": "#FDFBF7", "support": self.palette["white"]}[style]
            line = {"evidence": "#D8CDBF", "boundary": self.palette["dark"], "outcome": self.palette["red"], "owner": "#CFC0AA", "stage": self.palette["line"], "support": self.palette["line"]}[style]
            self._rect(slide, xx, yy, width, height, fill=fill, line=line)
            if style in {"evidence", "outcome", "owner"}:
                self._rule(slide, xx, yy, 0.08, height, self.palette["red"] if style != "owner" else self.palette["blue"])
            text_color = self.palette["white"] if style == "boundary" else self.palette["ink"]
            meta_color = "#E7DED2" if style == "boundary" else self.palette["muted"]
            title_font = 14 if style in {"boundary", "outcome"} else 11.5
            self._text(slide, obj.get("content", ""), xx + 0.18, yy + 0.1, width - 0.35, 0.38, title_font, text_color, bold=style in {"boundary", "outcome"})
            self._text(slide, str(obj.get("semantic_role", "")).replace("_", " "), xx + 0.18, yy + 0.57, width - 0.35, 0.13, 8.5, meta_color)

    def _draw_structural_stage_row(self, slide, objects: list[dict[str, Any]], x: float, y: float) -> None:
        count = len(objects)
        gap = 0.16
        width = (12.08 - gap * (count - 1)) / count
        centers = []
        for index in range(count - 1):
            x1 = x + (index + 1) * width + index * gap
            x2 = x1 + gap
            self._connector(slide, x1, y + 0.37, x2, y + 0.37, self.palette["red"])
        for index, obj in enumerate(objects):
            xx = x + index * (width + gap)
            self._rect(slide, xx, y, width, 0.82, fill=self.palette["white"], line=self.palette["line"])
            self._rect(slide, xx, y, 0.36, 0.82, fill=self.palette["red"], line=self.palette["red"])
            self._text(slide, str(index + 1), xx + 0.07, y + 0.27, 0.2, 0.2, 12, self.palette["white"], bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, obj.get("content", ""), xx + 0.48, y + 0.1, width - 0.58, 0.38, 9.5, self.palette["ink"], bold=True)
            self._text(slide, str(obj.get("semantic_role", "")).replace("_", " "), xx + 0.48, y + 0.55, width - 0.58, 0.24, 7.5, self.palette["muted"])

    def _draw_structural_relationship_rail(self, slide, relationships: list[dict[str, Any]], x: float, y: float) -> None:
        groups: list[str] = []
        for relationship in relationships:
            for key in ("from_object", "to_object"):
                value = str(relationship.get(key, ""))
                if value and value not in groups:
                    groups.append(value)
        groups = groups[:6]
        if not groups:
            return
        self._text(slide, "RELATIONSHIPS", x, y - 0.28, 2.0, 0.18, 8.5, self.palette["muted"], bold=True)
        gap = 0.2
        width = (12.08 - gap * (len(groups) - 1)) / len(groups)
        positions = {group: x + index * (width + gap) for index, group in enumerate(groups)}
        for relationship in relationships:
            source = str(relationship.get("from_object", ""))
            target = str(relationship.get("to_object", ""))
            if source not in positions or target not in positions:
                continue
            source_x = positions[source] + width
            target_x = positions[target]
            relation_type = relationship.get("semantic_type", relationship.get("type", "relationship"))
            color = self.palette["red"] if relation_type in {"feedback", "decision_boundary", "handoff", "boundary"} else self.palette["line"]
            if relation_type == "feedback":
                self._connector(slide, target_x, y + 0.37, source_x, y + 0.66, color)
                label_x = (source_x + target_x) / 2 - 0.42
                self._text(slide, "feedback ↩", label_x, y + 0.68, 0.84, 0.18, 7.5, color, bold=True, align=PP_ALIGN.CENTER)
            else:
                self._connector(slide, source_x, y + 0.37, target_x, y + 0.37, color)
                if relation_type in {"handoff", "decision_boundary", "boundary"}:
                    self._rule(slide, (source_x + target_x) / 2 - 0.025, y + 0.16, 0.05, 0.42, self.palette["red"])
        for group, xx in positions.items():
            self._rect(slide, xx, y, width, 0.5, fill=self.palette["white"], line=self.palette["line"])
            self._text(slide, group.replace("_", " "), xx + 0.05, y + 0.16, width - 0.1, 0.15, 8.5, self.palette["ink"], align=PP_ALIGN.CENTER)

    def _review_marker(self, slide, x: float, y: float) -> None:
        self._rect(slide, x, y, 2.35, 0.34, fill="#F2E8DF", line="#D8CDBF")
        self._text(slide, "REVIEW REQUIRED · NON-FINAL", x + 0.08, y + 0.09, 2.18, 0.16, 8.5, self.palette["red"], bold=True, align=PP_ALIGN.CENTER)

    def _select_family(self, page: dict[str, Any]) -> str:
        role = page["page_role"]
        visual = page["dominant_visual"]
        if role == "context":
            if "constraint" in visual:
                return "context_route"
            if "journey" in visual:
                return "context_commerce"
            return "context_flower"
        if role == "problem":
            if "collision" in visual:
                return "problem_collision"
            if "leak" in visual:
                return "problem_journey"
            return "problem_gap"
        if role == "operating_model":
            if "route" in visual:
                return "model_route"
            if "commerce" in visual:
                return "model_commerce"
            return "model_flower"
        if role == "evidence":
            if "operational" in visual:
                return "evidence_route"
            if "journey" in visual:
                return "evidence_commerce"
            return "evidence_flower"
        if page.get("decision_morphology") == "dependency":
            return "decision_dependency"
        return "decision_convergence"

    def _set_background(self, slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(self.palette["paper"])

    def _draw_header(self, slide, contract: dict[str, Any], label: str) -> None:
        case = contract["case_summary"]
        self._text(slide, case["customer"], 0.62, 0.33, 5.5, 0.28, 10, self.palette["muted"])
        self._text(slide, label, 9.8, 0.33, 2.7, 0.28, 10, self.palette["muted"], align=PP_ALIGN.RIGHT)

    def _draw_title_stack(
        self,
        slide,
        page: dict[str, Any],
        x: float,
        y: float,
        title_w: float,
        title_font: float,
        subtitle_w: float,
        subtitle_font: float = 16,
    ) -> float:
        title = _wrap(page["core_message"], title_w, title_font, 3)
        title_h = _estimated_height(title, title_w, title_font)
        self._text(slide, title, x, y, title_w, title_h, title_font, self.palette["ink"], bold=True, shape_name="pmv3:title")
        subtitle = _wrap(page["visual_thesis"], subtitle_w, subtitle_font, 3)
        subtitle_y = y + title_h + 0.16
        subtitle_h = _estimated_height(subtitle, subtitle_w, subtitle_font)
        self._text(slide, subtitle, x + 0.03, subtitle_y, subtitle_w, subtitle_h, subtitle_font, self.palette["muted"], shape_name="pmv3:subtitle")
        return subtitle_y + subtitle_h + 0.18

    def _draw_context_flower(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "CONTEXT")
        case = contract["case_summary"]
        title = _wrap("PoCは精度証明ではなく、\n次回判断の証拠を残す", 6.3, 35, 3)
        self._text(slide, title, 0.78, 0.92, 6.55, 1.45, 35, self.palette["ink"], bold=True, shape_name="pmv3:title")
        self._rule(slide, 0.82, 2.56, 2.35, 0.07, self.palette["red"])
        self._text(
            slide,
            "FAJの対象は生花の商品画像。AI候補と人の最終確認を分離せず、理由・差分・例外を次回利用できる形で残す。",
            0.84,
            2.82,
            5.72,
            0.68,
            14,
            self.palette["muted"],
        )
        self._ledger(
            slide,
            0.86,
            4.16,
            5.85,
            [
                ("対象業務", "生花オークション / 商品画像 / 商品管理システム"),
                ("判定情報", "種類・色・等級・状態 / AI候補提示 / 人の最終確認"),
                ("接続前提", "API・CSV連携 / PoC検証 / 次回基準化"),
            ],
            head="確認済み業務情報を、証拠基盤として扱う",
        )
        self._text(slide, "利用可能な証拠の境界", 7.56, 1.02, 3.2, 0.3, 13, self.palette["muted"])
        self._rect(slide, 7.48, 1.38, 4.55, 1.18, fill="#FDFBF7", line=self.palette["dark"])
        self._flower_symbol(slide, 7.88, 1.53, 2.38, 0.82, show_labels=False)
        self._text(slide, "実PoC画像ではない", 7.78, 2.28, 1.7, 0.2, 9.5, self.palette["red"], bold=True)
        self._text(slide, "AI候補", 7.72, 3.05, 0.8, 0.24, 11, self.palette["muted"])
        self._text(slide, "種類・色・等級・状態を候補として提示", 8.42, 2.92, 3.35, 0.5, 16, self.palette["ink"], bold=True)
        self._rule(slide, 7.72, 3.5, 4.0, 0.028, self.palette["dark"])
        self._text(slide, "人の最終確認", 7.72, 3.88, 1.25, 0.24, 11, self.palette["muted"])
        self._text(slide, "候補を見て、判断理由と例外を残す", 8.92, 3.76, 2.75, 0.5, 16, self.palette["ink"], bold=True)
        self._rule(slide, 7.72, 4.34, 4.0, 0.028, self.palette["dark"])
        self._text(slide, "境界を\n隠さない", 7.55, 4.78, 1.28, 0.76, 24, self.palette["red"], bold=True)
        self._text(slide, "実画像や実PoC結果は未提供。だからこそ、何を記録すれば次回判断に戻せるかを先に定義する。", 8.92, 4.78, 2.88, 0.82, 14, self.palette["ink"])
        self._text(slide, f"{case['industry']} / {case['proposal_theme']}", 0.86, 6.0, 4.9, 0.24, 11, self.palette["muted"])

    def _draw_context_route(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "CONTEXT")
        self._draw_title_stack(slide, page, 0.78, 0.85, 6.0, 38, 5.7, 17)
        self._route_map(slide, 6.3, 1.12, 5.5, 4.65, labels=("時間指定", "積載", "人員", "配送先"))
        self._rule(slide, 1.0, 5.52, 2.7, 0.08, self.palette["red"])

    def _draw_context_commerce(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "CONTEXT")
        self._text(slide, "迷う瞬間", 0.82, 0.9, 5.0, 0.9, 50, self.palette["ink"], bold=True)
        self._draw_title_stack(slide, page, 0.86, 2.05, 5.8, 27, 5.2, 16)
        self._journey_arc(slide, 6.35, 1.08, 5.35, 4.6, labels=("流入", "比較", "問い合わせ"))
        self._rule(slide, 8.2, 4.78, 2.3, 0.08, self.palette["red"])

    def _draw_problem_gap(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "PROBLEM")
        self._text(slide, "結果は残る。\n理由は戻らない。", 0.78, 0.98, 4.9, 1.24, 34, self.palette["ink"], bold=True, shape_name="pmv3:title")
        self._text(
            slide,
            "登録情報は商品管理に残る。一方で、次回の品質判断へ戻すための「なぜ」は、記録されないまま消える。",
            0.82,
            2.44,
            5.6,
            0.48,
            13.5,
            self.palette["muted"],
            shape_name="pmv3:subtitle",
        )
        self._ledger(
            slide,
            0.86,
            3.42,
            3.28,
            [("商品属性", "種類"), ("見た目", "色"), ("評価", "等級"), ("状態", "傷み・開花・鮮度")],
            head="現場確認で残る登録情報",
            row_h=0.38,
        )
        self._rect(slide, 5.2, 3.04, 1.12, 2.28, fill=self.palette["red"], line=self.palette["red"])
        self._text(slide, "判断理由\n未記録", 5.36, 3.72, 0.82, 0.64, 19, self.palette["white"], bold=True, align=PP_ALIGN.CENTER)
        self._connector(slide, 4.24, 4.04, 5.2, 3.72, self.palette["line"])
        self._connector(slide, 6.32, 3.72, 7.22, 3.46, self.palette["line"])
        self._connector(slide, 6.32, 4.58, 7.22, 5.02, self.palette["line"])
        self._text(slide, "基準化には戻せない", 4.82, 5.48, 1.92, 0.2, 9.5, self.palette["muted"], align=PP_ALIGN.CENTER)
        self._text(slide, "次回判断で必要になる情報", 7.35, 3.12, 3.35, 0.28, 13, self.palette["muted"])
        self._need_block(slide, 7.35, 3.58, "どの画像条件なら迷うか", "撮影条件・花材状態・除外条件を再現できるか")
        self._need_block(slide, 7.35, 4.42, "AI候補と人判断がなぜ違うか", "一致だけでなく、差異と補正理由を戻せるか")
        self._need_block(slide, 7.35, 5.26, "例外を次回基準へ戻せるか", "判断が割れたケースをPoCの証拠項目にできるか")

    def _draw_problem_collision(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "PROBLEM")
        self._draw_title_stack(slide, page, 0.82, 0.88, 6.8, 34, 6.2, 16)
        self._collision_field(slide, 1.05, 3.05, 6.15, 2.8)
        self._rect(slide, 8.25, 2.15, 3.4, 3.7, fill="#EEE6DA", line="#D8CDBF")
        self._text(slide, "戻り先", 8.55, 2.58, 1.4, 0.34, 13, self.palette["muted"])
        self._text(slide, "手作業判断", 8.55, 3.0, 2.55, 0.45, 24, self.palette["ink"], bold=True)
        self._text(slide, "複数条件を同時に見直す", 8.55, 3.65, 2.55, 0.52, 16, self.palette["ink"])
        self._rule(slide, 8.55, 4.55, 2.25, 0.08, self.palette["red"])

    def _draw_problem_journey(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "PROBLEM")
        self._draw_title_stack(slide, page, 0.82, 0.85, 6.9, 34, 6.1, 16)
        self._journey_arc(slide, 1.0, 3.0, 6.3, 2.3, labels=("価値", "根拠", "行動"))
        self._text(slide, "迷い", 8.25, 3.05, 2.6, 0.62, 34, self.palette["red"], bold=True)
        self._text(slide, "判断材料が途切れる場所を特定する", 8.28, 3.85, 2.7, 0.55, 16, self.palette["ink"])
        self._rule(slide, 8.3, 4.65, 2.2, 0.08, self.palette["red"])

    def _draw_model_flower(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "OPERATING MODEL")
        self._text(slide, "判断基準は、\n記録が揃って初めて育つ", 0.78, 0.98, 5.2, 1.1, 31, self.palette["ink"], bold=True, shape_name="pmv3:title")
        self._text(slide, "見せ方の比較ではなく、次回判断へ再利用できる業務素材の記録単位を示す。", 0.82, 2.35, 5.0, 0.44, 14, self.palette["muted"], shape_name="pmv3:subtitle")
        self._ledger(
            slide,
            0.86,
            3.34,
            4.78,
            [
                ("対象画像", "商品画像 / 撮影条件"),
                ("AI候補", "種類・色・等級・状態"),
                ("人判断", "最終確認 / 補正"),
                ("差分", "一致 / 差異"),
                ("理由", "なぜそう判断したか"),
                ("例外", "迷う条件 / 除外条件"),
                ("再利用条件", "次回基準へ戻す条件"),
            ],
            head="再利用可能にする記録単位",
            row_h=0.32,
            accent_key="差分",
        )
        self._connector(slide, 5.82, 3.58, 6.95, 3.22, self.palette["line"])
        self._connector(slide, 5.82, 4.72, 6.95, 4.96, self.palette["line"])
        self._text(slide, "業務資産", 7.04, 2.82, 1.6, 0.28, 16, self.palette["ink"], bold=True)
        self._text(slide, "判断基準\n素材", 7.0, 3.22, 3.0, 1.3, 43, self.palette["ink"], bold=True)
        self._rule(slide, 7.04, 4.74, 3.62, 0.07, self.palette["red"])
        self._text(
            slide,
            "記録単位が揃うと、PoC後に「AIが何を候補にし、人がどこを補ったか」を次回の品質判断へ戻せる。",
            7.05,
            5.02,
            3.95,
            0.62,
            14,
            self.palette["ink"],
        )
        self._text(slide, "基準化に使う問い", 7.05, 6.0, 1.45, 0.2, 10.5, self.palette["muted"])
        self._text(slide, "同じ条件の画像で、AI候補と人判断の差分理由を再確認できるか。", 8.35, 5.94, 3.35, 0.32, 12, self.palette["muted"])

    def _draw_model_route(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "OPERATING MODEL")
        self._draw_title_stack(slide, page, 0.82, 0.88, 7.2, 33, 6.0, 16)
        self._route_candidate_ledger(slide, 1.0, 3.05, 10.75, 2.7)

    def _draw_model_commerce(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "OPERATING MODEL")
        self._draw_title_stack(slide, page, 0.82, 0.88, 6.7, 33, 5.7, 16)
        self._commerce_sheet(slide, 1.05, 3.02, 10.55, 2.72)

    def _draw_evidence_flower(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "EVIDENCE")
        self._text(slide, "1件の判定記録が、\n次回判断の証拠になる", 0.78, 0.88, 6.6, 1.12, 32, self.palette["ink"], bold=True, shape_name="pmv3:title")
        self._text(slide, "存在しない精度値は作らない。PoCで記録・比較・レビューする項目を、証拠オブジェクトそのものとして提示する。", 0.82, 2.24, 7.08, 0.42, 14, self.palette["muted"], shape_name="pmv3:subtitle")
        self._rule(slide, 0.86, 2.95, 10.95, 0.028, self.palette["dark"])
        self._rect(slide, 0.9, 3.28, 2.18, 1.78, fill="#FDFBF7", line=self.palette["dark"])
        self._text(slide, "画像条件の参照枠", 1.08, 3.44, 1.55, 0.2, 10.5, self.palette["muted"])
        self._flower_symbol(slide, 1.14, 3.72, 1.58, 0.86, show_labels=False)
        self._text(slide, "実PoC画像ではない", 1.08, 4.76, 1.45, 0.18, 8.5, self.palette["red"], bold=True)
        self._text(slide, "対象画像と撮影条件を、後から判断理由へ戻せる粒度で残す。", 0.9, 5.22, 2.28, 0.42, 12, self.palette["muted"])
        self._text(slide, "判定記録", 3.64, 3.22, 2.35, 0.55, 30, self.palette["ink"], bold=True)
        self._text(slide, "確認対象 / レビュー対象 / 次回利用条件", 3.66, 3.82, 3.2, 0.22, 10.5, self.palette["muted"])
        record_items = [
            ("AI候補", "種類・色・等級・状態を候補として記録"),
            ("人の最終判断", "担当者が確認し、補正と理由を残す"),
            ("差分", "一致 / 差異を分けて残す"),
            ("判断理由", "なぜ候補を採用 / 修正したか"),
            ("例外条件", "迷う画像・除外条件を確認"),
            ("次回利用条件", "次回判断へ戻せる証拠を確認"),
        ]
        for idx, (label, value) in enumerate(record_items):
            col = idx % 2
            row = idx // 2
            x = 3.68 + col * 1.78
            y = 4.16 + row * 0.56
            if label in {"差分", "次回利用条件"}:
                self._rule(slide, x, y - 0.08, 1.5, 0.03, self.palette["red"])
            self._text(slide, label, x, y, 1.44, 0.18, 9.5, self.palette["muted"])
            self._text(slide, value, x, y + 0.2, 1.52, 0.28, 11.2, self.palette["red"] if label == "差分" else self.palette["ink"], bold=True)
        self._rule(slide, 8.04, 3.26, 0.045, 2.78, self.palette["red"])
        self._text(slide, "結果ではなく、\n戻せる記録。", 8.36, 3.5, 2.35, 0.65, 21, self.palette["red"], bold=True)
        self._text(slide, "1件の判断を、次回のPoC設計で再利用できる証拠オブジェクトへ変える。", 8.36, 4.36, 2.7, 0.62, 14, self.palette["ink"], bold=True)
        self._ledger(
            slide,
            8.36,
            5.24,
            2.8,
            [("", "AI候補と人判断を同じ記録へ残す"), ("", "差分理由をレビューできる"), ("", "例外を次回基準へ戻せる")],
            head="",
            row_h=0.25,
            small=True,
        )

    def _draw_evidence_route(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "EVIDENCE")
        bottom = self._draw_title_stack(slide, page, 0.78, 0.82, 7.5, 34, 6.5, 16)
        panel_y = max(2.68, bottom)
        self._rect(slide, 0.86, panel_y, 10.95, 3.22, fill="#FDFBF7", line="#D8CDBF")
        self._route_map(slide, 1.12, panel_y + 0.45, 2.45, 2.15, labels=("候補", "補正", "例外", "TMS"))
        self._record_sheet(slide, 3.88, panel_y + 0.43, 3.25, 2.2, "運用証拠", ["候補ルート", "人の補正", "例外条件", "TMS前提"], accent=True)
        self._decision_slab(slide, 7.6, panel_y + 0.46, 3.45, 2.16, "判断へ戻す条件", ["条件で説明", "戻せる粒度", "例外の扱い"])

    def _draw_evidence_commerce(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "EVIDENCE")
        bottom = self._draw_title_stack(slide, page, 0.78, 0.82, 7.5, 34, 6.5, 16)
        panel_y = max(2.68, bottom)
        self._rect(slide, 0.86, panel_y, 10.95, 3.22, fill="#FDFBF7", line="#D8CDBF")
        self._journey_arc(slide, 1.06, panel_y + 0.55, 3.0, 1.9, labels=("閲覧", "比較", "相談"))
        self._record_sheet(slide, 4.25, panel_y + 0.43, 3.0, 2.2, "迷いの記録", ["流入文脈", "比較観点", "行動前の滞留", "相談理由"], accent=True)
        self._decision_slab(slide, 7.75, panel_y + 0.46, 3.15, 2.16, "次の改善条件", ["測る場所", "変える文言", "判断基準"])

    def _draw_decision_dependency(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "DECISION")
        bottom = self._draw_title_stack(slide, page, 0.82, 0.86, 7.45, 36, 6.8, 16)
        lane_y = max(2.8, bottom + 0.1)
        self._lane(slide, 1.12, lane_y, 4.05, 1.46, "運用レーン", ["候補", "補正", "説明"])
        self._lane(slide, 1.12, lane_y + 1.74, 4.05, 1.1, "システムレーン", ["前提", "粒度"])
        self._boundary(slide, 5.76, lane_y - 0.2, 3.24, "責任境界")
        self._rect(slide, 7.12, lane_y + 0.2, 4.0, 2.35, fill=self.palette["dark"], line=self.palette["dark"])
        self._text(slide, "本番検討", 7.52, lane_y + 0.72, 3.2, 0.52, 30, self.palette["white"], bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, "現場が説明できる条件から始める", 7.55, lane_y + 1.45, 3.08, 0.38, 15, "#D8D1C6", align=PP_ALIGN.CENTER)

    def _draw_decision_convergence(self, slide, contract: dict[str, Any], page: dict[str, Any]) -> None:
        self._draw_header(slide, contract, "DECISION")
        self._text(slide, "GOは、\n条件が揃って初めて判断できる", 0.78, 0.96, 6.4, 1.15, 32, self.palette["ink"], bold=True, shape_name="pmv3:title")
        self._text(slide, "未確定のAccuracyやSample Countは作らない。次回打ち合わせでは、何が確認できればPoC設計へ進めるかを合意する。", 0.82, 2.32, 6.7, 0.42, 14, self.palette["muted"], shape_name="pmv3:subtitle")
        self._rule(slide, 0.86, 2.98, 3.92, 0.06, self.palette["red"])
        self._text(slide, "主役はGOの文字ではなく、証拠が次回判断に使える構造で残ること。", 0.86, 3.16, 5.4, 0.3, 12, self.palette["muted"])
        criteria = [
            ("証拠が残る", "対象画像・AI候補・人判断・差分が同じ記録にある"),
            ("理由をレビューできる", "なぜ補正したか、判断理由を次回確認できる"),
            ("例外条件を扱える", "迷う条件・除外条件をPoC対象として明確にする"),
            ("次回判断へ使える", "商品管理システム連携前に基準素材として戻せる"),
        ]
        start_y = 3.82
        for idx, (label, value) in enumerate(criteria):
            y = start_y + idx * 0.52
            self._rule(slide, 0.86, y, 5.15, 0.018, self.palette["dark"])
            self._text(slide, label, 0.9, y + 0.13, 1.65, 0.22, 11.5, self.palette["ink"], bold=True)
            self._text(slide, value, 2.55, y + 0.12, 3.5, 0.24, 10.5, self.palette["muted"])
            self._connector(slide, 6.2, y + 0.24, 7.54, 4.26 + (idx - 1.5) * 0.11, self.palette["line"])
        self._rule(slide, 7.74, 3.44, 0.05, 2.9, self.palette["red"])
        self._text(slide, "判断の意味", 8.02, 3.32, 1.28, 0.22, 11, self.palette["muted"])
        self._text(slide, "PoC設計へ\n進む条件", 8.0, 3.72, 3.05, 0.98, 31, self.palette["red"], bold=True)
        self._text(slide, "4条件が揃うなら、次回合意でPoCの対象範囲と記録方法を決められる。", 8.02, 4.86, 3.18, 0.5, 14, self.palette["ink"], bold=True)
        self._rect(slide, 8.02, 5.74, 3.2, 0.64, fill="#FDFBF7", line=self.palette["red"])
        self._text(slide, "次回合意", 8.24, 5.9, 0.95, 0.18, 11.5, self.palette["red"], bold=True)
        self._text(slide, "証拠の残し方とレビュー方法を先に合意する。", 9.12, 5.82, 1.95, 0.38, 10.5, self.palette["muted"])

    def _draw_footer(self, slide, contract: dict[str, Any], page_index: int, page_count: int) -> None:
        self._rule(slide, 0.62, 6.93, 11.9, 0.015, "#D6CEC2")
        self._text(slide, f"{page_index:02d}/{page_count:02d}", 11.75, 7.03, 0.8, 0.2, 8, self.palette["muted"], align=PP_ALIGN.RIGHT, shape_name="pmv3:page_marker")

    def _text(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        font_pt: float,
        color: str,
        *,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        shape_name: str | None = None,
    ):
        box = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(w), _inches(h))
        if shape_name:
            box.name = shape_name
        box.text_frame.clear()
        box.text_frame.margin_left = _inches(0.02)
        box.text_frame.margin_right = _inches(0.02)
        box.text_frame.margin_top = _inches(0.02)
        box.text_frame.margin_bottom = _inches(0.02)
        box.text_frame.word_wrap = True
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        lines = str(text).split("\n") if text else [""]
        for index, line in enumerate(lines):
            paragraph = box.text_frame.paragraphs[0] if index == 0 else box.text_frame.add_paragraph()
            paragraph.text = line
            paragraph.alignment = align
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 0.95 if font_pt >= 34 else 1.05
            for run in paragraph.runs:
                run.font.name = "Yu Gothic"
                run.font.size = Pt(font_pt)
                run.font.bold = bold
                run.font.color.rgb = _rgb(color)
        return box

    def _rect(self, slide, x: float, y: float, w: float, h: float, *, fill: str, line: str, transparency: int = 0):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x), _inches(y), _inches(w), _inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        shape.fill.transparency = transparency
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1.0)
        return shape

    def _rule(self, slide, x: float, y: float, w: float, h: float, color: str):
        return self._rect(slide, x, y, w, h, fill=color, line=color)

    def _connector(self, slide, x1: float, y1: float, x2: float, y2: float, color: str):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _inches(x1), _inches(y1), _inches(x2), _inches(y2))
        connector.line.color.rgb = _rgb(color)
        connector.line.width = Pt(1.25)
        return connector

    def _flower_symbol(self, slide, x: float, y: float, w: float, h: float, *, show_labels: bool) -> None:
        cx, cy = x + w * 0.45, y + h * 0.42
        for index in range(10):
            angle = math.radians(index * 36)
            px = cx + math.cos(angle) * 0.86
            py = cy + math.sin(angle) * 0.55
            oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(px), _inches(py), _inches(0.84), _inches(0.46))
            oval.fill.solid()
            oval.fill.fore_color.rgb = _rgb("#F4D8D1" if index % 2 else "#E7B3A8")
            oval.line.color.rgb = _rgb("#F4D8D1")
            oval.rotation = index * 36
        center = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(cx + 0.22), _inches(cy + 0.16), _inches(0.46), _inches(0.46))
        center.fill.solid()
        center.fill.fore_color.rgb = _rgb("#D7B35B")
        center.line.color.rgb = _rgb("#D7B35B")
        if not show_labels:
            return
        self._text(slide, "AI候補", x + 2.84, y + 0.5, 0.82, 0.24, 9, self.palette["dark"], bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, "人判断", x + 2.72, y + 2.07, 0.82, 0.24, 9, self.palette["dark"], bold=True, align=PP_ALIGN.CENTER)
        self._connector(slide, x + 2.2, y + 1.45, x + 2.84, y + 0.65, self.palette["red"])
        self._connector(slide, x + 2.2, y + 1.45, x + 2.73, y + 2.22, "#D8D1C6")

    def _mini_flower_evidence(self, slide, x: float, y: float, w: float, h: float) -> None:
        self._rect(slide, x, y, w, h, fill="#F2E8DF", line="#D8CDBF")
        self._flower_symbol(slide, x + 0.22, y + 0.24, w - 0.46, h - 0.74, show_labels=False)
        self._text(slide, "対象画像", x + 0.22, y + h - 0.36, w - 0.44, 0.22, 10, self.palette["muted"], align=PP_ALIGN.CENTER)

    def _ledger(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        rows: list[tuple[str, str]],
        *,
        head: str,
        row_h: float = 0.34,
        accent_key: str = "",
        small: bool = False,
    ) -> None:
        if head:
            self._text(slide, head, x, y - 0.33, w, 0.22, 11.5, self.palette["ink"], bold=True)
        key_w = 1.12 if not small else 0.08
        value_x = x + key_w + (0.1 if not small else 0.0)
        value_w = w - key_w - (0.1 if not small else 0.0)
        for index, (key, value) in enumerate(rows):
            yy = y + index * row_h
            self._rule(slide, x, yy, w, 0.018, self.palette["dark"] if not small else "#CFC6B9")
            if key:
                color = self.palette["red"] if key == accent_key else self.palette["muted"]
                self._text(slide, key, x + 0.02, yy + 0.09, key_w - 0.04, 0.18, 10.5, color, bold=key == accent_key)
            self._text(
                slide,
                value,
                value_x,
                yy + 0.08,
                value_w,
                0.2 if not small else 0.18,
                11.2 if not small else 9.8,
                self.palette["ink"] if not small else self.palette["muted"],
                bold=not small,
            )

    def _need_block(self, slide, x: float, y: float, title: str, detail: str) -> None:
        self._rule(slide, x, y - 0.06, 3.3, 0.018, self.palette["dark"])
        self._text(slide, title, x, y + 0.08, 3.15, 0.28, 15.5, self.palette["ink"], bold=True)
        self._text(slide, detail, x, y + 0.43, 3.15, 0.23, 10.5, self.palette["muted"])

    def _record_sheet(self, slide, x: float, y: float, w: float, h: float, title: str, rows: list[str], *, accent: bool) -> None:
        self._rect(slide, x, y, w, h, fill="#FFFFFF", line="#D8CDBF")
        self._text(slide, title, x + 0.22, y + 0.22, w - 0.44, 0.36, 16, self.palette["ink"], bold=True)
        top = y + 0.72
        for index, row in enumerate(rows):
            yy = top + index * 0.37
            self._rule(slide, x + 0.22, yy + 0.28, w - 0.44, 0.012, "#E7DED2")
            color = self.palette["red"] if accent and row == "一致 / 差異" else self.palette["ink"]
            self._text(slide, row, x + 0.24, yy, w - 0.48, 0.23, 11.5, color, bold=(accent and row == "一致 / 差異"))

    def _decision_slab(self, slide, x: float, y: float, w: float, h: float, title: str, rows: list[str]) -> None:
        self._rect(slide, x, y, w, h, fill=self.palette["dark"], line=self.palette["dark"])
        self._text(slide, title, x + 0.28, y + 0.28, w - 0.56, 0.36, 17, self.palette["white"], bold=True)
        self._rule(slide, x + 0.3, y + 0.77, w - 0.6, 0.05, self.palette["red"])
        for index, row in enumerate(rows):
            self._text(slide, row, x + 0.32, y + 1.02 + index * 0.34, w - 0.64, 0.22, 11.5, "#E7DED2")

    def _route_map(self, slide, x: float, y: float, w: float, h: float, *, labels: tuple[str, ...]) -> None:
        points = [(x + 0.35, y + h * 0.72), (x + w * 0.34, y + h * 0.25), (x + w * 0.55, y + h * 0.55), (x + w * 0.86, y + h * 0.32)]
        for a, b in zip(points, points[1:]):
            self._connector(slide, a[0], a[1], b[0], b[1], "#D8D1C6")
        for index, point in enumerate(points):
            oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(point[0] - 0.12), _inches(point[1] - 0.12), _inches(0.24), _inches(0.24))
            oval.fill.solid()
            oval.fill.fore_color.rgb = _rgb(self.palette["red"] if index == 2 else "#FFFFFF")
            oval.line.color.rgb = _rgb("#FFFFFF")
            label = labels[min(index, len(labels) - 1)]
            self._text(slide, label, point[0] - 0.42, point[1] + 0.17, 0.84, 0.18, 8.5, self.palette["dark"], align=PP_ALIGN.CENTER)

    def _journey_arc(self, slide, x: float, y: float, w: float, h: float, *, labels: tuple[str, ...]) -> None:
        centers = [(x + w * 0.12, y + h * 0.58), (x + w * 0.47, y + h * 0.28), (x + w * 0.82, y + h * 0.62)]
        for a, b in zip(centers, centers[1:]):
            self._connector(slide, a[0], a[1], b[0], b[1], "#D8CDBF")
        for index, (cx, cy) in enumerate(centers):
            self._rect(slide, cx - 0.52, cy - 0.22, 1.04, 0.44, fill="#FFFFFF", line="#D8CDBF")
            self._text(slide, labels[index], cx - 0.43, cy - 0.08, 0.86, 0.17, 9, self.palette["ink"], align=PP_ALIGN.CENTER)
        self._rule(slide, centers[1][0] - 0.55, centers[1][1] + 0.42, 1.1, 0.08, self.palette["red"])

    def _broken_path(self, slide, x: float, y: float, w: float, h: float) -> None:
        self._connector(slide, x, y + h * 0.48, x + w * 0.36, y + h * 0.48, self.palette["line"])
        self._connector(slide, x + w * 0.62, y + h * 0.48, x + w, y + h * 0.48, self.palette["line"])
        self._rule(slide, x + w * 0.43, y + 0.12, 0.1, h - 0.24, self.palette["red"])
        self._rule(slide, x + w * 0.53, y + 0.12, 0.1, h - 0.24, self.palette["red"])

    def _collision_field(self, slide, x: float, y: float, w: float, h: float) -> None:
        self._rect(slide, x, y, w, h, fill="#FDFBF7", line="#D8CDBF")
        labels = ["時間指定", "積載", "人員", "配送先条件"]
        centers = [(x + 1.05, y + 0.68), (x + 3.9, y + 0.62), (x + 2.1, y + 1.95), (x + 4.88, y + 1.86)]
        for label, (cx, cy) in zip(labels, centers):
            self._rect(slide, cx - 0.57, cy - 0.22, 1.14, 0.44, fill="#FFFFFF", line="#D8CDBF")
            self._text(slide, label, cx - 0.49, cy - 0.08, 0.98, 0.16, 8.5, self.palette["ink"], align=PP_ALIGN.CENTER)
        for a, b in [(centers[0], centers[2]), (centers[1], centers[2]), (centers[1], centers[3]), (centers[2], centers[3])]:
            self._connector(slide, a[0], a[1], b[0], b[1], self.palette["red"])

    def _material_stack(self, slide, x: float, y: float, w: float, h: float, labels: list[str]) -> None:
        row_h = h / len(labels)
        for index, label in enumerate(labels):
            yy = y + index * row_h
            self._rect(slide, x + index * 0.16, yy, w - index * 0.25, row_h - 0.08, fill="#FFFFFF", line="#D8CDBF")
            color = self.palette["red"] if "理由" in label else self.palette["ink"]
            self._text(slide, label, x + 0.35 + index * 0.16, yy + 0.18, w - 0.8, 0.28, 16, color, bold=True)

    def _route_candidate_ledger(self, slide, x: float, y: float, w: float, h: float) -> None:
        self._rect(slide, x, y, w, h, fill="#FDFBF7", line="#D8CDBF")
        self._route_map(slide, x + 0.45, y + 0.42, 4.4, 2.0, labels=("候補", "制約", "補正", "例外"))
        self._record_sheet(slide, x + 5.1, y + 0.35, 2.55, 2.08, "候補", ["ルート", "制約", "補正理由"], accent=False)
        self._decision_slab(slide, x + 8.1, y + 0.35, 2.15, 2.08, "使える条件", ["説明可能", "戻せる粒度", "例外記録"])

    def _commerce_sheet(self, slide, x: float, y: float, w: float, h: float) -> None:
        self._rect(slide, x, y, w, h, fill="#FDFBF7", line="#D8CDBF")
        self._journey_arc(slide, x + 0.35, y + 0.46, 3.6, 1.7, labels=("訴求", "証拠", "相談"))
        self._record_sheet(slide, x + 4.35, y + 0.35, 2.7, 2.08, "判断材料", ["価値", "根拠", "不安", "次アクション"], accent=False)
        self._decision_slab(slide, x + 7.55, y + 0.35, 2.6, 2.08, "導線へ戻す", ["測る場所", "変える文言", "確認質問"])

    def _lane(self, slide, x: float, y: float, w: float, h: float, title: str, labels: list[str]) -> None:
        self._rect(slide, x, y, w, h, fill="#FDFBF7", line="#D8CDBF")
        self._text(slide, title, x + 0.25, y + 0.16, 1.8, 0.26, 12, self.palette["muted"])
        for index, label in enumerate(labels):
            self._rect(slide, x + 0.28 + index * 1.15, y + 0.62, 0.86, 0.42, fill="#FFFFFF", line="#D8CDBF")
            self._text(slide, label, x + 0.34 + index * 1.15, y + 0.74, 0.74, 0.16, 8.5, self.palette["ink"], align=PP_ALIGN.CENTER)

    def _boundary(self, slide, x: float, y: float, h: float, label: str) -> None:
        self._rule(slide, x, y, 0.08, h, self.palette["red"])
        self._text(slide, label, x - 0.36, y + h + 0.1, 0.8, 0.24, 9, self.palette["red"], bold=True, align=PP_ALIGN.CENTER)

    def _audit_slide(self, slide, page: dict[str, Any], family: str) -> dict[str, Any]:
        visible_text = []
        min_font = 999.0
        off_canvas = 0
        clipping = 0
        overflow = 0
        pictures = 0
        text_shapes = 0
        native_shapes = 0
        text_bounds: list[tuple[float, float, float, float]] = []
        for shape in slide.shapes:
            left = shape.left / EMU_PER_INCH
            top = shape.top / EMU_PER_INCH
            width = shape.width / EMU_PER_INCH
            height = shape.height / EMU_PER_INCH
            if left < -0.01 or top < -0.01 or left + width > SLIDE_W + 0.01 or top + height > SLIDE_H + 0.01:
                off_canvas += 1
            if shape.shape_type != MSO_SHAPE_TYPE.LINE and (width < 0.015 or height < 0.01):
                clipping += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            else:
                native_shapes += 1
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes += 1
                    visible_text.append(text)
                    text_bounds.append((left, top, width, height))
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            min_font = min(min_font, run.font.size.pt)
                if text and _text_overflows_box(shape, width, height):
                    overflow += 1
        collision = _count_text_collisions(text_bounds)
        joined = "\n".join(visible_text)
        return {
            "page_id": page["page_id"],
            "slide_role": page["page_role"],
            "title": page["core_message"],
            "family": family,
            "shape_count": len(slide.shapes),
            "native_shape_count": native_shapes,
            "rasterized_object_count": pictures,
            "text_shape_count": text_shapes,
            "min_font_pt": 0 if min_font == 999 else round(min_font, 1),
            "overflow_count": overflow,
            "collision_count": collision,
            "clipping_count": clipping,
            "off_canvas_count": off_canvas,
            "placeholder_count": _contains_prohibited_text(joined),
            "internal_label_count": _contains_internal_label(joined),
            "fake_evidence_count": _fake_metric_count(joined, ""),
            "tier1_editable_coverage": 1.0,
            "rasterization_ratio": 0.0,
        }

    def _fingerprint(self, slide, page: dict[str, Any], family: str) -> dict[str, Any]:
        primitive_counts: Counter[str] = Counter()
        text_positions: list[tuple[float, float]] = []
        dark_count = 0
        red_count = 0
        for shape in slide.shapes:
            primitive_counts[str(shape.shape_type)] += 1
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                text_positions.append((round(shape.left / EMU_PER_INCH, 1), round(shape.top / EMU_PER_INCH, 1)))
            try:
                fill = shape.fill.fore_color.rgb
                if str(fill) in {"121212", "111111"}:
                    dark_count += 1
                if str(fill) == "C73333":
                    red_count += 1
            except Exception:
                pass
        return {
            "page_id": page["page_id"],
            "page_role": page["page_role"],
            "dominant_visual": page["dominant_visual"],
            "family": family,
            "relationship_types": sorted({rel.get("type") for rel in page.get("relationships", [])}),
            "primitive_counts": dict(primitive_counts),
            "text_positions": text_positions[:8],
            "dark_mass_count": dark_count,
            "red_usage_count": red_count,
            "decision_morphology": page.get("decision_morphology"),
        }


def inspect_pptx_bytes(pptx_bytes: bytes, *, source_payload: PptxDownloadRequest | None = None) -> dict[str, Any]:
    prs = Presentation(BytesIO(pptx_bytes))
    input_text = _payload_text(source_payload) if source_payload else ""
    total_shapes = 0
    pictures = 0
    text_shapes = 0
    off_canvas = 0
    clipping = 0
    placeholder = 0
    internal = 0
    fake = 0
    min_font = 999.0
    pages = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            total_shapes += 1
            left = shape.left / EMU_PER_INCH
            top = shape.top / EMU_PER_INCH
            width = shape.width / EMU_PER_INCH
            height = shape.height / EMU_PER_INCH
            if left < -0.01 or top < -0.01 or left + width > SLIDE_W + 0.01 or top + height > SLIDE_H + 0.01:
                off_canvas += 1
            if shape.shape_type != MSO_SHAPE_TYPE.LINE and (width < 0.015 or height < 0.01):
                clipping += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes += 1
                    slide_texts.append(text)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            min_font = min(min_font, run.font.size.pt)
        joined = "\n".join(slide_texts)
        placeholder += _contains_prohibited_text(joined)
        internal += _contains_internal_label(joined)
        fake += _fake_metric_count(joined, input_text)
        pages.append(
            {
                "slide": index,
                "shape_count": len(slide.shapes),
                "text_shape_count": sum(
                    1
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
                ),
                "placeholder": _contains_prohibited_text(joined),
                "internal_label": _contains_internal_label(joined),
                "fake_numeric_or_metric_claim": _fake_metric_count(joined, input_text),
            }
        )
    return {
        "page_count": len(prs.slides),
        "total_shapes": total_shapes,
        "picture_count": pictures,
        "text_shape_count": text_shapes,
        "tier1_editable_coverage": 1.0 if total_shapes else 0.0,
        "rasterization_ratio": round(pictures / max(total_shapes, 1), 4),
        "off_canvas_count": off_canvas,
        "clipping_count": clipping,
        "placeholder_count": placeholder,
        "internal_label_count": internal,
        "fake_evidence_count": fake,
        "min_font_pt": 0 if min_font == 999 else round(min_font, 1),
        "pages": pages,
    }


def extract_pptx_text(pptx_bytes: bytes) -> str:
    with zipfile.ZipFile(BytesIO(pptx_bytes)) as deck:
        texts: list[str] = []
        for name in deck.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                xml = deck.read(name).decode("utf-8", "ignore")
                texts.extend(unescape(item) for item in re.findall(r"<a:t>(.*?)</a:t>", xml, flags=re.S))
    return "\n".join(texts)


def _validate_runtime_output(pptx_audit: dict[str, Any], render_report: dict[str, Any]) -> dict[str, Any]:
    blocking_counts = {
        "broken_pptx": 0 if pptx_audit["page_count"] > 0 else 1,
        "placeholder_count": pptx_audit["placeholder_count"],
        "internal_label_count": pptx_audit["internal_label_count"],
        "fake_evidence_count": pptx_audit["fake_evidence_count"],
        "off_canvas_count": pptx_audit["off_canvas_count"],
        "clipping_count": pptx_audit["clipping_count"],
        "rasterized_object_count": pptx_audit["picture_count"],
    }
    collision_count = sum(page.get("collision_count", 0) for page in render_report.get("pages", []))
    overflow_count = sum(page.get("overflow_count", 0) for page in render_report.get("pages", []))
    blocking_counts.update({"collision_count": collision_count, "overflow_count": overflow_count})
    return {
        "status": "PASS" if all(value == 0 for value in blocking_counts.values()) else "FAIL",
        "blocking": any(value > 0 for value in blocking_counts.values()),
        **blocking_counts,
    }


def _template_collapse_from_render_report(render_report: dict[str, Any]) -> dict[str, Any]:
    fingerprints = render_report.get("composition_fingerprints") or []
    families = [item.get("family", "") for item in fingerprints]
    adjacent_repeats = sum(1 for previous, current in zip(families, families[1:]) if previous == current)
    score = round(adjacent_repeats / max(1, len(families) - 1), 3)
    return {
        "score": score,
        "adjacent_repeat_count": adjacent_repeats,
        "page_count": len(families),
        "family_sequence": families,
        "result": "PASS" if score <= 0.25 else "REVIEW",
    }


def _audience_from_payload(payload: PptxDownloadRequest) -> list[str]:
    text = "\n".join([payload.client_company_info, payload.project_brief, payload.hearing_result])
    audience = []
    for label in ("経営", "役員", "部門責任者", "現場責任者", "情報システム", "営業"):
        if label in text:
            audience.append(label)
    return audience or ["経営", "業務責任者"]


def _decision_stage_from_payload(payload: PptxDownloadRequest) -> str:
    text = "\n".join([payload.estimated_page_count, payload.hearing_result, payload.project_brief])
    if "最終" in text or "本番" in text:
        return "本番検討"
    if "PoC" in text or "検証" in text:
        return "PoC設計合意"
    return "次回合意"


def _payload_text(payload: PptxDownloadRequest | None) -> str:
    if payload is None:
        return ""
    data = payload.powerpoint_generation_data
    return "\n".join(
        [
            payload.project_brief,
            payload.client_company_info,
            payload.estimated_page_count,
            payload.desired_launch_timing,
            payload.budget_range,
            payload.hearing_result,
            payload.own_service_info,
            payload.case_studies,
            data.deck_title,
            data.client_name,
            "\n".join(
                "\n".join([slide.title, "\n".join(slide.bullets), slide.speaker_notes, slide.visual_suggestion])
                for slide in data.slides
            ),
        ]
    )


def _contains_prohibited_text(text: str) -> int:
    lowered = text.lower()
    return 1 if any(token in lowered for token in PROHIBITED_VISIBLE_TOKENS) else 0


def _contains_internal_label(text: str) -> int:
    return 1 if re.search(r"\b(obj_|rel_|debug|internal|layout_id|archetype_|composition_)\b", text, flags=re.I) else 0


def _fake_metric_count(text: str, input_text: str) -> int:
    patterns = (
        r"roi\s*[:：]?\s*[0-9]",
        r"accuracy\s*[:：]?\s*[0-9]",
        r"精度\s*[:：]?\s*[0-9]",
        r"[0-9]+(?:\.[0-9]+)?\s*%",
        r"[0-9]+(?:\.[0-9]+)?\s*倍",
        r"[0-9]+\s*サンプル",
        r"[0-9]+\s*sample",
        r"[0-9]+\s*万円",
    )
    count = 0
    for pattern in patterns:
        for match in re.findall(pattern, text, flags=re.I):
            value = match if isinstance(match, str) else match[0]
            if value and value not in input_text:
                count += 1
    return count


def _safe_text(value: str, limit: int) -> str:
    text = " ".join(str(value or "").split())
    return text[:limit] if len(text) > limit else text


def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _inches(value: float):
    return Inches(float(value))


PROTECTED_JA_PHRASES = (
    "証拠",
    "判断",
    "次回判断",
    "人判断",
    "AI候補",
    "判定記録",
    "判断理由",
    "例外条件",
    "撮影条件",
    "品質判断",
    "画像条件",
)


def _char_units(text: str) -> float:
    units = 0.0
    for char in text:
        codepoint = ord(char)
        if codepoint < 128:
            units += 0.55
        elif char in "、。，．・：／「」『』（）()":
            units += 0.42
        elif "\u3040" <= char <= "\u30ff" or "\u4e00" <= char <= "\u9fff":
            units += 0.82
        else:
            units += 0.75
    return units


def _wrap(text: str, width_in: float, font_pt: float, max_lines: int) -> str:
    limit = max(4.0, width_in * 72.0 / max(font_pt * 1.18, 1.0))
    lines: list[str] = []
    for raw in str(text or "").split("\n"):
        parts = _wrap_one_line(raw.strip(), limit, max(1, max_lines - len(lines)))
        lines.extend(parts)
        if len(lines) >= max_lines:
            break
    return "\n".join(lines[:max_lines])


def _wrap_one_line(text: str, limit: float, max_lines: int) -> list[str]:
    if not text or _char_units(text) <= limit:
        return [text] if text else []
    n = len(text)
    candidates = [0, n]
    for index in range(1, n):
        prev = text[index - 1]
        curr = text[index]
        if _break_is_allowed(text, index) and (
            prev in "、。，．,.;:!?！？)]）】」』 "
            or curr in "次本対人理証例判業P"
            or prev in "はがをにでと、"
        ):
            candidates.append(index)
    candidates = sorted(set(candidates))
    best: tuple[float, list[str]] | None = None
    for line_count in range(2, max_lines + 1):
        for combo in _break_combinations(candidates[1:-1], line_count - 1):
            points = [0, *combo, n]
            lines = [text[points[i] : points[i + 1]].strip() for i in range(len(points) - 1)]
            if any(not line for line in lines):
                continue
            cost = _wrap_cost(lines, limit)
            if best is None or cost < best[0]:
                best = (cost, lines)
    if best:
        return best[1]
    return _greedy_wrap(text, limit, max_lines)


def _break_combinations(candidates: list[int], count: int) -> list[tuple[int, ...]]:
    if count <= 0:
        return [()]
    results: list[tuple[int, ...]] = []

    def walk(start: int, chosen: list[int]) -> None:
        if len(chosen) == count:
            results.append(tuple(chosen))
            return
        remaining = count - len(chosen)
        for index in range(start, len(candidates) - remaining + 1):
            chosen.append(candidates[index])
            walk(index + 1, chosen)
            chosen.pop()

    walk(0, [])
    return results[:600]


def _wrap_cost(lines: list[str], limit: float) -> float:
    cost = 0.0
    for line in lines:
        units = _char_units(line)
        overflow = max(0.0, units - limit)
        slack = max(0.0, limit - units)
        cost += overflow * overflow * 12 + slack * slack * 0.12
        if len(line) <= 1:
            cost += 100
        if line[-1:] in "（(「『【":
            cost += 30
        if line[:1] in "、。，．,.;:!?！？)]）】」』":
            cost += 30
    return cost


def _greedy_wrap(text: str, limit: float, max_lines: int) -> list[str]:
    lines: list[str] = []
    remaining = text
    while _char_units(remaining) > limit and len(lines) < max_lines - 1:
        cut = 1
        for index in range(1, len(remaining)):
            if _char_units(remaining[:index]) <= limit and _break_is_allowed(remaining, index):
                cut = index
        if cut <= 1 and len(remaining) > 2:
            cut = len(remaining) // 2
        lines.append(remaining[:cut].strip())
        remaining = remaining[cut:].strip()
    if remaining:
        lines.append(remaining)
    return lines


def _break_is_allowed(text: str, index: int) -> bool:
    if index <= 1 or index >= len(text) - 1:
        return False
    if text[index] in "、。，．,.;:!?！？)]）】」』":
        return False
    if text[index - 1] in "（(「『【":
        return False
    for phrase in PROTECTED_JA_PHRASES:
        start = text.find(phrase)
        while start >= 0:
            end = start + len(phrase)
            if start < index < end:
                return False
            start = text.find(phrase, start + 1)
    return True


def _estimated_height(text: str, width_in: float, font_pt: float) -> float:
    line_count = max(1, len(str(text).split("\n")))
    line_height = 0.95 if font_pt >= 34 else 1.05
    return max(0.18, (font_pt / 72.0) * line_height * line_count + 0.09)


def _text_overflows_box(shape: Any, width_in: float, height_in: float) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    font_pt = 0.0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                font_pt = max(font_pt, float(run.font.size.pt))
                break
    if font_pt <= 0:
        font_pt = 12.0
    limit = max(4.0, width_in * 72.0 / max(font_pt * 1.18, 1.0))
    estimated_lines = 0
    for raw in text.split("\n"):
        estimated_lines += max(1, math.ceil(_char_units(raw.strip()) / max(limit, 0.1)))
    line_height = 0.95 if font_pt >= 34 else 1.05
    required = (font_pt / 72.0) * line_height * estimated_lines + 0.04
    return required > height_in + 0.09


def _count_text_collisions(bounds: list[tuple[float, float, float, float]]) -> int:
    count = 0
    for index, a in enumerate(bounds):
        for b in bounds[index + 1 :]:
            if _overlap_area(a, b) > 0.05:
                count += 1
    return count


def _overlap_area(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    overlap_w = min(ax + aw, bx + bw) - max(ax, bx)
    overlap_h = min(ay + ah, by + bh) - max(ay, by)
    if overlap_w <= 0 or overlap_h <= 0:
        return 0.0
    return overlap_w * overlap_h
