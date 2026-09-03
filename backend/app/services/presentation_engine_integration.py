from __future__ import annotations

from dataclasses import dataclass
import hashlib
import logging
from time import perf_counter
from typing import Any, TYPE_CHECKING

from app.config import settings
from app.models import PptxDownloadRequest
from app.services.pptx_service import build_pptx_context, build_pptx_result

if TYPE_CHECKING:
    from app.strategy_engine.models import HumanReviewReport, PresentationContext


logger = logging.getLogger(__name__)

ENGINE_MODE_LEGACY = "legacy"
ENGINE_MODE_STRATEGY_V1 = "strategy_v1"
ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1 = "presentation_director_v10_1"
ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1 = "presentation_design_master_v1"
ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP = "presentation_master_v3_renderer_mvp"
SUPPORTED_ENGINE_MODES = {ENGINE_MODE_LEGACY, ENGINE_MODE_STRATEGY_V1}
PRESENTATION_DIRECTOR_V10_1_FLAG = "PRESENTATION_DESIGN_AI_V10_ENABLED"
PRESENTATION_DESIGN_MASTER_FLAG = "PRESENTATION_DESIGN_AI_MASTER_ENABLED"
PRESENTATION_MASTER_V3_RENDERER_MVP_FLAG = "PRESENTATION_MASTER_V3_RENDERER_MVP_ENABLED"
PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_FLAG = "PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_ENABLED"
PRESENTATION_MASTER_V3_RENDERER_MVP_CANARY_FLAG = "PRESENTATION_MASTER_V3_RENDERER_MVP_CANARY_ENABLED"
PRESENTATION_MASTER_V3_RENDERER_MVP_AUTO_FALLBACK_FLAG = "PRESENTATION_MASTER_V3_RENDERER_MVP_AUTO_FALLBACK_ENABLED"
PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_MAX_WORKERS = "PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_MAX_WORKERS"
PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_MAX_PENDING = "PRESENTATION_MASTER_V3_RENDERER_MVP_SHADOW_MAX_PENDING"



@dataclass(frozen=True)
class PresentationEngineResult:
    pptx_bytes: bytes
    engine_mode: str
    presentation_context: Any | None = None
    quality_report: dict[str, Any] | None = None


class RendererMvpInternalCanaryDisabled(RuntimeError):
    reason_code = "renderer_mvp_internal_canary_disabled"
    failure_stage = "routing"


class RendererMvpInternalCanaryError(RuntimeError):
    def __init__(self, reason_code: str, *, failure_stage: str, fallback_category: str):
        self.reason_code = reason_code
        self.failure_stage = failure_stage
        self.fallback_category = fallback_category
        super().__init__(reason_code)


def resolve_engine_mode(value: str | None = None) -> str:
    mode = (value or settings.presentation_engine_mode or ENGINE_MODE_LEGACY).strip().lower()
    if mode in SUPPORTED_ENGINE_MODES:
        return mode
    logger.warning("presentation_engine_mode_invalid", extra={"engine_mode": mode})
    return ENGINE_MODE_LEGACY


def build_pptx_bytes_for_engine(
    payload: PptxDownloadRequest,
    *,
    engine_mode: str | None = None,
    shadow_master: bool = False,
    renderer_mvp_canary: bool = False,
    request_id: str | None = None,
    project_id: str | None = None,
) -> PresentationEngineResult:
    """Build the authoritative Primary result, then optionally observe it in Shadow."""
    primary = _build_primary_pptx_bytes_for_engine(
        payload,
        engine_mode=engine_mode,
        shadow_master=shadow_master,
        renderer_mvp_canary=renderer_mvp_canary,
        request_id=request_id,
        project_id=project_id,
    )
    return _submit_production_shadow_after_primary(
        primary,
        payload,
        request_id=request_id,
        project_id=project_id,
    )


def _build_primary_pptx_bytes_for_engine(
    payload: PptxDownloadRequest,
    *,
    engine_mode: str | None = None,
    shadow_master: bool = False,
    renderer_mvp_canary: bool = False,
    request_id: str | None = None,
    project_id: str | None = None,
) -> PresentationEngineResult:
    from app.services.presentation_master import MASTER_MODE_ENABLED, MASTER_MODE_SHADOW, resolve_master_runtime_mode

    started = perf_counter()
    if getattr(settings, "presentation_master_v3_renderer_mvp_enabled", False):
        return _build_renderer_mvp_with_fallback(
            payload,
            request_id=request_id,
            project_id=project_id,
            routing_mode="enabled",
            semantic_gate=True,
        )

    if getattr(settings, "presentation_master_v3_renderer_mvp_canary_enabled", False) and renderer_mvp_canary:
        return _build_renderer_mvp_with_fallback(
            payload,
            request_id=request_id,
            project_id=project_id,
            routing_mode="canary",
        )

    master_mode = resolve_master_runtime_mode(
        enabled=settings.presentation_design_ai_master_enabled,
        shadow_enabled=getattr(settings, "presentation_design_ai_master_shadow_enabled", False),
        explicit_shadow=shadow_master,
    )
    if master_mode == MASTER_MODE_ENABLED:
        result = _build_master_with_fallback(payload, request_id=request_id, project_id=project_id)
        return result
    shadow_enabled = master_mode == MASTER_MODE_SHADOW

    if settings.presentation_design_ai_v10_enabled:
        result = _build_v10_1_with_fallback(payload)
        engine_result = _attach_shadow_result(
            result,
            payload,
            shadow_master=shadow_enabled,
            request_id=request_id,
            project_id=project_id,
        )
        return engine_result

    mode = resolve_engine_mode(engine_mode)
    if mode == ENGINE_MODE_LEGACY:
        _log_engine_selection(mode)
        result = _build_legacy_pptx_result(
            payload,
            request_id=request_id,
            project_id=project_id,
            requested_version=ENGINE_MODE_LEGACY,
            started=started,
        )
        engine_result = PresentationEngineResult(
            pptx_bytes=result.pptx_bytes,
            engine_mode=mode,
            quality_report=result.quality_report.to_dict(),
        )
        engine_result = _attach_shadow_result(
            engine_result,
            payload,
            shadow_master=shadow_enabled,
            request_id=request_id,
            project_id=project_id,
        )
        return engine_result

    presentation_context = _presentation_context_from_review_report(payload.strategy_review_report)
    _log_engine_selection(mode, presentation_context)
    result = build_pptx_result(payload, presentation_context=presentation_context)
    engine_result = PresentationEngineResult(
        pptx_bytes=result.pptx_bytes,
        engine_mode=mode,
        presentation_context=presentation_context,
        quality_report=result.quality_report.to_dict(),
    )
    engine_result = _attach_shadow_result(
        engine_result,
        payload,
        shadow_master=shadow_enabled,
        request_id=request_id,
        project_id=project_id,
    )
    return engine_result


_PRODUCTION_SHADOW_CONTROLLER: Any | None = None
_READINESS_CLASSES = frozenset(
    {
        "READY",
        "READY_WITH_VALID_BINDINGS",
        "REVIEW_REQUIRED",
        "NOT_READY",
        "NO_MATCH",
        "INVALID_INPUT",
        "ADAPTER_ERROR",
    }
)
_INVALID_INPUT_REASONS = frozenset(
    {
        "PRODUCTION_REQUEST_INVALID",
        "CONFIRMATION_TRANSPORT_INVALID",
        "SEMANTIC_SUPPLY_INVALID",
        "SEMANTIC_PREPARATION_INVALID",
        "MASTER_SELECTION_INVALID",
        "COMPOSITION_INVALID",
        "RENDER_PREPARATION_INVALID",
        "ADAPTER_VALIDATION_ERROR",
    }
)
_SEMANTIC_SUPPLY_INVALID_REASONS = frozenset(
    {
        "NO_CANDIDATES",
        "REVIEWED_BUT_AUTHORITY_NOT_ADMISSIBLE",
        "ADMISSIBILITY_STATE_UNCLASSIFIED",
    }
)
_BOUNDED_READINESS_STAGES = frozenset(
    {"INPUT_ADAPTER", "SEMANTIC_ADAPTER", "SEMANTIC_RESOLUTION", "MASTER_SELECTION", "COMPOSITION", "RENDER_PREP"}
)
_BOUNDED_SELECTION_STATES = frozenset({"selected", "review_required", "no_match"})
_BOUNDED_COMPOSITION_STATES = frozenset({"VALID", "DEGRADED", "REVIEW_REQUIRED", "INVALID", "NOT_READY"})
_BOUNDED_DIAGNOSTIC_CODES = _INVALID_INPUT_REASONS | _SEMANTIC_SUPPLY_INVALID_REASONS | {"renderer_preparation"}


def _submit_production_shadow_after_primary(
    primary: PresentationEngineResult,
    payload: PptxDownloadRequest,
    *,
    request_id: str | None,
    project_id: str | None,
) -> PresentationEngineResult:
    """Best-effort, default-off observation seam; Primary remains authoritative."""
    if not getattr(settings, "presentation_master_v3_renderer_mvp_shadow_enabled", False):
        return primary
    if payload.summary:
        return primary
    raw_candidates = getattr(payload, "semantic_candidates", None)
    candidate_items = raw_candidates.get("candidates") if isinstance(raw_candidates, dict) else None
    candidate_count = len(candidate_items) if isinstance(candidate_items, list) else 0
    candidate_state = "OMITTED" if raw_candidates is None else ("NONEMPTY" if candidate_count else "EMPTY")
    correlation_id = getattr(payload, "candidate_boundary_correlation_id", None) or (hashlib.sha256(request_id.encode()).hexdigest()[:16] if request_id else "missing")
    try:
        _log_shadow_metadata(
            f"presentation_candidate_boundary_backend semantic_candidates_state={candidate_state} candidate_count={candidate_count} correlation_id={correlation_id}",
            semantic_candidates_state=candidate_state,
            candidate_count=candidate_count,
            correlation_id=correlation_id,
        )
    except Exception:
        pass
    decision_emitted = False

    def emit_decision(
        decision: str,
        reason: str,
        readiness_class: str | None = None,
        invalid_input_reason: str | None = None,
        semantic_supply_invalid_reason: str | None = None,
        readiness_metadata: dict[str, Any] | None = None,
    ) -> None:
        nonlocal decision_emitted
        if not decision_emitted:
            _log_shadow_eligibility_decision(
                request_id,
                decision,
                reason,
                readiness_class=readiness_class,
                invalid_input_reason=invalid_input_reason,
                semantic_supply_invalid_reason=semantic_supply_invalid_reason,
                readiness_metadata=readiness_metadata,
            )
            decision_emitted = True

    if not request_id:
        emit_decision("INELIGIBLE", "MISSING_REQUEST_ID")
        return primary
    hook_correlation_id = hashlib.sha256(request_id.encode()).hexdigest()[:16]
    try:
        _log_shadow_metadata(f"presentation_shadow_hook_entered correlation_id={hook_correlation_id}", correlation_id=hook_correlation_id)
    except Exception:
        pass
    if not isinstance(primary.pptx_bytes, bytes) or not primary.pptx_bytes:
        emit_decision("INELIGIBLE", "INVALID_PRIMARY_BYTES")
        return primary
    try:
        from app.services.presentation_master.integration import (
            ShadowController,
            ShadowProcessWorkload,
            ShadowJob,
            build_candidate_state_bridge,
            prepare_pmv3,
        )

        if payload.semantic_candidates is None or payload.semantic_confirmation_state is None:
            emit_decision("INELIGIBLE", "MISSING_CANDIDATE_STATE")
            return primary
        try:
            context = build_candidate_state_bridge(payload, request_id=request_id)
        except Exception:
            emit_decision("INELIGIBLE", "CANDIDATE_BINDING_FAILED")
            return primary
        if context is None:
            emit_decision("INELIGIBLE", "CANDIDATE_BINDING_FAILED")
            return primary
        try:
            prepared = prepare_pmv3(payload, semantic_candidates=context.binding.candidates)
        except Exception:
            emit_decision("INELIGIBLE", "PREPARATION_FAILED")
            return primary
        readiness_metadata = _bounded_readiness_metadata(prepared, candidate_count=candidate_count)
        try:
            eligibility = ShadowController.eligibility(
                summary=payload.summary,
                confirmation_state_present=True,
                prepared_status=prepared.status.value,
                selected_master=prepared.selected_master_id or "",
                composition_status=prepared.composition_readiness,
            )
        except Exception:
            emit_decision("INELIGIBLE", "INTERNAL_PRE_ADMISSION_ERROR")
            return primary
        invalid_input_reason = None
        semantic_supply_invalid_reason = None
        diagnostics = getattr(prepared, "diagnostics", {})
        if prepared.status.value == "INVALID_INPUT" and isinstance(diagnostics, dict):
            candidate_reason = diagnostics.get("invalid_input_reason")
            if candidate_reason in _INVALID_INPUT_REASONS:
                invalid_input_reason = candidate_reason
            supply_reason = diagnostics.get("semantic_supply_invalid_reason")
            if supply_reason in _SEMANTIC_SUPPLY_INVALID_REASONS:
                semantic_supply_invalid_reason = supply_reason
        if not eligibility.eligible:
            reason_map = {
                "SEMANTIC_REVIEW_REQUIRED": "READINESS_NOT_ELIGIBLE",
                "NOT_ELIGIBLE": "READINESS_NOT_ELIGIBLE",
                "SELECTION_NO_MATCH": "MASTER_NOT_M48",
                "COMPOSITION_INVALID": "COMPOSITION_INVALID",
            }
            emit_decision(
                "INELIGIBLE",
                reason_map.get(eligibility.reason, "INTERNAL_PRE_ADMISSION_ERROR"),
                readiness_class=prepared.status.value,
                invalid_input_reason=invalid_input_reason,
                semantic_supply_invalid_reason=semantic_supply_invalid_reason,
                readiness_metadata=readiness_metadata,
            )
            return primary
        global _PRODUCTION_SHADOW_CONTROLLER
        if _PRODUCTION_SHADOW_CONTROLLER is None:
            _PRODUCTION_SHADOW_CONTROLLER = ShadowController(enabled=True, event_logger=_log_shadow_metadata)
        stop_event = getattr(_PRODUCTION_SHADOW_CONTROLLER, "_stop", None)
        if not getattr(_PRODUCTION_SHADOW_CONTROLLER, "enabled", True) or (
            stop_event is not None and stop_event.is_set()
        ):
            emit_decision("INELIGIBLE", "CONTROLLER_DISABLED")
            return primary
        job = ShadowJob(
            request_id=request_id,
            primary_engine=primary.engine_mode,
            semantic_readiness=prepared.semantic_readiness,
            selected_master=prepared.selected_master_id or "",
            composition_status=prepared.composition_readiness,
            workload=ShadowProcessWorkload(payload=payload, binding=context.binding),
        )
        emit_decision("ELIGIBLE", "ELIGIBLE", readiness_class=prepared.status.value, readiness_metadata=readiness_metadata)
        submitted = _PRODUCTION_SHADOW_CONTROLLER.submit(job, eligibility=eligibility)
        del submitted
    except Exception:
        emit_decision("INELIGIBLE", "INTERNAL_PRE_ADMISSION_ERROR")
    return primary


def _log_shadow_eligibility_decision(
    request_id: str | None,
    decision: str,
    reason: str,
    *,
    readiness_class: str | None = None,
    invalid_input_reason: str | None = None,
    semantic_supply_invalid_reason: str | None = None,
    readiness_metadata: dict[str, Any] | None = None,
) -> None:
    """Emit one bounded, opaque eligibility decision without affecting Primary."""
    try:
        fields: dict[str, Any] = {"decision": decision, "reason": reason}
        message = f"presentation_shadow_eligibility_decision decision={decision} reason={reason}"
        if readiness_class in _READINESS_CLASSES:
            fields["readiness_class"] = readiness_class
            message += f" readiness_class={readiness_class}"
        if readiness_class == "INVALID_INPUT" and invalid_input_reason in _INVALID_INPUT_REASONS:
            fields["invalid_input_reason"] = invalid_input_reason
            message += f" invalid_input_reason={invalid_input_reason}"
        if (
            readiness_class == "INVALID_INPUT"
            and invalid_input_reason == "SEMANTIC_SUPPLY_INVALID"
            and semantic_supply_invalid_reason in _SEMANTIC_SUPPLY_INVALID_REASONS
        ):
            fields["semantic_supply_invalid_reason"] = semantic_supply_invalid_reason
            message += f" semantic_supply_invalid_reason={semantic_supply_invalid_reason}"
        if request_id:
            correlation_id = hashlib.sha256(request_id.encode()).hexdigest()[:16]
            fields["correlation_id"] = correlation_id
            message += f" correlation_id={correlation_id}"
        if isinstance(readiness_metadata, dict):
            for key, value in readiness_metadata.items():
                if value is not None:
                    fields[key] = value
        _log_shadow_metadata(message, **fields)
    except Exception:
        return


def _bounded_readiness_metadata(prepared: Any, *, candidate_count: int) -> dict[str, Any]:
    """Extract only existing, bounded adapter metadata for the eligibility event."""
    metadata: dict[str, Any] = {"candidate_count": max(0, min(int(candidate_count), 100000))}
    try:
        stage = getattr(getattr(prepared, "fallback_stage", None), "value", None)
        if stage in _BOUNDED_READINESS_STAGES:
            metadata["fallback_stage"] = stage

        selection_state = getattr(getattr(prepared, "selection", None), "state", None)
        if selection_state in _BOUNDED_SELECTION_STATES:
            metadata["selection_state"] = selection_state

        composition_status = getattr(prepared, "composition_readiness", None)
        if composition_status in _BOUNDED_COMPOSITION_STATES:
            metadata["composition_status"] = composition_status

        provenance = getattr(prepared, "provenance_summary", None)
        if isinstance(provenance, dict):
            bounded_provenance = {
                str(key): max(0, min(int(value), 100000))
                for key, value in provenance.items()
                if isinstance(key, str) and isinstance(value, int) and value >= 0
            }
            if bounded_provenance:
                metadata["provenance_counts"] = bounded_provenance

        diagnostics = getattr(prepared, "diagnostics", None)
        if isinstance(diagnostics, dict):
            codes = tuple(
                value
                for value in diagnostics.values()
                if isinstance(value, str) and value in _BOUNDED_DIAGNOSTIC_CODES
            )
            if codes:
                metadata["allowlisted_diagnostic_codes"] = codes[:16]
            metadata["diagnostic_count"] = min(len(diagnostics), 32)
            if diagnostics.get("semantic_supply_invalid_reason") in _SEMANTIC_SUPPLY_INVALID_REASONS:
                metadata["semantic_supply_status"] = "INVALID"
    except Exception:
        return {"candidate_count": metadata["candidate_count"]}
    return metadata


def _log_shadow_metadata(event: str, **fields: Any) -> None:
    """Application logger only; logger failures must never affect Primary."""
    try:
        logger.info(event, extra={key: value for key, value in fields.items() if key != "candidate_values"})
    except Exception:
        return


def build_renderer_mvp_internal_canary_pptx_bytes(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
) -> PresentationEngineResult:
    requested_version = ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP
    started = perf_counter()
    category = _category_from_payload(payload)
    if not getattr(settings, "presentation_master_v3_renderer_mvp_canary_enabled", False):
        logger.warning(
            "v3_internal_canary_disabled",
            extra={
                "requested_version": requested_version,
                "actual_version": "",
                "fallback_used": False,
                "fallback_reason": RendererMvpInternalCanaryDisabled.reason_code,
                "failure_stage": RendererMvpInternalCanaryDisabled.failure_stage,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        raise RendererMvpInternalCanaryDisabled()

    logger.info(
        "v3_internal_canary_requested",
        extra={
            "requested_version": requested_version,
            "actual_version": "",
            "fallback_used": False,
            "fallback_reason": "",
            "category": category,
            "request_id": request_id or "",
            "project_id": project_id or "",
        },
    )
    try:
        result = _build_renderer_mvp_pptx_result(payload, request_id=request_id, project_id=project_id)
        if result.engine_mode != requested_version:
            from app.services.presentation_master.renderer_mvp import RendererMvpIntegrationError

            raise RendererMvpIntegrationError(
                "renderer_mvp_unexpected_engine",
                failure_stage="internal_canary_validation",
            )
        if result.pptx_bytes[:2] != b"PK":
            from app.services.presentation_master.renderer_mvp import RendererMvpIntegrationError

            raise RendererMvpIntegrationError(
                "renderer_mvp_malformed_pptx",
                failure_stage="pptx_validation",
            )
        duration_ms = round((perf_counter() - started) * 1000)
        quality_report = dict(result.quality_report or {})
        pptx_audit = quality_report.get("pptx_audit") or {}
        quality_report.update(
            {
                "requested_version": requested_version,
                "actual_version": requested_version,
                "routing_mode": "internal_canary",
                "internal_canary": True,
                "canary_success": True,
                "fallback_used": False,
                "fallback_reason": "",
                "fallback_category": "",
                "failure_stage": "",
                "feature_flag": PRESENTATION_MASTER_V3_RENDERER_MVP_CANARY_FLAG,
                "request_id": request_id or "",
                "project_id": project_id or "",
            }
        )
        logger.info(
            "v3_internal_canary_success",
            extra={
                "requested_version": requested_version,
                "actual_version": requested_version,
                "fallback_used": False,
                "fallback_reason": "",
                "category": category,
                "slide_count": int(pptx_audit.get("page_count") or quality_report.get("page_count") or 0),
                "generation_time_ms": duration_ms,
                "renderer_latency_ms": duration_ms,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        return PresentationEngineResult(
            pptx_bytes=result.pptx_bytes,
            engine_mode=requested_version,
            presentation_context=result.presentation_context,
            quality_report=quality_report,
        )
    except Exception as exc:
        from app.services.presentation_master import fallback_category_for_exception

        reason = getattr(exc, "reason_code", exc.__class__.__name__)
        failure_stage = getattr(exc, "failure_stage", "internal_canary_generation")
        fallback_category = fallback_category_for_exception(exc)
        duration_ms = round((perf_counter() - started) * 1000)
        logger.warning(
            "v3_internal_canary_failure",
            extra={
                "requested_version": requested_version,
                "actual_version": "",
                "fallback_used": False,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": failure_stage,
                "category": category,
                "generation_time_ms": duration_ms,
                "renderer_latency_ms": duration_ms,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        _log_renderer_mvp_block_event(
            reason=reason,
            failure_stage=failure_stage,
            request_id=request_id,
            project_id=project_id,
            routing_mode="internal_canary",
            shadow_enabled=False,
        )
        raise RendererMvpInternalCanaryError(
            reason,
            failure_stage=failure_stage,
            fallback_category=fallback_category,
        ) from exc


def _build_legacy_pptx_result(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None,
    project_id: str | None,
    requested_version: str,
    started: float,
) -> Any:
    try:
        result = build_pptx_result(payload)
        total_request_latency_ms = round((perf_counter() - started) * 1000)
        logger.info(
            "legacy_success",
            extra={
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "total_request_latency_ms": total_request_latency_ms,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        return result
    except Exception:
        total_request_latency_ms = round((perf_counter() - started) * 1000)
        logger.exception(
            "legacy_failure",
            extra={
                "requested_version": requested_version,
                "actual_version": "",
                "total_request_latency_ms": total_request_latency_ms,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        raise


def _build_renderer_mvp_with_fallback(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
    routing_mode: str,
    semantic_gate: bool = False,
) -> PresentationEngineResult:
    requested_version = ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP
    started = perf_counter()
    try:
        result = _build_renderer_mvp_pptx_result(
            payload,
            request_id=request_id,
            project_id=project_id,
            semantic_gate=semantic_gate,
        )
        if routing_mode == "canary":
            logger.info(
                "v3_canary_success",
                extra={
                    "requested_version": requested_version,
                    "actual_version": requested_version,
                    "fallback_used": False,
                    "fallback_reason": "",
                    "request_id": request_id or "",
                    "project_id": project_id or "",
                },
            )
            quality_report = dict(result.quality_report or {})
            quality_report.update({"routing_mode": "canary", "canary_success": True})
            return PresentationEngineResult(
                pptx_bytes=result.pptx_bytes,
                engine_mode=result.engine_mode,
                presentation_context=result.presentation_context,
                quality_report=quality_report,
            )
        return result
    except Exception as exc:
        from app.services.presentation_master import fallback_category_for_exception

        if not getattr(settings, "presentation_master_v3_renderer_mvp_auto_fallback_enabled", True):
            raise
        reason = getattr(exc, "reason_code", exc.__class__.__name__)
        failure_stage = getattr(exc, "failure_stage", "renderer_mvp_generation")
        fallback_category = (
            "unsupported"
            if reason == "summary_deck_uses_existing_renderer"
            else fallback_category_for_exception(exc)
        )
        logger.warning(
            "v3_canary_fallback" if routing_mode == "canary" else "presentation_master_v3_renderer_mvp_fallback",
            extra={
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": failure_stage,
                "request_id": request_id or "",
                "project_id": project_id or "",
                "routing_mode": routing_mode,
            },
        )
        _log_renderer_mvp_block_event(
            reason=reason,
            failure_stage=failure_stage,
            request_id=request_id,
            project_id=project_id,
            routing_mode=routing_mode,
            shadow_enabled=False,
        )
        result = _build_legacy_pptx_result(
            payload,
            request_id=request_id,
            project_id=project_id,
            requested_version=requested_version,
            started=started,
        )
        quality_report = result.quality_report.to_dict()
        quality_report.update(
            {
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": failure_stage,
                "request_id": request_id or "",
                "project_id": project_id or "",
                "routing_mode": routing_mode,
                "canary_success": False if routing_mode == "canary" else None,
            }
        )
        return PresentationEngineResult(
            pptx_bytes=result.pptx_bytes,
            engine_mode=ENGINE_MODE_LEGACY,
            quality_report=quality_report,
        )


def _build_renderer_mvp_pptx_result(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
    semantic_gate: bool = False,
) -> PresentationEngineResult:
    if semantic_gate:
        from app.services.presentation_master.integration import (
            build_candidate_state_bridge,
            prepare_pmv3,
            render_pmv3,
        )

        context = build_candidate_state_bridge(payload, request_id=request_id or "")
        if context is None:
            raise RendererMvpIntegrationError(
                "semantic_readiness_not_ready",
                failure_stage="semantic_readiness",
            )
        prepared = prepare_pmv3(payload, semantic_candidates=context.binding.candidates)
        if prepared.status.value not in {"READY", "READY_WITH_VALID_BINDINGS"}:
            raise RendererMvpIntegrationError(
                "semantic_readiness_not_ready",
                failure_stage="semantic_readiness",
                details={"readiness_status": prepared.status.value},
            )
        rendered = render_pmv3(prepared)
        return PresentationEngineResult(
            pptx_bytes=rendered.pptx_bytes,
            engine_mode=ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
            quality_report={
                "requested_version": ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
                "actual_version": ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
                "fallback_used": False,
                "fallback_reason": "",
                "feature_flag": PRESENTATION_MASTER_V3_RENDERER_MVP_FLAG,
                "request_id": request_id or "",
                "project_id": project_id or "",
                "validation_status": rendered.validation_status,
                "slide_count": rendered.slide_count,
                "rasterization_ratio": rendered.rasterization_ratio,
                "clipping_count": rendered.clipping_count,
                "overflow_count": rendered.overflow_count,
                "off_canvas_count": rendered.off_canvas_count,
            },
        )

    from app.services.presentation_master.renderer_mvp import build_renderer_mvp_pptx

    renderer_result = build_renderer_mvp_pptx(payload, request_id=request_id, project_id=project_id)
    logger.info(
        "presentation_master_v3_renderer_mvp_selected",
        extra={
            "requested_version": ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
            "actual_version": ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
            "fallback_used": False,
            "fallback_reason": "",
            "request_id": request_id or "",
            "project_id": project_id or "",
        },
    )
    return PresentationEngineResult(
        pptx_bytes=renderer_result.pptx_bytes,
        engine_mode=ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
        quality_report=renderer_result.quality_report,
    )


def _log_renderer_mvp_block_event(
    *,
    reason: str,
    failure_stage: str,
    request_id: str | None,
    project_id: str | None,
    routing_mode: str,
    shadow_enabled: bool,
) -> None:
    event_name = _renderer_mvp_block_event_name(reason, failure_stage)
    if not event_name:
        return
    logger.warning(
        event_name,
        extra={
            "requested_version": ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
            "actual_version": "",
            "shadow_enabled": shadow_enabled,
            "fallback_reason": reason,
            "failure_stage": failure_stage,
            "request_id": request_id or "",
            "project_id": project_id or "",
            "routing_mode": routing_mode,
        },
    )


def _renderer_mvp_block_event_name(reason: str, failure_stage: str) -> str:
    lowered = f"{reason} {failure_stage}".lower()
    if "evidence" in lowered:
        return "v3_evidence_block"
    if "qa_block" in lowered or failure_stage == "visual_qa":
        return "v3_qa_block"
    return ""


def _compare_legacy_and_v3(
    payload: PptxDownloadRequest,
    baseline_result: PresentationEngineResult,
    v3_result: PresentationEngineResult,
) -> dict[str, Any]:
    from app.services.presentation_master.renderer_mvp import extract_pptx_text

    customer = _customer_name_from_payload(payload)
    legacy_text = extract_pptx_text(baseline_result.pptx_bytes)
    v3_text = extract_pptx_text(v3_result.pptx_bytes)
    legacy_slide_count = _pptx_slide_count(baseline_result.pptx_bytes)
    v3_slide_count = _pptx_slide_count(v3_result.pptx_bytes)
    return {
        "information_loss_candidate": v3_slide_count < min(legacy_slide_count, 5),
        "slide_count_delta": v3_slide_count - legacy_slide_count,
        "customer_name_loss": bool(customer and customer in legacy_text and customer not in v3_text),
        "summary_loss": payload.summary and v3_result.engine_mode == ENGINE_MODE_PRESENTATION_MASTER_V3_RENDERER_MVP,
        "estimate_kpi_handling": "no_fake_numeric_claims",
        "download_compatibility": v3_result.pptx_bytes[:2] == b"PK",
    }


def _pptx_slide_count(pptx_bytes: bytes) -> int:
    try:
        from io import BytesIO
        from pptx import Presentation

        return len(Presentation(BytesIO(pptx_bytes)).slides)
    except Exception:
        return 0


def _qa_blocking_count(quality_report: dict[str, Any]) -> int:
    return sum(
        int(quality_report.get(key) or 0)
        for key in (
            "architecture_deviation_count",
            "fake_evidence_count",
            "placeholder_internal_label_count",
        )
    )


def _render_count(quality_report: dict[str, Any], key: str) -> int:
    pages = (quality_report.get("render_report") or {}).get("pages") or []
    return sum(int(page.get(key) or 0) for page in pages if isinstance(page, dict))


def _case_identifier(payload: PptxDownloadRequest) -> str:
    raw = "|".join(
        [
            _customer_name_from_payload(payload),
            _category_from_payload(payload),
            payload.powerpoint_generation_data.deck_title,
        ]
    )
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:12]


def _build_master_with_fallback(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
) -> PresentationEngineResult:
    requested_version = ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1
    try:
        return _build_master_pptx_result(payload, request_id=request_id, project_id=project_id)
    except Exception as exc:
        from app.services.presentation_master import fallback_category_for_exception

        reason = getattr(exc, "reason_code", exc.__class__.__name__)
        failure_stage = getattr(exc, "failure_stage", "master_generation")
        fallback_category = fallback_category_for_exception(exc)
        logger.warning(
            "presentation_design_master_v1_fallback",
            extra={
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": failure_stage,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        result = build_pptx_result(payload)
        quality_report = result.quality_report.to_dict()
        quality_report.update(
            {
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": failure_stage,
                "request_id": request_id or "",
                "project_id": project_id or "",
            }
        )
        return PresentationEngineResult(
            pptx_bytes=result.pptx_bytes,
            engine_mode=ENGINE_MODE_LEGACY,
            quality_report=quality_report,
        )


def _build_master_pptx_result(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None = None,
    project_id: str | None = None,
) -> PresentationEngineResult:
    from app.services.presentation_master import build_presentation_master

    master_result = build_presentation_master(
        payload,
        core_builder=_build_v10_1_pptx_result,
        request_id=request_id,
        project_id=project_id,
    )
    logger.info(
        "presentation_design_master_v1_selected",
        extra={
            "requested_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
            "actual_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
            "fallback_used": False,
            "fallback_reason": "",
            "request_id": request_id or "",
            "project_id": project_id or "",
        },
    )
    return PresentationEngineResult(
        pptx_bytes=master_result.pptx_bytes,
        engine_mode=ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
        quality_report=master_result.quality_report,
    )


def _attach_shadow_result(
    result: PresentationEngineResult,
    payload: PptxDownloadRequest,
    *,
    shadow_master: bool,
    request_id: str | None,
    project_id: str | None,
) -> PresentationEngineResult:
    if not shadow_master:
        return result

    shadow_report = _run_master_shadow(payload, request_id=request_id, project_id=project_id)
    quality_report = dict(result.quality_report or {})
    quality_report["shadow_result"] = shadow_report
    return PresentationEngineResult(
        pptx_bytes=result.pptx_bytes,
        engine_mode=result.engine_mode,
        presentation_context=result.presentation_context,
        quality_report=quality_report,
    )


def _run_master_shadow(
    payload: PptxDownloadRequest,
    *,
    request_id: str | None,
    project_id: str | None,
) -> dict[str, Any]:
    started = perf_counter()
    try:
        result = _build_master_pptx_result(payload, request_id=request_id, project_id=project_id)
        generation_time_ms = round((perf_counter() - started) * 1000)
        report = _shadow_report_from_quality(
            payload,
            quality_report=result.quality_report or {},
            request_id=request_id,
            project_id=project_id,
            success=True,
            generation_time_ms=generation_time_ms,
        )
        logger.info(
            "presentation_design_master_v1_shadow_success",
            extra={
                "requested_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
                "actual_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
                "shadow_enabled": True,
                "shadow_success": True,
                "fallback_used": False,
                "fallback_reason": "",
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        return report
    except Exception as exc:
        from app.services.presentation_master import fallback_category_for_exception

        generation_time_ms = round((perf_counter() - started) * 1000)
        reason = exc.__class__.__name__
        fallback_category = fallback_category_for_exception(exc)
        logger.warning(
            "presentation_design_master_v1_shadow_failure",
            extra={
                "requested_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
                "actual_version": "",
                "shadow_enabled": True,
                "shadow_success": False,
                "fallback_used": False,
                "fallback_reason": reason,
                "fallback_category": fallback_category,
                "failure_stage": "shadow_master_generation",
                "failure_reason": reason,
                "request_id": request_id or "",
                "project_id": project_id or "",
            },
        )
        return _shadow_report_from_quality(
            payload,
            quality_report={},
            request_id=request_id,
            project_id=project_id,
            success=False,
            generation_time_ms=generation_time_ms,
            failure_stage="shadow_master_generation",
            failure_reason=reason,
            fallback_category=fallback_category,
        )


def _shadow_report_from_quality(
    payload: PptxDownloadRequest,
    *,
    quality_report: dict[str, Any],
    request_id: str | None,
    project_id: str | None,
    success: bool,
    generation_time_ms: int,
    failure_stage: str = "",
    failure_reason: str = "",
    fallback_category: str = "",
) -> dict[str, Any]:
    forms = _visual_forms_from_quality(quality_report)
    return {
        "request_id": request_id or "",
        "project_id": project_id or "",
        "customer": _customer_name_from_payload(payload),
        "category": _category_from_payload(payload),
        "audience": _decision_maker_hint(payload),
        "sales_stage": "production_shadow",
        "deck_objective": _deck_objective_from_payload(payload),
        "requested_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
        "actual_version": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1 if success else "",
        "shadow_enabled": True,
        "shadow_success": success,
        "fallback_used": False,
        "fallback_reason": "",
        "fallback_category": "" if success else (fallback_category or "unexpected_error"),
        "mode": "shadow",
        "engine_requested": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1,
        "engine_used": ENGINE_MODE_PRESENTATION_DESIGN_MASTER_V1 if success else "",
        "page_count": int(quality_report.get("director_slide_count") or 0),
        "story_strategy": "presentation_director_master_v1",
        "selected_visual_forms": forms,
        "composition_fingerprints": _composition_fingerprints_from_quality(quality_report),
        "template_repetition_score": _template_repetition_score(forms),
        "editable_tier1_coverage": 1.0 if success else 0.0,
        "editable_tier2_coverage": 1.0 if success else 0.0,
        "visual_qa_result": "PASS" if success else "FAIL",
        "generation_time_ms": generation_time_ms,
        "failure_stage": failure_stage,
        "failure_reason": failure_reason,
    }


def _visual_forms_from_quality(quality_report: dict[str, Any]) -> list[str]:
    pages = (quality_report.get("render_report") or {}).get("pages") or []
    forms: list[str] = []
    for page in pages:
        if isinstance(page, dict):
            forms.append(str(page.get("slide_role") or page.get("layout") or page.get("title") or "unknown"))
    return forms


def _composition_fingerprints_from_quality(quality_report: dict[str, Any]) -> list[dict[str, Any]]:
    pages = (quality_report.get("render_report") or {}).get("pages") or []
    fingerprints: list[dict[str, Any]] = []
    for index, page in enumerate(pages, start=1):
        if isinstance(page, dict):
            fingerprints.append(
                {
                    "page": index,
                    "form": str(page.get("slide_role") or page.get("layout") or "unknown"),
                    "title_length": len(str(page.get("title") or "")),
                }
            )
    return fingerprints


def _template_repetition_score(forms: list[str]) -> int:
    if not forms:
        return 0
    repeats = sum(1 for previous, current in zip(forms, forms[1:]) if previous == current)
    return round((repeats / max(1, len(forms) - 1)) * 100)


def _customer_name_from_payload(payload: PptxDownloadRequest) -> str:
    data = payload.powerpoint_generation_data
    context = build_pptx_context(payload)
    return _safe_text(data.client_name or context.client_name or payload.client_company_info)


def _category_from_payload(payload: PptxDownloadRequest) -> str:
    data = payload.powerpoint_generation_data
    context = build_pptx_context(payload)
    return _safe_text(context.proposal_category or data.deck_title or "Proposal")


def _deck_objective_from_payload(payload: PptxDownloadRequest) -> str:
    data = payload.powerpoint_generation_data
    return _safe_text(_first_non_empty(payload.project_brief, payload.hearing_result, data.deck_title), 120)


def _build_v10_1_with_fallback(payload: PptxDownloadRequest) -> PresentationEngineResult:
    requested_version = ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1
    try:
        if payload.summary:
            raise ValueError("presentation_director_v10_1 does not handle summary decks.")
        return _build_v10_1_pptx_result(payload)
    except Exception as exc:
        reason = exc.__class__.__name__
        logger.warning(
            "presentation_director_v10_1_fallback",
            extra={
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
            },
        )
        result = build_pptx_result(payload)
        quality_report = result.quality_report.to_dict()
        quality_report.update(
            {
                "requested_version": requested_version,
                "actual_version": ENGINE_MODE_LEGACY,
                "fallback_used": True,
                "fallback_reason": reason,
            }
        )
        return PresentationEngineResult(
            pptx_bytes=result.pptx_bytes,
            engine_mode=ENGINE_MODE_LEGACY,
            quality_report=quality_report,
        )


def _build_v10_1_pptx_result(payload: PptxDownloadRequest) -> PresentationEngineResult:
    from app.presentation_composer import render_director_plan_to_pptx
    from app.presentation_design_ai import design_presentation_deck
    from app.presentation_director import build_directed_presentation_plan, direct_case, validate_director_plan

    case = _case_context_from_payload(payload)
    design_deck = design_presentation_deck(case, force_enabled=True)
    director_plan = direct_case(case, presentation_time_minutes=_presentation_time_minutes(payload))
    validation_errors = validate_director_plan(director_plan)
    if validation_errors:
        raise ValueError(f"presentation_director_v10_1 validation failed: {validation_errors[0]}")
    presentation_plan = build_directed_presentation_plan(case, director_plan)
    pptx_bytes, render_report = render_director_plan_to_pptx(presentation_plan, director_plan)
    quality_report = {
        "requested_version": ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1,
        "actual_version": ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1,
        "fallback_used": False,
        "fallback_reason": "",
        "feature_flag": PRESENTATION_DIRECTOR_V10_1_FLAG,
        "director_recommended_page_count": director_plan.recommended_page_count,
        "director_slide_count": len(director_plan.slide_sequence),
        "design_contract_count": len(design_deck.slide_contracts),
        "render_report": render_report,
    }
    logger.info(
        "presentation_director_v10_1_selected",
        extra={
            "requested_version": ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1,
            "actual_version": ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1,
            "fallback_used": False,
            "fallback_reason": "",
        },
    )
    return PresentationEngineResult(
        pptx_bytes=pptx_bytes,
        engine_mode=ENGINE_MODE_PRESENTATION_DIRECTOR_V10_1,
        quality_report=quality_report,
    )


def _case_context_from_payload(payload: PptxDownloadRequest):
    from app.presentation_composer import CaseContext

    data = payload.powerpoint_generation_data
    context = build_pptx_context(payload)
    slides_text = " ".join(
        " ".join([slide.title, " ".join(slide.bullets), slide.visual_suggestion])
        for slide in data.slides
    )
    project_summary = _first_non_empty(payload.project_brief, payload.hearing_result, slides_text, data.deck_title)
    pain_points = _items_from_text(
        _first_non_empty(payload.hearing_result, payload.project_brief, slides_text),
        fallback=("現状課題を整理", "判断材料を明確化", "次アクションを合意"),
        limit=3,
    )
    expected_outcomes = _items_from_text(
        _first_non_empty(payload.own_service_info, payload.case_studies, slides_text),
        fallback=("合意形成を短縮", "提案品質を安定化", "運用負荷を抑制"),
        limit=3,
    )
    if payload.win_probability:
        expected_outcomes = tuple(
            dict.fromkeys(
                list(expected_outcomes)
                + payload.win_probability.positive_factors[:1]
                + payload.win_probability.recommended_next_actions[:1]
            )
        )[:3]
    return CaseContext(
        case_id="production_v10_1",
        case_name=_safe_text(data.deck_title or "提案書"),
        client_name=_safe_text(data.client_name or context.client_name or "顧客企業"),
        industry=_safe_text(_industry_hint(payload) or context.proposal_label or "提案対象"),
        category=_safe_text(context.proposal_category or data.deck_title or "Proposal"),
        project_summary=_safe_text(project_summary),
        pain_points=tuple(_safe_text(item) for item in pain_points),
        expected_outcomes=tuple(_safe_text(item) for item in expected_outcomes),
        budget=_safe_text(payload.budget_range),
        timeline=_safe_text(payload.desired_launch_timing or payload.estimated_page_count),
        decision_maker=_safe_text(_decision_maker_hint(payload)),
        competitor=_safe_text(payload.competitor_company_name or payload.competitor_site_url),
    )


def _presentation_time_minutes(payload: PptxDownloadRequest) -> int:
    text = " ".join([payload.estimated_page_count, payload.hearing_result, payload.project_brief])
    for pattern in (r"(\d+)\s*分", r"(\d+)\s*min", r"(\d+)\s*minutes"):
        match = __import__("re").search(pattern, text, flags=__import__("re").IGNORECASE)
        if match:
            return max(5, min(60, int(match.group(1))))
    return 30


def _decision_maker_hint(payload: PptxDownloadRequest) -> str:
    text = " ".join([payload.client_company_info, payload.hearing_result, payload.project_brief])
    for word in ("経営者", "社長", "役員", "部門責任者", "部長", "情報システム", "IT責任者", "現場責任者"):
        if word in text:
            return word
    return "経営者 / 部門責任者"


def _industry_hint(payload: PptxDownloadRequest) -> str:
    import re

    text = "\n".join([payload.client_company_info, payload.project_brief, payload.hearing_result])
    for pattern in (r"業界[:：]\s*([^\n\r]+)", r"業種[:：]\s*([^\n\r]+)"):
        match = re.search(pattern, text)
        if match:
            return match.group(1).strip()
    return ""


def _items_from_text(text: str, *, fallback: tuple[str, ...], limit: int) -> tuple[str, ...]:
    import re

    candidates = [
        item.strip(" ・-:：\t\r\n")
        for item in re.split(r"[\n。;；、]", text)
        if item.strip(" ・-:：\t\r\n")
    ]
    cleaned = [_safe_text(item, 42) for item in candidates if len(item.strip()) >= 3]
    return tuple(dict.fromkeys(cleaned[:limit] or list(fallback)[:limit]))


def _first_non_empty(*values: str) -> str:
    return next((value for value in values if value and value.strip()), "")


def _safe_text(value: str, limit: int = 80) -> str:
    text = " ".join(str(value or "").split())
    return text[:limit] if len(text) > limit else text


def _presentation_context_from_review_report(value: dict | None):
    report = _parse_review_report(value)
    if report is None:
        raise ValueError("strategy_v1 requires an approved strategy_review_report.")
    from app.strategy_engine import adapter as strategy_adapter

    return strategy_adapter.adapt_review_report_to_presentation_context(report)


def _parse_review_report(value: Any) -> Any | None:
    if value is None:
        return None
    from app.strategy_engine.models import HumanReviewReport

    if isinstance(value, HumanReviewReport):
        return value
    return HumanReviewReport(**value)


def _log_engine_selection(mode: str, presentation_context: Any | None = None) -> None:
    extra: dict[str, str | None] = {
        "engine_mode": mode,
        "strategy_version": None,
        "presentation_context_version": None,
        "presentation_pack": None,
        "story": None,
        "persona": None,
    }
    if presentation_context is not None:
        extra.update(
            {
                "strategy_version": presentation_context.source_strategy_schema_version,
                "presentation_context_version": presentation_context.schema_version,
                "presentation_pack": str(presentation_context.presentation_pack),
                "story": presentation_context.story_type,
                "persona": presentation_context.persona,
            }
        )
    logger.info("presentation_engine_selected", extra=extra)
