from __future__ import annotations

from dataclasses import asdict, dataclass, field
from time import perf_counter
import re
from typing import Any, Literal

from pptx import Presentation

from app.models import PowerPointSlide
from app.services.pptx_parts.models import PptxContext
from app.services.pptx_theme import SLIDE_HEIGHT, SLIDE_WIDTH
from app.services.pptx_design_system.typography import (
    find_internal_label_leaks,
    normalize_customer_facing_text,
    normalize_customer_facing_title,
    text_density_score,
)

Severity = Literal["info", "warning", "critical"]
DiagramType = Literal["comparison", "timeline", "roadmap", "kpi", "flow", "matrix"]

RULE_TITLE = "PPT-TITLE-001"
RULE_BODY = "PPT-BODY-001"
RULE_BULLET = "PPT-BULLET-001"
RULE_OVERFLOW = "PPT-OVERFLOW-001"
RULE_LAYOUT = "PPT-LAYOUT-001"
RULE_DIAGRAM = "PPT-DIAGRAM-001"
RULE_COMPARE = "PPT-COMPARE-001"
RULE_NUMERIC = "PPT-NUMERIC-001"
RULE_RENDER = "PPT-RENDER-001"
RULE_PLACEHOLDER = "PPT-CGV3-LABEL-001"
RULE_DENSITY = "PPT-CGV3-DENSITY-001"
RULE_CUSTOMER_NAME = "PPT-CGV3-CUSTOMER-001"


@dataclass(frozen=True)
class TemplateQualityRule:
    title_max_chars: int
    body_max_chars: int
    bullet_max_items: int
    min_font_size: int
    margin: float
    max_colors: int
    max_cards: int
    max_table_columns: int
    preferred_diagrams: tuple[DiagramType, ...]
    forbidden_layouts: tuple[str, ...] = ()
    closing_style: str = "next_action"


@dataclass(frozen=True)
class QualityFinding:
    rule_id: str
    category: str
    severity: Severity
    message: str
    recommendation: str
    slide_no: int | None = None
    slide_title: str = ""
    auto_fixable: bool = False
    before: str = ""
    after: str = ""
    confidence: float = 0.8
    human_review_required: bool = False


@dataclass(frozen=True)
class PptxQualityReport:
    overall_score: int
    category_scores: dict[str, int]
    findings: list[QualityFinding] = field(default_factory=list)
    auto_fixes_applied: list[QualityFinding] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    human_review_required: bool = False
    slide_count_before: int = 0
    slide_count_after: int = 0
    template: str = "corporate_clean"
    generation_duration: float = 0.0
    layout_decisions: list[dict[str, Any]] = field(default_factory=list)
    layout_fallbacks: list[dict[str, Any]] = field(default_factory=list)
    preview_pptx_differences: list[dict[str, Any]] = field(default_factory=list)
    predicted_score: int | None = None
    rendered_score: int | None = None
    score_delta: int | None = None
    unsupported_layouts: list[str] = field(default_factory=list)
    numeric_integrity: dict[str, Any] = field(default_factory=dict)
    template_token_application: dict[str, Any] = field(default_factory=dict)
    human_review_items: list[str] = field(default_factory=list)
    customer_ready_status: str = ""
    customer_ready_score: int | None = None
    customer_ready_reasons: list[str] = field(default_factory=list)
    customer_ready_blockers: list[str] = field(default_factory=list)
    customer_ready_auto_fixes: list[str] = field(default_factory=list)
    customer_ready_excluded_internal_items: list[str] = field(default_factory=list)
    customer_ready_sales_summary: list[str] = field(default_factory=list)
    customer_ready_expected_questions: list[dict[str, str]] = field(default_factory=list)
    customer_ready_rubric: dict[str, int] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        payload = asdict(self)
        payload["findings"] = payload["findings"][:30]
        payload["auto_fixes_applied"] = payload["auto_fixes_applied"][:30]
        payload["warnings"] = payload["warnings"][:30]
        return payload


@dataclass(frozen=True)
class QualityPipelineResult:
    slides: list[PowerPointSlide]
    report: PptxQualityReport


TEMPLATE_QUALITY_RULES: dict[str, TemplateQualityRule] = {
    "corporate_clean": TemplateQualityRule(38, 460, 5, 16, 0.82, 5, 4, 4, ("comparison", "flow", "kpi")),
    "modern_dark": TemplateQualityRule(34, 460, 5, 15, 0.82, 4, 3, 4, ("kpi", "timeline", "flow")),
    "creative_agency": TemplateQualityRule(32, 430, 5, 15, 0.85, 5, 3, 3, ("roadmap", "comparison", "flow")),
    "executive_minimal": TemplateQualityRule(34, 380, 5, 16, 0.92, 4, 3, 3, ("kpi", "comparison", "timeline")),
    "data_driven": TemplateQualityRule(34, 450, 5, 14, 0.75, 5, 4, 5, ("kpi", "matrix", "comparison")),
    "warm_professional": TemplateQualityRule(36, 480, 5, 14, 0.78, 5, 4, 4, ("flow", "roadmap", "comparison")),
    "japanese_business": TemplateQualityRule(40, 480, 5, 16, 0.78, 5, 4, 5, ("timeline", "comparison", "flow")),
    "bold_vision": TemplateQualityRule(30, 390, 4, 16, 0.92, 4, 3, 3, ("kpi", "roadmap", "matrix")),
}

SUMMARY_RULE = TemplateQualityRule(32, 340, 4, 16, 0.9, 4, 3, 3, ("kpi", "comparison", "timeline"))


def run_pptx_quality_pipeline(
    slides: list[PowerPointSlide],
    context: PptxContext,
    *,
    summary_mode: bool,
    presentation_quality_state: dict[str, Any] | None = None,
) -> QualityPipelineResult:
    started = perf_counter()
    template = context.design_template or "corporate_clean"
    rule = _rules_for(template, summary_mode=summary_mode)
    original_count = len(slides)
    findings: list[QualityFinding] = []
    auto_fixes: list[QualityFinding] = []

    normalized = [_normalize_slide(slide, rule, findings, auto_fixes) for slide in slides]
    findings.extend(_analyze_quality(normalized, rule, summary_mode=summary_mode))
    repaired = _split_dense_slides(normalized, rule, summary_mode=summary_mode, findings=findings, auto_fixes=auto_fixes)
    resolved = _resolve_layouts(repaired, rule, findings=findings, auto_fixes=auto_fixes)
    findings.extend(_find_layout_repetition(resolved))
    findings.extend(_find_numeric_issues(resolved))
    findings.extend(_state_findings(presentation_quality_state))

    report = _build_report(
        findings=findings,
        auto_fixes=auto_fixes,
        slide_count_before=original_count,
        slide_count_after=len(resolved),
        template=template,
        generation_duration=perf_counter() - started,
    )
    return QualityPipelineResult(slides=resolved, report=report)


def validate_rendered_pptx(prs: Presentation, report: PptxQualityReport) -> PptxQualityReport:
    findings = list(report.findings)
    auto_fixes = list(report.auto_fixes_applied)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    if len(prs.slides) != report.slide_count_after:
        findings.append(
            QualityFinding(
                rule_id=RULE_RENDER,
                category="post_render_validation",
                severity="warning",
                message="PPTX生成後のスライド数が品質パイプラインの想定と異なります。",
                recommendation="生成後のスライド構造を確認してください。",
                confidence=0.9,
                human_review_required=True,
            )
        )

    for slide_index, slide in enumerate(prs.slides, start=1):
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        visible_text = [shape.text.strip() for shape in text_shapes if shape.text.strip()]
        if not visible_text:
            findings.append(_render_finding(slide_index, "空スライドまたはテキスト欠落があります。", "タイトルまたは本文を確認してください。", "critical"))
        if not visible_text[0:1]:
            findings.append(_render_finding(slide_index, "タイトル候補のテキストが見つかりません。", "スライドタイトルを確認してください。", "warning"))
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
                findings.append(_render_finding(slide_index, "スライド外にはみ出す要素があります。", "要素位置とサイズを調整してください。", "warning"))
            if getattr(shape, "has_text_frame", False) and not shape.text.strip():
                findings.append(_render_finding(slide_index, "空のテキストボックスがあります。", "空要素を削除または内容を追加してください。", "info"))
            if getattr(shape, "has_table", False):
                table = shape.table
                if len(table.columns) > 5:
                    findings.append(_render_finding(slide_index, "表の列数が多く読みづらい可能性があります。", "列を減らすか複数表に分けてください。", "warning"))
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None and run.font.size.pt < 9:
                            findings.append(_render_finding(slide_index, "極端に小さい文字があります。", "最低文字サイズを維持してください。", "warning"))
                            break

    return _build_report(
        findings=findings,
        auto_fixes=auto_fixes,
        slide_count_before=report.slide_count_before,
        slide_count_after=report.slide_count_after,
        template=report.template,
        generation_duration=report.generation_duration,
        layout_decisions=report.layout_decisions,
        layout_fallbacks=report.layout_fallbacks,
        preview_pptx_differences=report.preview_pptx_differences,
        predicted_score=report.predicted_score,
        unsupported_layouts=report.unsupported_layouts,
        numeric_integrity=report.numeric_integrity,
        template_token_application=report.template_token_application,
        human_review_items=report.human_review_items,
    )


def merge_layout_integration_report(
    report: PptxQualityReport,
    *,
    findings: list[QualityFinding],
    layout_decisions: list[dict[str, Any]],
    layout_fallbacks: list[dict[str, Any]],
    preview_pptx_differences: list[dict[str, Any]],
    predicted_score: int | None,
    unsupported_layouts: list[str],
    numeric_integrity: dict[str, Any],
    template_token_application: dict[str, Any],
    human_review_items: list[str],
) -> PptxQualityReport:
    return _build_report(
        findings=[*report.findings, *findings],
        auto_fixes=report.auto_fixes_applied,
        slide_count_before=report.slide_count_before,
        slide_count_after=report.slide_count_after,
        template=report.template,
        generation_duration=report.generation_duration,
        layout_decisions=layout_decisions,
        layout_fallbacks=layout_fallbacks,
        preview_pptx_differences=preview_pptx_differences,
        predicted_score=predicted_score,
        unsupported_layouts=unsupported_layouts,
        numeric_integrity=numeric_integrity,
        template_token_application=template_token_application,
        human_review_items=human_review_items,
    )


def _rules_for(template: str, *, summary_mode: bool) -> TemplateQualityRule:
    if summary_mode:
        return SUMMARY_RULE
    return TEMPLATE_QUALITY_RULES.get(template, TEMPLATE_QUALITY_RULES["corporate_clean"])


def _normalize_slide(
    slide: PowerPointSlide,
    rule: TemplateQualityRule,
    findings: list[QualityFinding],
    auto_fixes: list[QualityFinding],
) -> PowerPointSlide:
    title = normalize_customer_facing_title(_clean_title(slide.title), limit=44)
    bullets = [normalize_customer_facing_text(_clean_body_item(item), limit=92) for item in slide.bullets if _clean_body_item(item)]
    update: dict[str, Any] = {}
    if title != slide.title:
        finding = QualityFinding(
            rule_id=RULE_TITLE,
            category="title",
            severity="info",
            slide_no=slide.slide_no,
            slide_title=slide.title,
            message="タイトルの不要な改行または末尾記号を整理しました。",
            recommendation="意味を変えない範囲の整形のみ自動適用しています。",
            auto_fixable=True,
            before=slide.title,
            after=title,
            confidence=0.95,
        )
        auto_fixes.append(finding)
        update["title"] = title
    if bullets != slide.bullets:
        finding = QualityFinding(
            rule_id=RULE_BODY,
            category="content_fit",
            severity="info",
            slide_no=slide.slide_no,
            slide_title=title,
            message="本文の空行や余分な空白を整理しました。",
            recommendation="本文の意味は変えず、PPTX内で読みやすい形へ整えています。",
            auto_fixable=True,
            before="\n".join(slide.bullets),
            after="\n".join(bullets),
            confidence=0.95,
        )
        auto_fixes.append(finding)
        update["bullets"] = bullets
    if len(title) > rule.title_max_chars:
        findings.append(
            QualityFinding(
                rule_id=RULE_TITLE,
                category="title",
                severity="warning",
                slide_no=slide.slide_no,
                slide_title=title,
                message=f"タイトルが{len(title)}文字でテンプレート上限の目安を超えています。",
                recommendation="意味を変える短縮が必要なため、自動適用せず人が確認してください。",
                auto_fixable=False,
                before=title,
                confidence=0.88,
                human_review_required=True,
            )
        )
    return slide.copy(update=update) if update else slide


def _analyze_quality(slides: list[PowerPointSlide], rule: TemplateQualityRule, *, summary_mode: bool) -> list[QualityFinding]:
    findings: list[QualityFinding] = []
    for slide in slides:
        body = "\n".join(slide.bullets)
        bullet_count = len([item for item in slide.bullets if item.strip()])
        if len(body) > rule.body_max_chars:
            findings.append(
                QualityFinding(
                    rule_id=RULE_BODY,
                    category="content_fit",
                    severity="warning",
                    slide_no=slide.slide_no,
                    slide_title=slide.title,
                    message=f"本文量が{len(body)}文字で、読みやすい目安を超えています。",
                    recommendation="空行整理、箇条書き整理、スライド分割候補の順で改善してください。",
                    auto_fixable=bullet_count > rule.bullet_max_items + 1,
                    before=body,
                    confidence=0.86,
                    human_review_required=not summary_mode,
                )
            )
        if bullet_count > rule.bullet_max_items:
            findings.append(
                QualityFinding(
                    rule_id=RULE_BULLET,
                    category="content_fit",
                    severity="warning",
                    slide_no=slide.slide_no,
                    slide_title=slide.title,
                    message=f"箇条書きが{bullet_count}件あり、1ページの情報量が多い状態です。",
                    recommendation=f"{rule.bullet_max_items}件以内を目安に整理するか、スライド分割してください。",
                    auto_fixable=True,
                    before=body,
                    confidence=0.9,
                )
            )
        if len(body) > rule.body_max_chars * 1.35:
            findings.append(
                QualityFinding(
                    rule_id=RULE_OVERFLOW,
                    category="overflow",
                    severity="warning",
                    slide_no=slide.slide_no,
                    slide_title=slide.title,
                    message="テキストボックス内でオーバーフローする可能性があります。",
                    recommendation="文字サイズを極端に下げず、分割または図解化してください。",
                    auto_fixable=bullet_count >= 4,
                    before=body,
                    confidence=0.84,
                    human_review_required=True,
                )
            )
        diagram = recommend_diagram_type(slide)
        if diagram and _is_plain_content_layout(slide):
            findings.append(
                QualityFinding(
                    rule_id=RULE_DIAGRAM,
                    category="diagram",
                    severity="info",
                    slide_no=slide.slide_no,
                    slide_title=slide.title,
                    message=f"{diagram}図解に適した内容です。",
                    recommendation="安全に構造化できる場合のみ、編集可能なShape/Tableへ変換します。",
                    auto_fixable=True,
                    before=slide.layout,
                    after=f"quality_{diagram}",
                    confidence=0.82,
                )
            )
        if diagram == "comparison":
            findings.append(
                QualityFinding(
                    rule_id=RULE_COMPARE,
                    category="diagram",
                    severity="info",
                    slide_no=slide.slide_no,
                    slide_title=slide.title,
                    message="比較表候補を検出しました。",
                    recommendation="比較軸が明確な場合、編集可能な比較表へ変換します。",
                    auto_fixable=True,
                    confidence=0.82,
                )
            )
    return findings


def _split_dense_slides(
    slides: list[PowerPointSlide],
    rule: TemplateQualityRule,
    *,
    summary_mode: bool,
    findings: list[QualityFinding],
    auto_fixes: list[QualityFinding],
) -> list[PowerPointSlide]:
    result: list[PowerPointSlide] = []
    max_slides = 12 if summary_mode else 28
    for slide in slides:
        bullets = _expand_long_bullets([item for item in slide.bullets if item.strip()], max_chars=max(120, rule.body_max_chars // 3))
        body_length = len("\n".join(bullets))
        should_split = len(bullets) > rule.bullet_max_items + 2 or (body_length > rule.body_max_chars * 1.35 and len(bullets) >= 2)
        if not should_split or len(result) + 2 > max_slides:
            result.append(slide.copy(update={"bullets": bullets}) if bullets != slide.bullets else slide)
            continue
        chunks = _chunk_items(bullets, max(3, rule.bullet_max_items))
        if len(chunks) <= 1:
            result.append(slide)
            continue
        for index, chunk in enumerate(chunks):
            title = slide.title if index == 0 else _subtitle_for(slide.title, chunk)
            result.append(slide.copy(update={"title": title, "bullets": chunk}))
        finding = QualityFinding(
            rule_id=RULE_OVERFLOW,
            category="slide_split",
            severity="info",
            slide_no=slide.slide_no,
            slide_title=slide.title,
            message="長文スライドを情報欠落なしで分割しました。",
            recommendation="分割後の各スライドの順序と意味を確認してください。",
            auto_fixable=True,
            before="\n".join(bullets),
            after=f"{len(chunks)} slides",
            confidence=0.82,
            human_review_required=True,
        )
        auto_fixes.append(finding)
        findings.append(finding)
    return result


def _expand_long_bullets(items: list[str], *, max_chars: int) -> list[str]:
    expanded: list[str] = []
    for item in items:
        if len(item) <= max_chars:
            expanded.append(item)
            continue
        parts = [part.strip() for part in re.split(r"(?<=[。.!?！？])\s*", item) if part.strip()]
        if len(parts) <= 1:
            parts = [item[index : index + max_chars].strip() for index in range(0, len(item), max_chars)]
        current = ""
        for part in parts:
            if current and len(current) + len(part) <= max_chars:
                current = f"{current}{part}"
                continue
            if current:
                expanded.append(current)
            current = part
        if current:
            expanded.append(current)
    return expanded


def _resolve_layouts(
    slides: list[PowerPointSlide],
    rule: TemplateQualityRule,
    *,
    findings: list[QualityFinding],
    auto_fixes: list[QualityFinding],
) -> list[PowerPointSlide]:
    resolved: list[PowerPointSlide] = []
    for slide in slides:
        diagram = recommend_diagram_type(slide)
        if not diagram or diagram not in rule.preferred_diagrams or not _is_plain_content_layout(slide):
            resolved.append(slide)
            continue
        layout = f"quality_{diagram}"
        updated = slide.copy(update={"layout": layout, "visual_suggestion": _diagram_label(diagram)})
        resolved.append(updated)
        auto_fixes.append(
            QualityFinding(
                rule_id=RULE_DIAGRAM,
                category="diagram",
                severity="info",
                slide_no=slide.slide_no,
                slide_title=slide.title,
                message=f"{_diagram_label(diagram)}へレイアウトを解決しました。",
                recommendation="文章の意味を変えず、編集可能なPowerPoint要素へ変換します。",
                auto_fixable=True,
                before=slide.layout,
                after=layout,
                confidence=0.82,
            )
        )
    return resolved


def recommend_diagram_type(slide: PowerPointSlide) -> DiagramType | None:
    text = f"{slide.title}\n{slide.visual_suggestion}\n" + "\n".join(slide.bullets)
    if re.search(r"比較|競合|差別|AとB|Before|After|現状と改善後|自社|プラン|選択肢", text, re.IGNORECASE):
        return "comparison"
    if re.search(r"月|週|Phase|Step|時系列|スケジュール|フェーズ|段階", text, re.IGNORECASE):
        return "timeline"
    if re.search(r"短期|中期|長期|Phase 1|Phase1|ロードマップ|段階導入", text, re.IGNORECASE):
        return "roadmap"
    if re.search(r"%|％|円|万円|件|CVR|ROI|工数|期間|時間|KPI", text, re.IGNORECASE):
        return "kpi"
    if re.search(r"入力から出力|申込から完了|現状から改善|連携|承認|登録|API|CSV|ステップ", text, re.IGNORECASE):
        return "flow"
    if re.search(r"重要度|緊急度|効果とコスト|難易度|価値|2軸|評価", text, re.IGNORECASE):
        return "matrix"
    return None


def extract_numbers(text: str) -> list[str]:
    return re.findall(r"\d+(?:\.\d+)?\s*(?:%|％|万円|円|件|時間|分|日|週間|か月|ヶ月|人|社|ページ|P)?", text)


def _find_layout_repetition(slides: list[PowerPointSlide]) -> list[QualityFinding]:
    findings: list[QualityFinding] = []
    for index in range(2, len(slides)):
        previous = slides[index - 2 : index + 1]
        if all(_is_exempt_repeated_slide(slide) for slide in previous):
            continue
        if len({slide.layout for slide in previous}) == 1:
            findings.append(
                QualityFinding(
                    rule_id=RULE_LAYOUT,
                    category="layout",
                    severity="warning",
                    slide_no=slides[index].slide_no,
                    slide_title=slides[index].title,
                    message="同じLayout IDが3ページ以上連続しています。",
                    recommendation="Storyの順序を保ったまま、比較、KPI、フローなどへ一部を変更してください。",
                    auto_fixable=False,
                    confidence=0.8,
                    human_review_required=True,
                )
            )
    return findings


def _find_numeric_issues(slides: list[PowerPointSlide]) -> list[QualityFinding]:
    findings: list[QualityFinding] = []
    for slide in slides:
        numbers = extract_numbers("\n".join(slide.bullets))
        if not numbers:
            continue
        if recommend_diagram_type(slide) == "kpi":
            continue
        findings.append(
            QualityFinding(
                rule_id=RULE_NUMERIC,
                category="numeric",
                severity="info",
                slide_no=slide.slide_no,
                slide_title=slide.title,
                message=f"主要数値候補を{len(numbers)}件検出しました。",
                recommendation="数値の意味を判断できる場合のみKPIカード化してください。数値は改変しません。",
                auto_fixable=False,
                before=", ".join(numbers[:6]),
                confidence=0.78,
                human_review_required=True,
            )
        )
    return findings


def _state_findings(state: dict[str, Any] | None) -> list[QualityFinding]:
    if not state:
        return []
    findings: list[QualityFinding] = []
    rejected = state.get("rejected_fixes") or []
    if rejected:
        findings.append(
            QualityFinding(
                rule_id="PPT-HUMAN-001",
                category="human_review",
                severity="info",
                message=f"ユーザーが却下したAuto Fixが{len(rejected)}件あります。",
                recommendation="却下済みの改善はBackendで意味変更を伴う自動適用をしません。",
                auto_fixable=False,
                confidence=1.0,
            )
        )
    return findings


def _build_report(
    *,
    findings: list[QualityFinding],
    auto_fixes: list[QualityFinding],
    slide_count_before: int,
    slide_count_after: int,
    template: str,
    generation_duration: float,
    layout_decisions: list[dict[str, Any]] | None = None,
    layout_fallbacks: list[dict[str, Any]] | None = None,
    preview_pptx_differences: list[dict[str, Any]] | None = None,
    predicted_score: int | None = None,
    unsupported_layouts: list[str] | None = None,
    numeric_integrity: dict[str, Any] | None = None,
    template_token_application: dict[str, Any] | None = None,
    human_review_items: list[str] | None = None,
) -> PptxQualityReport:
    category_scores: dict[str, int] = {}
    categories = sorted({finding.category for finding in findings} | {"title", "content_fit", "content_density", "customer_facing_copy", "diagram", "layout", "numeric", "post_render_validation"})
    for category in categories:
        critical = sum(1 for finding in findings if finding.category == category and finding.severity == "critical")
        warning = sum(1 for finding in findings if finding.category == category and finding.severity == "warning")
        info = sum(1 for finding in findings if finding.category == category and finding.severity == "info")
        category_scores[category] = max(35, 100 - critical * 25 - warning * 12 - info * 3)
    overall = round(sum(category_scores.values()) / max(1, len(category_scores)))
    rendered_score = overall
    score_delta = rendered_score - predicted_score if predicted_score is not None else None
    warnings = [finding.message for finding in findings if finding.severity in {"warning", "critical"}]
    human_review_required = any(finding.human_review_required or finding.severity == "critical" for finding in findings)
    human_items = list(human_review_items or [])
    if layout_fallbacks:
        human_items.append("layout_fallback_review")
    if unsupported_layouts:
        human_items.append("unsupported_layout_review")
    return PptxQualityReport(
        overall_score=overall,
        category_scores=category_scores,
        findings=findings,
        auto_fixes_applied=auto_fixes,
        warnings=warnings,
        human_review_required=human_review_required or bool(human_items),
        slide_count_before=slide_count_before,
        slide_count_after=slide_count_after,
        template=template,
        generation_duration=round(generation_duration, 3),
        layout_decisions=layout_decisions or [],
        layout_fallbacks=layout_fallbacks or [],
        preview_pptx_differences=preview_pptx_differences or [],
        predicted_score=predicted_score,
        rendered_score=rendered_score,
        score_delta=score_delta,
        unsupported_layouts=unsupported_layouts or [],
        numeric_integrity=numeric_integrity or {},
        template_token_application=template_token_application or {},
        human_review_items=human_items,
    )


def _render_finding(slide_index: int, message: str, recommendation: str, severity: Severity) -> QualityFinding:
    return QualityFinding(
        rule_id=RULE_RENDER,
        category="post_render_validation",
        severity=severity,
        slide_no=slide_index,
        message=message,
        recommendation=recommendation,
        confidence=0.86,
        human_review_required=severity != "info",
    )


def _clean_title(value: str) -> str:
    title = re.sub(r"\s+", " ", value.replace("\n", " ")).strip()
    title = re.sub(r"[。．.]+$", "", title)
    return title


def _clean_body_item(value: str) -> str:
    return re.sub(r"[ \t]+", " ", value.replace("\r\n", "\n")).strip()


def _chunk_items(items: list[str], size: int) -> list[list[str]]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def _subtitle_for(title: str, bullets: list[str]) -> str:
    label = _clean_title(bullets[0])[:18] if bullets else "補足"
    return f"{title} - {label}" if label and label not in title else f"{title} - 補足"


def _is_plain_content_layout(slide: PowerPointSlide) -> bool:
    normalized = (slide.layout or "").strip().lower()
    return normalized in {"", "content", "summary", "three points", "bullet", "text", "proposal", "case study"}


def _is_exempt_repeated_slide(slide: PowerPointSlide) -> bool:
    text = f"{slide.title} {slide.layout}"
    return bool(re.search(r"見積|費用|Estimate|章|Chapter", text, re.IGNORECASE))


def _diagram_label(diagram: DiagramType) -> str:
    labels = {
        "comparison": "比較表",
        "timeline": "タイムライン",
        "roadmap": "ロードマップ",
        "kpi": "KPIカード",
        "flow": "フロー",
        "matrix": "マトリクス",
    }
    return labels[diagram]
