from __future__ import annotations

from io import BytesIO
import logging
import re

from pptx import Presentation
from pptx.util import Inches

from app.models import PowerPointData, PowerPointSlide, PptxDownloadRequest
from app.proposal_profiles import proposal_profile_for_text
from app.services.pptx_parts.models import PptxContext
from app.services.pptx_parts.content import (
    _fallback_slide,
    _looks_generic_client_name,
    _trim,
    build_case_triplets_from_items,
    case_action_hint,
    case_current_hint,
    case_result_hint,
    concept_statement,
    derive_competitor_rows,
    derive_concept,
    derive_content_items,
    derive_current_understanding,
    derive_estimate_summary,
    derive_journey_points,
    derive_kpi_rows,
    derive_kpi_targets,
    derive_sitemap_items,
    derive_target_user_rows,
    derive_web_strategy_items,
    derive_winning_strategy,
    display_case_title,
    estimate_flag,
    extract_budget_amount,
    extract_case_result_phrase,
    extract_case_study_items,
    extract_case_triplets,
    extract_client_name,
    extract_confirmation_items,
    extract_cost_points,
    extract_current_segment,
    ensure_items,
    extract_labeled_segment,
    extract_page_count,
    extract_project_points,
    extract_schedule_points,
    extract_service_points,
    extract_solution_points,
    extract_team_points,
    extract_text_items,
    extract_unlabeled_case_result,
    has_competitor_context,
    looks_like_case_result,
    normalize_number_text,
    projected_probability_for,
    rank_probability_for,
    remove_case_title,
    risk_label_for,
    risk_score_for_probability,
    sanitize_proposal_text,
    score_label,
    sitemap_note,
    strip_probability_label,
    unique_items,
)
from app.services.pptx_parts.slides import (
    add_budget_fit_slide,
    add_case_studies_slide,
    add_competitor_slide,
    add_concept_slide,
    add_content_design_slide,
    add_cost_slide,
    add_cover_slide,
    add_current_understanding_slide,
    add_customer_journey_slide,
    add_designed_slide,
    add_estimate_priority_slide,
    add_estimate_slide,
    add_generic_slide,
    add_issues_slide,
    add_kpi_slide,
    add_next_steps_slide,
    add_process_slide,
    add_proposal_summary_slide,
    add_schedule_slide,
    add_sitemap_slide,
    add_solution_slide,
    add_summary_slide,
    add_target_user_slide,
    add_team_slide,
    add_understanding_slide,
    add_web_strategy_slide,
    add_win_probability_slide,
    resolve_slide_kind,
)
from app.services.pptx_design.validators import validate_premium_deck
from app.services.pptx_theme import MEDIA_TYPE, SLIDE_HEIGHT, SLIDE_WIDTH


logger = logging.getLogger(__name__)


def build_pptx_bytes(payload: PptxDownloadRequest) -> bytes:
    return _build_pptx_bytes(payload, summary_mode=payload.summary)


def build_summary_pptx_bytes(payload: PptxDownloadRequest) -> bytes:
    return _build_pptx_bytes(payload, summary_mode=True)


def _build_pptx_bytes(payload: PptxDownloadRequest, *, summary_mode: bool) -> bytes:
    try:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)

        data = payload.powerpoint_generation_data
        context = build_pptx_context(payload)
        slides = data.slides or [_fallback_slide(data)]
        slides = insert_competitor_analysis_slide(slides, context)
        slides = insert_case_studies_slide(slides, context)
        slides = insert_estimate_slides(slides, context)
        slides = insert_win_probability_slide(slides, context)
        if summary_mode:
            slides = build_summary_slides(slides, context)
        else:
            slides = normalize_detailed_slides(slides)

        display_slide_no = 1
        for index, slide_data in enumerate(slides):
            numbered_slide = slide_data.copy(update={"slide_no": display_slide_no})
            add_designed_slide(prs, numbered_slide, data, index, context)
            display_slide_no += 1

        issues = validate_premium_deck(prs)
        if issues:
            logger.info(
                "pptx_design_validation_issues",
                extra={"issue_count": len(issues), "issue_codes": [issue.code for issue in issues[:8]]},
            )

        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
    except Exception:
        logger.exception("Failed to build proposal PowerPoint deck.")
        raise


def build_pptx_context(payload: PptxDownloadRequest) -> PptxContext:
    data = payload.powerpoint_generation_data
    all_input = "\n".join(
        [
            payload.project_brief,
            payload.client_company_info,
            payload.competitor_site_url,
            payload.competitor_company_name,
            payload.estimated_page_count,
            payload.cms_required,
            payload.contact_form_required,
            payload.special_function_required,
            payload.seo_required,
            payload.content_creation_required,
            payload.desired_launch_timing,
            payload.budget_range,
            payload.hearing_result,
            payload.own_service_info,
            payload.past_proposal_template,
            payload.case_studies,
        ]
    )
    profile = proposal_profile_for_text(all_input)
    concept = derive_concept(all_input)
    estimate = derive_estimate_summary(payload, all_input)
    return PptxContext(
        client_name=extract_client_name(payload.client_company_info, data.client_name),
        proposal_category=profile.category,
        proposal_label=profile.label,
        concept=concept,
        current_understanding=derive_current_understanding(payload.project_brief, concept),
        journey_points=derive_journey_points(concept),
        sitemap_items=derive_sitemap_items(all_input),
        kpi_rows=derive_kpi_rows(all_input, concept),
        kpi_targets=derive_kpi_targets(all_input, concept),
        competitor_rows=derive_competitor_rows(all_input),
        competitor_site_url=payload.competitor_site_url.strip(),
        competitor_company_name=payload.competitor_company_name.strip(),
        winning_strategy=derive_winning_strategy(all_input),
        target_user_rows=derive_target_user_rows(all_input, concept),
        content_items=derive_content_items(all_input, concept),
        web_strategy_items=derive_web_strategy_items(all_input, concept),
        case_studies=extract_case_study_items(payload.case_studies),
        case_triplets=extract_case_triplets(payload.case_studies),
        project_points=extract_project_points(payload.project_brief),
        solution_points=extract_solution_points(all_input),
        service_points=extract_service_points(payload.own_service_info),
        schedule_points=extract_schedule_points(all_input),
        team_points=extract_team_points(payload.own_service_info),
        cost_points=extract_cost_points(all_input),
        estimate=estimate,
        confirmation_items=extract_confirmation_items(all_input),
        win_probability=payload.win_probability,
    )


def build_pptx_filename(data: PowerPointData, client_company_info: str = "", *, summary_mode: bool = False) -> str:
    base_name = data.deck_title or "ready-crew-proposal"
    client_name = extract_client_name(client_company_info, data.client_name)
    if client_name != "提案先企業" and client_name not in base_name:
        base_name = f"{client_name}_{base_name}"
    if summary_mode:
        suffix = "_要約版"
        sanitized_base = sanitize_filename(base_name, max_length=80 - len(suffix))
        return f"{sanitized_base}{suffix}.pptx"
    return f"{sanitize_filename(base_name)}.pptx"


def sanitize_filename(value: str, *, max_length: int = 80) -> str:
    sanitized = re.sub(r'[\\/:*?"<>|]', "-", value)
    sanitized = re.sub(r"\s+", " ", sanitized).strip()
    return sanitized[:max_length].strip() or "ready-crew-proposal"


def insert_case_studies_slide(slides: list[PowerPointSlide], context: PptxContext) -> list[PowerPointSlide]:
    if not context.case_triplets and not context.case_studies:
        return slides
    if any("実績" in slide.title for slide in slides):
        return slides

    next_slide_no = max((slide.slide_no for slide in slides), default=0) + 1
    case_slide = PowerPointSlide(
        slide_no=next_slide_no,
        layout="content",
        title="関連実績",
        bullets=context.case_studies or [case["title"] for case in context.case_triplets],
        speaker_notes="入力された成功事例データをもとに、現状、施策、成果の順で実績を紹介します。",
        visual_suggestion="現状、施策、成果を整理した実績カード",
    )

    insert_at = next((idx for idx, slide in enumerate(slides) if "体制" in slide.title), len(slides))
    return [*slides[:insert_at], case_slide, *slides[insert_at:]]


def insert_competitor_analysis_slide(slides: list[PowerPointSlide], context: PptxContext) -> list[PowerPointSlide]:
    if not has_competitor_context(context):
        return slides
    if any("競合比較分析" in slide.title for slide in slides):
        return slides

    competitor_name = context.competitor_company_name or ("競合サイト" if context.proposal_category == "web" else "比較対象")
    comparison_note = (
        "デザイン、SEO、導線設計、コンテンツ量、更新性、CTAの6観点で比較"
        if context.proposal_category == "web"
        else f"{context.proposal_label}に必要な比較軸で改善優先度を整理"
    )
    bullets = ensure_items(
        [
            f"比較対象: {competitor_name}",
            f"参考URL: {context.competitor_site_url}" if context.competitor_site_url else "",
            f"勝ち筋: {context.winning_strategy}",
            comparison_note,
        ],
        ["競合比較をもとに、勝てる訴求軸と改善優先度を明確化"],
        4,
    )
    next_slide_no = max((slide.slide_no for slide in slides), default=0) + 1
    competitor_slide = PowerPointSlide(
        slide_no=next_slide_no,
        layout="content",
        title="競合比較分析",
        bullets=bullets,
        speaker_notes="比較観点と勝ち筋を提示し、提案内容の優先順位に反映します。",
        visual_suggestion="6観点の競合比較表と勝ち筋の強調表示",
    )

    insert_at = next(
        (
            idx + 1
            for idx, slide in enumerate(slides)
            if ("市場" in slide.title or "競合" in slide.title) and "競合比較分析" not in slide.title
        ),
        None,
    )
    if insert_at is None:
        insert_at = next((idx + 1 for idx, slide in enumerate(slides) if "課題" in slide.title), min(5, len(slides)))
    return [*slides[:insert_at], competitor_slide, *slides[insert_at:]]


def insert_estimate_slides(slides: list[PowerPointSlide], context: PptxContext) -> list[PowerPointSlide]:
    existing_titles = {slide.title for slide in slides}
    additions: list[PowerPointSlide] = []
    next_slide_no = max((slide.slide_no for slide in slides), default=0) + 1

    if "概算見積" not in existing_titles:
        additions.append(
            PowerPointSlide(
                slide_no=next_slide_no,
                layout="content",
                title="概算見積",
                bullets=[
                    f"合計概算: {context.estimate.total_label}",
                    context.estimate.scope_label,
                    "固定金額ではなく、要件確定前のレンジとして提示",
                ],
                speaker_notes="案件条件から概算見積レンジを整理します。",
                visual_suggestion="見積内訳と合計レンジの表",
            )
        )
        next_slide_no += 1

    if "予算適合判定" not in existing_titles:
        additions.append(
            PowerPointSlide(
                slide_no=next_slide_no,
                layout="content",
                title="予算適合判定",
                bullets=[
                    f"判定: {context.estimate.budget_fit}",
                    f"予算感: {context.estimate.budget_label}",
                    f"概算見積: {context.estimate.total_label}",
                ],
                speaker_notes="予算感と概算見積レンジの差分を示します。",
                visual_suggestion="予算適合判定カード",
            )
        )
        next_slide_no += 1

    if "必須・推奨・オプション対応" not in existing_titles:
        additions.append(
            PowerPointSlide(
                slide_no=next_slide_no,
                layout="content",
                title="必須・推奨・オプション対応",
                bullets=[
                    f"必須対応: {'、'.join(context.estimate.required[:4])}",
                    f"推奨対応: {'、'.join(context.estimate.recommended[:4]) or '次回確認'}",
                    f"オプション対応: {'、'.join(context.estimate.optional[:4]) or '次回確認'}",
                ],
                speaker_notes="予算内で優先すべき対応範囲を分類します。",
                visual_suggestion="必須・推奨・オプションの3カラム",
            )
        )

    if not additions:
        return slides

    insert_at = next((idx + 1 for idx, slide in enumerate(slides) if "費用" in slide.title or "概算" in slide.title), None)
    if insert_at is None:
        insert_at = next((idx for idx, slide in enumerate(slides) if "今後" in slide.title or "進め方" in slide.title), len(slides))
    return [*slides[:insert_at], *additions, *slides[insert_at:]]


def insert_win_probability_slide(slides: list[PowerPointSlide], context: PptxContext) -> list[PowerPointSlide]:
    win = context.win_probability
    if win is None:
        return slides
    if any("受注確率" in slide.title or "案件ランク" in slide.title for slide in slides):
        return slides

    probability = win.probability or rank_probability_for(win.rank)
    risk_score = max(1, min(5, win.risk_score or risk_score_for_probability(probability, len(win.risk_factors))))
    risk_label = win.risk_label or risk_label_for(risk_score)
    projected = win.projected_probability_after_actions or projected_probability_for(probability, risk_score, len(win.improvement_actions))
    next_slide_no = max((slide.slide_no for slide in slides), default=0) + 1
    win_slide = PowerPointSlide(
        slide_no=next_slide_no,
        layout="content",
        title="受注確率判定",
        bullets=[
            f"受注確率: {probability}%",
            f"受注リスク: {risk_label}",
            f"リスク要因: {'、'.join(unique_items(win.risk_factors, 3))}",
            f"改善アクション: {'、'.join(unique_items(win.improvement_actions or win.recommended_next_actions, 3))}",
            f"受注確率向上予測: {probability}% → {projected}%",
        ],
        speaker_notes="受注確率、リスク要因、改善アクション、向上予測を営業判断材料として提示します。",
        visual_suggestion="受注確率カード、リスク要因、改善アクション、向上予測の3カラム",
    )

    insert_at = next((idx for idx, slide in enumerate(slides) if "今後" in slide.title or "進め方" in slide.title), len(slides))
    return [*slides[:insert_at], win_slide, *slides[insert_at:]]


def normalize_detailed_slides(slides: list[PowerPointSlide], *, max_slides: int = 25) -> list[PowerPointSlide]:
    compacted: list[PowerPointSlide] = []
    seen_titles: set[str] = set()
    for slide in slides:
        key = re.sub(r"\s+", "", slide.title or "").lower()
        if key and key in seen_titles:
            continue
        if key:
            seen_titles.add(key)
        compacted.append(slide)

    if len(compacted) <= max_slides:
        return compacted

    priority_kinds = [
        "cover",
        "proposal_summary",
        "current_understanding",
        "issues",
        "proposal_concept",
        "solution",
        "customer_journey",
        "sitemap",
        "kpi",
        "process",
        "schedule",
        "team",
        "cost",
        "estimate",
        "budget_fit",
        "estimate_priority",
        "win_probability",
        "next_steps",
    ]
    selected: list[PowerPointSlide] = []
    for kind in priority_kinds:
        for index, slide in enumerate(compacted):
            if slide in selected:
                continue
            if resolve_slide_kind(slide, index) == kind:
                selected.append(slide)
                break
        if len(selected) >= max_slides:
            return selected[:max_slides]

    for slide in compacted:
        if slide not in selected:
            selected.append(slide)
        if len(selected) >= max_slides:
            break
    return selected[:max_slides]


def build_summary_slides(slides: list[PowerPointSlide], context: PptxContext) -> list[PowerPointSlide]:
    def existing_bullets(*keywords: str) -> list[str]:
        slide = find_slide_by_keywords(slides, *keywords)
        return slide.bullets if slide else []

    def summary_slide(title: str, bullets: list[str], fallback: list[str], visual: str = "") -> PowerPointSlide:
        return PowerPointSlide(
            slide_no=0,
            layout="content",
            title=title,
            bullets=ensure_items([bullet for bullet in bullets if bullet], fallback, 3),
            speaker_notes=f"{title}を要約版向けに簡潔に説明します。",
            visual_suggestion=visual or f"{title}を図解する要約レイアウト",
        )

    cover_bullets = slides[0].bullets if slides else []
    strategy_title = "Web戦略" if context.proposal_category == "web" else "導入戦略"
    structure_title = "サイトマップ" if context.proposal_category == "web" else "導入構成"
    strategy_fallback = (
        ["集客から問い合わせまでの流れを設計", "比較検討に必要な情報を整備", "公開後の改善運用を前提に構築"]
        if context.proposal_category == "web"
        else ["現状整理から導入までの流れを設計", "必要データと連携条件を整理", "運用改善を前提に構築"]
    )
    structure_fallback = ["トップ", "サービス", "お問い合わせ"] if context.proposal_category == "web" else context.sitemap_items[:3]
    kpi_fallback = (
        ["問い合わせ数: 月20件", "CV率: 2.8%", "自然検索流入: 月3,500セッション"]
        if context.proposal_category == "web"
        else [f"{name}: {value}" for name, value in context.kpi_rows[:3]]
    )
    concept_fallback = (
        ["サイトの目的を1つの提案軸に集約", "成果導線を中心に設計", "運用改善まで含めて推進"]
        if context.proposal_category == "web"
        else ["案件目的を1つの提案軸に集約", "導入範囲と効果測定を中心に設計", "運用改善まで含めて推進"]
    )
    selected = [
        PowerPointSlide(
            slide_no=1,
            layout="title",
            title="",
            bullets=ensure_items(
                cover_bullets,
                [f"{context.client_name}向け提案", concept_statement(context.concept), "重要論点を短時間で共有"],
                3,
            ),
            speaker_notes="要約版の表紙です。",
            visual_suggestion="表紙",
        ),
        summary_slide(
            "提案サマリー",
            [*existing_bullets("提案サマリー"), concept_statement(context.concept), *(context.project_points[:2] or context.solution_points[:2])],
            ["提案の狙いを端的に共有", "成果につながる導入範囲と施策を整理", "導入後の運用改善まで見据えて設計"],
            "提案価値を3点で示すサマリーカード",
        ),
        summary_slide(
            "現状理解",
            [*existing_bullets("現状理解"), *list(context.current_understanding.values())],
            ["現状業務と導入条件に改善余地", "利用者と意思決定者の論点を再整理", "効果測定につながる運用を強化"],
            "現状・課題・機会を並べるカード",
        ),
        summary_slide(
            "主要課題",
            [*existing_bullets("課題"), *context.project_points, *context.solution_points],
            ["対象範囲を分かりやすく再整理", "施策と実行条件を整理", "運用と改善サイクルを設計"],
            "優先度順の課題カード",
        ),
        summary_slide(
            "提案コンセプト",
            [
                f"提案コンセプト: {context.concept}",
                concept_statement(context.concept),
                "導入範囲、施策設計、効果測定、運用改善を一体で実施",
            ],
            concept_fallback,
            "提案コンセプトを中心にしたメッセージボード",
        ),
        summary_slide(
            strategy_title,
            [*existing_bullets(strategy_title), *context.web_strategy_items],
            strategy_fallback,
            "導入ステップと改善サイクルをつなぐ戦略図",
        ),
        summary_slide(
            structure_title,
            [*existing_bullets(structure_title), *context.sitemap_items],
            structure_fallback,
            "主要構成を階層で示す構成図",
        ),
        summary_slide(
            "KPI設計",
            [*existing_bullets("KPI"), *[f"{name}: {value}" for name, value in context.kpi_rows]],
            kpi_fallback,
            "主要KPIを比較しやすいカードで表示",
        ),
        summary_slide(
            "スケジュール",
            [*existing_bullets("スケジュール"), *context.schedule_points],
            ["要件整理・設計", "実装・検証", "運用開始"],
            "短期共有向けの簡易タイムライン",
        ),
        summary_slide(
            "費用概算",
            [
                f"合計概算: {context.estimate.total_label}",
                f"予算適合: {context.estimate.budget_fit}",
                f"必須対応: {'、'.join(context.estimate.required[:3])}",
                *existing_bullets("費用"),
                *context.cost_points,
            ],
            ["必須範囲とオプションを分けて提示", "導入範囲と運用支援の範囲を整理", "次回確認後に見積精度を更新"],
            "必須範囲とオプションを分けた費用表",
        ),
        summary_slide(
            "今後の進め方",
            [*existing_bullets("今後"), *context.confirmation_items],
            ["次回ヒアリングで要件を確定", "提案範囲と優先順位を合意", "概算費用とスケジュールを更新"],
            "次回アクションのステップカード",
        ),
    ]

    return [slide.copy(update={"slide_no": index}) for index, slide in enumerate(selected, start=1)]


def find_slide_by_keywords(slides: list[PowerPointSlide], *keywords: str) -> PowerPointSlide | None:
    for slide in slides:
        if all(keyword in slide.title for keyword in keywords):
            return slide
    return None
