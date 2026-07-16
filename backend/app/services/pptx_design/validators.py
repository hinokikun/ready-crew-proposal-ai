from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass(frozen=True)
class PptxDesignIssue:
    code: str
    message: str


def validate_premium_deck(prs: Presentation) -> list[PptxDesignIssue]:
    issues: list[PptxDesignIssue] = []
    if not prs.slides:
        return [PptxDesignIssue("blank_deck", "No slides were generated.")]

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if not texts:
            issues.append(PptxDesignIssue("blank_slide", f"Slide {slide_index} has no readable text."))
        if not any("ProposalPilot" in text for text in texts):
            issues.append(PptxDesignIssue("missing_brand", f"Slide {slide_index} does not show ProposalPilot."))
        if len(slide.shapes) < 6:
            issues.append(PptxDesignIssue("low_visual_density", f"Slide {slide_index} has too few editable shapes."))
        if _is_text_only(slide):
            issues.append(PptxDesignIssue("text_only_slide", f"Slide {slide_index} has text but no editable visual elements."))

    return issues


def count_external_relationships(prs: Presentation) -> int:
    total = 0
    for slide in prs.slides:
        for relationship in slide.part.rels.values():
            if getattr(relationship, "is_external", False):
                total += 1
    return total


def shape_type_counts(prs: Presentation) -> Counter[str]:
    counts: Counter[str] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            counts[str(shape.shape_type)] += 1
    return counts


def _is_text_only(slide) -> bool:
    has_text = any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes)
    has_visual = any(
        shape.shape_type
        in {
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.PICTURE,
        }
        for shape in slide.shapes
    )
    return has_text and not has_visual
