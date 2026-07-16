from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from app.services.pptx_theme import COLORS, FONT_FACE, FOOTER_Y, HEADER_Y, MARGIN_X, SECTION_COLORS
from app.services.pptx_parts.content import _fit_font_size, _trim, unique_items


def add_factor_column(slide, title: str, items: list[str], x: float, y: float, accent: str) -> None:
    add_card(slide, title, "", x, y, 3.05, 1.58, accent, COLORS["white"])
    add_bullet_list(slide, unique_items(items, 2), x + 0.22, y + 0.58, 2.6, 0.78, max_items=2, size=10)


def add_case_line(slide, label: str, body: str, x: float, y: float, accent: str) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.66, 0.28, fill=accent, line=accent)
    add_text(slide, label, x + 0.08, y + 0.08, 0.5, 0.1, size=8, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, _trim(body, 42), x + 0.84, y - 0.02, 2.18, 0.48, size=11, color=COLORS["text"], bold=True)


def add_icon_badge(slide, label: str, x: float, y: float, accent: str, *, size: float = 0.82) -> None:
    add_shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=accent, line=accent)
    add_text(slide, label, x + 0.06, y + size * 0.42, size - 0.12, 0.12, size=10 if len(label) <= 3 else 8, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)


def chapter_icon(title: str) -> str:
    if "競合" in title:
        return "VS"
    if "KPI" in title:
        return "KPI"
    if "サイト" in title:
        return "MAP"
    if "戦略" in title:
        return "STR"
    if "体制" in title:
        return "TEAM"
    return "PLAN"


def chapter_message(title: str, context: PptxContext) -> str:
    if "提案サマリー" in title:
        if context.proposal_category == "web":
            return f"{context.concept}を軸に、成果につながるWebサイト改善の全体像を提示します。"
        return f"{context.concept}を軸に、導入範囲、効果、運用までの全体像を提示します。"
    if "現状理解" in title:
        return "案件概要から読み取れる現状、事業課題、改善機会を整理します。"
    if "課題" in title:
        return "成果に影響する課題を優先度順に絞り込み、提案の論点を明確にします。"
    if "競合" in title:
        return f"競合標準を超えるための改善余地を6つの観点で可視化し、{context.winning_strategy}を提案します。"
    if "ターゲット" in title:
        return "主要ユーザーのニーズと不安を整理し、必要な情報接点を定義します。"
    if "ジャーニー" in title or "業務フロー" in title:
        if context.proposal_category == "web":
            return "認知から問い合わせまでの行動を分解し、離脱を減らす接点を設計します。"
        return "導入前後の業務を分解し、例外処理と改善接点を設計します。"
    if "Web戦略" in title or "WEB戦略" in title or "導入戦略" in title:
        if context.proposal_category == "web":
            return f"{context.concept}を実現するための集客、比較検討、CV、運用改善を設計します。"
        return f"{context.concept}を実現するための現状整理、導入設計、検証、運用改善を設計します。"
    if "サイトマップ" in title or "サイト構成" in title or "導入構成" in title:
        if context.proposal_category == "web":
            return "ユーザーが必要情報へ迷わず進み、運用もしやすいサイト構成を提示します。"
        return "対象業務、連携、運用、効果測定の構成を提示します。"
    if "コンテンツ" in title or "施策設計" in title:
        if context.proposal_category == "web":
            return "サービス理解、比較検討、不安解消、問い合わせを支えるコンテンツを設計します。"
        return "導入判断、実行、効果測定、運用定着を支える施策を設計します。"
    if "KPI" in title:
        if context.proposal_category == "web":
            return "問い合わせ数、CV率、自然検索流入、資料DL数を目標化します。"
        return "カテゴリに合う効果指標を目標化します。"
    if "制作方針" in title or "方針" in title:
        return "設計、デザイン、実装、検証を段階化し、品質と進行の見通しを示します。"
    if "スケジュール" in title:
        return "公開希望時期から逆算し、要件整理からリリースまでの進行を提示します。"
    if "関連実績" in title or "実績" in title:
        return "入力された成功事例を、現状、施策、成果の流れで提案の裏付けとして示します。"
    if "体制" in title:
        return "プロジェクトを安定して進める役割分担と支援範囲を明確にします。"
    if "見積" in title:
        return f"想定範囲と機能条件をもとに、概算見積レンジ{context.estimate.total_label}を提示します。"
    if "予算適合" in title:
        return f"入力予算と概算見積を比較し、{context.estimate.budget_fit}として調整方針を整理します。"
    if "必須" in title and "推奨" in title and "オプション" in title:
        return "予算内で成果を出すため、対応範囲を必須・推奨・オプションに分類します。"
    if "費用" in title or "概算" in title:
        return "必須範囲とオプションを分け、予算判断しやすい費用構成を提示します。"
    if "今後" in title:
        return "意思決定に必要な確認事項と、実行開始までの進め方を明確にします。"
    return "本提案の実行内容と判断材料を、章ごとに具体化します。"


def content_title(item: str, index: int) -> str:
    if ":" in item:
        return item.split(":", 1)[0]
    if "：" in item:
        return item.split("：", 1)[0]
    return ["サービス訴求", "実績・事例", "FAQ", "CV導線"][index % 4]


def add_header(slide, title: str, section: str, accent: str = COLORS["teal"]) -> None:
    add_section_label(slide, section, MARGIN_X, HEADER_Y, fill=accent, color=COLORS["white"])
    add_title(slide, title, MARGIN_X, 0.8, 10.2, 0.48, size=30, color=COLORS["navy"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, MARGIN_X, 1.32, 11.76, 0.025, fill=COLORS["line"], line=COLORS["line"])


def add_footer(slide, slide_no: int) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, MARGIN_X, 6.76, 11.88, 0.02, fill=COLORS["line"], line=COLORS["line"])
    add_text(slide, "ProposalPilot / AI営業秘書", MARGIN_X, FOOTER_Y, 3.4, 0.18, size=9, color=COLORS["muted"])
    add_text(slide, f"{slide_no:02}", 11.58, 6.86, 0.76, 0.22, size=10, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_title(slide, text: str, x: float, y: float, w: float, h: float, *, size: int, color: str) -> None:
    add_text(slide, _trim(text, 42), x, y, w, h, size=_fit_font_size(text, size), color=color, bold=True)


def add_card(
    slide,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    accent: str,
    fill: str,
    number: str | None = None,
) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=COLORS["line"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, fill=accent, line=accent)
    if number:
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.22, y + 0.22, 0.36, 0.36, fill=accent, line=accent)
        add_text(slide, number, x + 0.22, y + 0.29, 0.36, 0.12, size=10, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        title_x = x + 0.72
        title_w = w - 0.95
    else:
        title_x = x + 0.28
        title_w = w - 0.54
    add_text(slide, _trim(title, 28), title_x, y + 0.22, title_w, 0.28, size=14, color=accent, bold=True)
    if body:
        add_text(slide, _trim(body, 56), x + 0.28, y + 0.72, w - 0.54, h - 0.88, size=14, color=COLORS["text"])


def add_section_label(slide, text: str, x: float, y: float, *, fill: str, color: str) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.55, 0.32, fill=fill, line=fill)
    add_text(slide, text, x + 0.13, y + 0.08, 1.3, 0.12, size=8, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_bullet_list(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    max_items: int = 5,
    size: int = 14,
) -> None:
    clean_items = unique_items(items, max_items) or ["次回確認事項として整理します。"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    set_frame_margins(frame, 0)

    for index, item in enumerate(clean_items[:max_items]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {_trim(item, 52)}"
        paragraph.space_after = Pt(6)
        style_paragraph(paragraph, size=size, color=COLORS["text"], bold=False)


def add_table(slide, headers: list[str], rows: list[list[str]], x: float, y: float, w: float, h: float, col_widths: list[float]) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table

    for col_index, width in enumerate(col_widths):
        table.columns[col_index].width = Inches(width)

    for col_index, header in enumerate(headers):
        set_cell(table.cell(0, col_index), header, fill=COLORS["navy"], color=COLORS["white"], size=12, bold=True)

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            fill = COLORS["canvas"] if row_index % 2 == 0 else COLORS["white"]
            set_cell(table.cell(row_index, col_index), _trim(value, 52), fill=fill, color=COLORS["text"], size=12, bold=col_index == 0)


def add_step_flow(slide, steps: list[str], x: float, y: float, w: float, h: float) -> None:
    step_count = min(len(steps), 4)
    gap = 0.18
    step_w = (w - gap * (step_count - 1)) / step_count
    for index, step in enumerate(steps[:step_count]):
        sx = x + index * (step_w + gap)
        accent = SECTION_COLORS[index % len(SECTION_COLORS)]
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, sx, y, step_w, h, fill=COLORS["white"], line=COLORS["line"])
        add_shape(slide, MSO_SHAPE.OVAL, sx + 0.28, y + 0.28, 0.46, 0.46, fill=accent, line=accent)
        add_text(slide, str(index + 1), sx + 0.28, y + 0.38, 0.46, 0.14, size=10, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, _trim(step, 36), sx + 0.32, y + 0.95, step_w - 0.64, 0.55, size=14, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
        if index < step_count - 1:
            add_shape(slide, MSO_SHAPE.CHEVRON, sx + step_w - 0.05, y + 0.72, 0.32, 0.42, fill=COLORS["line_dark"], line=COLORS["line_dark"])


def add_visual_frame(slide, label: str, x: float, y: float, w: float, h: float) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=COLORS["canvas"], line=COLORS["line"])
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.35, y + 0.4, 0.55, 0.55, fill=COLORS["teal_light"], line=COLORS["teal_light"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 1.05, y + 0.54, w - 1.55, 0.08, fill=COLORS["line_dark"], line=COLORS["line_dark"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.35, y + 1.35, w - 0.7, 0.08, fill=COLORS["line"], line=COLORS["line"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.35, y + 1.78, w - 1.15, 0.08, fill=COLORS["line"], line=COLORS["line"])
    add_text(slide, _trim(label, 56), x + 0.38, y + h - 0.82, w - 0.76, 0.36, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_insight_band(slide, title: str, body: str, x: float, y: float, w: float, h: float) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=COLORS["navy"], line=COLORS["navy"])
    add_text(slide, title, x + 0.34, y + 0.2, 2.4, 0.25, size=13, color=COLORS["teal_light"], bold=True)
    add_text(slide, _trim(body, 82), x + 2.65, y + 0.19, w - 3.0, 0.3, size=13, color=COLORS["white"])


def add_side_panel(slide, title: str, items: list[str], x: float, y: float, w: float, h: float, accent: str) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=COLORS["navy"], line=COLORS["navy"])
    add_text(slide, title, x + 0.32, y + 0.34, w - 0.64, 0.3, size=17, color=COLORS["white"], bold=True)
    for idx, item in enumerate(items):
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.36, y + 1.1 + idx * 0.78, 0.22, 0.22, fill=accent, line=accent)
        add_text(slide, item, x + 0.72, y + 1.05 + idx * 0.78, w - 1.0, 0.24, size=13, color=COLORS["white"], bold=True)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(COLORS["white"])


def add_shape(slide, shape_type, x: float, y: float, w: float, h: float, *, fill: str, line: str):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int,
    color: str,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    set_frame_margins(frame, 0)

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    style_paragraph(paragraph, size=size, color=color, bold=bold)


def set_cell(cell, text: str, *, fill: str, color: str, size: int, bold: bool) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    set_frame_margins(frame, 0.06)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    style_paragraph(paragraph, size=size, color=color, bold=bold)


def set_frame_margins(frame, value: float) -> None:
    margin = Inches(value)
    frame.margin_left = margin
    frame.margin_right = margin
    frame.margin_top = margin
    frame.margin_bottom = margin


def style_paragraph(paragraph, *, size: int, color: str, bold: bool) -> None:
    for run in paragraph.runs:
        run.font.name = FONT_FACE
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)

def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)
