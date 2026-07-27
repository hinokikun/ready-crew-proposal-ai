from __future__ import annotations

import csv
from io import BytesIO, StringIO
from typing import Any

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response, StreamingResponse

from app.auth import require_roles
from app.db import get_db
from app.models import ProposalAgentMemoryRequest
from app.organization import get_user_workspace_context
from app.repositories import create_audit_log

router = APIRouter(prefix="/api/proposal-agent", tags=["proposal-agent"])


@router.get("/dashboard")
async def get_proposal_agent_dashboard(user: dict = Depends(require_roles("admin", "manager", "member", "viewer"))) -> dict:
    with get_db() as db:
        dashboard = _load_dashboard(db, user)
        create_audit_log(db, int(user["id"]), "view_proposal_agent_dashboard", "proposal_agent", "", "success", "sanitized=true")
    return {"dashboard": dashboard}


@router.get("/dashboard/export")
async def export_proposal_agent_dashboard(
    format: str = "markdown",
    user: dict = Depends(require_roles("admin", "manager", "member", "viewer")),
) -> Response:
    normalized_format = format.strip().lower()
    with get_db() as db:
        dashboard = _load_dashboard(db, user)
        create_audit_log(db, int(user["id"]), "export_proposal_agent_dashboard", "proposal_agent", "", "success", f"sanitized=true;format={normalized_format}")
    if normalized_format in {"md", "markdown"}:
        content = _render_dashboard_markdown(dashboard)
        return Response(
            content=f"\ufeff{content}",
            media_type="text/markdown; charset=utf-8",
            headers={"Content-Disposition": "attachment; filename=proposal-intelligence-dashboard.md"},
        )
    if normalized_format == "csv":
        content = _render_dashboard_csv(dashboard)
        return Response(
            content=f"\ufeff{content}",
            media_type="text/csv; charset=utf-8",
            headers={"Content-Disposition": "attachment; filename=proposal-intelligence-dashboard.csv"},
        )
    if normalized_format == "pdf":
        return StreamingResponse(
            BytesIO(_render_dashboard_pdf(dashboard)),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=proposal-intelligence-dashboard.pdf"},
        )
    if normalized_format in {"ppt", "pptx", "powerpoint"}:
        return StreamingResponse(
            BytesIO(_render_dashboard_pptx(dashboard)),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=proposal-intelligence-dashboard.pptx"},
        )
    raise HTTPException(status_code=400, detail="export format must be markdown, csv, pdf, or pptx")


@router.post("/memory")
async def save_proposal_agent_memory(
    payload: ProposalAgentMemoryRequest,
    user: dict = Depends(require_roles("admin", "manager", "member")),
) -> dict:
    with get_db() as db:
        context = get_user_workspace_context(db, int(user["id"]))
        organization_id = int(context["organization_id"])
        workspace_id = int(context["workspace_id"])
        if payload.project_id and not _can_access_project(db, user, int(payload.project_id), organization_id, workspace_id):
            raise HTTPException(status_code=403, detail="この案件メモリを保存する権限がありません。")
        cursor = db.execute(
            """
            INSERT INTO proposal_agent_memories
            (
                user_id,
                project_id,
                project_name,
                hearing_notes,
                confirmation_items,
                proposal_content,
                competitor_analysis,
                improvement_history,
                organization_id,
                workspace_id
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                int(user["id"]),
                payload.project_id,
                payload.project_name[:200],
                payload.hearing_notes[:2000],
                payload.confirmation_items[:2000],
                payload.proposal_content[:2000],
                payload.competitor_analysis[:2000],
                payload.improvement_history[:2000],
                organization_id,
                workspace_id,
            ),
        )
        memory = _get_memory(db, int(cursor.lastrowid))
        create_audit_log(db, int(user["id"]), "save_proposal_agent_memory", "proposal_agent_memories", str(cursor.lastrowid), "success", "sanitized=true")
    return {"memory": memory}


def _list_projects(db: Any, user: dict, organization_id: int, workspace_id: int) -> list[dict[str, Any]]:
    params: list[Any] = [organization_id, workspace_id]
    member_clause = ""
    if str(user.get("role") or "") not in {"admin", "manager"}:
        member_clause = """
            AND EXISTS (
                SELECT 1 FROM project_lifecycle_events e
                WHERE e.project_id = p.id AND e.user_id = ?
            )
        """
        params.append(int(user["id"]))
    rows = db.execute(
        f"""
        SELECT p.*, COALESCE(c.company_name, '') AS customer_name
        FROM projects p
        LEFT JOIN customers c ON c.id = p.customer_id
        WHERE p.organization_id = ? AND p.workspace_id = ?
        {member_clause}
        ORDER BY p.updated_at DESC, p.id DESC
        LIMIT 50
        """,
        tuple(params),
    ).fetchall()
    return [dict(row) for row in rows]


def _list_histories(db: Any, user: dict, organization_id: int, workspace_id: int, project_ids: list[int]) -> list[dict[str, Any]]:
    params: list[Any] = [organization_id, workspace_id]
    member_clause = ""
    if str(user.get("role") or "") not in {"admin", "manager"}:
        member_clause = "AND h.user_id = ?"
        params.append(int(user["id"]))
    rows = db.execute(
        f"""
        SELECT h.*
        FROM proposal_histories h
        WHERE h.organization_id = ? AND h.workspace_id = ? AND COALESCE(h.is_demo, 0) = 0
        {member_clause}
        ORDER BY h.created_at DESC, h.id DESC
        LIMIT 200
        """,
        tuple(params),
    ).fetchall()
    histories = [dict(row) for row in rows]
    if project_ids:
        return histories
    return histories


def _list_memories(db: Any, user: dict, organization_id: int, workspace_id: int) -> list[dict[str, Any]]:
    params: list[Any] = [organization_id, workspace_id]
    member_clause = ""
    if str(user.get("role") or "") not in {"admin", "manager"}:
        member_clause = "AND m.user_id = ?"
        params.append(int(user["id"]))
    rows = db.execute(
        f"""
        SELECT
            m.*,
            COALESCE(u.display_name, '') AS created_by_name,
            COALESCE(u.email, '') AS created_by_email
        FROM proposal_agent_memories m
        LEFT JOIN users u ON u.id = m.user_id
        WHERE m.organization_id = ? AND m.workspace_id = ?
        {member_clause}
        ORDER BY m.updated_at DESC, m.id DESC
        LIMIT 50
        """,
        tuple(params),
    ).fetchall()
    return [dict(row) for row in rows]


def _get_memory(db: Any, memory_id: int) -> dict[str, Any]:
    row = db.execute(
        """
        SELECT
            m.*,
            COALESCE(u.display_name, '') AS created_by_name,
            COALESCE(u.email, '') AS created_by_email
        FROM proposal_agent_memories m
        LEFT JOIN users u ON u.id = m.user_id
        WHERE m.id = ?
        """,
        (memory_id,),
    ).fetchone()
    return dict(row) if row else {}


def _can_access_project(db: Any, user: dict, project_id: int, organization_id: int, workspace_id: int) -> bool:
    row = db.execute(
        "SELECT organization_id, workspace_id FROM projects WHERE id = ?",
        (project_id,),
    ).fetchone()
    if not row:
        return False
    if int(row["organization_id"] or 0) != organization_id or int(row["workspace_id"] or 0) != workspace_id:
        return False
    if str(user.get("role") or "") in {"admin", "manager"}:
        return True
    event = db.execute("SELECT 1 FROM project_lifecycle_events WHERE project_id = ? AND user_id = ? LIMIT 1", (project_id, int(user["id"]))).fetchone()
    return bool(event)


def _load_dashboard(db: Any, user: dict) -> dict[str, Any]:
    context = get_user_workspace_context(db, int(user["id"]))
    organization_id = int(context["organization_id"])
    workspace_id = int(context["workspace_id"])
    projects = _list_projects(db, user, organization_id, workspace_id)
    project_ids = [int(project["id"]) for project in projects]
    histories = _list_histories(db, user, organization_id, workspace_id, project_ids)
    memories = _list_memories(db, user, organization_id, workspace_id)
    business_reports = _list_business_improvement_reports(db, user, organization_id, workspace_id)
    timelines = _build_timelines(db, project_ids)
    status_cards = _build_status_cards(projects, histories)
    scores = [_score_project(project, histories, memories) for project in projects[:8]]
    priorities = _build_priorities(projects[:8], scores, histories, memories)
    win_probabilities = _build_win_probabilities(projects[:8], scores, histories, memories)
    health = _build_health(projects[:8], scores, priorities, win_probabilities)
    competitors = _build_competitors(projects[:8], memories)
    actions = _build_sales_actions(projects[:8], priorities, win_probabilities, health, memories)
    todo = _build_todo(status_cards, scores, memories)
    review = _build_review(todo, scores)
    kpi = _build_kpi(histories, scores, win_probabilities, business_reports)
    insights = _build_insights(projects, scores, priorities, win_probabilities, competitors, todo)
    summaries = _build_summaries(projects, status_cards, scores, todo, review)
    return {
        "status_cards": status_cards,
        "todo": todo,
        "scores": scores,
        "timeline": timelines,
        "memories": memories,
        "review": review,
        "summaries": summaries,
        "priorities": priorities,
        "win_probabilities": win_probabilities,
        "competitors": competitors,
        "sales_actions": actions,
        "health": health,
        "kpi": kpi,
        "insights": insights,
    }


def _list_business_improvement_reports(db: Any, user: dict, organization_id: int, workspace_id: int) -> list[dict[str, Any]]:
    params: list[Any] = [organization_id, workspace_id]
    member_clause = ""
    if str(user.get("role") or "") not in {"admin", "manager"}:
        member_clause = "AND user_id = ?"
        params.append(int(user["id"]))
    rows = db.execute(
        f"""
        SELECT *
        FROM business_improvement_reports
        WHERE organization_id = ? AND workspace_id = ? AND COALESCE(is_demo, 0) = 0
        {member_clause}
        ORDER BY created_at DESC, id DESC
        LIMIT 500
        """,
        tuple(params),
    ).fetchall()
    return [dict(row) for row in rows]


def _build_status_cards(projects: list[dict[str, Any]], histories: list[dict[str, Any]]) -> list[dict[str, Any]]:
    project_ids = {int(project["id"]) for project in projects}
    proposal_projects = _project_ids_by_history(histories, {"markdown", "markdown+pptx-data"})
    estimate_projects = _project_ids_by_history(histories, {"estimate-pdf"})
    beautiful_projects = _project_ids_by_history(histories, {"beautiful-ai"})
    completed_projects = {int(project["id"]) for project in projects if str(project.get("status") or "") in {"提出済み", "商談中", "受注", "完了"}}
    return [
        {"key": "proposal_waiting", "label": "提案待ち", "count": max(len(project_ids - proposal_projects), 0), "tone": "warn"},
        {"key": "proposal_building", "label": "提案書作成中", "count": _count_failed_or_pending(histories), "tone": "info"},
        {"key": "proposal_done", "label": "提案完了", "count": len(proposal_projects), "tone": "ok"},
        {"key": "estimate_waiting", "label": "見積作成待ち", "count": max(len(proposal_projects - estimate_projects), 0), "tone": "warn"},
        {"key": "beautiful_waiting", "label": "Beautiful.ai生成待ち", "count": max(len(proposal_projects - beautiful_projects), 0), "tone": "warn"},
        {"key": "customer_waiting", "label": "顧客送付待ち", "count": max(len(proposal_projects - completed_projects), 0), "tone": "info"},
    ]


def _project_ids_by_history(histories: list[dict[str, Any]], output_types: set[str]) -> set[int]:
    ids: set[int] = set()
    for item in histories:
        project_id = int(item.get("project_id") or 0)
        if project_id and str(item.get("output_type") or "") in output_types and str(item.get("status") or "") == "success":
            ids.add(project_id)
    return ids


def _count_failed_or_pending(histories: list[dict[str, Any]]) -> int:
    return sum(1 for item in histories if str(item.get("status") or "") not in {"success", "completed"})


def _score_project(project: dict[str, Any], histories: list[dict[str, Any]], memories: list[dict[str, Any]]) -> dict[str, Any]:
    project_id = int(project.get("id") or 0)
    project_histories = [item for item in histories if int(item.get("project_id") or 0) == project_id]
    project_memory = next((item for item in memories if int(item.get("project_id") or 0) == project_id), {})
    criteria = [
        _criterion("課題整理", 80 if len(str(project.get("summary") or "")) >= 30 or project_memory.get("hearing_notes") else 50, "案件概要とヒアリング内容を補強します。"),
        _criterion("提案内容", 85 if _has_output(project_histories, {"markdown", "markdown+pptx-data"}) or project_memory.get("proposal_content") else 45, "提案書の初稿を作成します。"),
        _criterion("競合分析", 80 if project_memory.get("competitor_analysis") else 45, "競合比較を追加します。"),
        _criterion("見積", 82 if _has_output(project_histories, {"estimate-pdf"}) else 50, "見積PDFまたは見積条件を確認します。"),
        _criterion("ストーリー性", 78 if _has_output(project_histories, {"pptx", "summary-pptx", "beautiful-ai"}) else 55, "スライド構成と説明順を整えます。"),
        _criterion("確認事項", 82 if project_memory.get("confirmation_items") else 50, "未確認事項をチェックリスト化します。"),
    ]
    score = round(sum(item["score"] for item in criteria) / len(criteria))
    improvements = [item["improvement"] for item in criteria if item["score"] < 70]
    return {
        "project_id": project_id,
        "project_name": project.get("name") or project.get("project_name") or "案件",
        "customer_name": project.get("customer_name") or "",
        "score": score,
        "criteria": criteria,
        "improvements": improvements[:5],
    }


def _criterion(label: str, score: int, improvement: str) -> dict[str, Any]:
    return {"label": label, "score": score, "improvement": improvement}


def _has_output(histories: list[dict[str, Any]], output_types: set[str]) -> bool:
    return any(str(item.get("output_type") or "") in output_types and str(item.get("status") or "") == "success" for item in histories)


def _build_todo(status_cards: list[dict[str, Any]], scores: list[dict[str, Any]], memories: list[dict[str, Any]]) -> list[dict[str, Any]]:
    count_by_key = {item["key"]: int(item["count"]) for item in status_cards}
    tasks = [
        {"label": "ヒアリング不足", "reason": "確認事項やヒアリング内容をAgent Memoryへ保存してください。", "priority": "高", "checked": bool(memories)},
        {"label": "予算確認が必要", "reason": "見積作成前に予算条件を確認してください。", "priority": "高", "checked": count_by_key.get("estimate_waiting", 0) == 0},
        {"label": "競合比較を追加", "reason": "競合分析が弱い案件の説得力を補強します。", "priority": "中", "checked": not any("競合比較を追加します。" in score.get("improvements", []) for score in scores)},
        {"label": "見積確認", "reason": "見積PDFが未作成の案件を確認します。", "priority": "中", "checked": count_by_key.get("estimate_waiting", 0) == 0},
        {"label": "スライド生成", "reason": "PowerPointまたはBeautiful.aiの出力待ちを確認します。", "priority": "中", "checked": count_by_key.get("beautiful_waiting", 0) == 0},
        {"label": "顧客へ送付", "reason": "提案完了後、送付待ち案件を進めます。", "priority": "低", "checked": count_by_key.get("customer_waiting", 0) == 0},
    ]
    return tasks


def _build_priorities(
    projects: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    histories: list[dict[str, Any]],
    memories: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    score_by_project = {int(item["project_id"]): item for item in scores}
    priorities = []
    for project in projects:
        project_id = int(project.get("id") or 0)
        memory = _memory_for_project(memories, project_id)
        text = _project_text(project, memory)
        project_histories = [item for item in histories if int(item.get("project_id") or 0) == project_id]
        proposal_score = int(score_by_project.get(project_id, {}).get("score") or 0)
        criteria = [
            _priority_criterion("予算", _has_any(text, ["予算", "万円", "budget"]) or int(project.get("win_probability") or 0) >= 60, "予算条件が見えている案件は優先度を上げます。"),
            _priority_criterion("納期", _has_any(text, ["納期", "導入", "公開", "月", "202"]), "納期や導入時期がある案件は次アクションを急ぎます。"),
            _priority_criterion("業種", _has_any(text, ["AI", "OCR", "SaaS", "IT", "DX", "CRM", "RPA", "製造", "医療"]), "業種やテーマが明確です。"),
            _priority_criterion("案件規模", len(text) >= 120 or len(project_histories) >= 2, "入力情報や履歴が多く、提案化しやすい案件です。"),
            _priority_criterion("競合状況", bool(memory.get("competitor_analysis")), "競合状況が整理されています。"),
            _priority_criterion("過去受注率", int(project.get("win_probability") or 0) >= 60, "CRM上の受注確度が高めです。"),
            _priority_criterion("提案難易度", proposal_score >= 70, "提案品質スコアが一定以上です。"),
        ]
        priority_score = round(sum(item["score"] for item in criteria) / len(criteria))
        grade = _priority_grade(priority_score)
        priorities.append(
            {
                "project_id": project_id,
                "project_name": project.get("name") or "案件",
                "priority_score": priority_score,
                "grade": grade,
                "stars": _priority_stars(priority_score),
                "criteria": criteria,
                "reasons": [item["reason"] for item in criteria if item["score"] >= 70][:4],
            }
        )
    return priorities


def _priority_criterion(label: str, passed: bool, reason: str) -> dict[str, Any]:
    return {"label": label, "score": 85 if passed else 45, "reason": reason if passed else f"{label}の確認が必要です。"}


def _priority_grade(score: int) -> str:
    if score >= 85:
        return "A"
    if score >= 70:
        return "B"
    if score >= 55:
        return "C"
    if score >= 40:
        return "D"
    return "E"


def _priority_stars(score: int) -> str:
    count = 5 if score >= 85 else 4 if score >= 70 else 3 if score >= 55 else 2 if score >= 40 else 1
    return "★" * count + "☆" * (5 - count)


def _build_win_probabilities(
    projects: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    histories: list[dict[str, Any]],
    memories: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    score_by_project = {int(item["project_id"]): item for item in scores}
    items = []
    for project in projects:
        project_id = int(project.get("id") or 0)
        memory = _memory_for_project(memories, project_id)
        project_histories = [item for item in histories if int(item.get("project_id") or 0) == project_id]
        proposal_score = int(score_by_project.get(project_id, {}).get("score") or 0)
        probability = int(project.get("win_probability") or 45)
        reasons = []
        if memory.get("hearing_notes"):
            probability += 10
            reasons.append("ヒアリング十分")
        if _has_output(project_histories, {"estimate-pdf"}):
            probability += 8
            reasons.append("見積作成済み")
        if memory.get("competitor_analysis"):
            probability += 7
            reasons.append("競合状況を把握")
        if proposal_score >= 75:
            probability += 8
            reasons.append("Proposal Scoreが高い")
        if _has_any(_project_text(project, memory), ["予算", "万円", "budget"]):
            probability += 7
            reasons.append("予算条件が見えている")
        probability = max(5, min(95, probability))
        if not reasons:
            reasons = ["案件情報がまだ不足しています。"]
        items.append(
            {
                "project_id": project_id,
                "project_name": project.get("name") or "案件",
                "probability": probability,
                "reasons": reasons[:4],
            }
        )
    return items


def _build_competitors(projects: list[dict[str, Any]], memories: list[dict[str, Any]]) -> list[dict[str, Any]]:
    cards: list[dict[str, Any]] = []
    for project in projects:
        project_id = int(project.get("id") or 0)
        memory = _memory_for_project(memories, project_id)
        competitor_text = str(memory.get("competitor_analysis") or "").strip()
        if not competitor_text:
            continue
        names = _competitor_names(competitor_text)
        for name in names[:2]:
            cards.append(
                {
                    "project_id": project_id,
                    "project_name": project.get("name") or "案件",
                    "competitor_name": name,
                    "strengths": ["既存接点や価格訴求で比較される可能性があります。"],
                    "weaknesses": ["業務理解や運用定着まで踏み込めない可能性があります。"],
                    "differentiation": ["提案書、見積、Beautiful.ai、改善測定まで一気通貫で支援する点を訴求します。"],
                    "cautions": ["価格だけで比較されないよう、導入後の運用成果を説明してください。"],
                }
            )
    if cards:
        return cards[:6]
    return [
        {
            "project_id": 0,
            "project_name": "共通",
            "competitor_name": "価格重視の競合",
            "strengths": ["初期費用を低く見せやすい。"],
            "weaknesses": ["提案品質や導入後の運用支援が弱くなりやすい。"],
            "differentiation": ["時間短縮、品質管理、提出物の一貫性を数値で説明します。"],
            "cautions": ["価格比較になった場合は、削減時間と確認品質をセットで示してください。"],
        },
        {
            "project_id": 0,
            "project_name": "共通",
            "competitor_name": "既存ベンダー",
            "strengths": ["顧客業務をすでに理解している。"],
            "weaknesses": ["新しいAI活用や提案改善のスピードが出にくい。"],
            "differentiation": ["Proposal Agentが次アクション、リスク、レビューまで継続支援する点を訴求します。"],
            "cautions": ["既存ベンダーとの置き換えではなく、補完・高度化として説明すると進めやすいです。"],
        },
    ]


def _competitor_names(text: str) -> list[str]:
    normalized = text.replace("、", "\n").replace(",", "\n").replace("・", "\n")
    names = [line.strip(" 　:：。") for line in normalized.splitlines() if line.strip()]
    return names[:4] or ["競合候補"]


def _build_sales_actions(
    projects: list[dict[str, Any]],
    priorities: list[dict[str, Any]],
    win_probabilities: list[dict[str, Any]],
    health: list[dict[str, Any]],
    memories: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    priority_by_project = {int(item["project_id"]): item for item in priorities}
    win_by_project = {int(item["project_id"]): item for item in win_probabilities}
    health_by_project = {int(item["project_id"]): item for item in health}
    actions = []
    for project in projects[:8]:
        project_id = int(project.get("id") or 0)
        memory = _memory_for_project(memories, project_id)
        action = "メール送信"
        reason = "提案内容を整理し、次回打ち合わせの確認を送ります。"
        if not memory.get("hearing_notes"):
            action = "ヒアリング追加"
            reason = "提案前に課題、予算、決裁者を確認してください。"
        elif not memory.get("competitor_analysis"):
            action = "競合調査"
            reason = "競合比較が不足しているため差別化ポイントを補強します。"
        elif int(win_by_project.get(project_id, {}).get("probability") or 0) < 55:
            action = "電話する"
            reason = "受注確率が低めのため、懸念点を直接確認します。"
        elif str(health_by_project.get(project_id, {}).get("status") or "") != "Healthy":
            action = "見積修正"
            reason = "案件健康度に注意があるため、見積と提案範囲を確認します。"
        actions.append(
            {
                "project_id": project_id,
                "project_name": project.get("name") or "案件",
                "action": action,
                "priority": priority_by_project.get(project_id, {}).get("grade", "C"),
                "reason": reason,
            }
        )
    return actions


def _build_health(
    projects: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    priorities: list[dict[str, Any]],
    win_probabilities: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    score_by_project = {int(item["project_id"]): int(item["score"]) for item in scores}
    priority_by_project = {int(item["project_id"]): item for item in priorities}
    win_by_project = {int(item["project_id"]): int(item["probability"]) for item in win_probabilities}
    items = []
    for project in projects:
        project_id = int(project.get("id") or 0)
        score = score_by_project.get(project_id, 0)
        probability = win_by_project.get(project_id, 0)
        priority = str(priority_by_project.get(project_id, {}).get("grade") or "C")
        if score < 55 or probability < 35 or priority in {"D", "E"}:
            status = "Critical"
            reason = "提案品質、受注確率、優先度のいずれかが低く、早めの見直しが必要です。"
        elif score < 70 or probability < 60 or priority == "C":
            status = "Warning"
            reason = "提案は進められますが、確認事項や競合比較の補強が必要です。"
        else:
            status = "Healthy"
            reason = "提案準備は順調です。顧客送付や次回商談へ進められます。"
        items.append({"project_id": project_id, "project_name": project.get("name") or "案件", "status": status, "reason": reason})
    return items


def _build_kpi(
    histories: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    win_probabilities: list[dict[str, Any]],
    business_reports: list[dict[str, Any]],
) -> dict[str, Any]:
    proposal_histories = [item for item in histories if str(item.get("output_type") or "") in {"markdown", "markdown+pptx-data"}]
    success_count = sum(1 for item in proposal_histories if str(item.get("status") or "") == "success")
    proposal_count = len(proposal_histories)
    total_duration = [int(item.get("total_generation_duration_ms") or 0) for item in histories if int(item.get("total_generation_duration_ms") or 0) > 0]
    return {
        "proposal_count": proposal_count,
        "proposal_success_rate": round((success_count / proposal_count) * 100, 1) if proposal_count else 0,
        "average_proposal_score": round(sum(int(item.get("score") or 0) for item in scores) / len(scores), 1) if scores else 0,
        "average_win_probability": round(sum(int(item.get("probability") or 0) for item in win_probabilities) / len(win_probabilities), 1) if win_probabilities else 0,
        "average_generation_time_seconds": round((sum(total_duration) / len(total_duration)) / 1000, 1) if total_duration else 0,
        "total_saved_minutes": round(sum(float(item.get("saved_minutes") or 0) for item in business_reports), 1),
        "beautiful_ai_count": sum(1 for item in histories if str(item.get("output_type") or "") == "beautiful-ai" and str(item.get("status") or "") == "success"),
    }


def _build_insights(
    projects: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    priorities: list[dict[str, Any]],
    win_probabilities: list[dict[str, Any]],
    competitors: list[dict[str, Any]],
    todo: list[dict[str, Any]],
) -> list[str]:
    insights: list[str] = []
    missing_competition = sum(1 for score in scores if "競合比較を追加します。" in score.get("improvements", []))
    hearing_todo = next((item for item in todo if item["label"] == "ヒアリング不足"), {})
    high_win_ai = [
        item
        for item in win_probabilities
        if int(item.get("probability") or 0) >= 70 and _has_any(_project_name(projects, int(item.get("project_id") or 0)), ["AI", "OCR", "DX", "CRM", "RPA"])
    ]
    if missing_competition:
        insights.append(f"競合比較不足が{missing_competition}件あります。差別化ポイントを先に補強すると提案品質が上がります。")
    if hearing_todo and not hearing_todo.get("checked"):
        insights.append("ヒアリング不足が見られます。予算、決裁者、導入時期を確認してください。")
    if high_win_ai:
        insights.append("AI/DX系案件は受注確率が高めです。PoC条件と評価指標を早めに固めると進めやすいです。")
    if competitors and competitors[0].get("project_id") != 0:
        insights.append("競合情報が蓄積されています。商談では価格ではなく運用定着と成果測定を訴求してください。")
    if any(str(item.get("grade") or "") in {"A", "B"} for item in priorities):
        insights.append("優先度A/Bの案件があります。提案書、見積、顧客送付まで一気に進める候補です。")
    return insights or ["案件データが増えると、最近の傾向や改善ポイントを自動で表示します。"]


def _project_name(projects: list[dict[str, Any]], project_id: int) -> str:
    project = next((item for item in projects if int(item.get("id") or 0) == project_id), {})
    return str(project.get("name") or "")


def _memory_for_project(memories: list[dict[str, Any]], project_id: int) -> dict[str, Any]:
    return next((item for item in memories if int(item.get("project_id") or 0) == project_id), {})


def _project_text(project: dict[str, Any], memory: dict[str, Any]) -> str:
    return "\n".join(
        [
            str(project.get("name") or ""),
            str(project.get("summary") or ""),
            str(project.get("next_action") or ""),
            str(project.get("customer_name") or ""),
            str(memory.get("hearing_notes") or ""),
            str(memory.get("confirmation_items") or ""),
            str(memory.get("proposal_content") or ""),
            str(memory.get("competitor_analysis") or ""),
        ]
    )


def _has_any(text: str, keywords: list[str]) -> bool:
    normalized = text.lower()
    return any(keyword.lower() in normalized for keyword in keywords)


def _build_timelines(db: Any, project_ids: list[int]) -> list[dict[str, Any]]:
    if not project_ids:
        return []
    placeholders = ",".join("?" for _ in project_ids[:20])
    rows = db.execute(
        f"""
        SELECT e.project_id, p.name AS project_name, e.event_type, e.from_status, e.to_status, e.note, e.created_at
        FROM project_lifecycle_events e
        INNER JOIN projects p ON p.id = e.project_id
        WHERE e.project_id IN ({placeholders})
        ORDER BY e.created_at DESC, e.id DESC
        LIMIT 80
        """,
        tuple(project_ids[:20]),
    ).fetchall()
    events = [
        {
            "project_id": int(row["project_id"]),
            "project_name": row["project_name"],
            "label": _timeline_label(str(row["event_type"] or ""), str(row["to_status"] or "")),
            "detail": row["note"] or f"{row['from_status']} → {row['to_status']}",
            "created_at": row["created_at"],
        }
        for row in rows
    ]
    return events


def _timeline_label(event_type: str, to_status: str) -> str:
    if to_status:
        return to_status
    labels = {
        "created": "案件登録",
        "analysis": "AI分析",
        "proposal": "提案生成",
        "beautiful_ai": "Beautiful.ai生成",
        "sent": "顧客送付",
        "won": "受注",
    }
    return labels.get(event_type, event_type or "更新")


def _build_review(todo: list[dict[str, Any]], scores: list[dict[str, Any]]) -> dict[str, list[str]]:
    improvements = [task["label"] for task in todo if not task["checked"]][:5]
    low_score_projects = [score for score in scores if int(score.get("score") or 0) < 70]
    risks = [f"{score['project_name']}の提案品質スコアが{score['score']}点です。" for score in low_score_projects[:3]]
    missing = []
    for score in scores:
        missing.extend(score.get("improvements", [])[:2])
    return {
        "improvements": improvements or ["現在のチェックリストは概ね完了しています。"],
        "risks": risks or ["重大なリスクは検出されていません。"],
        "missing_information": list(dict.fromkeys(missing))[:5] or ["不足情報はありません。"],
    }


def _build_summaries(
    projects: list[dict[str, Any]],
    status_cards: list[dict[str, Any]],
    scores: list[dict[str, Any]],
    todo: list[dict[str, Any]],
    review: dict[str, list[str]],
) -> dict[str, str]:
    project_count = len(projects)
    average_score = round(sum(int(item.get("score") or 0) for item in scores) / len(scores)) if scores else 0
    waiting = next((item["count"] for item in status_cards if item["key"] == "proposal_waiting"), 0)
    next_task = next((item["label"] for item in todo if not item["checked"]), "顧客送付準備")
    executive = f"現在の案件は{project_count}件、平均提案スコアは{average_score}点です。提案待ちは{waiting}件で、次に優先すべき作業は「{next_task}」です。"
    sales = (
        f"営業向けには、まず提案待ち{waiting}件を確認し、次に{next_task}を進めます。"
        f"平均スコアは{average_score}点のため、低い項目は{', '.join(review['missing_information'][:2])}を中心に補強します。"
        "提案書、見積、Beautiful.ai、顧客送付の順に状態を確認すれば、次の商談準備まで迷わず進められます。"
    )
    detail = (
        f"案件総数: {project_count}件\n"
        f"平均提案スコア: {average_score}点\n"
        f"次アクション: {next_task}\n"
        f"改善点: {', '.join(review['improvements'])}\n"
        f"リスク: {', '.join(review['risks'])}\n"
        f"不足情報: {', '.join(review['missing_information'])}"
    )
    return {"executive_30s": executive, "sales_3m": sales, "detail": detail}


def _render_dashboard_markdown(dashboard: dict[str, Any]) -> str:
    lines = [
        "# Proposal Intelligence Dashboard",
        "",
        "## KPI",
    ]
    for key, label in _kpi_labels().items():
        lines.append(f"- {label}: {dashboard['kpi'].get(key, 0)}")
    lines.extend(["", "## Priority Engine"])
    for item in dashboard.get("priorities", []):
        lines.append(f"- {item['project_name']}: {item['grade']} / {item['stars']} / {item['priority_score']}点")
        for reason in item.get("reasons", []):
            lines.append(f"  - {reason}")
    lines.extend(["", "## 受注確率"])
    for item in dashboard.get("win_probabilities", []):
        lines.append(f"- {item['project_name']}: {item['probability']}%（{', '.join(item.get('reasons', []))}）")
    lines.extend(["", "## 案件健康度"])
    for item in dashboard.get("health", []):
        lines.append(f"- {item['project_name']}: {item['status']} - {item['reason']}")
    lines.extend(["", "## AI Insights"])
    for insight in dashboard.get("insights", []):
        lines.append(f"- {insight}")
    lines.extend(["", "## 営業アクション"])
    for item in dashboard.get("sales_actions", []):
        lines.append(f"- {item['project_name']}: {item['action']}（{item['reason']}）")
    lines.extend(["", "## Executive Summary", dashboard.get("summaries", {}).get("executive_30s", "")])
    return "\n".join(lines)


def _render_dashboard_csv(dashboard: dict[str, Any]) -> str:
    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(["section", "project", "metric", "value", "reason"])
    for key, label in _kpi_labels().items():
        writer.writerow(["KPI", "", label, dashboard["kpi"].get(key, 0), ""])
    for item in dashboard.get("priorities", []):
        writer.writerow(["Priority", item["project_name"], item["grade"], item["priority_score"], " / ".join(item.get("reasons", []))])
    for item in dashboard.get("win_probabilities", []):
        writer.writerow(["Win Probability", item["project_name"], "受注確率", f"{item['probability']}%", " / ".join(item.get("reasons", []))])
    for item in dashboard.get("health", []):
        writer.writerow(["Health", item["project_name"], item["status"], "", item["reason"]])
    for item in dashboard.get("sales_actions", []):
        writer.writerow(["Sales Action", item["project_name"], item["action"], item["priority"], item["reason"]])
    for insight in dashboard.get("insights", []):
        writer.writerow(["Insight", "", "AI Insight", "", insight])
    return output.getvalue()


def _render_dashboard_pdf(dashboard: dict[str, Any]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    from reportlab.lib import colors

    from app.services import pdf_service

    pdf_service.register_japanese_fonts()
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=16 * mm, rightMargin=16 * mm, topMargin=16 * mm, bottomMargin=14 * mm)
    title = ParagraphStyle("AgentTitle", fontName=pdf_service.FONT_GOTHIC, fontSize=18, leading=24, textColor=colors.HexColor("#10233F"))
    body = ParagraphStyle("AgentBody", fontName=pdf_service.FONT_GOTHIC, fontSize=10, leading=15, textColor=colors.HexColor("#1B2430"))
    story: list[Any] = [Paragraph("Proposal Intelligence Dashboard", title), Spacer(1, 8)]
    story.append(Paragraph(dashboard.get("summaries", {}).get("executive_30s", ""), body))
    story.append(Spacer(1, 10))
    kpi_rows = [["指標", "値"]] + [[label, str(dashboard["kpi"].get(key, 0))] for key, label in _kpi_labels().items()]
    table = Table(kpi_rows, colWidths=[80 * mm, 80 * mm])
    table.setStyle(
        TableStyle(
            [
                ("FONTNAME", (0, 0), (-1, -1), pdf_service.FONT_GOTHIC),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EAF2FF")),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D8DEE8")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 10))
    story.append(Paragraph("AI Insights", title))
    for insight in dashboard.get("insights", []):
        story.append(Paragraph(f"・{insight}", body))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Priority Engine", title))
    for item in dashboard.get("priorities", [])[:8]:
        story.append(Paragraph(f"{item['project_name']}: {item['grade']} / {item['stars']} / {item['priority_score']}点", body))
    doc.build(story)
    return buffer.getvalue()


def _render_dashboard_pptx(dashboard: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide: Any, title: str, subtitle: str = "") -> None:
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.5), Inches(0.55))
        frame = title_box.text_frame
        frame.text = title
        frame.paragraphs[0].font.size = Pt(26)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = RGBColor(16, 35, 63)
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.03), Inches(10), Inches(0.4))
            sub.text_frame.text = subtitle
            sub.text_frame.paragraphs[0].font.size = Pt(12)
            sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(85, 96, 116)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Proposal Intelligence Platform", dashboard.get("summaries", {}).get("executive_30s", ""))
    kpi_items = list(_kpi_labels().items())
    for index, (key, label) in enumerate(kpi_items):
        x = Inches(0.65 + (index % 4) * 3.1)
        y = Inches(1.75 + (index // 4) * 1.35)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.75), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
        shape.line.color.rgb = RGBColor(191, 219, 254)
        tf = shape.text_frame
        tf.text = f"{label}\n{dashboard['kpi'].get(key, 0)}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[1].font.size = Pt(22)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Priority / Win Probability", "優先順位と受注確率")
    rows = dashboard.get("priorities", [])[:6]
    for index, item in enumerate(rows):
        win = next((prob for prob in dashboard.get("win_probabilities", []) if prob["project_id"] == item["project_id"]), {})
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35 + index * 0.72), Inches(11.7), Inches(0.55))
        box.text_frame.text = f"{item['project_name']}  {item['grade']} / {item['stars']}  受注確率 {win.get('probability', 0)}%"
        box.text_frame.paragraphs[0].font.size = Pt(16)
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(16, 35, 63)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "AI Insights / Next Actions", "最近の傾向と次アクション")
    lines = [f"・{item}" for item in dashboard.get("insights", [])]
    lines += [f"・{item['project_name']}: {item['action']}" for item in dashboard.get("sales_actions", [])[:5]]
    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.7), Inches(5.2))
    body.text_frame.text = "\n".join(lines)
    for paragraph in body.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(16, 35, 63)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _kpi_labels() -> dict[str, str]:
    return {
        "proposal_count": "提案数",
        "proposal_success_rate": "提案成功率",
        "average_proposal_score": "平均Proposal Score",
        "average_win_probability": "平均受注確率",
        "average_generation_time_seconds": "平均作成時間(秒)",
        "total_saved_minutes": "累計削減時間(分)",
        "beautiful_ai_count": "Beautiful.ai生成数",
    }
