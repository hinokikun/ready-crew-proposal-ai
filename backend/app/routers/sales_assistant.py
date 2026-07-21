from datetime import datetime
from io import BytesIO
from typing import Any, Literal
import unicodedata
from urllib.parse import quote
from zipfile import BadZipFile, ZipFile

from fastapi import APIRouter, Depends, Header, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pydantic import BaseModel, Field, validator
from pydantic.fields import ModelField

from app.auth import require_roles
from app.beautiful_ai.schemas import BeautifulAiPresentationRequest, BeautifulAiPresentationResponse
from app.config import settings
from app.db import get_db
from app.models import AnalysisResponse, ProposalRequest, PptxDownloadRequest
from app.rate_limit import rate_limit_dependency
from app.sales_assistant.generator import generate_sales_assistant_brief
from app.sales_assistant.models import SalesAssistantBrief, SalesAssistantInput
from app.services.beautiful_ai_service import BeautifulAiServiceError, create_beautiful_ai_presentation
from app.services.openai_service import OpenAIServiceError, generate_proposal
from app.services.pptx_service import MEDIA_TYPE, build_pptx_filename
from app.services.presentation_engine_integration import build_pptx_bytes_for_engine
from app.strategy_engine.enums import BudgetType
from app.strategy_engine.evaluator import evaluate_strategy
from app.strategy_engine.models import ProposalStrategyInput, StrategyBrief


router = APIRouter(prefix="/api/sales-assistant", tags=["sales-assistant"])

MAX_REQUEST_BYTES = 64_000
MAX_EXPORT_REQUEST_BYTES = 256_000
MAX_LIST_ITEMS = {
    "known_requirements": 20,
    "known_constraints": 20,
    "previous_interactions": 20,
    "evidence_items": 30,
}
MAX_LIST_ITEM_LENGTH = {
    "known_requirements": 2_000,
    "known_constraints": 2_000,
    "previous_interactions": 3_000,
    "evidence_items": 3_000,
}


def _clean_text(value: Any) -> str:
    text = str(value or "").strip()
    return "".join(
        char
        for char in unicodedata.normalize("NFKC", text)
        if char in {"\n", "\t"} or unicodedata.category(char)[0] != "C"
    ).strip()


class SalesAssistantGenerateRequest(BaseModel):
    project_title: str = Field("", max_length=200)
    project_summary: str = Field("", max_length=10_000)
    client_name: str = Field("", max_length=200)
    known_requirements: list[str] = Field(default_factory=list)
    known_constraints: list[str] = Field(default_factory=list)
    budget_information: str = Field("", max_length=2_000)
    schedule_information: str = Field("", max_length=2_000)
    meeting_stage: str = "preparation"
    previous_interactions: list[str] = Field(default_factory=list)
    evidence_items: list[str] = Field(default_factory=list)
    industry: str = Field("", max_length=200)
    business_goals: list[str] = Field(default_factory=list)
    current_problems: list[str] = Field(default_factory=list)
    proposed_solution: str = Field("", max_length=5_000)
    expected_deliverables: list[str] = Field(default_factory=list)
    integrations: list[str] = Field(default_factory=list)
    expected_kpis: list[str] = Field(default_factory=list)
    risks: list[str] = Field(default_factory=list)
    stakeholders: list[str] = Field(default_factory=list)

    @validator(
        "project_title",
        "project_summary",
        "client_name",
        "budget_information",
        "schedule_information",
        "industry",
        "proposed_solution",
        pre=True,
        always=True,
        allow_reuse=True,
    )
    def normalize_text(cls, value: Any) -> str:
        return _clean_text(value)

    @validator(
        "known_requirements",
        "known_constraints",
        "previous_interactions",
        "evidence_items",
        "business_goals",
        "current_problems",
        "expected_deliverables",
        "integrations",
        "expected_kpis",
        "risks",
        "stakeholders",
        pre=True,
        always=True,
        allow_reuse=True,
    )
    def normalize_list(cls, value: Any, field: ModelField) -> list[str]:
        raw_items = value if isinstance(value, list) else ([] if value is None else [value])
        limit = MAX_LIST_ITEMS.get(field.name, 20)
        item_limit = MAX_LIST_ITEM_LENGTH.get(field.name, 2_000)
        if len(raw_items) > limit:
            raise ValueError(f"{field.name} must contain no more than {limit} items")
        normalized = []
        for item in raw_items:
            text = _clean_text(item)
            if not text:
                continue
            if len(text) > item_limit:
                raise ValueError(f"{field.name} item must be {item_limit} characters or fewer")
            normalized.append(text)
        return normalized

    @validator("meeting_stage", pre=True, always=True, allow_reuse=True)
    def normalize_meeting_stage(cls, value: Any) -> str:
        return _clean_text(value or "preparation")


class SalesAssistantProposalPreviewRequest(BaseModel):
    source_request: SalesAssistantGenerateRequest
    sales_assistant_brief: SalesAssistantBrief
    strategy_brief: StrategyBrief | None = None


ReviewStatus = Literal["unreviewed", "reviewed", "needs_revision", "regenerate_recommended", "approved", "exportable"]
ExportType = Literal["powerpoint", "beautiful_ai"]


class SalesAssistantExportRequest(BaseModel):
    export_type: ExportType
    source_request: SalesAssistantGenerateRequest
    sales_assistant_brief: SalesAssistantBrief
    proposal_preview: dict[str, Any]
    proposal_response: AnalysisResponse
    strategy_brief: StrategyBrief | None = None
    human_review_status: ReviewStatus = "unreviewed"
    human_review_required: bool = False
    project_id: str = Field("", max_length=120)
    force_new: bool = False


def _ensure_enabled() -> None:
    if not settings.sales_assistant_enabled:
        raise HTTPException(
            status_code=404,
            detail={
                "error_type": "sales_assistant_disabled",
                "message": "AI Sales Assistant is disabled by feature flag.",
            },
        )


def _ensure_proposal_enabled() -> None:
    _ensure_enabled()
    if not settings.sales_assistant_proposal_enabled:
        raise HTTPException(
            status_code=404,
            detail={
                "error_type": "sales_assistant_proposal_disabled",
                "message": "Sales Assistant to Proposal integration is disabled by feature flag.",
            },
        )


def _ensure_export_enabled() -> None:
    _ensure_proposal_enabled()
    if not settings.proposal_export_enabled:
        raise HTTPException(
            status_code=404,
            detail={
                "error_type": "proposal_export_disabled",
                "message": "Proposal Export integration is disabled by feature flag.",
            },
        )


def _check_content_length(content_length: int | None, max_bytes: int = MAX_REQUEST_BYTES) -> None:
    if content_length is not None and content_length > max_bytes:
        raise HTTPException(
            status_code=413,
            detail={
                "error_type": "sales_assistant_request_too_large",
                "message": "Request body is too large.",
            },
        )


@router.get("/status")
async def get_sales_assistant_status(_: dict = Depends(require_roles("admin"))) -> dict:
    return {
        "enabled": settings.sales_assistant_enabled,
        "version": "50",
        "requires_admin": True,
        "persistence_enabled": False,
        "external_ai_enabled": False,
        "proposal_preview_enabled": settings.sales_assistant_proposal_enabled,
        "proposal_export_enabled": settings.proposal_export_enabled,
        "beautiful_ai_export_enabled": settings.beautiful_ai_enabled or settings.beautiful_ai_mock,
    }


@router.post("/generate", dependencies=[Depends(rate_limit_dependency("admin"))])
async def generate_sales_assistant(
    payload: SalesAssistantGenerateRequest,
    _: dict = Depends(require_roles("admin")),
    content_length: int | None = Header(default=None),
) -> dict:
    _check_content_length(content_length)
    _ensure_enabled()
    if not payload.project_title or not payload.project_summary:
        raise HTTPException(
            status_code=422,
            detail={
                "error_type": "sales_assistant_input_error",
                "message": "project_title and project_summary are required.",
            },
        )
    try:
        strategy_input = _build_strategy_input(payload)
        strategy_brief = evaluate_strategy(strategy_input)
        assistant_input = SalesAssistantInput(
            strategy_brief=strategy_brief,
            project_title=payload.project_title,
            project_summary=payload.project_summary,
            client_name=payload.client_name,
            known_requirements=payload.known_requirements,
            known_constraints=payload.known_constraints,
            budget_information=payload.budget_information,
            schedule_information=payload.schedule_information,
            meeting_stage=payload.meeting_stage,
            previous_interactions=payload.previous_interactions,
            evidence_items=payload.evidence_items,
        )
        sales_assistant_brief = generate_sales_assistant_brief(assistant_input)
    except ValueError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "error_type": "sales_assistant_input_error",
                "message": str(exc),
            },
        ) from exc
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "sales_assistant_generation_error",
                "message": "AI Sales Assistant could not generate a brief.",
            },
        ) from exc

    return {
        "sales_assistant_brief": sales_assistant_brief.dict(),
        "strategy_brief": strategy_brief.dict(),
        "strategy_brief_summary": {
            "schema_version": strategy_brief.schema_version,
            "category": strategy_brief.project_category,
            "project_category": strategy_brief.project_category,
            "persona": strategy_brief.primary_persona,
            "primary_persona": strategy_brief.primary_persona,
            "decision_maker": strategy_brief.decision_maker,
            "strategy": strategy_brief.primary_strategy,
            "primary_strategy": strategy_brief.primary_strategy,
            "story": strategy_brief.story_type,
            "story_type": strategy_brief.story_type,
            "primary_pack": strategy_brief.primary_pack,
            "kpi_pack": strategy_brief.kpi_pack,
            "estimate_pack": strategy_brief.estimate_pack,
            "confidence": strategy_brief.confidence,
            "human_review_required": strategy_brief.human_review_required,
            "human_review_reasons": strategy_brief.human_review_reasons,
        },
        "warnings": sales_assistant_brief.generation_metadata.warnings,
        "human_review_required": sales_assistant_brief.summary.human_review_required,
        "human_review_reasons": sales_assistant_brief.summary.human_review_reasons,
        "generation_metadata": sales_assistant_brief.generation_metadata.dict(),
    }


@router.post("/proposal-preview", dependencies=[Depends(rate_limit_dependency("admin"))])
async def generate_sales_assistant_proposal_preview(
    payload: SalesAssistantProposalPreviewRequest,
    _: dict = Depends(require_roles("admin")),
    content_length: int | None = Header(default=None),
) -> dict:
    _check_content_length(content_length)
    _ensure_proposal_enabled()
    strategy_brief = payload.strategy_brief or evaluate_strategy(_build_strategy_input(payload.source_request))
    proposal_request = _build_proposal_request_for_preview(
        payload.source_request,
        payload.sales_assistant_brief,
        strategy_brief,
    )
    try:
        proposal_response = await generate_proposal(proposal_request)
    except OpenAIServiceError as exc:
        raise HTTPException(
            status_code=exc.status_code,
            detail={
                "error_type": "sales_assistant_proposal_generation_error",
                "message": "Proposal Preview could not be generated.",
            },
        ) from exc
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "sales_assistant_proposal_generation_error",
                "message": "Proposal Preview could not be generated.",
            },
        ) from exc

    preview = _build_proposal_preview(
        proposal_response,
        payload.source_request,
        payload.sales_assistant_brief,
        strategy_brief,
    )
    return {
        "proposal_preview": preview,
        "proposal_response": proposal_response.dict(),
        "human_review_required": preview["human_review_required"],
        "human_review_reasons": preview["human_review_reasons"],
        "generation_metadata": {
            "schema_version": "sales_assistant_proposal_preview_v1",
            "proposal_generator": "app.services.openai_service.generate_proposal",
            "strategy_brief_version": strategy_brief.schema_version,
            "sales_assistant_version": payload.sales_assistant_brief.generation_metadata.generator_version,
            "persistence_enabled": False,
            "pptx_enabled": False,
            "beautiful_ai_enabled": False,
        },
    }


@router.post("/export", dependencies=[Depends(rate_limit_dependency("admin"))])
async def export_sales_assistant_proposal(
    payload: SalesAssistantExportRequest,
    user: dict = Depends(require_roles("admin")),
    content_length: int | None = Header(default=None),
) -> dict:
    _check_content_length(content_length, MAX_EXPORT_REQUEST_BYTES)
    _ensure_export_enabled()
    strategy_brief = payload.strategy_brief or evaluate_strategy(_build_strategy_input(payload.source_request))
    _ensure_human_review_allows_export(payload, strategy_brief)
    pptx_payload = _build_pptx_download_request(payload, strategy_brief)
    request_json_safe = {
        "export_type": payload.export_type,
        "human_review_status": payload.human_review_status,
        "human_review_required": _is_human_review_required(payload, strategy_brief),
        "project_id": _project_id(payload),
        "slides": len(pptx_payload.powerpoint_generation_data.slides),
        "deck_title": pptx_payload.powerpoint_generation_data.deck_title[:120],
        "client_name": pptx_payload.powerpoint_generation_data.client_name[:120],
    }

    if payload.export_type == "powerpoint":
        return _export_powerpoint(pptx_payload, request_json_safe)

    return await _export_beautiful_ai(payload, pptx_payload, request_json_safe, int(user["id"]))


@router.post("/export/download", dependencies=[Depends(rate_limit_dependency("admin"))])
async def download_sales_assistant_export(
    payload: SalesAssistantExportRequest,
    _: dict = Depends(require_roles("admin")),
    content_length: int | None = Header(default=None),
) -> StreamingResponse:
    _check_content_length(content_length, MAX_EXPORT_REQUEST_BYTES)
    _ensure_export_enabled()
    if payload.export_type != "powerpoint":
        raise HTTPException(
            status_code=422,
            detail={
                "error_type": "proposal_export_download_type_invalid",
                "message": "Only PowerPoint exports can be downloaded from this endpoint.",
            },
        )

    # Download uses the same validation as metadata export because Version54
    # intentionally avoids DB-backed artifact IDs.
    strategy_brief = payload.strategy_brief or evaluate_strategy(_build_strategy_input(payload.source_request))
    _ensure_human_review_allows_export(payload, strategy_brief)
    pptx_payload = _build_pptx_download_request(payload, strategy_brief)
    pptx_bytes, filename, _engine_mode = _build_powerpoint_export_file(pptx_payload)
    return StreamingResponse(
        BytesIO(pptx_bytes),
        media_type=MEDIA_TYPE,
        headers={
            "Content-Disposition": _content_disposition(filename),
            "X-Content-Type-Options": "nosniff",
            "Cache-Control": "no-store",
        },
    )


def _build_strategy_input(payload: SalesAssistantGenerateRequest) -> ProposalStrategyInput:
    source_text = " ".join(
        item
        for item in [
            payload.project_title,
            payload.project_summary,
            payload.client_name,
            payload.budget_information,
            payload.schedule_information,
            " ".join(payload.known_requirements),
            " ".join(payload.known_constraints),
            " ".join(payload.previous_interactions),
            " ".join(payload.evidence_items),
        ]
        if item
    )
    return ProposalStrategyInput(
        project_title=payload.project_title,
        project_summary=payload.project_summary,
        industry=payload.industry,
        business_goals=payload.business_goals or payload.known_requirements,
        current_problems=payload.current_problems or payload.known_constraints,
        proposed_solution=payload.proposed_solution or payload.project_summary,
        expected_deliverables=payload.expected_deliverables,
        integrations=payload.integrations,
        schedule=payload.schedule_information,
        budget=payload.budget_information,
        budget_type=BudgetType.CONFIRMED if payload.budget_information else BudgetType.UNKNOWN,
        expected_kpis=payload.expected_kpis,
        risks=payload.risks,
        stakeholders=payload.stakeholders,
        constraints=payload.known_constraints,
        source_text=source_text,
    )


def _build_proposal_request_for_preview(
    payload: SalesAssistantGenerateRequest,
    assistant_brief: SalesAssistantBrief,
    strategy_brief: StrategyBrief,
) -> ProposalRequest:
    summary = assistant_brief.summary
    meeting_plan = assistant_brief.meeting_plan
    objections = [
        f"{item.expected_objection}: {item.recommended_response}"
        for item in assistant_brief.objection_handling[:5]
    ]
    next_actions = [item.action for item in assistant_brief.next_actions[:6]]
    customer_pain = payload.current_problems or payload.known_constraints or assistant_brief.evidence_guidance.evidence_gaps
    project_brief = _join_lines(
        [
            f"案件名: {payload.project_title}",
            f"案件概要: {payload.project_summary}",
            f"Strategy: {strategy_brief.primary_strategy}",
            f"Story: {strategy_brief.story_type}",
            f"Persona: {strategy_brief.primary_persona}",
            f"Decision Maker: {strategy_brief.decision_maker}",
            f"Meeting Goal: {meeting_plan.objective}",
            f"Customer Pain: {_join_items(customer_pain)}",
            f"Recommended Positioning: {summary.recommended_positioning}",
            f"Objections: {_join_items(objections)}",
            f"Next Actions: {_join_items(next_actions)}",
            f"Main Message: {summary.primary_message}",
        ]
    )
    return ProposalRequest(
        project_brief=project_brief[:8000],
        client_company_info=payload.client_name,
        desired_launch_timing=payload.schedule_information,
        budget_range=payload.budget_information,
        hearing_result=_join_lines(
            [
                f"Known requirements: {_join_items(payload.known_requirements)}",
                f"Known constraints: {_join_items(payload.known_constraints)}",
                f"Previous interactions: {_join_items(payload.previous_interactions)}",
                f"Human review reasons: {_join_items(assistant_brief.summary.human_review_reasons)}",
            ]
        ),
        own_service_info=_join_lines(
            [
                f"Strategy Brief main message: {strategy_brief.main_message}",
                f"Priority messages: {_join_items(strategy_brief.priority_messages)}",
                f"Risk messages: {_join_items(strategy_brief.risk_messages)}",
                f"KPI pack: {strategy_brief.kpi_pack}",
                f"Estimate pack: {strategy_brief.estimate_pack}",
            ]
        ),
        past_proposal_template=_join_lines(
            [
                f"Story type: {strategy_brief.story_type}",
                f"Presentation pack: {strategy_brief.primary_pack}",
                f"Required slides: {_join_items(strategy_brief.required_slide_types)}",
            ]
        ),
        case_studies=_join_items(payload.evidence_items),
    )


def _ensure_human_review_allows_export(payload: SalesAssistantExportRequest, strategy_brief: StrategyBrief) -> None:
    if payload.human_review_status in {"needs_revision", "regenerate_recommended"}:
        raise HTTPException(
            status_code=409,
            detail={
                "error_type": "proposal_export_review_not_approved",
                "message": "Human Review indicates the proposal is not ready for export.",
            },
        )
    if _is_human_review_required(payload, strategy_brief) and payload.human_review_status not in {"approved", "exportable"}:
        raise HTTPException(
            status_code=409,
            detail={
                "error_type": "proposal_export_review_required",
                "message": "Human Review approval is required before export.",
            },
        )


def _is_human_review_required(payload: SalesAssistantExportRequest, strategy_brief: StrategyBrief) -> bool:
    preview_required = bool(payload.proposal_preview.get("human_review_required")) if isinstance(payload.proposal_preview, dict) else False
    return (
        bool(payload.human_review_required)
        or preview_required
        or payload.sales_assistant_brief.summary.human_review_required
        or strategy_brief.human_review_required
    )


def _build_pptx_download_request(
    payload: SalesAssistantExportRequest,
    strategy_brief: StrategyBrief,
) -> PptxDownloadRequest:
    response = payload.proposal_response
    return PptxDownloadRequest(
        powerpoint_generation_data=response.powerpoint_generation_data,
        win_probability=response.analysis.win_probability,
        strategy_review_report=_strategy_review_report(payload, strategy_brief),
        project_brief=payload.source_request.project_summary,
        client_company_info=payload.source_request.client_name,
        desired_launch_timing=payload.source_request.schedule_information,
        budget_range=payload.source_request.budget_information,
        hearing_result=_join_lines(
            [
                f"Known requirements: {_join_items(payload.source_request.known_requirements)}",
                f"Known constraints: {_join_items(payload.source_request.known_constraints)}",
                f"Human review status: {payload.human_review_status}",
            ]
        ),
        own_service_info=_join_lines(
            [
                f"Strategy: {strategy_brief.primary_strategy}",
                f"Story: {strategy_brief.story_type}",
                f"Persona: {strategy_brief.primary_persona}",
                f"Sales positioning: {payload.sales_assistant_brief.summary.recommended_positioning}",
            ]
        ),
        past_proposal_template=_join_lines(
            [
                f"Presentation pack: {strategy_brief.primary_pack}",
                f"Required slides: {_join_items(strategy_brief.required_slide_types)}",
            ]
        ),
        case_studies=_join_items(payload.source_request.evidence_items),
    )


def _strategy_review_report(payload: SalesAssistantExportRequest, strategy_brief: StrategyBrief) -> dict[str, Any]:
    return {
        "review_result": "approve" if payload.human_review_status in {"approved", "exportable", "reviewed"} else payload.human_review_status,
        "strategy_brief": strategy_brief.dict(),
        "review_comment": "Sales Assistant Proposal Preview export review.",
    }


def _export_powerpoint(pptx_payload: PptxDownloadRequest, request_json_safe: dict[str, Any]) -> dict:
    pptx_bytes, filename, engine_mode = _build_powerpoint_export_file(pptx_payload)
    return {
        "export_type": "powerpoint",
        "status": "success",
        "message": "PowerPoint export generated successfully.",
        "artifact": {
            "filename": filename,
            "content_type": MEDIA_TYPE,
            "byte_size": len(pptx_bytes),
            "download_url": "/api/sales-assistant/export/download",
            "download_method": "POST",
            "expires_on_refresh": True,
        },
        "request_json_safe": request_json_safe,
        "response_json_safe": {
            "filename": filename,
            "byte_size": len(pptx_bytes),
            "engine_mode": engine_mode,
            "download_url": "/api/sales-assistant/export/download",
            "download_method": "POST",
        },
    }


def _build_powerpoint_export_file(pptx_payload: PptxDownloadRequest) -> tuple[bytes, str, str]:
    try:
        engine_result = build_pptx_bytes_for_engine(pptx_payload)
        pptx_bytes = engine_result.pptx_bytes
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_powerpoint_failed",
                "message": "PowerPoint export failed.",
            },
        ) from exc
    _validate_pptx_bytes(pptx_bytes)
    filename = _build_export_filename(pptx_payload)
    return pptx_bytes, filename, engine_result.engine_mode


def _validate_pptx_bytes(pptx_bytes: bytes) -> None:
    if not pptx_bytes:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_empty",
                "message": "PowerPoint export generated an empty file.",
            },
        )
    try:
        with ZipFile(BytesIO(pptx_bytes)) as archive:
            broken_member = archive.testzip()
            names = set(archive.namelist())
    except BadZipFile as exc:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_invalid",
                "message": "PowerPoint export generated an invalid file.",
            },
        ) from exc
    if broken_member:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_invalid",
                "message": "PowerPoint export generated a broken file.",
            },
        )
    required = {"[Content_Types].xml", "ppt/presentation.xml"}
    if not required.issubset(names) or not any(name.startswith("ppt/slides/slide") and name.endswith(".xml") for name in names):
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_invalid",
                "message": "PowerPoint export is missing required presentation files.",
            },
        )
    try:
        presentation = Presentation(BytesIO(pptx_bytes))
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_invalid",
                "message": "PowerPoint export could not be reopened for validation.",
            },
        ) from exc
    if len(presentation.slides) < 1:
        raise HTTPException(
            status_code=500,
            detail={
                "error_type": "proposal_export_pptx_empty",
                "message": "PowerPoint export contains no slides.",
            },
        )


def _build_export_filename(pptx_payload: PptxDownloadRequest) -> str:
    base_filename = build_pptx_filename(pptx_payload.powerpoint_generation_data, pptx_payload.client_company_info)
    base = base_filename.rsplit(".", 1)[0]
    cleaned = _safe_filename_part(base)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M")
    return f"ProposalPilot_{cleaned}_{timestamp}.pptx"


def _safe_filename_part(value: str) -> str:
    normalized = unicodedata.normalize("NFKC", value or "proposal")
    allowed = []
    for char in normalized:
        category = unicodedata.category(char)
        if char.isalnum() or char in {"-", "_"}:
            allowed.append(char)
        elif char.isspace():
            allowed.append("_")
        elif category.startswith("L") or category.startswith("N"):
            allowed.append(char)
    safe = "".join(allowed).strip("._-")
    while "__" in safe:
        safe = safe.replace("__", "_")
    return (safe or "proposal")[:80]


def _content_disposition(filename: str) -> str:
    return f"attachment; filename=\"ProposalPilot_proposal.pptx\"; filename*=UTF-8''{quote(filename)}"


async def _export_beautiful_ai(
    payload: SalesAssistantExportRequest,
    pptx_payload: PptxDownloadRequest,
    request_json_safe: dict[str, Any],
    user_id: int,
) -> dict:
    request = BeautifulAiPresentationRequest(
        project_id=_project_id(payload),
        powerpoint_generation_data=pptx_payload.powerpoint_generation_data,
        win_probability=pptx_payload.win_probability,
        project_brief=pptx_payload.project_brief,
        client_company_info=pptx_payload.client_company_info,
        desired_launch_timing=pptx_payload.desired_launch_timing,
        budget_range=pptx_payload.budget_range,
        own_service_info=pptx_payload.own_service_info,
        past_proposal_template=pptx_payload.past_proposal_template,
        case_studies=pptx_payload.case_studies,
        force_new=payload.force_new,
    )
    try:
        with get_db() as db:
            response = await create_beautiful_ai_presentation(db, request=request, user_id=user_id)
    except BeautifulAiServiceError as exc:
        raise HTTPException(
            status_code=exc.status_code,
            detail={
                "error_type": exc.error_type,
                "message": exc.message,
                "fallback_available": True,
            },
        ) from exc
    return {
        "export_type": "beautiful_ai",
        "status": "success",
        "message": "Beautiful.ai export generated successfully.",
        "artifact": {
            "presentation_id": response.presentation_id,
            "editor_url": response.editor_url,
            "player_url": response.player_url,
            "title": response.title,
        },
        "request_json_safe": {
            **request_json_safe,
            "project_id": request.project_id,
            "force_new": request.force_new,
        },
        "response_json_safe": response.dict(),
    }


def _project_id(payload: SalesAssistantExportRequest) -> str:
    if payload.project_id.strip():
        return payload.project_id.strip()[:120]
    seed = f"{payload.source_request.project_title}-{payload.source_request.client_name}".strip("-") or "sales-assistant-proposal"
    normalized = unicodedata.normalize("NFKC", seed)
    safe = "".join(char.lower() if char.isalnum() else "-" for char in normalized)
    safe = "-".join(part for part in safe.split("-") if part)
    return f"sales-assistant-{safe}"[:120] or "sales-assistant-proposal"


def _build_proposal_preview(
    proposal_response: AnalysisResponse,
    payload: SalesAssistantGenerateRequest,
    assistant_brief: SalesAssistantBrief,
    strategy_brief: StrategyBrief,
) -> dict:
    analysis = proposal_response.analysis
    issues = [
        {
            "issue": item.issue,
            "background": item.background,
            "proposed_response": "",
            "confidence": item.confidence,
        }
        for item in analysis.assumed_customer_issues[:5]
    ]
    for item in analysis.issue_priorities[:5]:
        issues.append(
            {
                "issue": item.issue,
                "background": item.reason,
                "proposed_response": item.proposed_response,
                "confidence": "",
            }
        )
    slide_outline = [
        {
            "slide_no": slide.slide_no,
            "title": slide.title,
            "bullets": slide.bullets[:5],
            "visual_suggestion": slide.visual_suggestion,
        }
        for slide in proposal_response.powerpoint_generation_data.slides[:12]
    ]
    kpis = payload.expected_kpis or strategy_brief.priority_messages[:3] or [strategy_brief.kpi_pack]
    human_review_reasons = _unique_strings(
        list(assistant_brief.summary.human_review_reasons)
        + list(strategy_brief.human_review_reasons)
        + [analysis.quality_check.human_review_notes]
    )
    return {
        "proposal_summary": analysis.project_summary,
        "issues": issues[:8],
        "proposal_story": analysis.proposal_story,
        "proposal_policy": analysis.proposal_policy,
        "slide_outline": slide_outline,
        "kpis": kpis,
        "estimate_summary": _estimate_summary(payload, strategy_brief),
        "deck_title": proposal_response.powerpoint_generation_data.deck_title,
        "client_name": proposal_response.powerpoint_generation_data.client_name,
        "human_review_required": assistant_brief.summary.human_review_required or strategy_brief.human_review_required,
        "human_review_reasons": human_review_reasons,
        "source_versions": {
            "strategy_brief": strategy_brief.schema_version,
            "sales_assistant": assistant_brief.generation_metadata.generator_version,
        },
    }


def _estimate_summary(payload: SalesAssistantGenerateRequest, strategy_brief: StrategyBrief) -> str:
    budget = payload.budget_information or "予算情報は未確認です。"
    return (
        f"{budget} Estimate Packは {strategy_brief.estimate_pack} を使用します。"
        "正式見積は提案範囲、連携条件、スケジュール確認後に確定します。"
    )


def _join_lines(items: list[str]) -> str:
    return "\n".join(item for item in items if item and item.strip())


def _join_items(items: list[str]) -> str:
    return " / ".join(item for item in items if item and item.strip())


def _unique_strings(items: list[str]) -> list[str]:
    seen = set()
    result = []
    for item in items:
        text = str(item or "").strip()
        if text and text not in seen:
            seen.add(text)
            result.append(text)
    return result
