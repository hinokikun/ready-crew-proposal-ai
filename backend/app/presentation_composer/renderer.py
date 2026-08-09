"""Native PPTX renderer for Consulting Design System pages."""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.design_system import CONSULTING_DESIGN_SYSTEM, Palette

from .models import PageSpec, PresentationPlan


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _in(value: float):
    return Inches(value)


def _add_box(slide, x: float, y: float, w: float, h: float, fill: str, line: str = "", radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, _in(x), _in(y), _in(w), _in(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_icon(slide, x: float, y: float, label: str, palette: Palette, size: float = 0.36):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y), _in(size), _in(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(palette.accent)
    shape.line.fill.background()
    _add_text(slide, label[:1], x, y + 0.05, size, size * 0.55, size=9, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def _add_connector(slide, x1: float, y1: float, x2: float, y2: float, color: str, width: float = 1.4, arrow: bool = False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _in(x1), _in(y1), _in(x2), _in(y2))
    connector.line.color.rgb = _rgb(color)
    connector.line.width = Pt(width)
    return connector


def _add_slide_chrome(slide, plan: PresentationPlan, page: PageSpec, palette: Palette, dark: bool = False) -> None:
    fg = palette.ink if not dark else "#FFFFFF"
    muted = palette.muted if not dark else "#BFD3EA"
    _add_text(slide, page.action_title, 0.54, 0.28, 8.1, 0.48, size=25, color=fg, bold=True)
    _add_text(slide, page.conclusion, 0.56, 0.82, 7.9, 0.28, size=12, color=muted)
    _add_text(slide, f"{page.slide_no:02d}", 12.34, 0.34, 0.42, 0.18, size=8, color=muted, align=PP_ALIGN.RIGHT)
    _add_text(slide, plan.case.client_name, 9.1, 0.32, 3.0, 0.22, size=8, color=muted, align=PP_ALIGN.RIGHT)
    _add_text(slide, page.next_action, 0.56, 6.96, 8.8, 0.22, size=8, color=muted)
    _add_connector(slide, 0.54, 6.82, 12.72, 6.82, palette.line if not dark else "#314862", width=0.7)


def _render_hero(slide, plan: PresentationPlan, page: PageSpec, palette: Palette) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("#071426")
    _add_box(slide, 8.8, 0.0, 4.6, 7.5, "#0E2A47", "")
    _add_text(slide, plan.case.case_name, 0.72, 1.34, 7.3, 0.74, size=42, color="#FFFFFF", bold=True)
    _add_text(slide, plan.case.client_name, 0.78, 2.2, 5.4, 0.26, size=15, color="#BFD3EA")
    for i, label in enumerate(("背景", "課題", "提案", "効果")):
        _add_icon(slide, 8.95 + i * 0.94, 3.2 + (i % 2) * 0.86, label, palette, size=0.62)
        _add_connector(slide, 9.25 + i * 0.94, 3.52 + (i % 2) * 0.86, 9.78 + i * 0.94, 3.52 + ((i + 1) % 2) * 0.86, "#6DA9FF", width=1.8)
    _add_text(slide, "Consulting Design System Output", 0.78, 6.76, 3.5, 0.24, size=9, color="#BFD3EA")


def _render_issue_tree(slide, page: PageSpec, palette: Palette) -> None:
    root = _add_box(slide, 0.86, 2.8, 2.2, 0.82, palette.primary, "")
    _add_text(slide, page.diagram_labels[0] if page.diagram_labels else "課題", 1.0, 3.03, 1.9, 0.26, size=14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(page.diagram_labels[1:4] or ("原因A", "原因B", "原因C")):
        y = 1.52 + i * 1.34
        _add_connector(slide, 3.08, 3.21, 4.18, y + 0.36, palette.secondary, width=1.7)
        _add_box(slide, 4.18, y, 2.25, 0.72, "#FFFFFF", palette.line)
        _add_text(slide, label, 4.36, y + 0.22, 1.85, 0.2, size=13, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        _add_box(slide, 7.0, y - 0.06, 4.55, 0.84, "#EEF5FF", "")
        _add_text(slide, page.evidence, 7.18, y + 0.18, 4.1, 0.25, size=11, color=palette.ink)


def _render_split(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels) + ["現状", "改善後", "効果"]
    for i, title in enumerate(("現状", "転換点", "目指す姿")):
        x = 0.78 + i * 4.08
        fill = ["#F7FAFC", "#EAF4FF", "#ECFDF5"][i]
        _add_box(slide, x, 1.62, 3.42, 4.42, fill, palette.line)
        _add_text(slide, title, x + 0.22, 1.88, 2.8, 0.28, size=16, color=palette.primary, bold=True)
        _add_text(slide, labels[i], x + 0.24, 2.58, 2.75, 0.32, size=18, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        _add_icon(slide, x + 1.44, 3.36, title, palette, size=0.58)
        _add_text(slide, labels[i + 1], x + 0.38, 4.48, 2.5, 0.35, size=11, color=palette.muted, align=PP_ALIGN.CENTER)
    _add_connector(slide, 4.1, 3.7, 4.72, 3.7, palette.accent, width=2.1, arrow=True)
    _add_connector(slide, 8.18, 3.7, 8.8, 3.7, palette.accent, width=2.1, arrow=True)


def _render_matrix(slide, page: PageSpec, palette: Palette) -> None:
    x0, y0, w, h = 1.35, 1.55, 9.9, 4.82
    _add_connector(slide, x0, y0 + h / 2, x0 + w, y0 + h / 2, palette.line, width=1.1)
    _add_connector(slide, x0 + w / 2, y0, x0 + w / 2, y0 + h, palette.line, width=1.1)
    fills = ["#F8FAFC", "#EEF6FF", "#EEFDF8", "#FFF7ED"]
    labels = list(page.diagram_labels) + ["低", "高", "短期", "中期"]
    for i in range(4):
        x = x0 + (i % 2) * (w / 2 + 0.08)
        y = y0 + (i // 2) * (h / 2 + 0.08)
        _add_box(slide, x, y, w / 2 - 0.08, h / 2 - 0.08, fills[i], "")
        _add_text(slide, labels[i], x + 0.28, y + 0.38, w / 2 - 0.56, 0.28, size=15, color=palette.primary, bold=True)
        _add_text(slide, page.evidence, x + 0.28, y + 1.16, w / 2 - 0.56, 0.26, size=10, color=palette.muted)
    _add_text(slide, "実行しやすさ", 4.9, 6.48, 2.5, 0.2, size=9, color=palette.muted, align=PP_ALIGN.CENTER)
    _add_text(slide, "効果", 0.68, 3.62, 0.44, 0.2, size=9, color=palette.muted)


def _render_flow(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:5] or ["入力", "処理", "判断", "実行", "改善"]
    count = min(5, len(labels))
    for i, label in enumerate(labels[:count]):
        x = 0.78 + i * 2.38
        y = 3.0 + (0.42 if i % 2 else 0)
        _add_box(slide, x, y, 1.62, 0.92, "#FFFFFF", palette.line, radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_icon(slide, x + 0.18, y + 0.28, str(i + 1), palette, size=0.36)
        _add_text(slide, label, x + 0.58, y + 0.32, 0.86, 0.22, size=12, color=palette.ink, bold=True)
        if i < count - 1:
            _add_connector(slide, x + 1.62, y + 0.46, x + 2.3, 3.46 + (0.42 if (i + 1) % 2 else 0), palette.accent, width=1.8, arrow=True)


def _render_timeline(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:5] or ["要件", "設計", "PoC", "評価", "本番"]
    _add_connector(slide, 1.0, 3.62, 11.8, 3.62, palette.secondary, width=2.0)
    for i, label in enumerate(labels[:5]):
        x = 1.0 + i * 2.55
        _add_icon(slide, x, 3.34, str(i + 1), palette, size=0.55)
        _add_text(slide, label, x - 0.42, 4.08, 1.42, 0.25, size=12, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"Phase {i + 1}", x - 0.42, 2.72, 1.42, 0.18, size=8, color=palette.muted, align=PP_ALIGN.CENTER)


def _render_dashboard(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["時間", "品質", "精度", "負荷"]
    values = ["30%", "90%", "20%", "8週"]
    for i, label in enumerate(labels[:4]):
        x = 0.88 + i * 3.05
        _add_box(slide, x, 1.78, 2.42, 3.88, "#FFFFFF", palette.line)
        _add_text(slide, values[i], x + 0.2, 2.22, 2.0, 0.52, size=30, color=palette.primary, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, label, x + 0.25, 3.02, 1.92, 0.28, size=13, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        _add_connector(slide, x + 0.42, 4.2, x + 2.0, 4.2, palette.accent, width=4.5)
        _add_text(slide, page.evidence, x + 0.24, 4.68, 1.94, 0.26, size=9, color=palette.muted, align=PP_ALIGN.CENTER)


def _render_waterfall(slide, page: PageSpec, palette: Palette) -> None:
    labels = ["投資", "削減", "品質", "回収", "拡張"]
    heights = [1.1, 1.55, 1.0, 1.75, 1.25]
    x, baseline = 1.0, 5.65
    for i, label in enumerate(labels):
        fill = palette.primary if i in {0, 3} else palette.accent
        _add_box(slide, x + i * 2.1, baseline - heights[i], 1.32, heights[i], fill, "")
        _add_text(slide, label, x + i * 2.1 - 0.05, baseline + 0.18, 1.44, 0.22, size=10, color=palette.ink, align=PP_ALIGN.CENTER)
        _add_text(slide, ["費用", "時間", "精度", "回収", "将来"][i], x + i * 2.1 + 0.1, baseline - heights[i] + 0.18, 1.12, 0.2, size=11, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            _add_connector(slide, x + i * 2.1 + 1.32, baseline - heights[i], x + (i + 1) * 2.1, baseline - heights[i + 1], palette.line, width=1.0)


def _render_risk(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["精度", "定着", "連携", "運用"]
    colors = ["#ECFDF5", "#FEF3C7", "#FFEDD5", "#FEE2E2"]
    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            x, y = 2.0 + col * 4.2, 1.74 + row * 2.22
            _add_box(slide, x, y, 3.58, 1.72, colors[idx], palette.line)
            _add_text(slide, labels[idx], x + 0.24, y + 0.28, 2.8, 0.26, size=15, color=palette.primary, bold=True)
            _add_text(slide, "確認方法と対策をセットで管理", x + 0.24, y + 0.9, 2.9, 0.22, size=10, color=palette.muted)


def _render_tree(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["判断", "進む", "止める", "広げる"]
    _add_box(slide, 1.0, 3.0, 1.75, 0.78, palette.primary, "")
    _add_text(slide, labels[0], 1.16, 3.24, 1.4, 0.18, size=13, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(labels[1:4]):
        x = 5.2
        y = 1.72 + i * 1.42
        _add_connector(slide, 2.75, 3.39, x, y + 0.38, palette.accent, width=1.5)
        _add_box(slide, x, y, 4.2, 0.76, "#FFFFFF", palette.line)
        _add_text(slide, label, x + 0.24, y + 0.23, 3.4, 0.2, size=13, color=palette.ink, bold=True)


def _render_pyramid(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["結論", "効果", "根拠", "実行"]
    levels = [
        (5.25, 1.55, 2.8, 0.78, palette.primary, "#FFFFFF"),
        (4.55, 2.62, 4.2, 0.82, palette.accent, "#FFFFFF"),
        (3.75, 3.76, 5.8, 0.88, "#EAF4FF", palette.ink),
        (2.8, 4.96, 7.7, 0.92, "#F8FAFC", palette.ink),
    ]
    for i, (x, y, w, h, fill, text_color) in enumerate(levels):
        _add_box(slide, x, y, w, h, fill, "")
        _add_text(slide, labels[i % len(labels)], x + 0.18, y + 0.27, w - 0.36, 0.22, size=13, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, page.evidence, 4.1, 6.18, 5.1, 0.2, size=10, color=palette.muted, align=PP_ALIGN.CENTER)


def _render_fishbone(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:5] or ["課題", "人", "業務", "データ", "仕組み"]
    _add_connector(slide, 1.4, 3.6, 10.7, 3.6, palette.primary, width=2.0)
    _add_box(slide, 10.65, 3.15, 1.58, 0.9, palette.primary, "")
    _add_text(slide, labels[0], 10.78, 3.45, 1.28, 0.2, size=12, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    bones = [(2.5, 2.05), (4.3, 4.88), (6.1, 2.05), (7.9, 4.88)]
    for i, (x, y) in enumerate(bones, start=1):
        _add_connector(slide, x + 0.76, y + (0.72 if y < 3.6 else -0.1), x + 1.35, 3.6, palette.accent, width=1.5)
        _add_box(slide, x, y, 1.76, 0.64, "#FFFFFF", palette.line)
        _add_text(slide, labels[i % len(labels)], x + 0.12, y + 0.2, 1.5, 0.18, size=11, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)


def _render_architecture_layers(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:5] or ["入力", "保存", "AI", "確認", "連携"]
    for i, label in enumerate(labels[:5]):
        y = 1.48 + i * 0.9
        fill = palette.primary if i == 2 else "#FFFFFF"
        text = "#FFFFFF" if i == 2 else palette.ink
        _add_box(slide, 2.0 + i * 0.28, y, 8.5 - i * 0.56, 0.58, fill, palette.line)
        _add_text(slide, label, 2.28 + i * 0.28, y + 0.18, 7.7 - i * 0.56, 0.16, size=11, color=text, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "既存業務へ薄く重ね、深い連携はPoC後に判断", 3.1, 6.18, 6.9, 0.24, size=11, color=palette.muted, align=PP_ALIGN.CENTER)


def _render_landscape(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["自社", "競合A", "競合B", "代替案"]
    for i, label in enumerate(labels):
        x = 1.1 + (i % 2) * 5.2
        y = 1.72 + (i // 2) * 2.18
        fill = palette.primary if i == 0 else "#FFFFFF"
        text = "#FFFFFF" if i == 0 else palette.ink
        _add_box(slide, x, y, 4.28, 1.55, fill, palette.line)
        _add_text(slide, label, x + 0.22, y + 0.22, 3.3, 0.24, size=14, color=text, bold=True)
        _add_text(slide, "評価軸: 効果 / 速度 / 安心感", x + 0.22, y + 0.84, 3.5, 0.2, size=9, color=text if i == 0 else palette.muted)


def _render_blueprint(slide, page: PageSpec, palette: Palette) -> None:
    lanes = ["顧客接点", "業務処理", "データ", "管理"]
    steps = list(page.diagram_labels)[:4] or ["入力", "確認", "登録", "改善"]
    for i, lane in enumerate(lanes):
        y = 1.5 + i * 1.08
        _add_text(slide, lane, 0.86, y + 0.24, 1.15, 0.2, size=10, color=palette.primary, bold=True)
        _add_connector(slide, 2.08, y + 0.48, 12.1, y + 0.48, palette.line, width=0.8)
        for j, step in enumerate(steps):
            _add_box(slide, 2.25 + j * 2.34, y + 0.12, 1.5, 0.72, "#FFFFFF", palette.line)
            _add_text(slide, step, 2.34 + j * 2.34, y + 0.36, 1.28, 0.18, size=9, color=palette.ink, align=PP_ALIGN.CENTER)


def _render_value_chain(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:5] or ["入力", "処理", "判断", "実行", "成果"]
    for i, label in enumerate(labels[:5]):
        x = 0.9 + i * 2.36
        _add_box(slide, x, 2.44, 1.72, 1.32, palette.primary if i in {0, 4} else "#FFFFFF", palette.line)
        color = "#FFFFFF" if i in {0, 4} else palette.ink
        _add_text(slide, label, x + 0.14, 2.92, 1.42, 0.22, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            _add_connector(slide, x + 1.72, 3.1, x + 2.26, 3.1, palette.accent, width=2.0, arrow=True)
    _add_box(slide, 1.24, 4.68, 10.25, 0.62, "#EEF6FF", "")
    _add_text(slide, page.evidence, 1.46, 4.9, 9.7, 0.18, size=10, color=palette.primary, align=PP_ALIGN.CENTER)


def _render_org(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["責任者", "業務", "IT", "外部"]
    _add_box(slide, 5.25, 1.5, 2.1, 0.72, palette.primary, "")
    _add_text(slide, labels[0], 5.42, 1.73, 1.72, 0.18, size=13, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(labels[1:4]):
        x = 2.0 + i * 3.2
        _add_connector(slide, 6.3, 2.22, x + 1.05, 3.34, palette.accent, width=1.2)
        _add_box(slide, x, 3.34, 2.1, 0.82, "#FFFFFF", palette.line)
        _add_text(slide, label, x + 0.14, 3.6, 1.78, 0.18, size=12, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        _add_box(slide, x + 0.24, 4.62, 1.6, 0.44, "#EEF6FF", "")


def _render_cycle(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["計画", "実行", "確認", "改善"]
    coords = [(4.2, 1.55), (7.2, 2.8), (5.9, 5.0), (2.9, 3.75)]
    for i, (x, y) in enumerate(coords):
        _add_icon(slide, x, y, str(i + 1), palette, size=0.72)
        _add_text(slide, labels[i], x - 0.38, y + 0.86, 1.5, 0.22, size=12, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)
        nx, ny = coords[(i + 1) % len(coords)]
        _add_connector(slide, x + 0.62, y + 0.34, nx + 0.1, ny + 0.34, palette.accent, width=1.6, arrow=True)
    _add_text(slide, page.evidence, 5.0, 3.38, 2.4, 0.28, size=13, color=palette.primary, bold=True, align=PP_ALIGN.CENTER)


def _render_governance(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:4] or ["方針", "承認", "監査", "改善"]
    sizes = [(4.75, 1.28, 3.7, 3.7), (5.25, 1.78, 2.7, 2.7), (5.78, 2.31, 1.64, 1.64)]
    for i, (x, y, w, h) in enumerate(sizes):
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y), _in(w), _in(h))
        oval.fill.solid()
        oval.fill.fore_color.rgb = _rgb(["#EAF4FF", "#F8FAFC", palette.primary][i])
        oval.line.color.rgb = _rgb(palette.line)
    _add_text(slide, labels[0], 5.98, 2.98, 1.25, 0.2, size=12, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(labels[1:4]):
        _add_box(slide, 1.1 + i * 3.7, 5.55, 2.4, 0.58, "#FFFFFF", palette.line)
        _add_text(slide, label, 1.26 + i * 3.7, 5.76, 2.05, 0.16, size=10, color=palette.ink, bold=True, align=PP_ALIGN.CENTER)


def _render_next_action(slide, page: PageSpec, palette: Palette) -> None:
    labels = list(page.diagram_labels)[:3] or ["確認", "合意", "開始"]
    _add_box(slide, 1.1, 1.7, 10.8, 3.85, "#FFFFFF", palette.line)
    _add_text(slide, "次に決めること", 1.62, 2.18, 3.6, 0.34, size=21, color=palette.primary, bold=True)
    for i, label in enumerate(labels[:3]):
        x = 1.62 + i * 3.1
        _add_icon(slide, x, 3.22, str(i + 1), palette, size=0.62)
        _add_text(slide, label, x + 0.82, 3.38, 1.9, 0.22, size=13, color=palette.ink, bold=True)
        if i < 2:
            _add_connector(slide, x + 2.55, 3.54, x + 3.04, 3.54, palette.accent, width=2.2, arrow=True)
    _add_text(slide, page.next_action, 1.62, 4.66, 8.0, 0.24, size=12, color=palette.muted)


def _render_generic(slide, page: PageSpec, palette: Palette) -> None:
    visual = page.visual_type
    if visual in {"hero"}:
        return
    if visual in {"pyramid"}:
        _render_pyramid(slide, page, palette)
    elif visual in {"fishbone"}:
        _render_fishbone(slide, page, palette)
    elif visual in {"issue_tree"}:
        _render_issue_tree(slide, page, palette)
    elif visual in {"current_future", "before_after", "comparison", "feature_comparison"}:
        _render_split(slide, page, palette)
    elif visual in {"matrix", "priority_matrix", "swot", "pest", "three_c", "heatmap", "capability_map", "radar"}:
        _render_matrix(slide, page, palette)
    elif visual in {"architecture"}:
        _render_architecture_layers(slide, page, palette)
    elif visual in {"flow", "process", "sitemap", "implementation", "support", "operation"}:
        _render_flow(slide, page, palette)
    elif visual in {"timeline", "roadmap", "phase_gate", "milestone", "schedule", "gantt"}:
        _render_timeline(slide, page, palette)
    elif visual in {"kpi_dashboard", "kpi_cards"}:
        _render_dashboard(slide, page, palette)
    elif visual in {"value_chain", "business_model"}:
        _render_value_chain(slide, page, palette)
    elif visual in {"waterfall", "breakdown", "investment"}:
        _render_waterfall(slide, page, palette)
    elif visual in {"governance", "security"}:
        _render_governance(slide, page, palette)
    elif visual in {"risk_matrix"}:
        _render_risk(slide, page, palette)
    elif visual in {"next_action", "closing", "contact"}:
        _render_next_action(slide, page, palette)
    elif visual in {"decision_tree", "faq", "appendix", "funnel"}:
        _render_tree(slide, page, palette)
    elif visual in {"landscape"}:
        _render_landscape(slide, page, palette)
    elif visual in {"service_blueprint", "journey"}:
        _render_blueprint(slide, page, palette)
    elif visual in {"organization", "team", "stakeholder"}:
        _render_org(slide, page, palette)
    elif visual in {"cycle"}:
        _render_cycle(slide, page, palette)
    else:
        _render_flow(slide, page, palette)


def render_plan_to_pptx(plan: PresentationPlan) -> tuple[bytes, dict[str, Any]]:
    ds = CONSULTING_DESIGN_SYSTEM
    palette = ds.palette_for_category(plan.case.category)
    prs = Presentation()
    prs.slide_width = _in(ds.grid.slide_width)
    prs.slide_height = _in(ds.grid.slide_height)
    blank = prs.slide_layouts[6]
    render_pages: list[dict[str, Any]] = []
    for page in plan.pages:
        slide = prs.slides.add_slide(blank)
        dark = page.slide_no == 1
        if page.slide_no != 1:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _rgb(palette.background if page.slide_no % 5 else "#FFFFFF")
            _add_slide_chrome(slide, plan, page, palette, dark=False)
            _render_generic(slide, page, palette)
        else:
            _render_hero(slide, plan, page, palette)
        render_pages.append(
            {
                "slide_no": page.slide_no,
                "component_id": page.component_id,
                "component_name": page.component_name,
                "visual_type": page.visual_type,
                "layout_family": page.layout_family,
                "diagram_ratio": page.diagram_ratio,
                "text_ratio": page.text_ratio,
            }
        )
    output = BytesIO()
    prs.save(output)
    return output.getvalue(), {
        "provider": plan.provider,
        "slide_count": plan.slide_count,
        "distinct_layout_count": plan.distinct_layout_count,
        "average_diagram_ratio": plan.average_diagram_ratio,
        "average_text_ratio": plan.average_text_ratio,
        "pages": render_pages,
    }


V91_NAVY = "#071426"
V91_NAVY_2 = "#0B1F36"
V91_INK = "#111827"
V91_MUTED = "#475467"
V91_LINE = "#64748B"
V91_SOFT_LINE = "#CBD5E1"
V91_BLUE = "#2563EB"
V91_CYAN = "#00A6C8"
V91_GREEN = "#12B76A"
V91_AMBER = "#F59E0B"
V91_RED = "#EF4444"
V91_PAPER = "#F8FAFC"
V91_PANEL = "#EEF4FF"


def _add_v91_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str = "",
    *,
    line_width: float = 1.5,
    radius=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(radius, _in(x), _in(y), _in(w), _in(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_v91_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_v91_connector(slide, x1: float, y1: float, x2: float, y2: float, color: str, width: float = 2.0):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _in(x1), _in(y1), _in(x2), _in(y2))
    connector.line.color.rgb = _rgb(color)
    connector.line.width = Pt(width)
    return connector


def _add_v91_badge(slide, text: str, x: float, y: float, w: float, fill: str, color: str = "#FFFFFF"):
    _add_v91_box(slide, x, y, w, 0.34, fill, "", line_width=0.0)
    _add_v91_text(slide, text, x + 0.12, y + 0.09, w - 0.24, 0.12, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def _render_v91_cover(slide, case: CaseContext, proposal_date: str, proposer_name: str) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(V91_NAVY)
    _add_v91_box(slide, 8.45, 0.0, 4.9, 7.5, V91_NAVY_2, "")
    _add_v91_box(slide, 8.98, 0.68, 3.5, 5.95, "#102C4A", "#1E4F7A", line_width=1.2)
    _add_v91_text(slide, "Confidential", 0.74, 0.56, 1.45, 0.2, size=11, color="#BFD3EA", bold=True)
    _add_v91_text(slide, proposal_date, 2.28, 0.56, 1.6, 0.2, size=11, color="#BFD3EA")
    _add_v91_text(slide, case.client_name, 0.72, 1.14, 6.4, 0.42, size=22, color="#E6F0FF", bold=True)
    _add_v91_text(slide, "画像認識AI導入提案", 0.70, 1.92, 7.25, 1.03, size=44, color="#FFFFFF", bold=True)
    _add_v91_text(
        slide,
        "検品・判定業務を、AIで再現性ある判断プロセスへ変える",
        0.76,
        3.18,
        6.8,
        0.42,
        size=19,
        color="#BFD3EA",
        bold=True,
    )
    _add_v91_box(slide, 0.74, 4.22, 6.35, 1.18, "#0C2038", "#2B5B8A", line_width=1.4)
    _add_v91_text(slide, "提案の焦点", 1.0, 4.48, 1.35, 0.22, size=12, color=V91_CYAN, bold=True)
    _add_v91_text(slide, "目視判定の属人化を抑え、品質・速度・教育負荷を同時に改善する", 2.28, 4.42, 4.4, 0.38, size=17, color="#FFFFFF", bold=True)
    _add_v91_text(slide, proposer_name, 0.76, 6.66, 3.2, 0.26, size=13, color="#E6F0FF", bold=True)
    _add_v91_text(slide, "ProposalPilot / Ready Crew", 0.76, 6.96, 3.2, 0.18, size=10, color="#91A8C2")

    stages = [("目視判定", 9.28, 1.36), ("画像認識AI", 10.02, 3.04), ("業務判断", 10.78, 4.72)]
    for idx, (label, x, y) in enumerate(stages):
        _add_v91_box(slide, x, y, 1.68, 0.72, "#FFFFFF", "#5C7DA3", line_width=1.2)
        _add_v91_text(slide, label, x + 0.14, y + 0.23, 1.38, 0.18, size=12, color=V91_NAVY, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_box(slide, x + 0.54, y + 0.86, 0.58, 0.58, V91_CYAN if idx == 1 else "#2563EB", "")
    _add_v91_connector(slide, 10.11, 2.08, 10.84, 3.03, "#69D2E7", 2.4)
    _add_v91_connector(slide, 10.85, 3.76, 11.58, 4.71, "#69D2E7", 2.4)
    for x, y, r in [(9.36, 5.9, 0.26), (11.76, 1.02, 0.22), (12.0, 6.22, 0.18), (8.86, 3.98, 0.18)]:
        _add_v91_box(slide, x, y, r, r, "#69D2E7", "", radius=MSO_SHAPE.OVAL)
    return {
        "slide": "cover",
        "art_direction": "dark hero with customer identity and an editable image-recognition transformation visual",
        "takeaway": "画像認識AIの提案であること、顧客名、日付、提案者が一目で分かる表紙へ変更",
    }


def _render_v91_executive_summary(slide, case: CaseContext) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("#FFFFFF")
    _add_v91_text(slide, "Executive Summary", 0.62, 0.34, 2.6, 0.24, size=12, color=V91_CYAN, bold=True)
    _add_v91_text(slide, "目視依存の判定業務を、AIで再現性ある運用へ移行", 0.62, 0.68, 9.8, 0.52, size=28, color=V91_INK, bold=True)
    _add_v91_text(slide, case.client_name, 10.15, 0.43, 2.3, 0.18, size=10, color=V91_MUTED, align=PP_ALIGN.RIGHT)

    _add_v91_box(slide, 0.72, 1.48, 3.28, 4.62, V91_NAVY, "")
    _add_v91_text(slide, "今回の結論", 1.02, 1.84, 1.5, 0.22, size=13, color=V91_CYAN, bold=True)
    _add_v91_text(slide, "AIを小さく試し、判定基準を業務資産に変える", 1.0, 2.34, 2.56, 1.02, size=25, color="#FFFFFF", bold=True)
    _add_v91_text(slide, "投資判断は、精度だけでなく教育負荷・処理速度・判断ログの残り方で確認します。", 1.02, 4.1, 2.44, 0.72, size=15, color="#D8E8F9")
    _add_v91_badge(slide, "次回判断: PoC範囲と評価指標", 1.0, 5.32, 2.62, V91_CYAN)

    xs = [4.45, 6.85, 9.25]
    titles = ["Current", "Proposal", "Outcome"]
    subtitles = ["人の経験に依存", "AIで判定を補助", "判断を標準化"]
    fills = ["#F8FAFC", "#EEF4FF", "#ECFDF5"]
    accents = [V91_LINE, V91_BLUE, V91_GREEN]
    for i, x in enumerate(xs):
        _add_v91_box(slide, x, 2.0, 1.92, 2.74, fills[i], accents[i], line_width=1.8)
        _add_v91_text(slide, titles[i], x + 0.18, 2.26, 1.48, 0.18, size=12, color=accents[i], bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, subtitles[i], x + 0.22, 3.02, 1.42, 0.48, size=18, color=V91_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_box(slide, x + 0.62, 3.82, 0.66, 0.66, accents[i], "", radius=MSO_SHAPE.OVAL)
    _add_v91_connector(slide, 6.36, 3.36, 6.84, 3.36, V91_BLUE, 2.7)
    _add_v91_connector(slide, 8.76, 3.36, 9.24, 3.36, V91_GREEN, 2.7)

    metric_items = [("品質", "判定ばらつき低減"), ("速度", "確認時間の短縮"), ("教育", "新人教育を軽くする")]
    for i, (head, body) in enumerate(metric_items):
        x = 4.45 + i * 2.4
        _add_v91_box(slide, x, 5.15, 1.92, 0.72, "#FFFFFF", V91_SOFT_LINE, line_width=1.2)
        _add_v91_text(slide, head, x + 0.14, 5.32, 0.52, 0.16, size=12, color=V91_BLUE, bold=True)
        _add_v91_text(slide, body, x + 0.68, 5.31, 1.0, 0.18, size=11, color=V91_INK, bold=True, align=PP_ALIGN.CENTER)

    _add_v91_box(slide, 0.72, 6.42, 11.84, 0.55, "#F1F5F9", "")
    _add_v91_text(slide, "Takeaway: まずPoCで評価条件を合意し、精度・工数・運用負荷を確認して本導入判断へ進みます。", 0.98, 6.62, 10.9, 0.16, size=13, color=V91_INK, bold=True)
    _add_v91_text(slide, "数値は顧客提供データ確認後に確定 / 現時点は評価設計として提示", 8.72, 6.98, 3.5, 0.14, size=10, color=V91_MUTED, align=PP_ALIGN.RIGHT)
    return {
        "slide": "executive_summary",
        "art_direction": "large conclusion block plus current-proposal-outcome transformation diagram and decision takeaway",
        "takeaway": "6枚カードではなく投資判断の流れが見える1枚へ変更",
    }


def _render_v91_problem_structure(slide, case: CaseContext) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(V91_PAPER)
    _add_v91_text(slide, "現状と課題", 0.62, 0.34, 1.8, 0.24, size=12, color=V91_CYAN, bold=True)
    _add_v91_text(slide, "根本問題は、判定基準が人に閉じて業務資産になっていないことです", 0.62, 0.68, 10.6, 0.5, size=27, color=V91_INK, bold=True)
    _add_v91_text(slide, case.client_name, 10.15, 0.43, 2.3, 0.18, size=10, color=V91_MUTED, align=PP_ALIGN.RIGHT)

    columns = [
        ("症状", V91_BLUE, [("目視確認が多い", "判定が人ごとに揺れる"), ("繁忙時に滞留", "処理量の増減に弱い")]),
        ("業務上の原因", V91_AMBER, [("基準が暗黙知", "教育と引継ぎに時間がかかる"), ("記録が分散", "改善材料が残りにくい")]),
        ("根本原因", V91_RED, [("判断データ未蓄積", "AIで再利用できる形になっていない"), ("例外判断が属人化", "品質管理の再現性が弱い")]),
        ("事業影響", V91_GREEN, [("品質リスク", "誤判定と手戻りの発生"), ("成長制約", "取扱量拡大時に人員依存が残る")]),
    ]
    for i, (title, color, items) in enumerate(columns):
        x = 0.72 + i * 3.03
        _add_v91_box(slide, x, 1.58, 2.55, 0.54, color, "")
        _add_v91_text(slide, title, x + 0.12, 1.76, 2.28, 0.14, size=12, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        for j, (head, body) in enumerate(items):
            y = 2.45 + j * 1.32
            fill = "#FFFFFF" if i != 2 else "#FEF2F2"
            _add_v91_box(slide, x, y, 2.55, 0.92, fill, color, line_width=1.7)
            _add_v91_text(slide, head, x + 0.17, y + 0.18, 2.12, 0.18, size=13, color=V91_INK, bold=True, align=PP_ALIGN.CENTER)
            _add_v91_text(slide, body, x + 0.17, y + 0.52, 2.12, 0.16, size=10, color=V91_MUTED, align=PP_ALIGN.CENTER)
        if i < len(columns) - 1:
            _add_v91_connector(slide, x + 2.55, 3.5, x + 3.01, 3.5, color, 2.4)

    _add_v91_box(slide, 3.76, 5.52, 5.92, 0.78, V91_NAVY, "")
    _add_v91_text(slide, "本提案が変える範囲", 4.0, 5.76, 1.8, 0.16, size=12, color=V91_CYAN, bold=True)
    _add_v91_text(slide, "判定ルールの可視化 → AI補助 → 判断ログ蓄積 → 継続改善", 5.78, 5.73, 3.45, 0.18, size=13, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 0.72, 6.55, 11.84, 0.44, "#E2E8F0", "")
    _add_v91_text(slide, "Takeaway: AI導入の目的は置き換えではなく、判定品質を再現できる業務プロセスへ変えることです。", 0.98, 6.69, 10.9, 0.14, size=12, color=V91_INK, bold=True)
    return {
        "slide": "problem_structure",
        "art_direction": "large four-stage causal chain with strong root-cause emphasis and proposal scope",
        "takeaway": "小さなIssue Treeから、症状・原因・根本原因・影響が読める大判構造図へ変更",
    }


def render_v91_gate_slide_to_pptx(
    case: CaseContext,
    slide_kind: str,
    *,
    proposal_date: str = "2026年8月2日",
    proposer_name: str = "Ready Crew",
) -> tuple[bytes, dict[str, Any]]:
    """Render one art-directed Version 9.1 gate slide without adding a new engine."""

    prs = Presentation()
    ds = CONSULTING_DESIGN_SYSTEM
    prs.slide_width = _in(ds.grid.slide_width)
    prs.slide_height = _in(ds.grid.slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if slide_kind == "cover":
        report = _render_v91_cover(slide, case, proposal_date, proposer_name)
    elif slide_kind == "executive_summary":
        report = _render_v91_executive_summary(slide, case)
    elif slide_kind == "problem_structure":
        report = _render_v91_problem_structure(slide, case)
    else:
        raise ValueError(f"Unsupported Version 9.1 gate slide: {slide_kind}")
    output = BytesIO()
    prs.save(output)
    report.update(
        {
            "provider": "presentation_design_ai_v9_1_gate",
            "slide_count": 1,
            "case_id": case.case_id,
            "client_name": case.client_name,
        }
    )
    return output.getvalue(), report


# Version 9.2 three-slide final gate renderer

V92_NAVY = "#071426"
V92_NAVY_2 = "#0B1F36"
V92_INK = "#111827"
V92_MUTED = "#475467"
V92_LINE = "#64748B"
V92_SOFT = "#F3F6FA"
V92_PANEL = "#EAF2FF"
V92_BLUE = "#2563EB"
V92_CYAN = "#00A6C8"


def _case_item(items: tuple[str, ...], index: int, fallback: str) -> str:
    return items[index] if index < len(items) and items[index] else fallback


def _short_copy(text: str, limit: int = 28) -> str:
    text = " ".join(str(text).split())
    return text if len(text) <= limit else text[: limit - 1] + "…"


def _category_label(case: CaseContext) -> str:
    text = f"{case.category} {case.case_name} {case.project_summary}".lower()
    if "ocr" in text or "画像認識" in text:
        return "画像認識AI"
    if "ec" in text:
        return "EC改善"
    if "dx" in text:
        return "業務DX"
    if "recruit" in text or "採用" in text:
        return "採用改善"
    if "sales" in text or "営業" in text:
        return "営業改善"
    if "marketing" in text or "マーケ" in text:
        return "顧客獲得"
    if "ai" in text:
        return "AI活用"
    return _short_copy(case.category.replace("_", " ").strip() or "業務改善", 12)


def _case_value_message(case: CaseContext) -> str:
    first = _case_item(case.expected_outcomes, 0, "成果を測定できる状態")
    return f"{_category_label(case)}で{_short_copy(first, 16)}へつなげる"


def _case_decision_theme(case: CaseContext) -> str:
    if "PoC" in case.budget or "PoC" in case.timeline or "PoC" in case.project_summary:
        return "PoC条件を合意"
    return "次フェーズを合意"


def _case_solution_actor(case: CaseContext) -> str:
    text = f"{case.case_name} {case.category} {case.project_summary}"
    if "AI" in text or "OCR" in text:
        return "AI"
    if "CRM" in text or "SFA" in text or "ポータル" in text or "予約" in text:
        return "システム"
    if "サイト" in text or "EC" in text or "LP" in text:
        return "Web導線"
    return "仕組み"


def _case_validation_mode(case: CaseContext) -> str:
    text = f"{case.budget} {case.timeline} {case.project_summary}"
    return "PoC" if "PoC" in text or "検証" in text else "初期検証"


def _add_v92_step(slide, number: str, title: str, x: float, y: float, w: float, h: float, *, fill: str, line: str = V92_BLUE):
    _add_v91_box(slide, x, y, w, h, fill, line, line_width=2.0)
    _add_v91_box(slide, x + 0.18, y + 0.22, 0.44, 0.44, line, "", radius=MSO_SHAPE.OVAL)
    _add_v91_text(slide, number, x + 0.30, y + 0.34, 0.20, 0.08, size=14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, title, x + 0.78, y + 0.26, w - 1.02, 0.24, size=16, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)


def _add_v92_takeaway(slide, text: str, x: float = 0.68, y: float = 6.45, w: float = 11.95, *, fill: str = V92_NAVY):
    _add_v91_box(slide, x, y, w, 0.56, fill, "")
    _add_v91_text(slide, text, x + 0.34, y + 0.17, w - 0.68, 0.16, size=18, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def _render_v92_cover(slide, case: CaseContext, proposal_date: str, proposer_name: str) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(V92_NAVY)
    _add_v91_box(slide, 8.25, 0.0, 5.08, 7.5, V92_NAVY_2, "")
    _add_v91_text(slide, "Confidential", 0.72, 0.54, 1.55, 0.18, size=14, color="#D6E4F2", bold=True)
    _add_v91_text(slide, proposal_date, 2.52, 0.54, 1.72, 0.18, size=14, color="#D6E4F2")
    _add_v91_text(slide, _short_copy(case.case_name, 18), 0.70, 1.42, 7.2, 0.82, size=46, color="#FFFFFF", bold=True)
    _add_v91_text(slide, case.client_name, 0.74, 2.42, 5.8, 0.25, size=18, color="#BFD3EA", bold=True)
    _add_v91_text(slide, _short_copy(case.project_summary, 34), 0.74, 3.08, 6.6, 0.42, size=23, color="#E6F0FF", bold=True)
    _add_v91_box(slide, 0.74, 4.08, 6.55, 0.74, "#0D2844", "#2E6EA3", line_width=1.6)
    _add_v91_text(slide, "価値メッセージ", 0.98, 4.34, 1.6, 0.13, size=14, color=V92_CYAN, bold=True)
    _add_v91_text(slide, _short_copy(_case_value_message(case), 18), 2.56, 4.28, 3.6, 0.18, size=18, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, proposer_name, 0.76, 6.55, 2.6, 0.2, size=16, color="#FFFFFF", bold=True)
    _add_v91_text(slide, "ProposalPilot", 0.76, 6.86, 2.2, 0.18, size=14, color="#BFD3EA")

    _add_v91_text(slide, "提案の進め方", 9.04, 0.82, 2.6, 0.18, size=16, color="#BFD3EA", bold=True, align=PP_ALIGN.CENTER)
    flow = [("1", "現状整理", 9.02, 1.28), ("2", "解決策", 9.46, 2.62), ("3", "効果確認", 9.02, 3.96), ("4", "次回合意", 9.46, 5.30)]
    for i, (number, title, x, y) in enumerate(flow):
        _add_v92_step(slide, number, title, x, y, 2.55, 0.82, fill="#FFFFFF")
        if i < len(flow) - 1:
            _add_v91_connector(slide, x + 1.28, y + 0.82, flow[i + 1][2] + 1.28, flow[i + 1][3], V92_CYAN, 3.0)
    _add_v91_box(slide, 8.70, 6.55, 3.64, 0.45, "#123657", "#2E6EA3", line_width=1.4)
    _add_v91_text(slide, _case_decision_theme(case), 9.08, 6.69, 2.86, 0.13, size=16, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    return {"slide": "cover", "art_direction": "dark executive cover with a four-step image recognition business flow", "takeaway": "装飾図を処理フローへ変え、提案価値を表紙だけで理解できるようにした", "speaker_notes_moved_count": 1}


def _render_v92_executive_summary(slide, case: CaseContext) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("#FFFFFF")
    _add_v91_text(slide, "提案要約", 0.70, 0.36, 2.5, 0.18, size=16, color=V92_CYAN, bold=True)
    _add_v91_text(slide, _short_copy(_case_value_message(case), 34), 0.70, 0.78, 11.1, 0.34, size=30, color=V92_INK, bold=True)
    _add_v91_box(slide, 0.76, 1.82, 3.42, 2.18, V92_NAVY, "")
    _add_v91_text(slide, "提案結論", 1.04, 2.08, 1.2, 0.16, size=16, color=V92_CYAN, bold=True)
    _add_v91_text(slide, f"{_short_copy(_category_label(case), 9)}\n{_short_copy(_case_item(case.expected_outcomes, 0, '成果創出'), 9)}", 1.04, 2.56, 2.74, 0.48, size=24, color="#FFFFFF", bold=True)
    _add_v91_box(slide, 1.02, 3.48, 2.72, 0.34, V92_CYAN, "")
    _add_v91_text(slide, _case_decision_theme(case), 1.20, 3.55, 2.36, 0.18, size=14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    x_positions = [4.75, 7.05, 9.35]
    titles = [
        ("現状", _short_copy(_case_item(case.pain_points, 0, "業務課題が残る"), 9)),
        ("提案", _short_copy(_category_label(case), 9)),
        ("成果", _short_copy(_case_item(case.expected_outcomes, 0, "成果を測定"), 9)),
    ]
    for i, (head, body) in enumerate(titles):
        x = x_positions[i]
        fill = ["#F8FAFC", V92_PANEL, "#F8FAFC"][i]
        line = [V92_LINE, V92_BLUE, V92_LINE][i]
        _add_v91_box(slide, x, 1.86, 1.92, 2.28, fill, line, line_width=2.2)
        _add_v91_text(slide, head, x + 0.22, 2.18, 1.44, 0.18, size=18, color=line if i == 1 else V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.20, 2.84, 1.46, 0.44, size=20, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_box(slide, x + 0.66, 3.46, 0.58, 0.58, line, "", radius=MSO_SHAPE.OVAL)
    _add_v91_connector(slide, 6.68, 2.92, 7.04, 2.92, V92_BLUE, 3.0)
    _add_v91_connector(slide, 8.98, 2.92, 9.34, 2.92, V92_BLUE, 3.0)

    effects = [
        ("効果1", _short_copy(_case_item(case.expected_outcomes, 0, "成果を測る"), 11)),
        ("効果2", _short_copy(_case_item(case.expected_outcomes, 1, "負荷を下げる"), 11)),
        ("効果3", _short_copy(_case_item(case.expected_outcomes, 2, "判断を早める"), 11)),
    ]
    for i, (head, body) in enumerate(effects):
        x = 0.96 + i * 3.72
        _add_v91_box(slide, x, 4.72, 3.02, 0.72, "#F8FAFC", V91_SOFT_LINE, line_width=1.3)
        _add_v91_text(slide, head, x + 0.22, 4.94, 0.7, 0.16, size=16, color=V92_BLUE, bold=True)
        _add_v91_text(slide, body, x + 0.92, 4.92, 1.94, 0.16, size=15, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v92_takeaway(slide, f"{_case_decision_theme(case)}し、{_short_copy(case.industry, 8)}の成果を測れる形で次へ進みます。")
    return {"slide": "executive_summary", "art_direction": "reduced executive decision page with one conclusion, one transformation diagram, and three effects", "takeaway": "情報量を削り、経営者が現状・提案・価値・合意事項を即時理解できる構成にした", "speaker_notes_moved_count": 3}


def _render_v92_problem_structure(slide, case: CaseContext) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(V92_SOFT)
    _add_v91_text(slide, "現状と課題", 0.70, 0.36, 1.8, 0.18, size=16, color=V92_CYAN, bold=True)
    _add_v91_text(slide, _short_copy(f"{_case_item(case.pain_points, 0, '主要課題')}が、成果創出のボトルネックです", 36), 0.70, 0.78, 11.2, 0.34, size=26, color=V92_INK, bold=True)
    _add_v91_box(slide, 0.82, 1.58, 2.28, 2.96, "#FFFFFF", V91_SOFT_LINE, line_width=1.4)
    _add_v91_text(slide, "症状", 1.08, 1.90, 0.8, 0.16, size=16, color=V92_BLUE, bold=True)
    _add_v91_text(slide, _short_copy(_case_item(case.pain_points, 0, "課題が残る"), 9), 1.04, 2.42, 1.72, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, _short_copy(_case_item(case.pain_points, 1, "確認が遅い"), 9), 1.02, 3.26, 1.78, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 3.82, 1.86, 2.12, 2.42, "#FFFFFF", V92_BLUE, line_width=2.0)
    _add_v91_text(slide, "業務上の原因", 4.08, 2.18, 1.58, 0.16, size=16, color=V92_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "情報が分断", 4.10, 2.84, 1.5, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "判断が遅い", 4.22, 3.50, 1.28, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 6.62, 1.46, 2.34, 3.22, V92_NAVY, "")
    _add_v91_text(slide, "根本原因", 7.02, 1.88, 1.5, 0.18, size=17, color=V92_CYAN, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "改善材料\n未活用", 6.92, 2.70, 1.74, 0.54, size=27, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 9.62, 1.58, 2.60, 2.96, "#FFFFFF", V91_SOFT_LINE, line_width=1.4)
    _add_v91_text(slide, "事業影響", 10.20, 1.90, 1.3, 0.16, size=16, color=V92_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, _short_copy(_case_item(case.expected_outcomes, 0, "成果が出にくい"), 9), 10.02, 2.42, 1.72, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, _short_copy(_case_item(case.expected_outcomes, 1, "負荷が残る"), 9), 9.94, 3.26, 1.90, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_connector(slide, 3.12, 3.0, 3.80, 3.0, V92_BLUE, 3.2)
    _add_v91_connector(slide, 5.96, 3.0, 6.60, 3.0, V92_BLUE, 3.2)
    _add_v91_connector(slide, 8.98, 3.0, 9.60, 3.0, V92_BLUE, 3.2)
    _add_v91_box(slide, 3.34, 5.24, 6.50, 0.54, "#FFFFFF", V92_BLUE, line_width=1.8)
    _add_v91_text(slide, f"変える範囲: 現状整理 → {_short_copy(_category_label(case), 10)} → 効果測定 → 次判断", 3.56, 5.42, 6.08, 0.15, size=15, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v92_takeaway(slide, f"焦点は、{_short_copy(case.industry, 8)}の課題を測定できる改善テーマへ変えることです。", y=6.38)
    return {"slide": "problem_structure", "art_direction": "causal chain centered on one dominant root cause and connected business impact", "takeaway": "4列カードから因果の流れが見える構造へ変え、根本原因を主役にした", "speaker_notes_moved_count": 2}


def render_v92_gate_slide_to_pptx(case: CaseContext, slide_kind: str, *, proposal_date: str = "2026年8月2日", proposer_name: str = "Ready Crew") -> tuple[bytes, dict[str, Any]]:
    """Render one Version 9.2 final gate slide using the existing renderer layer."""
    prs = Presentation()
    ds = CONSULTING_DESIGN_SYSTEM
    prs.slide_width = _in(ds.grid.slide_width)
    prs.slide_height = _in(ds.grid.slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if slide_kind == "cover":
        report = _render_v92_cover(slide, case, proposal_date, proposer_name)
    elif slide_kind == "executive_summary":
        report = _render_v92_executive_summary(slide, case)
    elif slide_kind == "problem_structure":
        report = _render_v92_problem_structure(slide, case)
    else:
        raise ValueError(f"Unsupported Version 9.2 gate slide: {slide_kind}")
    output = BytesIO()
    prs.save(output)
    report.update({"provider": "presentation_design_ai_v9_2_gate", "slide_count": 1, "case_id": case.case_id, "client_name": case.client_name})
    return output.getvalue(), report


V101_PALE = "#F5F7FB"
V101_BLUE = "#1D4ED8"
V101_CYAN = "#00A6C8"
V101_GREEN = "#16A34A"
V101_AMBER = "#D97706"
V101_RED = "#DC2626"


def _add_v101_header(slide, case: CaseContext, label: str, title: str, slide_no: int, *, background: str = "#FFFFFF") -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(background)
    _add_v91_text(slide, label, 0.70, 0.34, 2.4, 0.20, size=14, color=V92_CYAN, bold=True)
    _add_v91_text(slide, title, 0.70, 0.72, 11.0, 0.48, size=26, color=V92_INK, bold=True)
    _add_v91_text(slide, case.client_name, 9.55, 0.38, 2.62, 0.18, size=14, color=V92_MUTED, align=PP_ALIGN.RIGHT)
    _add_v91_text(slide, f"{slide_no:02d}", 12.34, 0.38, 0.36, 0.18, size=14, color=V92_MUTED, align=PP_ALIGN.RIGHT)


def _case_slide_title(case: CaseContext, role: str, fallback: str) -> str:
    pain = _case_item(case.pain_points, 0, "主要課題")
    outcome = _case_item(case.expected_outcomes, 0, "成果を測定")
    unit = _case_implementation_unit(case)
    titles = {
        "root_cause": f"原因は、{_short_copy(pain, 12)}を改善材料に変えられていないことです",
        "poc_scope": f"{unit}に絞り、評価条件を先に合意します",
        "solution": f"{_short_copy(_category_label(case), 12)}で現場判断を支援し、重要業務に集中します",
        "process": f"{_short_copy(_category_label(case), 12)}を既存業務フローへ無理なく組み込みます",
        "kpi": f"評価指標は{_short_copy(outcome, 12)}を含む3点に絞ります",
        "roadmap": f"{_short_copy(case.timeline, 12)}で判断材料を揃えます",
        "roi": "投資判断は費用ではなく実測効果で確認します",
        "risk": "リスクは開始前に確認事項へ分解します",
        "governance": f"体制は{_short_copy(case.decision_maker, 14)}を中心に分担します",
        "decision": f"次回は{unit}・評価指標・開始条件を合意します",
        "appendix_evidence": "補足資料: 評価条件の詳細",
        "appendix_faq": "補足資料: リスク確認とFAQ",
    }
    return _short_copy(titles.get(role, fallback), 40)


def _case_implementation_unit(case: CaseContext) -> str:
    if "サイト" in case.case_name or "Web" in case.category or "EC" in case.category:
        return "主要導線"
    if "AI" in case.category or "OCR" in case.category:
        return "対象業務"
    if "DX" in case.category or "Workflow" in case.category:
        return "一部業務"
    if "営業" in case.case_name or "Sales" in case.category:
        return "重点案件"
    return "優先範囲"


def _add_v101_takeaway(slide, text: str, *, y: float = 6.42, fill: str = V92_NAVY) -> None:
    _add_v91_box(slide, 0.72, y, 11.88, 0.58, fill, "")
    _add_v91_text(slide, text, 1.04, y + 0.17, 11.18, 0.18, size=17, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def _add_v101_node(slide, title: str, body: str, x: float, y: float, w: float, h: float, *, fill: str = "#FFFFFF", line: str = V91_SOFT_LINE, accent: str = V101_BLUE) -> None:
    _add_v91_box(slide, x, y, w, h, fill, line, line_width=1.6)
    _add_v91_text(slide, title, x + 0.22, y + 0.20, w - 0.44, 0.20, size=16, color=accent, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, body, x + 0.28, y + 0.70, w - 0.56, 0.36, size=15, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)


def _render_v101_root_cause(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "根本原因", _case_slide_title(case, "root_cause", page.action_title), page.slide_no, background=V101_PALE)
    center_x, center_y = 5.26, 2.18
    _add_v91_box(slide, center_x, center_y, 2.78, 1.42, V92_NAVY, "")
    _add_v91_text(slide, "業務資産化\nされていない", center_x + 0.24, center_y + 0.34, 2.30, 0.42, size=24, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    causes = [
        ("基準", "経験則に閉じる", 1.00, 1.58, V101_BLUE),
        ("記録", "改善材料が分散", 9.52, 1.58, V101_CYAN),
        ("教育", "引継ぎ負荷が残る", 1.00, 4.02, V101_AMBER),
        ("例外", "判断が属人化", 9.52, 4.02, V101_RED),
    ]
    for title, body, x, y, accent in causes:
        _add_v101_node(slide, title, body, x, y, 2.56, 1.28, fill="#FFFFFF", accent=accent)
        _add_v91_connector(slide, x + (2.56 if x < 5 else 0), y + 0.64, center_x + (0 if x < 5 else 2.78), center_y + 0.72, accent, 2.4)
    _add_v101_takeaway(slide, f"{_case_validation_mode(case)}では、課題を改善材料へ変えられるかを確認します。")
    return {"slide": "root_cause", "art_direction": "central root cause map with four surrounding causes", "takeaway": "原因を責任論ではなくPoC仮説へ接続した"}


def _render_v101_poc_scope(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "PoC仮説", _case_slide_title(case, "poc_scope", page.action_title), page.slide_no)
    quadrants = [
        ("対象範囲", f"{_case_implementation_unit(case)}\n一次対応まで", 0.92, 1.72, "#EEF6FF", V101_BLUE),
        ("対象外", "全範囲展開\n詳細連携の実装", 6.82, 1.72, "#F8FAFC", V92_LINE),
        ("評価内容", f"{_short_copy(_case_item(case.expected_outcomes, 0, '成果'), 8)}\n3指標で判断", 0.92, 4.06, "#ECFDF5", V101_GREEN),
        ("確認事項", "対象範囲・担当\n開始条件", 6.82, 4.06, "#FFF7ED", V101_AMBER),
    ]
    for title, body, x, y, fill, accent in quadrants:
        _add_v101_node(slide, title, body, x, y, 5.20, 1.50, fill=fill, line=V91_SOFT_LINE, accent=accent)
    _add_v91_connector(slide, 6.48, 1.74, 6.48, 5.56, V91_SOFT_LINE, 1.2)
    _add_v91_connector(slide, 0.92, 3.72, 12.02, 3.72, V91_SOFT_LINE, 1.2)
    _add_v101_takeaway(slide, f"今回決める範囲と決めない範囲を分け、{_case_validation_mode(case)}を始めやすくします。")
    return {"slide": "poc_scope", "art_direction": "scope boundary matrix separating in-scope, out-of-scope, evaluation, and confirmation", "takeaway": "PoC合意に必要な境界を1枚で明確化した"}


def _render_v101_solution(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "提案全体像", _case_slide_title(case, "solution", page.action_title), page.slide_no, background=V101_PALE)
    actor = _case_solution_actor(case)
    _add_v91_box(slide, 0.86, 1.72, 3.12, 3.70, "#FFFFFF", V91_SOFT_LINE, line_width=1.6)
    _add_v91_text(slide, f"{actor}が担うこと", 1.16, 2.04, 2.52, 0.22, size=20, color=V101_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, f"{_short_copy(_category_label(case), 8)}で\n一次対応を支援", 1.22, 2.80, 2.40, 0.48, size=22, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 5.04, 2.20, 3.22, 2.76, V92_NAVY, "")
    _add_v91_text(slide, "判断ログ", 5.58, 2.72, 2.12, 0.24, size=22, color=V92_CYAN, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "対応結果と修正を\n次の改善材料にする", 5.38, 3.42, 2.52, 0.50, size=20, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_box(slide, 9.22, 1.72, 3.12, 3.70, "#FFFFFF", V91_SOFT_LINE, line_width=1.6)
    _add_v91_text(slide, "人が担うこと", 9.52, 2.04, 2.52, 0.22, size=20, color=V101_GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "例外確認と\n最終判断に集中", 9.58, 2.80, 2.40, 0.48, size=22, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v91_connector(slide, 3.98, 3.58, 5.04, 3.58, V101_BLUE, 3.2)
    _add_v91_connector(slide, 8.26, 3.58, 9.22, 3.58, V101_GREEN, 3.2)
    _add_v101_takeaway(slide, f"{actor}は置き換えではなく、現場判断を支援する仕組みとして導入します。")
    return {"slide": "solution", "art_direction": "human-ai role split with decision log as the central operating asset", "takeaway": "AI・人・ログの役割分担を大きな構図へ変更した"}


def _render_v101_process(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "検証フロー", _case_slide_title(case, "process", page.action_title), page.slide_no)
    steps = [
        ("情報取得", "対象情報を集める"),
        (_case_solution_actor(case), "一次対応を支援"),
        ("人が確認", "例外と修正を記録"),
        ("ログ蓄積", "改善材料を残す"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.92 + i * 3.08
        y = 2.34 + (0.36 if i % 2 else 0)
        _add_v91_box(slide, x, y, 2.30, 1.26, "#FFFFFF", V91_SOFT_LINE, line_width=1.5)
        _add_v91_box(slide, x + 0.28, y + 0.34, 0.52, 0.52, V101_CYAN, "", radius=MSO_SHAPE.OVAL)
        _add_v91_text(slide, str(i + 1), x + 0.43, y + 0.48, 0.22, 0.10, size=16, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, title, x + 0.92, y + 0.28, 1.10, 0.20, size=17, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.36, y + 0.82, 1.60, 0.18, size=14, color=V92_MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            _add_v91_connector(slide, x + 2.30, y + 0.66, x + 3.02, 2.98 + (0.36 if (i + 1) % 2 else 0), V101_BLUE, 2.4)
    _add_v91_box(slide, 1.30, 5.18, 10.50, 0.56, "#EEF6FF", "")
    _add_v91_text(slide, f"確認ポイント: 詳細連携は{_case_validation_mode(case)}後の判断で扱います。", 1.56, 5.36, 9.86, 0.15, size=16, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v101_takeaway(slide, f"{_case_validation_mode(case)}中は、成果だけでなく運用に残せる記録を確認します。")
    return {"slide": "process", "art_direction": "large four-step operating flow with confirmation boundary", "takeaway": "小さなフローを営業説明しやすい4ステップへ整理した"}


def _render_v101_kpi(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "KPI", _case_slide_title(case, "kpi", page.action_title), page.slide_no, background=V101_PALE)
    metrics = [
        (_short_copy(_case_item(case.expected_outcomes, 0, "成果指標"), 7), "基準値", "現状値と比較"),
        (_short_copy(_case_item(case.expected_outcomes, 1, "改善幅"), 7), "目標値", "合意後に決定"),
        (_short_copy(_case_item(case.expected_outcomes, 2, "負荷削減"), 7), "測定", "実行前後で比較"),
        ("合格条件", "判定", f"{_case_validation_mode(case)}後に合意"),
    ]
    for i, (metric, tag, body) in enumerate(metrics):
        x = 0.82 + i * 3.06
        _add_v91_box(slide, x, 1.76, 2.42, 3.82, "#FFFFFF", V91_SOFT_LINE, line_width=1.4)
        _add_v91_text(slide, metric, x + 0.18, 2.12, 2.02, 0.24, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_box(slide, x + 0.54, 2.80, 1.34, 1.34, "#EEF6FF", V101_BLUE, line_width=1.8, radius=MSO_SHAPE.OVAL)
        _add_v91_text(slide, tag, x + 0.72, 3.28, 0.98, 0.14, size=15, color=V101_BLUE, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.30, 4.62, 1.82, 0.36, size=15, color=V92_MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_v101_takeaway(slide, "未確定値は置かず、何をどう測って合格判断するかを先に合意します。")
    return {"slide": "kpi", "art_direction": "SMART-style KPI cards without invented numeric targets", "takeaway": "30%などの仮値羅列を避け、測定設計として再構成した"}


def _render_v101_roadmap(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "スケジュール", _case_slide_title(case, "roadmap", page.action_title), page.slide_no)
    _add_v91_connector(slide, 1.18, 3.50, 11.82, 3.50, V101_BLUE, 3.0)
    phases = [
        ("Phase 1", "準備", "対象範囲を決める"),
        ("Phase 2", "検証", f"{_case_solution_actor(case)}と人の確認を回す"),
        ("Week 5-6", "評価", "KPIと運用負荷を見る"),
        ("Week 7-8", "判断", "本番投資の可否を決める"),
    ]
    for i, (week, title, body) in enumerate(phases):
        x = 0.90 + i * 3.04
        y = 2.36 if i % 2 == 0 else 3.88
        _add_v91_box(slide, x, y, 2.36, 1.02, "#FFFFFF", V91_SOFT_LINE, line_width=1.5)
        _add_v91_text(slide, week, x + 0.18, y + 0.20, 0.84, 0.16, size=14, color=V101_CYAN, bold=True)
        _add_v91_text(slide, title, x + 1.04, y + 0.18, 0.94, 0.18, size=18, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.22, y + 0.62, 1.88, 0.18, size=14, color=V92_MUTED, align=PP_ALIGN.CENTER)
        _add_v91_connector(slide, x + 1.18, 3.50, x + 1.18, y + (1.02 if y < 3.5 else 0), V101_CYAN, 1.8)
    _add_v101_takeaway(slide, f"{_short_copy(case.timeline, 16)}で判断に必要な材料だけを揃えます。")
    return {"slide": "roadmap", "art_direction": "eight-week decision roadmap with alternating phase cards", "takeaway": "単純な横棒ではなく、判断材料が揃う流れにした"}


def _render_v101_roi(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "ROI・投資判断", _case_slide_title(case, "roi", page.action_title), page.slide_no, background=V101_PALE)
    lanes = [
        ("投資", f"{_case_validation_mode(case)}・準備\n運用設計・教育", V101_BLUE),
        ("効果", f"{_short_copy(_case_item(case.expected_outcomes, 0, '効果'), 10)}\n負荷削減・品質安定", V101_GREEN),
        ("判断", f"{_case_validation_mode(case)}実測値で\n次投資可否を判断", V92_NAVY),
    ]
    for i, (title, body, color) in enumerate(lanes):
        x = 0.94 + i * 4.08
        _add_v91_box(slide, x, 2.08, 3.18, 2.34, "#FFFFFF", color, line_width=2.0)
        _add_v91_box(slide, x + 1.14, 1.50, 0.90, 0.90, color, "", radius=MSO_SHAPE.OVAL)
        _add_v91_text(slide, str(i + 1), x + 1.40, 1.78, 0.34, 0.12, size=18, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, title, x + 0.42, 2.68, 2.34, 0.22, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.34, 3.38, 2.50, 0.42, size=16, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            _add_v91_connector(slide, x + 3.18, 3.26, x + 3.96, 3.26, color, 3.0)
    _add_v101_takeaway(slide, f"金額は作らず、{_case_validation_mode(case)}実測値で投資判断できる形にします。")
    return {"slide": "roi", "art_direction": "investment-effect-decision path without fabricated money or payback", "takeaway": "ROIを未確定数値ではなく投資判断プロセスとして示した"}


def _render_v101_risk(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "リスク", _case_slide_title(case, "risk", page.action_title), page.slide_no)
    risks = [
        ("入力準備", "対象情報と件数を先に確認", "#ECFDF5", V101_GREEN),
        ("例外処理", "人が見る条件を決める", "#FEF3C7", V101_AMBER),
        ("現場負荷", f"{_case_validation_mode(case)}時の担当を明確化", "#FFF7ED", V101_AMBER),
        ("連携", "本番判断後に詳細化", "#FEE2E2", V101_RED),
    ]
    for i, (title, body, fill, accent) in enumerate(risks):
        x = 1.02 + (i % 2) * 5.72
        y = 1.68 + (i // 2) * 2.02
        _add_v101_node(slide, title, body, x, y, 4.72, 1.30, fill=fill, line=accent, accent=accent)
    _add_v91_box(slide, 3.20, 5.64, 6.90, 0.48, "#FFFFFF", V101_BLUE, line_width=1.5)
    _add_v91_text(slide, f"確認方法と対策をセットで合意すれば、{_case_validation_mode(case)}開始前の不安を減らせます。", 3.48, 5.80, 6.34, 0.14, size=15, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
    _add_v101_takeaway(slide, f"不安点は隠さず、{_case_validation_mode(case)}前に管理可能な確認事項へ変えます。")
    return {"slide": "risk", "art_direction": "risk response matrix with confirmation-oriented wording", "takeaway": "リスクを懸念ではなく事前確認事項として整理した"}


def _render_v101_governance(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    _add_v101_header(slide, case, "実施体制", _case_slide_title(case, "governance", page.action_title), page.slide_no, background=V101_PALE)
    roles = [
        ("顧客側", f"{_case_implementation_unit(case)}\n確認・判断"),
        ("管理側", "データ管理・連携方針\nリスク確認"),
        ("Ready Crew", f"{_short_copy(_category_label(case), 8)}設計\n分析・改善提案"),
    ]
    for i, (title, body) in enumerate(roles):
        x = 0.94 + i * 4.08
        _add_v91_box(slide, x, 1.86, 3.18, 3.12, "#FFFFFF", V91_SOFT_LINE, line_width=1.5)
        _add_v91_box(slide, x + 1.10, 2.24, 0.96, 0.96, [V101_GREEN, V101_BLUE, V92_NAVY][i], "", radius=MSO_SHAPE.OVAL)
        _add_v91_text(slide, title, x + 0.42, 3.56, 2.34, 0.22, size=20, color=V92_INK, bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, body, x + 0.36, 4.12, 2.46, 0.42, size=15, color=V92_MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_v101_takeaway(slide, "誰が何を確認するかを分け、次回合意後すぐ動ける体制にします。")
    return {"slide": "governance", "art_direction": "three-party responsibility map", "takeaway": "体制を抽象図から会議後に動ける責任分担へ変更した"}


def _render_v101_decision(slide, case: CaseContext, page: PageSpec) -> dict[str, Any]:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(V92_NAVY)
    _add_v91_text(slide, "次回合意事項", 0.78, 0.46, 2.2, 0.22, size=16, color=V92_CYAN, bold=True)
    _add_v91_text(slide, _case_slide_title(case, "decision", page.action_title), 0.78, 0.88, 10.8, 0.50, size=30, color="#FFFFFF", bold=True)
    items = [
        (_case_implementation_unit(case), "対象・範囲"),
        ("評価指標", _short_copy(_case_item(case.expected_outcomes, 0, "成果指標"), 11)),
        ("使用データ", "入力情報・件数"),
        ("実施体制", _short_copy(case.decision_maker, 13)),
        ("開始条件", "日程・担当・準備物"),
    ]
    for i, (title, body) in enumerate(items):
        x = 0.94 + (i % 3) * 4.02
        y = 2.02 + (i // 3) * 1.64
        _add_v91_box(slide, x, y, 3.20, 1.08, "#FFFFFF", "#2E6EA3", line_width=1.7)
        _add_v91_box(slide, x + 0.24, y + 0.30, 0.48, 0.48, V92_CYAN, "", radius=MSO_SHAPE.OVAL)
        _add_v91_text(slide, str(i + 1), x + 0.38, y + 0.43, 0.20, 0.10, size=14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _add_v91_text(slide, title, x + 0.86, y + 0.26, 1.86, 0.18, size=16, color=V92_INK, bold=True)
        _add_v91_text(slide, body, x + 0.86, y + 0.66, 1.92, 0.16, size=14, color=V92_MUTED)
    _add_v91_box(slide, 7.00, 5.36, 4.68, 0.62, V92_CYAN, "")
    _add_v91_text(slide, f"合意後: {_case_validation_mode(case)}準備を開始", 7.36, 5.56, 3.92, 0.18, size=18, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_v91_text(slide, "次回日程と確認担当を会議内で決めてください。", 0.96, 6.48, 7.0, 0.20, size=16, color="#D6E4F2", bold=True)
    return {"slide": "decision", "art_direction": "dark decision page with five concrete agreement items", "takeaway": "曖昧なお願いを、誰が何を決めるか分かる合意ページへ変えた"}


def _render_v101_appendix(slide, case: CaseContext, page: PageSpec, *, faq: bool = False) -> dict[str, Any]:
    label = "補足資料"
    _add_v101_header(slide, case, label, _case_slide_title(case, "appendix_faq" if faq else "appendix_evidence", page.action_title), page.slide_no, background="#FFFFFF")
    if faq:
        rows = [
            ("セキュリティ", "入力情報の扱いと保存範囲を確認"),
            ("既存システム連携", f"{_case_validation_mode(case)}では深い連携を前提にしない"),
            ("例外処理", "人が確認する条件を事前定義"),
            ("現場教育", f"運用手順は{_case_validation_mode(case)}中に確認"),
        ]
        takeaway = "詳細論点は本編を止めず、質問時に参照します。"
    else:
        rows = [
            ("測定方法", "現状値と改善後を比較"),
            ("サンプル条件", "対象範囲・件数を確認"),
            ("判定基準", f"合格条件は{_case_validation_mode(case)}開始前に合意"),
            ("担当", "現場確認者とIT確認者を明確化"),
        ]
        takeaway = "評価条件の詳細は、次回合意時の確認リストとして使います。"
    for i, (title, body) in enumerate(rows):
        x = 1.02 + (i % 2) * 5.62
        y = 1.76 + (i // 2) * 1.82
        _add_v101_node(slide, title, body, x, y, 4.62, 1.16, fill=V101_PALE, line=V91_SOFT_LINE, accent=V101_BLUE)
    _add_v101_takeaway(slide, takeaway)
    return {"slide": "appendix_faq" if faq else "appendix_evidence", "art_direction": "customer-facing appendix reference grid", "takeaway": takeaway}


def _render_v101_director_slide(slide, plan: PresentationPlan, page: PageSpec, role: str) -> dict[str, Any]:
    if page.slide_no == 1 and role == "cover":
        return _render_v92_cover(slide, plan.case, proposal_date="2026年8月2日", proposer_name="Ready Crew")
    if page.slide_no == 2 and role == "executive_summary":
        return _render_v92_executive_summary(slide, plan.case)
    if page.slide_no == 3 and role == "problem":
        return _render_v92_problem_structure(slide, plan.case)
    if role == "root_cause":
        return _render_v101_root_cause(slide, plan.case, page)
    if role == "poc_scope":
        return _render_v101_poc_scope(slide, plan.case, page)
    if role == "solution":
        return _render_v101_solution(slide, plan.case, page)
    if role == "process":
        return _render_v101_process(slide, plan.case, page)
    if role == "kpi":
        return _render_v101_kpi(slide, plan.case, page)
    if role == "roadmap":
        return _render_v101_roadmap(slide, plan.case, page)
    if role == "roi":
        return _render_v101_roi(slide, plan.case, page)
    if role == "risk":
        return _render_v101_risk(slide, plan.case, page)
    if role == "governance":
        return _render_v101_governance(slide, plan.case, page)
    if role == "decision":
        return _render_v101_decision(slide, plan.case, page)
    if role == "appendix_evidence":
        return _render_v101_appendix(slide, plan.case, page)
    if role == "appendix_faq":
        return _render_v101_appendix(slide, plan.case, page, faq=True)
    palette = CONSULTING_DESIGN_SYSTEM.palette_for_category(plan.case.category)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb("#FFFFFF")
    _add_slide_chrome(slide, plan, page, palette, dark=False)
    _render_generic(slide, page, palette)
    return {"slide": role, "art_direction": "fallback existing consulting renderer", "takeaway": page.action_title}


def _write_speaker_notes(slide, page: PageSpec) -> bool:
    notes = page.speaker_notes or {}
    if not notes:
        return False
    labels = {
        "slide_conclusion": "このページの結論",
        "talking_points": "話す順番",
        "confirmation_question": "顧客へ確認する質問",
        "transition": "次ページへのつなぎ",
        "caution": "注意事項",
    }
    lines = [f"{label}: {notes[key]}" for key, label in labels.items() if notes.get(key)]
    if not lines:
        return False
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)
    return True


def render_director_plan_to_pptx(plan: PresentationPlan, director_plan: Any) -> tuple[bytes, dict[str, Any]]:
    """Render the Version 10 one-deck gate without replacing the renderer.

    The one-deck remediation keeps the Director plan intact while rendering
    every slide with the Version 9.2 art-direction baseline or equivalent
    audience-facing visuals.
    """

    ds = CONSULTING_DESIGN_SYSTEM
    palette = ds.palette_for_category(plan.case.category)
    prs = Presentation()
    prs.slide_width = _in(ds.grid.slide_width)
    prs.slide_height = _in(ds.grid.slide_height)
    blank = prs.slide_layouts[6]
    role_by_slide = {
        item.slide_no: item.slide_role
        for item in getattr(director_plan, "slide_sequence", ())
    }
    render_pages: list[dict[str, Any]] = []
    for page in plan.pages:
        slide = prs.slides.add_slide(blank)
        role = role_by_slide.get(page.slide_no, page.visual_type)
        detail = _render_v101_director_slide(slide, plan, page, role)
        notes_written = _write_speaker_notes(slide, page)
        render_pages.append(
            {
                "slide_no": page.slide_no,
                "slide_id": f"slide-{page.slide_no:02d}",
                "slide_role": role,
                "priority_level": next(
                    (item.priority_level for item in getattr(director_plan, "slide_sequence", ()) if item.slide_no == page.slide_no),
                    "support",
                ),
                "component_id": page.component_id,
                "visual_type": page.visual_type,
                "layout_family": page.layout_family,
                "speaker_notes_written": notes_written,
                "render_detail": detail,
            }
        )
    output = BytesIO()
    prs.save(output)
    return output.getvalue(), {
        "provider": "presentation_director_v10",
        "slide_count": plan.slide_count,
        "distinct_layout_count": plan.distinct_layout_count,
        "average_diagram_ratio": plan.average_diagram_ratio,
        "average_text_ratio": plan.average_text_ratio,
        "pages": render_pages,
    }
