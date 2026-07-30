from __future__ import annotations

from io import BytesIO
from typing import Any

import pytest
from pptx import Presentation

from app.beautiful_ai.presentation_mapper import map_to_beautiful_ai_payload
from app.beautiful_ai.schemas import BeautifulAiPresentationRequest
from app.models import PowerPointData, PowerPointSlide, PptxDownloadRequest
from app.services.customer_ready_judgement import release_judge_from_gate_status
from app.services.customer_ready_quality import (
    CustomerReadyBlockedError,
    run_customer_ready_quality_gate,
)
from app.services.pptx_service import build_pptx_context, build_pptx_result
from app.services.proposal_validation_engine import validate_proposal
from app.services.proposal_quality_upgrade import upgrade_slides_for_v11


def _golden_payload(
    *,
    client: str,
    deck_title: str,
    industry: str,
    category: str,
    brief: str,
    budget: str,
    schedule: str,
    solution: str,
    competitor: str,
) -> dict[str, Any]:
    return {
        "project_brief": brief,
        "client_company_info": f"{client}\n業界: {industry}\n意思決定者: 部長、現場責任者",
        "competitor_company_name": competitor,
        "desired_launch_timing": schedule,
        "budget_range": budget,
        "own_service_info": f"{category}の要件整理、設計、導入、運用定着を支援します。",
        "case_studies": "類似案件では初期検証、段階導入、運用改善を組み合わせて成果を確認しました。",
        "special_function_required": solution,
        "powerpoint_generation_data": {
            "deck_title": deck_title,
            "client_name": client,
            "slides": [
                {
                    "slide_no": 1,
                    "layout": "summary",
                    "title": "提案サマリー",
                    "bullets": [
                        brief,
                        f"解決策: {solution}",
                        f"予算: {budget}",
                        f"時期: {schedule}",
                    ],
                    "speaker_notes": "顧客の課題、提案の結論、期待効果を最初に説明します。",
                    "visual_suggestion": "3つの要点カード",
                }
            ],
        },
    }


GOLDEN_PAYLOADS = [
    _golden_payload(
        client="株式会社青葉不動産",
        deck_title="Webサイトリニューアル提案",
        industry="不動産",
        category="Web制作",
        brief="物件問い合わせの導線が弱く、更新作業も属人化しているため、Webサイトを営業活用できる状態へ改善したい。",
        budget="500万円から700万円",
        schedule="4か月以内に公開希望",
        solution="CMS刷新、問い合わせ導線改善、SEO初期設計、物件更新フロー整備",
        competitor="地域Web制作会社",
    ),
    _golden_payload(
        client="株式会社北斗物流",
        deck_title="AI業務効率化導入提案",
        industry="物流",
        category="AI/DX",
        brief="伝票確認と配車調整の確認作業に時間がかかり、繁忙期の処理遅延が発生している。",
        budget="1000万円以内",
        schedule="2027年5月の段階導入を想定",
        solution="AI候補提示、確認UI、既存システム連携、PoC評価",
        competitor="既存RPAベンダー",
    ),
    _golden_payload(
        client="東都製作所",
        deck_title="製造業DX改善提案",
        industry="製造業",
        category="製造業DX",
        brief="検査記録、品質確認、日報集計が紙とExcelに分散しており、品質分析まで時間がかかる。",
        budget="800万円前後",
        schedule="6か月で現場検証まで完了",
        solution="検査データ集約、品質KPI可視化、現場入力フロー改善",
        competitor="生産管理パッケージベンダー",
    ),
    _golden_payload(
        client="株式会社キャリアリンクス",
        deck_title="採用サイト改善・採用支援提案",
        industry="人材",
        category="採用",
        brief="応募数はあるが求職者とのミスマッチが多く、採用担当の説明工数が増えている。",
        budget="300万円から500万円",
        schedule="採用繁忙期前の3か月で改善",
        solution="採用サイト改善、職種別訴求、FAQ整理、応募導線改善",
        competitor="採用媒体代理店",
    ),
    _golden_payload(
        client="みどり市教育委員会",
        deck_title="教育機関向け問い合わせ対応改善提案",
        industry="自治体・教育",
        category="教育DX",
        brief="保護者や学校からの問い合わせが複数窓口に分散し、回答品質と対応時間にばらつきがある。",
        budget="600万円以内",
        schedule="新年度前にPoCを実施",
        solution="問い合わせ分類、ナレッジ検索、FAQ整備、職員確認フロー",
        competitor="自治体向けチャットボット事業者",
    ),
]


def test_customer_ready_gate_removes_internal_review_and_returns_sales_summary(sample_pptx_payload: dict[str, Any]) -> None:
    payload = PptxDownloadRequest(**sample_pptx_payload)
    context = build_pptx_context(payload)
    upgraded = upgrade_slides_for_v11(payload.powerpoint_generation_data.slides, context, summary_mode=False)

    result = run_customer_ready_quality_gate(upgraded, context, summary_mode=False)

    assert result.status in {"READY", "REVIEW_REQUIRED"}
    assert result.score >= 70
    assert len(result.expected_questions) == 10
    assert result.sales_summary
    assert any("提出前レビュー" in item for item in result.excluded_internal_items)
    assert "提出前レビュー" not in "\n".join(slide.title for slide in result.slides)


def test_customer_ready_gate_blocks_internal_or_generic_customer_deck(sample_pptx_payload: dict[str, Any]) -> None:
    payload = PptxDownloadRequest(
        **{
            **sample_pptx_payload,
            "client_company_info": "",
            "powerpoint_generation_data": {
                **sample_pptx_payload["powerpoint_generation_data"],
                "client_name": "Client",
                "slides": [
                    {
                        "slide_no": 1,
                        "layout": "summary",
                        "title": "Debug",
                        "bullets": ["schema_version=1.0", "api key is internal only"],
                        "speaker_notes": "internal only",
                        "visual_suggestion": "",
                    }
                ],
            },
        }
    )

    with pytest.raises(CustomerReadyBlockedError):
        build_pptx_result(payload)


def test_customer_ready_download_report_contains_v12_fields(sample_pptx_payload: dict[str, Any]) -> None:
    result = build_pptx_result(PptxDownloadRequest(**sample_pptx_payload))

    assert result.quality_report.customer_ready_status in {"READY", "REVIEW_REQUIRED"}
    assert result.quality_report.customer_ready_score is not None
    assert result.quality_report.customer_ready_reasons
    assert len(result.quality_report.customer_ready_expected_questions) == 10
    assert result.quality_report.customer_ready_rubric
    assert all(0 <= value <= 100 for value in result.quality_report.customer_ready_rubric.values())
    assert "story_flow" in result.quality_report.customer_ready_rubric


def test_customer_ready_gate_and_proposal_validation_share_judgement(sample_pptx_payload: dict[str, Any]) -> None:
    payload = PptxDownloadRequest(**sample_pptx_payload)
    context = build_pptx_context(payload)
    upgraded = upgrade_slides_for_v11(payload.powerpoint_generation_data.slides, context, summary_mode=False)

    gate = run_customer_ready_quality_gate(upgraded, context, summary_mode=False)
    validation = validate_proposal(
        PowerPointData(
            deck_title=payload.powerpoint_generation_data.deck_title,
            client_name=context.client_name,
            slides=gate.slides,
        ),
        {
            "client_name": context.client_name,
            "industry": context.proposal_label,
            "proposal_category": context.proposal_category,
            "decision_maker": "manager",
        },
    )

    assert release_judge_from_gate_status(gate.status) == validation.release_judge
    assert gate.score == validation.acceptance_scores.total_score
    assert gate.reasons
    assert validation.required_fixes == list(dict.fromkeys([*validation.required_fixes]))


def test_customer_ready_golden_fixtures_generate_customer_ready_pptx() -> None:
    for item in GOLDEN_PAYLOADS:
        result = build_pptx_result(PptxDownloadRequest(**item))
        prs = Presentation(BytesIO(result.pptx_bytes))
        report = result.quality_report

        assert len(prs.slides) >= 8
        assert report.customer_ready_status in {"READY", "REVIEW_REQUIRED"}
        assert report.customer_ready_score is not None and report.customer_ready_score >= 70
        assert report.customer_ready_expected_questions
        assert not report.customer_ready_blockers
        assert report.customer_ready_excluded_internal_items


def test_beautiful_ai_prompt_contains_v12_slide_design_metadata(sample_pptx_payload: dict[str, Any]) -> None:
    request = BeautifulAiPresentationRequest(
        project_id="v12-customer-ready",
        powerpoint_generation_data=sample_pptx_payload["powerpoint_generation_data"],
        project_brief=sample_pptx_payload["project_brief"],
        client_company_info=sample_pptx_payload["client_company_info"],
        competitor_company_name=sample_pptx_payload["competitor_company_name"],
        desired_launch_timing=sample_pptx_payload["desired_launch_timing"],
        budget_range=sample_pptx_payload["budget_range"],
        special_function_required=sample_pptx_payload["special_function_required"],
    )

    payload = map_to_beautiful_ai_payload(request)

    assert "Slide design metadata:" in payload.prompt
    assert "slide_number=" in payload.prompt
    assert "slide_purpose=" in payload.prompt
    assert "preferred_layout=" in payload.prompt
    assert "source_or_basis=" in payload.prompt
