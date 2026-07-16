from __future__ import annotations

from io import BytesIO
from typing import Any
from zipfile import ZipFile

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def _font_names(prs: Presentation) -> list[str]:
    names: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        names.add(run.font.name)
    return sorted(names)


def _broken_relationship_count(prs: Presentation) -> int:
    broken = 0
    for slide in prs.slides:
        for relationship in slide.part.rels.values():
            try:
                _ = relationship.target_ref
            except Exception:
                broken += 1
    return broken


def _external_relationship_count(content: bytes) -> int:
    with ZipFile(BytesIO(content)) as deck:
        return sum(
            deck.read(name).decode("utf-8", errors="ignore").count('TargetMode="External"')
            for name in deck.namelist()
            if name.endswith(".rels")
        )


def extract_pptx_structure(content: bytes) -> dict[str, Any]:
    prs = Presentation(BytesIO(content))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        text_values = [_shape_text(shape) for shape in slide.shapes]
        non_empty_text = [text for text in text_values if text]
        slides.append(
            {
                "index": index,
                "title": non_empty_text[0] if non_empty_text else "",
                "text": "\n".join(non_empty_text),
                "shape_count": len(slide.shapes),
                "text_shape_count": sum(1 for text in text_values if text),
                "table_count": sum(1 for shape in slide.shapes if getattr(shape, "has_table", False)),
                "chart_count": sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False)),
                "picture_count": sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE),
                "auto_shape_count": sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE),
            }
        )
    return {
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "total_shapes": sum(slide["shape_count"] for slide in slides),
        "total_tables": sum(slide["table_count"] for slide in slides),
        "total_charts": sum(slide["chart_count"] for slide in slides),
        "total_pictures": sum(slide["picture_count"] for slide in slides),
        "total_auto_shapes": sum(slide["auto_shape_count"] for slide in slides),
        "font_names": _font_names(prs),
        "broken_relationship_count": _broken_relationship_count(prs),
        "slides": slides,
    }


def _assert_premium_pptx(content: bytes, *, min_slides: int, max_slides: int) -> dict[str, Any]:
    assert content[:2] == b"PK"
    assert len(content) > 10_000
    actual = extract_pptx_structure(content)
    assert min_slides <= actual["slide_count"] <= max_slides
    assert actual["broken_relationship_count"] == 0
    assert _external_relationship_count(content) == 0
    assert actual["total_pictures"] == 0
    assert actual["total_auto_shapes"] >= actual["slide_count"] * 4
    assert "Noto Sans JP" in actual["font_names"]
    assert all(slide["text_shape_count"] > 0 for slide in actual["slides"])
    assert all("ProposalPilot" in slide["text"] for slide in actual["slides"])
    return actual


def test_detailed_pptx_uses_premium_editable_structure(
    client: TestClient,
    admin_headers: dict[str, str],
    sample_pptx_payload: dict[str, Any],
) -> None:
    response = client.post("/api/download-pptx", headers=admin_headers, json=sample_pptx_payload)

    assert response.status_code == 200
    assert response.headers["content-type"].startswith(PPTX_MEDIA_TYPE)
    actual = _assert_premium_pptx(response.content, min_slides=5, max_slides=12)
    assert actual["total_tables"] >= 1


def test_summary_pptx_stays_concise_and_visual(
    client: TestClient,
    admin_headers: dict[str, str],
    sample_pptx_payload: dict[str, Any],
) -> None:
    response = client.post(
        "/api/download-pptx",
        headers=admin_headers,
        json={**sample_pptx_payload, "summary": True},
    )

    assert response.status_code == 200
    assert response.headers["content-type"].startswith(PPTX_MEDIA_TYPE)
    _assert_premium_pptx(response.content, min_slides=8, max_slides=12)


def test_generated_image_recognition_detailed_deck_is_20_to_25_pages(
    client: TestClient,
    admin_headers: dict[str, str],
) -> None:
    payload = {
        "project_brief": "生花オークションの商品画像から花の種類、色、等級、状態をAI画像認識で判定し、商品管理システムへAPIまたはCSVで連携したい。PoCで精度評価してから本導入を検討する。",
        "client_company_info": "株式会社サンプル\n担当者: テスト担当",
        "competitor_site_url": "",
        "competitor_company_name": "",
        "estimated_page_count": "",
        "cms_required": "",
        "contact_form_required": "",
        "special_function_required": "",
        "seo_required": "",
        "content_creation_required": "",
        "desired_launch_timing": "3か月後",
        "budget_range": "1000万円以内",
        "hearing_result": "",
        "own_service_info": "",
        "past_proposal_template": "",
        "case_studies": "",
    }
    analyzed = client.post("/api/analyze", headers=admin_headers, json=payload)
    assert analyzed.status_code == 200, analyzed.text
    data = analyzed.json()
    response = client.post(
        "/api/download-pptx",
        headers=admin_headers,
        json={
            **payload,
            "powerpoint_generation_data": data["powerpoint_generation_data"],
            "win_probability": data["analysis"]["win_probability"],
        },
    )

    assert response.status_code == 200
    actual = _assert_premium_pptx(response.content, min_slides=20, max_slides=25)
    assert actual["total_tables"] >= 2
