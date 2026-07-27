from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation

from app.models import PowerPointSlide, PresentationLayoutDecisionContract, PptxDownloadRequest
from app.services.pptx_layout_integration import (
    REQUIRED_RENDERED_LAYOUT_IDS,
    apply_layout_decisions_to_slides,
    fallback_layout_for_slide_type,
)
from app.services.pptx_parts.slides import LAYOUT_RENDERER_REGISTRY
from app.services.pptx_service import build_pptx_result


def _decision(
    *,
    slide_index: int = 1,
    slide_type: str = "Problem",
    layout_id: str = "LAYOUT-003",
    status: str = "applied",
) -> PresentationLayoutDecisionContract:
    return PresentationLayoutDecisionContract(
        slide_id=f"slide-{slide_index}",
        slide_index=slide_index,
        slide_type=slide_type,
        selected_layout_id=layout_id,
        recommended_layout_ids=[layout_id, "LAYOUT-002"],
        selection_reason="Unit test layout decision.",
        expected_effect="Improves structural readability.",
        template_id="corporate_clean",
        design_token_id="corporate_clean:balanced",
        applied_by="user",
        status=status,
        confidence=0.9,
        human_review_required=False,
    )


def _payload_with_decision(sample_pptx_payload: dict[str, Any], decision: PresentationLayoutDecisionContract, *, summary: bool = False) -> PptxDownloadRequest:
    return PptxDownloadRequest(
        **{
            **sample_pptx_payload,
            "summary": summary,
            "presentation_layout_decisions": [decision.dict()],
        }
    )


def test_layout_decision_contract_accepts_optional_phase4_fields() -> None:
    decision = _decision(slide_type="Comparison", layout_id="LAYOUT-007")

    assert decision.slide_id == "slide-1"
    assert decision.slide_type == "Comparison"
    assert decision.selected_layout_id == "LAYOUT-007"
    assert decision.status == "applied"
    assert decision.applied_by == "user"


def test_required_layouts_have_registered_pptx_renderers() -> None:
    assert REQUIRED_RENDERED_LAYOUT_IDS <= set(LAYOUT_RENDERER_REGISTRY)


def test_applied_layout_decision_generates_openable_pptx_and_report(sample_pptx_payload: dict[str, Any]) -> None:
    payload = _payload_with_decision(sample_pptx_payload, _decision(slide_type="Comparison", layout_id="LAYOUT-007"))

    result = build_pptx_result(payload)
    prs = Presentation(BytesIO(result.pptx_bytes))

    assert len(prs.slides) >= 1
    assert result.quality_report.layout_decisions
    assert result.quality_report.layout_decisions[0]["selected_layout_id"] == "LAYOUT-007"
    assert result.quality_report.numeric_integrity["preserved"] is True
    assert result.quality_report.predicted_score is not None
    assert result.quality_report.rendered_score == result.quality_report.overall_score


def test_unsupported_layout_falls_back_without_changing_numbers() -> None:
    slides = [
        PowerPointSlide(
            slide_no=1,
            layout="content",
            title="AI-OCR PoC Estimate",
            bullets=["PoC budget 1000万円", "Target accuracy is 95%", "May 2027 launch"],
            speaker_notes="",
            visual_suggestion="",
        )
    ]

    result = apply_layout_decisions_to_slides(
        slides,
        [_decision(slide_type="Estimate", layout_id="LAYOUT-999")],
        template="corporate_clean",
        summary_mode=False,
        predicted_score=82,
    )

    assert result.layout_fallbacks
    assert result.unsupported_layouts == ["LAYOUT-999"]
    assert result.layout_decisions[0]["status"] == "backend_fallback"
    assert result.numeric_integrity["preserved"] is True
    assert result.slides[0].bullets == slides[0].bullets


def test_rejected_layout_is_reported_but_not_applied() -> None:
    slides = [
        PowerPointSlide(slide_no=1, layout="content", title="Risk", bullets=["確認項目"], speaker_notes="", visual_suggestion="")
    ]

    result = apply_layout_decisions_to_slides(
        slides,
        [_decision(slide_type="Risk", layout_id="LAYOUT-011", status="rejected")],
        template="corporate_clean",
        summary_mode=False,
        predicted_score=76,
    )

    assert result.slides[0].layout == "content"
    assert result.layout_decisions[0]["status"] == "rejected"
    assert result.human_review_items == ["slide-1: rejected by user"]


def test_summary_and_detailed_modes_use_different_safe_fallbacks() -> None:
    assert fallback_layout_for_slide_type("Estimate", summary_mode=False) == "LAYOUT-007"
    assert fallback_layout_for_slide_type("Estimate", summary_mode=True) == "LAYOUT-006"
    assert fallback_layout_for_slide_type("Next Action", summary_mode=False) == "LAYOUT-010"
    assert fallback_layout_for_slide_type("Next Action", summary_mode=True) == "LAYOUT-016"


def test_all_required_layouts_generate_openable_pptx(sample_pptx_payload: dict[str, Any]) -> None:
    for layout_id in sorted(REQUIRED_RENDERED_LAYOUT_IDS):
        payload = _payload_with_decision(sample_pptx_payload, _decision(layout_id=layout_id))
        result = build_pptx_result(payload)
        prs = Presentation(BytesIO(result.pptx_bytes))
        assert len(prs.slides) >= 1
        assert result.quality_report.overall_score > 0
