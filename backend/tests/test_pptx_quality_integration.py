from __future__ import annotations

from io import BytesIO
import json
from typing import Any
from urllib.parse import unquote

from fastapi.testclient import TestClient
from pptx import Presentation

from app.models import PowerPointSlide, PptxDownloadRequest
from app.services.pptx_quality import run_pptx_quality_pipeline
from app.services.pptx_service import build_pptx_context, build_pptx_result


def _report_from_header(response) -> dict[str, Any]:
    value = response.headers.get("X-Presentation-Quality-Report")
    assert value
    return json.loads(unquote(value))


def test_quality_pipeline_splits_dense_slides_and_recommends_diagrams(sample_pptx_payload: dict[str, Any]) -> None:
    payload = PptxDownloadRequest(**sample_pptx_payload)
    context = build_pptx_context(payload)
    dense_slide = PowerPointSlide(
        slide_no=1,
        layout="content",
        title="競合比較とKPIを含むAI導入提案スライド",
        bullets=[
            "AとBの比較を行う",
            "現状と改善後を比較する",
            "Phase 1でPoCを行う",
            "Phase 2でAPI連携する",
            "KPIは82%と1000万円を確認する",
            "申込から完了までのフローを整理する",
            "リスクと次のアクションを確認する",
            "運用体制を確認する",
            "関係者の確認観点を整理する",
            "導入後の改善サイクルを整理する",
        ],
        speaker_notes="確認用",
        visual_suggestion="",
    )

    result = run_pptx_quality_pipeline([dense_slide], context, summary_mode=False)

    assert result.report.slide_count_before == 1
    assert result.report.slide_count_after >= 2
    assert {finding.rule_id for finding in result.report.findings} >= {"PPT-BULLET-001", "PPT-DIAGRAM-001", "PPT-COMPARE-001"}
    assert "82%" in "\n".join(bullet for slide in result.slides for bullet in slide.bullets)
    assert "1000万円" in "\n".join(bullet for slide in result.slides for bullet in slide.bullets)


def test_download_pptx_returns_quality_report_header(
    client: TestClient,
    admin_headers: dict[str, str],
    sample_pptx_payload: dict[str, Any],
) -> None:
    response = client.post("/api/download-pptx", headers=admin_headers, json=sample_pptx_payload)

    assert response.status_code == 200
    assert response.content[:2] == b"PK"
    report = _report_from_header(response)
    assert report["overall_score"] > 0
    assert report["slide_count_after"] >= report["slide_count_before"]
    assert report["template"] == "corporate_clean"
    assert "findings" in report


def test_summary_pptx_uses_quality_rules_without_breaking_download(
    client: TestClient,
    admin_headers: dict[str, str],
    sample_pptx_payload: dict[str, Any],
) -> None:
    response = client.post("/api/download-pptx", headers=admin_headers, json={**sample_pptx_payload, "summary": True})

    assert response.status_code == 200
    report = _report_from_header(response)
    assert report["overall_score"] > 0
    assert report["slide_count_after"] <= 12
    assert report["category_scores"]


def test_all_v80_templates_generate_openable_pptx_with_quality_report(sample_pptx_payload: dict[str, Any]) -> None:
    templates = [
        "corporate_clean",
        "modern_dark",
        "creative_agency",
        "executive_minimal",
        "data_driven",
        "warm_professional",
        "japanese_business",
        "bold_vision",
    ]

    for template in templates:
        payload = PptxDownloadRequest(**{**sample_pptx_payload, "design_template": template})
        result = build_pptx_result(payload)
        prs = Presentation(BytesIO(result.pptx_bytes))
        assert len(prs.slides) >= 1
        assert result.quality_report.template == template
        assert result.quality_report.overall_score > 0


def test_quality_pipeline_preserves_numbers_when_recommending_kpi(sample_pptx_payload: dict[str, Any]) -> None:
    payload = PptxDownloadRequest(
        **{
            **sample_pptx_payload,
            "powerpoint_generation_data": {
                **sample_pptx_payload["powerpoint_generation_data"],
                "slides": [
                    {
                        "slide_no": 1,
                        "layout": "content",
                        "title": "KPIを含むEC改善提案",
                        "bullets": ["CVR 2.8%", "月間1000件", "費用300万円", "3か月で検証"],
                        "speaker_notes": "数値保持確認",
                        "visual_suggestion": "KPIカード",
                    }
                ],
            },
        }
    )

    context = build_pptx_context(payload)
    result = run_pptx_quality_pipeline(payload.powerpoint_generation_data.slides, context, summary_mode=False)
    text = "\n".join(bullet for slide in result.slides for bullet in slide.bullets)
    assert "2.8%" in text
    assert "1000" in text
    assert "300" in text
    assert result.report.overall_score > 0
