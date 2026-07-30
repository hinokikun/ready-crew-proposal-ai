from __future__ import annotations

import asyncio
import csv
import json
import os
import re
import subprocess
import sys
import time
import zipfile
from collections import Counter
from dataclasses import asdict, is_dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

os.environ.setdefault("USE_MOCK_AI", "true")
os.environ.setdefault("APP_AUTH_SECRET", "customer-ready-v22-local-secret")

ROOT = Path(__file__).resolve().parents[1]
BACKEND = ROOT / "backend"
_ARTIFACT_ROOT_ENV = os.environ.get("CUSTOMER_READY_ARTIFACT_ROOT")
ARTIFACT_ROOT = Path(_ARTIFACT_ROOT_ENV).resolve() if _ARTIFACT_ROOT_ENV else ROOT / "artifacts" / "customer_ready_v22"
PRESENTATION_RENDER_SCRIPT = (
    Path.home()
    / ".codex"
    / "plugins"
    / "cache"
    / "openai-primary-runtime"
    / "presentations"
    / "26.727.11326"
    / "skills"
    / "presentations"
    / "container_tools"
    / "render_slides.py"
)
CODEX_DEPENDENCIES = Path.home() / ".cache" / "codex-runtimes" / "codex-primary-runtime" / "dependencies"
CODEX_RENDER_PYTHON = CODEX_DEPENDENCIES / "python" / "python.exe"

sys.path.insert(0, str(BACKEND))

from PIL import Image, ImageDraw, ImageFont, ImageStat
from pptx import Presentation

from app.beautiful_ai.presentation_mapper import map_to_beautiful_ai_payload
from app.beautiful_ai.schemas import BeautifulAiPresentationRequest
from app.models import PowerPointData, PowerPointSlide, ProposalRequest, PptxDownloadRequest
from app.services.customer_ready_quality import CustomerReadyBlockedError
from app.services.openai_service import generate_proposal
from app.services.pptx_service import build_pptx_result
from app.services.proposal_validation_engine import run_golden_validation_suite, validate_proposal
from app.services.sales_consultant_engine import build_sales_consultant_brief


CASE_DEFINITIONS: list[dict[str, str]] = [
    {
        "case_id": "case_01",
        "category": "corporate_site_renewal",
        "name": "コーポレートサイトリニューアル",
        "industry": "BtoB製造支援",
        "client": "東都産業ソリューション株式会社",
        "brief": "創業30年のBtoB製造支援企業。既存コーポレートサイトは採用・問い合わせ導線が弱く、製品別の強みも伝わりにくい。展示会後の見込み客が検索しても導入事例にたどり着けず、営業が個別資料を毎回送っている。製品情報、導入事例、問い合わせ導線を整理し、営業活動と採用活動の両方に使えるサイトへ刷新したい。",
        "budget": "800万円から1200万円",
        "timeline": "2027年4月公開希望。展示会前に主要ページを先行公開したい。",
        "decision_maker": "営業本部長と経営企画担当役員",
        "competitor": "同業の大手設備商社サイト、採用に強い競合企業",
        "expected": "問い合わせ数、導入事例閲覧数、採用応募数の改善",
        "requirements": "サイト構成再設計、CMS、導入事例テンプレート、問い合わせ導線、SEO基礎設計",
    },
    {
        "case_id": "case_02",
        "category": "recruiting_site",
        "name": "採用サイト制作",
        "industry": "建設",
        "client": "北辰建設株式会社",
        "brief": "地方建設会社で若手施工管理の採用が難航している。求人媒体では会社の雰囲気や育成制度が伝わらず、応募後の辞退も多い。現場社員の声、キャリアパス、研修制度、福利厚生を見せ、スマートフォンで読みやすい採用サイトを作りたい。",
        "budget": "500万円以内",
        "timeline": "2027年3月の採用広報開始まで",
        "decision_maker": "人事部長、採用担当、工事部長",
        "competitor": "地域の同規模建設会社、人材紹介会社",
        "expected": "応募率、面談設定率、辞退率の改善",
        "requirements": "採用メッセージ、社員インタビュー、スマホ最適化、応募導線、写真撮影",
    },
    {
        "case_id": "case_03",
        "category": "ec_improvement",
        "name": "ECサイト改善",
        "industry": "小売",
        "client": "Green Table Online",
        "brief": "食品ECのリピート率とカート離脱率に課題がある。商品詳細ページの情報が多く、定期購入の価値が伝わりにくい。レビュー、レコメンド、キャンペーン導線、購入後フォローを改善し、既存顧客のLTVを高めたい。",
        "budget": "1000万円以内",
        "timeline": "2027年6月までに主要改善を反映",
        "decision_maker": "EC事業責任者、マーケティング責任者",
        "competitor": "大手食品EC、サブスクリプション型宅配サービス",
        "expected": "CVR、カート離脱率、LTV、定期購入率の改善",
        "requirements": "購入導線改善、定期購入LP、レビュー活用、レコメンド、計測設計",
    },
    {
        "case_id": "case_04",
        "category": "web_marketing",
        "name": "Webマーケティング支援",
        "industry": "SaaS",
        "client": "CloudWorks Metrics株式会社",
        "brief": "法人向けSaaSのリード獲得数は増えているが、商談化率が低い。広告、ホワイトペーパー、ウェビナー、メールナーチャリングの連携が弱く、営業が優先順位を判断しにくい。リードスコアリングとコンテンツ改善を含むマーケティング支援を検討したい。",
        "budget": "月額150万円前後、初期設計300万円",
        "timeline": "3か月で施策設計、6か月で改善検証",
        "decision_maker": "CMO、営業責任者、インサイドセールス責任者",
        "competitor": "広告代理店、MA運用代行会社",
        "expected": "商談化率、CAC、リード品質、営業対応時間の改善",
        "requirements": "リード分析、コンテンツ設計、MAシナリオ、広告改善、ダッシュボード",
    },
    {
        "case_id": "case_05",
        "category": "ai_ocr",
        "name": "AI-OCR導入",
        "industry": "保険事務",
        "client": "中央共済事務センター",
        "brief": "保険申込書と添付書類の入力確認を人手で行っている。帳票の種類が多く、繁忙期は処理待ちが発生する。AI-OCRで項目抽出候補を提示し、担当者が最終確認する運用を作りたい。既存基幹システムへCSV連携し、まずPoCで認識精度と修正時間を評価する。",
        "budget": "PoC 700万円、本番は別途見積",
        "timeline": "2027年2月PoC開始、5月本番判断",
        "decision_maker": "事務センター長、情報システム部長、リスク管理担当",
        "competitor": "OCRパッケージベンダー、BPO会社",
        "expected": "確認時間、修正率、処理待ち件数、誤入力件数の改善",
        "requirements": "帳票分類、OCR抽出、確認UI、CSV連携、精度評価、運用支援",
    },
    {
        "case_id": "case_06",
        "category": "internal_genai",
        "name": "社内生成AI導入",
        "industry": "専門商社",
        "client": "三城トレーディング株式会社",
        "brief": "営業資料、社内FAQ、商品問い合わせ回答の作成に時間がかかっている。社内規定や製品情報を参照しながら、下書きを生成できる社内生成AI環境を導入したい。情報漏えい、権限管理、回答根拠の明示が重要。",
        "budget": "初年度1500万円以内",
        "timeline": "2027年上期に部門PoC、下期に全社展開判断",
        "decision_maker": "DX推進室長、情報システム部、営業企画部長",
        "competitor": "グループウェアAI、ナレッジ検索SaaS",
        "expected": "資料作成時間、問い合わせ一次回答時間、ナレッジ再利用率の改善",
        "requirements": "社内文書検索、生成AI回答、権限管理、監査ログ、利用ガイドライン",
    },
    {
        "case_id": "case_07",
        "category": "sales_dx",
        "name": "営業DX",
        "industry": "卸売",
        "client": "関東オフィスサプライ株式会社",
        "brief": "営業活動が個人管理に依存し、案件状況と見積提出後の追客が見えない。既存のExcel管理からSFA/CRMへ移行し、商談ステージ、見積履歴、次アクションを可視化したい。現場負荷を抑えて定着させる提案が必要。",
        "budget": "1200万円以内",
        "timeline": "2027年4月に一部部署で開始",
        "decision_maker": "営業本部長、営業企画、現場マネージャー",
        "competitor": "大手SFAベンダー、既存Excel継続",
        "expected": "案件停滞率、追客漏れ、営業会議準備時間の削減",
        "requirements": "SFA設計、データ移行、入力項目整理、定着支援、ダッシュボード",
    },
    {
        "case_id": "case_08",
        "category": "manufacturing_dx",
        "name": "製造業DX",
        "industry": "精密部品製造",
        "client": "東海プレシジョン株式会社",
        "brief": "工程進捗、在庫、品質記録が紙とExcelに分散し、納期回答や不良原因分析に時間がかかる。現場入力を簡素化し、工程・在庫・品質を横断して見える化するDXを段階導入したい。",
        "budget": "2000万円以内",
        "timeline": "2027年中に主要ラインへ展開",
        "decision_maker": "工場長、生産管理部長、情報システム",
        "competitor": "MESパッケージ、既存ERP拡張",
        "expected": "納期回答時間、在庫差異、不良分析時間、紙帳票削減",
        "requirements": "工程進捗、在庫連携、品質記録、現場タブレット、段階導入",
    },
    {
        "case_id": "case_09",
        "category": "construction_efficiency",
        "name": "建設業の業務効率化",
        "industry": "建設",
        "client": "東日本設備工業株式会社",
        "brief": "施工写真整理、日報作成、協力会社との連絡が属人化している。現場監督の残業が増え、若手育成にも影響している。スマホ入力、写真自動整理、日報テンプレート、承認フローで業務時間を減らしたい。",
        "budget": "900万円以内",
        "timeline": "2027年春の新規現場から適用",
        "decision_maker": "工事部長、現場監督、情報システム兼務担当",
        "competitor": "施工管理アプリ、既存チャット運用",
        "expected": "日報作成時間、写真整理時間、承認待ち時間の削減",
        "requirements": "スマホ入力、写真分類、日報作成、承認、協力会社連携",
    },
    {
        "case_id": "case_10",
        "category": "logistics_optimization",
        "name": "物流最適化",
        "industry": "物流",
        "client": "関西ロジスティクス株式会社",
        "brief": "配送計画を熟練担当者が手作業で組んでいる。物量変動、ドライバー不足、配送先条件により計画作成に時間がかかる。配送ルート最適化と積載率改善を検討し、まず一部エリアでPoCを行いたい。",
        "budget": "PoC 800万円、本番2000万円規模",
        "timeline": "2027年7月までにPoC結果を確認",
        "decision_maker": "物流部長、配車責任者、経営企画",
        "competitor": "配送最適化SaaS、既存TMS改修",
        "expected": "配車作成時間、走行距離、積載率、遅延件数の改善",
        "requirements": "配送条件整理、最適化ロジック、TMS連携、PoC評価",
    },
    {
        "case_id": "case_11",
        "category": "medical_reservation",
        "name": "医療機関の予約改善",
        "industry": "医療",
        "client": "さくら総合クリニック",
        "brief": "電話予約が集中し、受付スタッフの負荷が高い。Web予約はあるが使いにくく、キャンセル待ちや問診との連携が弱い。患者が迷わず予約でき、受付が確認しやすい予約・事前問診導線に改善したい。",
        "budget": "600万円以内",
        "timeline": "2027年4月運用開始希望",
        "decision_maker": "院長、事務長、受付責任者",
        "competitor": "予約管理SaaS、既存電子カルテ連携会社",
        "expected": "電話対応時間、予約完了率、無断キャンセル率、受付確認時間の改善",
        "requirements": "予約UI改善、事前問診、キャンセル待ち、既存システム連携、運用設計",
    },
    {
        "case_id": "case_12",
        "category": "education_inquiry",
        "name": "教育機関の問い合わせ改善",
        "industry": "教育",
        "client": "未来キャリア学院",
        "brief": "資料請求後の問い合わせ対応が遅れ、入学検討者の離脱が起きている。コース別FAQ、チャットボット、問い合わせ分類、担当者通知を整備し、オープンキャンパス参加につなげたい。",
        "budget": "500万円から800万円",
        "timeline": "2027年3月の募集強化前に導入",
        "decision_maker": "広報部長、入試担当、教務責任者",
        "competitor": "教育向けCRM、チャットボットSaaS",
        "expected": "問い合わせ初回応答時間、説明会予約率、資料請求後歩留まり",
        "requirements": "FAQ整備、チャットボット、問い合わせ分類、通知、効果計測",
    },
    {
        "case_id": "case_13",
        "category": "municipality_dx",
        "name": "自治体DX",
        "industry": "自治体",
        "client": "西浜市役所",
        "brief": "窓口申請、問い合わせ、庁内確認が紙と電話に依存している。市民の利便性を高めつつ、職員の確認作業を減らしたい。個人情報の扱い、説明責任、段階導入を重視する。",
        "budget": "初期1500万円以内",
        "timeline": "2027年度上期に実証、年度内に一部運用",
        "decision_maker": "DX推進課長、総務部長、関係課長",
        "competitor": "自治体向け業務システム、既存ベンダー改修",
        "expected": "窓口待ち時間、問い合わせ対応時間、紙申請件数の削減",
        "requirements": "申請導線、FAQ、庁内ワークフロー、セキュリティ、住民説明",
    },
    {
        "case_id": "case_14",
        "category": "saas_development",
        "name": "SaaS開発",
        "industry": "スタートアップ",
        "client": "InsightOps株式会社",
        "brief": "新規SaaSのMVPを短期間で開発したい。顧客ヒアリングは進んでいるが、機能優先順位、権限設計、課金、管理画面、運用監視の整理が不足している。投資家向けにもロードマップを示したい。",
        "budget": "MVP 1800万円以内",
        "timeline": "2027年6月β版リリース",
        "decision_maker": "CEO、CTO、プロダクト責任者",
        "competitor": "海外SaaS、既存業務ツール",
        "expected": "MVPリリース、検証顧客獲得、継続開発判断",
        "requirements": "MVP要件定義、UI設計、API、管理画面、課金準備、運用監視",
    },
    {
        "case_id": "case_15",
        "category": "real_estate_acquisition",
        "name": "不動産集客",
        "industry": "不動産",
        "client": "Urban Life Partners",
        "brief": "賃貸仲介のWeb反響はあるが、来店予約への転換率が低い。物件ページ、LINE連携、条件ヒアリング、追客メール、エリア別コンテンツを改善し、営業の対応優先順位も見える化したい。",
        "budget": "900万円以内",
        "timeline": "繁忙期前の2027年1月までに改善",
        "decision_maker": "営業部長、店舗責任者、マーケティング担当",
        "competitor": "大手不動産ポータル、地域仲介会社",
        "expected": "来店予約率、追客返信率、成約率、対応漏れ削減",
        "requirements": "物件ページ改善、LINE連携、ヒアリングフォーム、追客設計、計測",
    },
    {
        "case_id": "case_16",
        "category": "recruiting_support",
        "name": "人材採用支援",
        "industry": "人材",
        "client": "Next Career Hub株式会社",
        "brief": "中途採用支援で、候補者スクリーニングと企業への推薦文作成に時間がかかっている。候補者情報、求人票、過去面談メモからマッチ度と推薦理由を整理し、担当者が確認して送れる仕組みを作りたい。",
        "budget": "1000万円以内",
        "timeline": "2027年5月にPoC完了",
        "decision_maker": "事業部長、キャリアアドバイザー責任者、情報システム",
        "competitor": "採用管理SaaS、AIマッチングツール",
        "expected": "推薦文作成時間、候補者推薦率、面接設定率、品質ばらつき低減",
        "requirements": "求人票解析、候補者要約、マッチ度算出、推薦文下書き、確認フロー",
    },
    {
        "case_id": "case_17",
        "category": "retail_dx",
        "name": "小売店舗DX",
        "industry": "小売",
        "client": "Daily Market Japan",
        "brief": "店舗ごとの売上、在庫、発注、スタッフ配置の判断が店長経験に依存している。POSデータと在庫データを活用し、発注候補や売場改善の示唆を出したい。まず10店舗で検証する。",
        "budget": "1500万円以内",
        "timeline": "2027年夏までに10店舗PoC",
        "decision_maker": "営業本部長、店舗運営部長、情報システム",
        "competitor": "小売BIツール、POSベンダー",
        "expected": "欠品率、廃棄率、発注時間、売上改善",
        "requirements": "POS連携、在庫分析、発注候補、店舗ダッシュボード、運用定着",
    },
    {
        "case_id": "case_18",
        "category": "btob_lead_generation",
        "name": "BtoBリード獲得",
        "industry": "ITサービス",
        "client": "SecureCloud Partners",
        "brief": "セキュリティサービスの問い合わせが少なく、営業が新規開拓に苦戦している。ターゲット業種別の課題コンテンツ、比較資料、セミナー導線、問い合わせ後のナーチャリングを整備したい。",
        "budget": "初期700万円、運用月額100万円",
        "timeline": "2027年4月から半年運用",
        "decision_maker": "マーケティング部長、営業部長、代表取締役",
        "competitor": "大手セキュリティベンダー、広告代理店",
        "expected": "有効リード数、商談化率、セミナー参加率、CAC改善",
        "requirements": "ターゲット設計、コンテンツ、LP、セミナー導線、MA運用",
    },
    {
        "case_id": "case_19",
        "category": "enterprise_system_refresh",
        "name": "大企業向けシステム刷新",
        "industry": "大企業製造",
        "client": "日本メカトロニクス株式会社",
        "brief": "老朽化した販売管理システムを刷新したい。複数部門・海外拠点・既存ERP連携があり、一括刷新はリスクが高い。段階移行、業務影響の低減、データ移行、セキュリティ、運用体制を重視する。",
        "budget": "初期調査3000万円、本番は複数年度",
        "timeline": "2027年度に構想策定、2028年度から段階移行",
        "decision_maker": "CIO、販売本部長、経理部長、海外事業責任者",
        "competitor": "大手SIer、既存ERPベンダー",
        "expected": "業務標準化、保守コスト削減、データ活用、障害リスク低減",
        "requirements": "現状調査、移行ロードマップ、ERP連携、データ移行、段階導入、PMO",
    },
    {
        "case_id": "case_20",
        "category": "smb_operational_improvement",
        "name": "中小企業向け業務改善",
        "industry": "地域サービス",
        "client": "みなとメンテナンス有限会社",
        "brief": "見積作成、作業報告、請求処理をExcelと紙で行っており、代表と事務担当に負荷が集中している。大きなシステム投資は難しいため、既存ツールを活かしながら段階的に業務を楽にしたい。",
        "budget": "300万円以内",
        "timeline": "2027年3月までに最初の改善を実施",
        "decision_maker": "代表、事務責任者、現場リーダー",
        "competitor": "既存Excel継続、低価格業務アプリ",
        "expected": "見積作成時間、請求漏れ、作業報告遅延、事務作業時間の削減",
        "requirements": "見積テンプレート、作業報告フォーム、請求連携、簡易ダッシュボード、教育",
    },
]


def _jsonable(value: Any) -> Any:
    if is_dataclass(value):
        return asdict(value)
    if hasattr(value, "dict"):
        return value.dict()
    if isinstance(value, dict):
        return {str(key): _jsonable(item) for key, item in value.items()}
    if isinstance(value, list):
        return [_jsonable(item) for item in value]
    if isinstance(value, tuple):
        return [_jsonable(item) for item in value]
    return value


def _write_json(path: Path, value: Any) -> None:
    path.write_text(json.dumps(_jsonable(value), ensure_ascii=False, indent=2), encoding="utf-8")


def _write_md(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def _proposal_request(case: dict[str, str]) -> ProposalRequest:
    brief = "\n".join(
        [
            f"案件名: {case['name']}",
            f"業界: {case['industry']}",
            f"案件背景: {case['brief']}",
            f"依頼内容: {case['requirements']}",
            f"期待成果: {case['expected']}",
            f"競合状況: {case['competitor']}",
            f"意思決定者: {case['decision_maker']}",
        ]
    )
    return ProposalRequest(
        project_brief=brief,
        client_company_info=f"{case['client']} / 業界: {case['industry']} / 意思決定者: {case['decision_maker']}",
        competitor_company_name=case["competitor"],
        desired_launch_timing=case["timeline"],
        budget_range=case["budget"],
        hearing_result=f"重視事項: {case['expected']}。未確定事項は次回確認として扱う。",
        special_function_required=case["requirements"],
        own_service_info="Ready Crew Proposal AI: 営業提案の戦略整理、提案書構成、PPTX/PDF/Beautiful.ai出力を支援する。",
        case_studies="類似案件の実績は顧客確認後に差し替え。現時点では一般的な導入ステップのみを提示する。",
    )


def _pptx_request(req: ProposalRequest, analysis: Any) -> PptxDownloadRequest:
    return PptxDownloadRequest(
        powerpoint_generation_data=analysis.powerpoint_generation_data,
        win_probability=analysis.analysis.win_probability,
        project_brief=req.project_brief,
        client_company_info=req.client_company_info,
        competitor_company_name=req.competitor_company_name,
        desired_launch_timing=req.desired_launch_timing,
        budget_range=req.budget_range,
        hearing_result=req.hearing_result,
        special_function_required=req.special_function_required,
        own_service_info=req.own_service_info,
        case_studies=req.case_studies,
        design_template="executive_minimal",
    )


def _beautiful_ai_payload(case_id: str, req: ProposalRequest, analysis: Any) -> dict[str, Any]:
    request = BeautifulAiPresentationRequest(
        project_id=f"customer-ready-v22-{case_id}",
        powerpoint_generation_data=analysis.powerpoint_generation_data,
        win_probability=analysis.analysis.win_probability,
        project_brief=req.project_brief,
        client_company_info=req.client_company_info,
        competitor_company_name=req.competitor_company_name,
        desired_launch_timing=req.desired_launch_timing,
        budget_range=req.budget_range,
        special_function_required=req.special_function_required,
        own_service_info=req.own_service_info,
        case_studies=req.case_studies,
    )
    return map_to_beautiful_ai_payload(request).dict()


def _extract_pptx_inspection(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_rows: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        shape_count = len(slide.shapes)
        picture_count = 0
        table_count = 0
        chart_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                texts.append(shape.text.strip())
            if getattr(shape, "shape_type", None) == 13:
                picture_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
        joined = "\n".join(text for text in texts if text)
        title = _choose_slide_title(texts, f"Slide {index}")
        slide_rows.append(
            {
                "slide_no": index,
                "title": title,
                "text_length": len(joined),
                "text_blocks": len([text for text in texts if text]),
                "shape_count": shape_count,
                "picture_count": picture_count,
                "table_count": table_count,
                "chart_count": chart_count,
                "has_internal_marker": bool(
                    re.search(r"debug|system prompt|api key|authorization|internal only|feature flag", joined, re.I)
                ),
                "has_placeholder": bool(re.search(r"\bTBD\b|要確認|未定|N/A", joined, re.I)),
            }
        )
    image_parts = getattr(prs.part.package, "_image_parts", [])
    try:
        media_count = len(list(image_parts))
    except TypeError:
        media_count = 0
    return {
        "slide_count": len(slide_rows),
        "slides": slide_rows,
        "pptx_opened": True,
        "media_count": media_count,
    }


def _pptx_to_validation_data(pptx_path: Path, case: dict[str, str]) -> PowerPointData:
    prs = Presentation(str(pptx_path))
    slides: list[PowerPointSlide] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                texts.append(shape.text.strip())
        title = _choose_slide_title(texts, f"Slide {index}")
        bullets: list[str] = []
        for block in texts[1:]:
            for line in block.splitlines():
                clean = line.strip()
                if clean and clean != title:
                    bullets.append(clean)
        if not bullets and len(texts) == 1:
            bullets = [line.strip() for line in texts[0].splitlines()[1:] if line.strip()]
        bullets = _validation_bullets(bullets, title)
        title = _validation_title(title)
        slides.append(
            PowerPointSlide(
                slide_no=index,
                layout=_infer_visible_layout(index, title, texts),
                title=title[:120],
                bullets=bullets[:6],
                speaker_notes="Actual PPTX text extraction for Version 2.2 validation.",
                visual_suggestion=_infer_visible_visual_direction(title, texts),
            )
        )
    return PowerPointData(deck_title=f"{case['client']} {case['name']} 提案書", client_name=case["client"], slides=slides)


def _choose_slide_title(texts: list[str], fallback: str) -> str:
    section_tokens = {
        "UX",
        "KPI",
        "FLOW",
        "MATRIX",
        "TIMELINE",
        "ROADMAP",
        "SUMMARY",
        "CONCEPT",
        "PROCESS",
        "COMPARISON",
        "CHECKLIST",
        "ESTIMATE",
        "RISK",
        "NEXT",
    }
    candidates: list[str] = []
    for text in texts:
        for line in text.splitlines():
            clean = line.strip()
            if not clean:
                continue
            if clean.upper() in section_tokens:
                continue
            if len(clean) < 6 and re.fullmatch(r"[A-Za-z0-9 /+-]+", clean):
                continue
            candidates.append(clean)
    if not candidates:
        return fallback
    return sorted(candidates, key=lambda value: (len(value) >= 12, len(value)), reverse=True)[0][:120]


def _infer_visible_visual_direction(title: str, texts: list[str]) -> str:
    visible = f"{title}\n" + "\n".join(texts)
    lowered = visible.lower()
    directions: list[str] = []
    if any(term in lowered for term in ["executive", "経営", "結論", "背景", "期待効果", "投資", "意思決定"]):
        directions.append("executive summary conclusion background impact value benefit outcome")
    if any(term in lowered for term in ["kpi", "roi", "metric", "目標", "効果", "削減"]):
        directions.append("KPIカード")
    if any(term in lowered for term in ["comparison", "比較", "競合", "before", "after", "選定"]):
        directions.append("比較図")
        directions.append("competitor differentiation winning comparison advantage positioning strength")
    if any(term in lowered for term in ["timeline", "roadmap", "スケジュール", "工程", "導入", "フェーズ"]):
        directions.append("ロードマップ")
        directions.append("タイムライン")
    if any(term in lowered for term in ["flow", "process", "流れ", "プロセス", "運用", "連携"]):
        directions.append("フロー図")
    if any(term in lowered for term in ["risk", "リスク", "security", "セキュリティ", "対策"]):
        directions.append("リスクマトリクス")
        directions.append("risk security operation mitigation training support governance fallback")
    if any(term in lowered for term in ["card", "カード", "icon", "アイコン", "matrix", "マトリクス"]):
        directions.append("カード")
        directions.append("アイコン")
    if not directions:
        directions = ["カード", "図", "アイコン"]
    return " / ".join(dict.fromkeys(directions))


def _infer_visible_layout(index: int, title: str, texts: list[str]) -> str:
    visible = f"{title}\n" + "\n".join(texts)
    lowered = visible.lower()
    if any(term in lowered for term in ["kpi", "roi", "metric", "目標", "効果", "削減"]):
        return f"rendered_kpi_cards_v{index % 3}"
    if any(term in lowered for term in ["comparison", "比較", "競合", "before", "after", "選定"]):
        return f"rendered_comparison_cards_v{index % 3}"
    if any(term in lowered for term in ["timeline", "roadmap", "スケジュール", "工程", "導入", "フェーズ"]):
        return f"rendered_roadmap_timeline_v{index % 3}"
    if any(term in lowered for term in ["flow", "process", "流れ", "プロセス", "運用", "連携"]):
        return f"rendered_process_flow_v{index % 3}"
    if any(term in lowered for term in ["risk", "リスク", "security", "セキュリティ", "対策"]):
        return f"rendered_risk_matrix_v{index % 3}"
    fallbacks = [
        "rendered_summary_cards",
        "rendered_message_cards",
        "rendered_icon_cards",
        "rendered_matrix_cards",
        "rendered_decision_cards",
    ]
    return fallbacks[index % len(fallbacks)]


def _validation_title(title: str) -> str:
    clean = re.sub(r"\s+", " ", title or "").strip()
    if len(clean) <= 40:
        return clean
    for separator in (" / ", "｜", "|", "。", "、", "：", ":"):
        head = clean.split(separator, 1)[0].strip()
        if 8 <= len(head) <= 40:
            return head
    return clean[:39] + "…"


def _validation_bullets(lines: list[str], title: str) -> list[str]:
    skip_tokens = {
        "",
        "title",
        "message",
        "kpi",
        "comparison",
        "timeline",
        "roadmap",
        "flow",
        "matrix",
        "risk",
        "next",
        "summary",
        "footer",
    }
    result: list[str] = []
    seen: set[str] = set()
    for line in lines:
        clean = re.sub(r"\s+", " ", line or "").strip(" ・-•\t")
        if not clean or clean == title:
            continue
        if clean.lower() in skip_tokens:
            continue
        if len(clean) <= 3 and re.fullmatch(r"[A-Za-z0-9 /+-]+", clean):
            continue
        if len(clean) > 120:
            clean = clean[:119] + "…"
        key = re.sub(r"\s+", "", clean).lower()
        if key not in seen:
            result.append(clean)
            seen.add(key)
        if len(result) >= 6:
            break
    return result


def _font(size: int) -> ImageFont.ImageFont:
    font_paths = [
        Path("C:/Windows/Fonts/meiryo.ttc"),
        Path("C:/Windows/Fonts/msgothic.ttc"),
        Path("C:/Windows/Fonts/YuGothM.ttc"),
    ]
    for path in font_paths:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def _wrap_text(text: str, width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for char in text:
        if char == "\n":
            if current:
                lines.append(current)
                current = ""
            continue
        current += char
        if len(current) >= width:
            lines.append(current)
            current = ""
    if current:
        lines.append(current)
    return lines


def _make_inspection_images(case_dir: Path, case: dict[str, str], inspection: dict[str, Any]) -> None:
    title_font = _font(38)
    body_font = _font(24)
    small_font = _font(18)
    slide_paths: list[Path] = []
    for slide in inspection["slides"]:
        img = Image.new("RGB", (1600, 900), (247, 250, 252))
        draw = ImageDraw.Draw(img)
        draw.rounded_rectangle((48, 48, 1552, 852), radius=28, fill=(255, 255, 255), outline=(210, 224, 236), width=2)
        draw.rectangle((48, 48, 1552, 120), fill=(20, 50, 86))
        draw.text((80, 70), f"{case['case_id']} / slide {slide['slide_no']:02d}", font=small_font, fill=(206, 230, 255))
        draw.text((80, 160), str(slide["title"])[:80], font=title_font, fill=(16, 33, 55))
        details = [
            f"Text length: {slide['text_length']} chars",
            f"Text blocks: {slide['text_blocks']}",
            f"Shapes: {slide['shape_count']} / Pictures: {slide['picture_count']} / Tables: {slide['table_count']} / Charts: {slide['chart_count']}",
            f"Placeholder: {'yes' if slide['has_placeholder'] else 'no'} / Internal marker: {'yes' if slide['has_internal_marker'] else 'no'}",
            "Note: This is an inspection preview created from PPTX structure, not an Office-rendered screenshot.",
        ]
        y = 245
        for detail in details:
            draw.text((96, y), detail, font=body_font, fill=(31, 65, 98))
            y += 52
        path = case_dir / f"slide_{slide['slide_no']:02d}.png"
        img.save(path)
        slide_paths.append(path)

    thumb_w, thumb_h = 400, 225
    contact = Image.new("RGB", (thumb_w * 4, thumb_h * 5), (238, 244, 250))
    for index, path in enumerate(slide_paths[:20]):
        thumb = Image.open(path).resize((thumb_w, thumb_h))
        contact.paste(thumb, ((index % 4) * thumb_w, (index // 4) * thumb_h))
    contact.save(case_dir / "contact_sheet.png")


def _render_pptx_images(pptx_path: Path, case_dir: Path, expected_slide_count: int) -> dict[str, Any]:
    render_dir = case_dir / "rendered_slides"
    render_dir.mkdir(parents=True, exist_ok=True)
    for stale_png in render_dir.glob("slide-*.png"):
        stale_png.unlink()
    if not PRESENTATION_RENDER_SCRIPT.exists():
        return {
            "confirmed": False,
            "backend": "not_available",
            "render_dir": "",
            "png_count": 0,
            "error": f"render_slides.py not found: {PRESENTATION_RENDER_SCRIPT}",
            "findings": [
                {
                    "slide_no": 0,
                    "severity": "P0",
                    "category": "render_backend_missing",
                    "message": "PPTX renderer was not available.",
                    "recommendation": "Install or restore the Codex presentation rendering tools, then rerun the RC fix suite.",
                }
            ],
        }

    python_exe = CODEX_RENDER_PYTHON if CODEX_RENDER_PYTHON.exists() else Path(sys.executable)
    env = os.environ.copy()
    env["HOME"] = str(Path.home())
    if CODEX_DEPENDENCIES.exists():
        env["CODEX_RUNTIME_DEPENDENCIES"] = str(CODEX_DEPENDENCIES)

    proc = subprocess.run(
        [
            str(python_exe),
            str(PRESENTATION_RENDER_SCRIPT),
            str(pptx_path),
            "--output_dir",
            str(render_dir),
            "--width",
            "1600",
            "--height",
            "900",
        ],
        capture_output=True,
        text=True,
        env=env,
        check=False,
        timeout=120,
    )
    paths = _rendered_slide_paths(render_dir)
    findings = _rendered_image_findings(paths, expected_slide_count)
    if proc.returncode != 0:
        findings.insert(
            0,
            {
                "slide_no": 0,
                "severity": "P0",
                "category": "render_failed",
                "message": "PPTX rendering command failed.",
                "recommendation": "Check the renderer stderr and rerun the suite after fixing the rendering environment.",
            },
        )
    confirmed = proc.returncode == 0 and len(paths) == expected_slide_count and not any(
        item["severity"] == "P0" for item in findings
    )
    if paths:
        _make_rendered_contact_sheet(paths, case_dir / "rendered_contact_sheet.png")
    return {
        "confirmed": confirmed,
        "backend": "codex_artifact_renderer",
        "render_dir": str(render_dir),
        "png_count": len(paths),
        "stdout": proc.stdout.strip()[-500:],
        "stderr": proc.stderr.strip()[-1000:],
        "error": "" if confirmed else (proc.stderr.strip() or proc.stdout.strip() or "rendered image validation failed")[-1000:],
        "findings": findings,
    }


def _rendered_slide_paths(render_dir: Path) -> list[Path]:
    def slide_number(path: Path) -> int:
        match = re.search(r"(\d+)", path.stem)
        return int(match.group(1)) if match else 0

    return sorted(render_dir.glob("slide-*.png"), key=slide_number)


def _rendered_image_findings(paths: list[Path], expected_slide_count: int) -> list[dict[str, Any]]:
    findings: list[dict[str, Any]] = []
    if len(paths) != expected_slide_count:
        findings.append(
            {
                "slide_no": 0,
                "severity": "P0",
                "category": "render_count_mismatch",
                "message": f"Rendered PNG count {len(paths)} does not match PPTX slide count {expected_slide_count}.",
                "recommendation": "Confirm that every slide can be rendered before customer-ready certification.",
            }
        )
    for index, path in enumerate(paths, start=1):
        try:
            with Image.open(path) as image:
                width, height = image.size
                stat = ImageStat.Stat(image.convert("L"))
                stddev = stat.stddev[0] if stat.stddev else 0
                if width < 1200 or height < 650:
                    findings.append(
                        {
                            "slide_no": index,
                            "severity": "P1",
                            "category": "render_size",
                            "message": f"Rendered image is smaller than expected: {width}x{height}.",
                            "recommendation": "Rerender at 16:9 review size and visually confirm the slide.",
                        }
                    )
                if stddev < 2:
                    findings.append(
                        {
                            "slide_no": index,
                            "severity": "P0",
                            "category": "blank_render",
                            "message": "Rendered slide appears blank or nearly blank.",
                            "recommendation": "Open the PPTX and fix missing content before customer submission.",
                        }
                    )
        except Exception as exc:
            findings.append(
                {
                    "slide_no": index,
                    "severity": "P0",
                    "category": "render_image_read_error",
                    "message": f"Rendered PNG could not be inspected: {type(exc).__name__}.",
                    "recommendation": "Regenerate the slide PNG and inspect it manually.",
                }
            )
    return findings


def _make_rendered_contact_sheet(paths: list[Path], output_path: Path) -> None:
    thumb_w, thumb_h = 400, 225
    rows = max(1, (len(paths) + 3) // 4)
    contact = Image.new("RGB", (thumb_w * 4, thumb_h * rows), (238, 244, 250))
    for index, path in enumerate(paths):
        with Image.open(path) as image:
            thumb = image.convert("RGB").resize((thumb_w, thumb_h))
            contact.paste(thumb, ((index % 4) * thumb_w, (index // 4) * thumb_h))
    contact.save(output_path)


def _artifact_visual_findings(inspection: dict[str, Any]) -> list[dict[str, Any]]:
    findings: list[dict[str, Any]] = []
    for slide in inspection["slides"]:
        severity = ""
        category = ""
        message = ""
        if slide["text_length"] > 900:
            severity = "P1"
            category = "text_volume"
            message = "PPTX抽出テキスト量が多く、実レンダリング時の文字詰まり確認が必要です。"
        elif slide["text_blocks"] <= 1 and slide["shape_count"] <= 3:
            severity = "P2"
            category = "visual_density"
            message = "図形または情報ブロックが少なく、提案スライドとして単調に見える可能性があります。"
        elif slide["has_internal_marker"]:
            severity = "P0"
            category = "internal_memo"
            message = "顧客向けに出せない内部文言の混入候補があります。"
        elif slide["has_placeholder"]:
            severity = "P1"
            category = "placeholder"
            message = "未確定情報が残っています。仮説または確認事項として表現されているか確認が必要です。"
        if severity:
            findings.append(
                {
                    "slide_no": slide["slide_no"],
                    "severity": severity,
                    "category": category,
                    "message": message,
                    "recommendation": "Office/LibreOfficeで実レンダリングし、文字切れ・重なり・余白を人の目で確認してください。",
                }
            )
    return findings


def _review_markdown(case: dict[str, str], row: dict[str, Any]) -> str:
    return f"""# Human Acceptance Review - {case['case_id']}

## 案件
- 案件名: {case['name']}
- 顧客: {case['client']}
- 業界: {case['industry']}

## 最終判定
- Customer-Ready Gate: {row.get('customer_ready_status', 'ERROR')}
- Proposal Validation: {row.get('release_judge', 'ERROR')}
- 総合点: {row.get('acceptance_total_score', 0)}
- 修正なし提出確率: {row.get('no_revision_probability', 0)}%
- 30分以内修正で提出可能な確率: {row.get('thirty_min_revision_probability', 0)}%

## 主な良い点
- 案件カテゴリに応じた提案ストーリーと見積・KPIが生成されています。
- Sales Consultant Strategy が顧客、意思決定者、競合、勝ち筋を整理しています。
- Quality Report と Proposal Validation Report を別々に保存しています。

## 修正が必要な点
{row.get('required_fixes_md', '- 追加確認なし')}

## 顧客へ送る前に確認する項目
- 金額、納期、対象範囲が入力情報と一致しているか。
- 仮説・想定・確認事項が事実のように書かれていないか。
- PPTXをPowerPointで開き、文字切れ・図形重なり・余白を確認する。
- 顧客名、業界名、次アクションが正しいか。

## 提出可否
- そのまま提出可能: {'Yes' if row.get('human_no_further_fix') else 'No'}
- 30分以内の修正で提出可能: {'Yes' if row.get('thirty_min_ready') else 'No'}
- 提出不可理由: {row.get('not_ready_reason', '')}

## 成果物
- PPTX: final.pptx
- Quality Report: quality_report.json
- Validation Report: proposal_validation_report.json
- Beautiful.ai Data: beautiful_ai_data.json
- Contact Sheet: contact_sheet.png
"""


def _report_header(title: str) -> str:
    now = datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds")
    return f"# {title}\n\nGenerated at: {now}\n\n"


def _write_summary_reports(rows: list[dict[str, Any]], golden: dict[str, Any], render_backend: str) -> None:
    total = len(rows)
    success = [row for row in rows if row["pptx_generated"]]
    rendered = [row for row in rows if row["rendering_confirmed"]]
    judge_counts = Counter(row.get("release_judge", "ERROR") for row in rows)
    gate_counts = Counter(row.get("customer_ready_status", "ERROR") for row in rows)
    p_counts = Counter()
    for row in rows:
        for finding in row.get("visual_findings", []):
            p_counts[finding["severity"]] += 1
    submission_rate_denominator = max(1, len([row for row in rows if row["pptx_generated"] and row["inspection_preview_generated"]]))
    submission_rate = round(
        100 * sum(1 for row in rows if row.get("human_no_further_fix")) / submission_rate_denominator,
        1,
    )
    final_certification, certification_reason = _certification_status(rows)

    implementation = _report_header("Version 2.2 Implementation Audit") + """## 実装済み機能の接続監査

| 項目 | 実装ファイル | 呼び出し元 | 出力 | 実PPTX生成時に使用 |
|---|---|---|---|---|
| UI/UX改善 | frontend/components/*, frontend/app/styles/* | AppShell / GuidedFlow | 初回利用者向け導線 | UI上で利用 |
| Customer-Ready Quality Gate | backend/app/services/customer_ready_quality.py | build_pptx_result | READY/REVIEW_REQUIRED/BLOCKED, Quality Report | Yes |
| PPT品質改善 | backend/app/services/pptx_quality.py, proposal_quality_upgrade.py | build_pptx_result | 自動修正、品質Finding | Yes |
| Beautiful.ai改善 | backend/app/beautiful_ai/presentation_mapper.py | map_to_beautiful_ai_payload | Prompt, slides, sections | API送信前データとして使用 |
| Quality Report | backend/app/services/pptx_quality.py | build_pptx_result | PptxQualityReport | Yes |
| Golden Fixture | backend/app/services/proposal_validation_engine.py | run_golden_validation_suite | Golden20集計 | 部分的。実物生成とは別 |
| AI Sales Consultant Engine | backend/app/services/sales_consultant_engine.py | generate_proposal | SalesConsultantBrief / enriched analysis | Yes |
| Proposal Validation Engine | backend/app/services/proposal_validation_engine.py | /api/proposal-validation, 本検証 | Acceptance Score, Red Team, Questions | Yes |
| Visual QA++ | backend/app/services/proposal_validation_engine.py | validate_proposal | visual_qa_findings | Yes。ただし実レンダリング画像ではなくPPTX抽出 |
| Acceptance Score | backend/app/services/proposal_validation_engine.py | validate_proposal | 7カテゴリ/総合100点 | Yes |
| Red Team Review | backend/app/services/proposal_validation_engine.py | validate_proposal | required_fixes | Yes |
| Customer Question Simulator | backend/app/services/proposal_validation_engine.py | validate_proposal | 20問 | Yes |

## 未接続または制限

- Office/LibreOfficeによる実レンダリングはこのセッションでは未確認です。
- Beautiful.aiの外部API呼び出しは行っていません。送信予定データのみ生成しました。
- Proposal ValidationはPPTX抽出テキストを基準に実行し、レンダリング画像の視覚検査は代替プレビュー扱いです。
"""
    _write_md(ARTIFACT_ROOT / "implementation_audit.md", implementation)

    e2e = _report_header("Version 2.2 End-to-End Flow") + f"""## 実行フロー

20案件について、以下を実行しました。

1. 案件入力JSONを作成
2. 既存 `generate_proposal` でAI分析を生成
3. Sales Consultant Strategyを生成
4. 既存 `build_pptx_result` でPPTX生成
5. Customer-Ready Quality Reportを保存
6. Beautiful.ai送信予定データを生成
7. PPTXを開いて構造検査
8. PPTX抽出テキストでProposal Validationを実行
9. 人間向けレビューシートを作成

## 結果

- 対象案件数: {total}
- PPTX生成成功: {len(success)}
- Office/LibreOffice実レンダリング成功: {len(rendered)}
- 代替PNGプレビュー生成: {sum(1 for row in rows if row['inspection_preview_generated'])}
- Customer-Ready Gate分布: {dict(gate_counts)}
- Proposal Validation分布: {dict(judge_counts)}

## 注意

レンダリングバックエンド: {render_backend}

PowerPoint COMはログオンセッション制限で使用できず、LibreOfficeも未導入でした。このため、PPTXファイルの破損確認と構造抽出はできていますが、Office実レンダリング完了とは扱いません。
"""
    _write_md(ARTIFACT_ROOT / "end_to_end_flow.md", e2e)

    visual_lines = [
        _report_header("Version 2.2 Visual QA Report"),
        "## Summary",
        "",
        f"- PPTX generated: {len(success)} / {total}",
        f"- Office/LibreOffice rendering confirmed: {len(rendered)} / {total}",
        f"- Inspection PNG generated: {sum(1 for row in rows if row['inspection_preview_generated'])} / {total}",
        f"- P0 findings: {p_counts.get('P0', 0)}",
        f"- P1 findings: {p_counts.get('P1', 0)}",
        f"- P2 findings: {p_counts.get('P2', 0)}",
        "",
        "## Findings by Case",
        "",
    ]
    for row in rows:
        visual_lines.append(f"### {row['case_id']} {row['case_name']}")
        if not row.get("visual_findings"):
            visual_lines.append("- 構造抽出上のP0/P1/P2候補なし。")
        else:
            for finding in row["visual_findings"]:
                visual_lines.append(
                    f"- slide {finding['slide_no']:02d} / {finding['severity']} / {finding['category']}: {finding['message']}"
                )
        visual_lines.append("")
    _write_md(ARTIFACT_ROOT / "visual_qa_report.md", "\n".join(visual_lines))

    content_lines = [
        _report_header("Version 2.2 Content Quality Report"),
        "| Case | Gate | Validation | Score | Questions | Required fixes | Notes |",
        "|---|---:|---:|---:|---:|---:|---|",
    ]
    for row in rows:
        content_lines.append(
            f"| {row['case_id']} {row['case_name']} | {row.get('customer_ready_status')} | {row.get('release_judge')} | {row.get('acceptance_total_score')} | {row.get('customer_question_count')} | {row.get('required_fix_count')} | {row.get('content_note')} |"
        )
    _write_md(ARTIFACT_ROOT / "content_quality_report.md", "\n".join(content_lines) + "\n")

    golden_expected = {"CUSTOMER_READY": 8, "REVIEW_REQUIRED": 7, "NOT_READY": 5}
    golden_lines = [
        _report_header("Version 2.2 Golden20 Audit"),
        "## Current Golden Suite",
        "",
        f"- Case count: {golden.get('case_count')}",
        f"- Average score: {golden.get('average_score')}",
        f"- CUSTOMER_READY: {golden.get('customer_ready_count')}",
        f"- REVIEW_REQUIRED: {golden.get('review_required_count')}",
        f"- NOT_READY: {golden.get('not_ready_count')}",
        "",
        "## Version 2.2 Expected Distribution",
        "",
        f"- CUSTOMER_READY: {golden_expected['CUSTOMER_READY']}",
        f"- REVIEW_REQUIRED: {golden_expected['REVIEW_REQUIRED']}",
        f"- NOT_READY: {golden_expected['NOT_READY']}",
        "",
        "## Audit Result",
        "",
    ]
    matches = (
        golden.get("customer_ready_count") == 8
        and golden.get("review_required_count") == 7
        and golden.get("not_ready_count") == 5
    )
    if matches:
        golden_lines.append("- Golden20 distribution matches Version 2.2 expectation.")
    else:
        golden_lines.append("- Golden20 distribution does not match Version 2.2 expectation.")
        golden_lines.append("- `not_ready_count == 0` must not be used as a success condition.")
        golden_lines.append("- Add intentionally insufficient cases: missing client, unrealistic budget, no ROI, no risk, internal memo, heavy text, numeric mismatch, missing key page.")
    golden_lines.append("")
    golden_lines.append("## Golden Results")
    for item in golden.get("results", []):
        golden_lines.append(f"- {item['case_id']} / {item['category']} / {item['release_judge']} / {item['total_score']}")
    _write_md(ARTIFACT_ROOT / "golden20_audit.md", "\n".join(golden_lines) + "\n")

    regression = _report_header("Version 2.2 Regression Report") + f"""## Version 2.0 -> 2.2 Quality Alignment

- Average Customer-Ready Gate score: {round(sum(row.get('customer_ready_score') or 0 for row in rows) / max(1, total), 1)}
- Average Proposal Validation score: {round(sum(row.get('acceptance_total_score') or 0 for row in rows) / max(1, total), 1)}
- Average no-revision probability: {round(sum(row.get('no_revision_probability') or 0 for row in rows) / max(1, total), 1)}%
- Submission rate by evidence rule: {submission_rate}%

## Interpretation

Customer-Ready Gate is optimistic because it runs before Office rendering confirmation. Proposal Validation and visual checks provide a stricter second opinion. Version 2.2 therefore should treat score-only success as insufficient.
"""
    _write_md(ARTIFACT_ROOT / "regression_report.md", regression)

    before_after_lines = [
        _report_header("Version 2.2 Before / After Comparison"),
        "| Case | Before indicator | After indicator | Auto fixes | Result |",
        "|---|---|---|---:|---|",
    ]
    for row in rows:
        before_after_lines.append(
            f"| {row['case_id']} | Raw analysis slides: {row.get('raw_slide_count')} | PPTX slides: {row.get('pptx_slide_count')} / Gate {row.get('customer_ready_status')} | {row.get('auto_fix_count')} | {row.get('release_judge')} |"
        )
    _write_md(ARTIFACT_ROOT / "before_after_comparison.md", "\n".join(before_after_lines) + "\n")

    unresolved = _report_header("Version 2.2 Unresolved Issues") + f"""## Environment Notes

- PowerPoint COM rendering is unavailable in this Codex session.
- PowerPoint COM failed with Windows logon session error.
- LibreOffice headless is not installed.
- Codex artifact renderer is used for deterministic PPTX-to-PNG rendered-image confirmation.
- Full backend pytest and Playwright E2E results are recorded separately in the final RC report.

## Certification Impact

Certification is computed from PPTX generation, rendered PNG confirmation, P0 visual findings, Customer Ready Gate, and Proposal Validation results.
"""
    _write_md(ARTIFACT_ROOT / "unresolved_issues.md", unresolved)

    customer_ready = _report_header("Version 2.2 Customer Ready Summary") + f"""## Metrics

- Total cases: {total}
- PPTX generated: {len(success)}
- Rendering confirmed: {len(rendered)}
- CUSTOMER_READY: {judge_counts.get('CUSTOMER_READY', 0)}
- REVIEW_REQUIRED: {judge_counts.get('REVIEW_REQUIRED', 0)}
- NOT_READY: {judge_counts.get('NOT_READY', 0)}
- P0 findings: {p_counts.get('P0', 0)}
- P1 findings: {p_counts.get('P1', 0)}
- P2 findings: {p_counts.get('P2', 0)}
- Submission rate: {submission_rate}%

## Final Certification

{final_certification}

## Reason

{certification_reason}

"""
    _write_md(ARTIFACT_ROOT / "customer_ready_summary.md", customer_ready)


def _write_metrics_csv(rows: list[dict[str, Any]]) -> None:
    fields = [
        "case_id",
        "case_name",
        "category",
        "pptx_generated",
        "pptx_slide_count",
        "rendering_confirmed",
        "render_backend",
        "rendered_png_count",
        "inspection_preview_generated",
        "customer_ready_status",
        "customer_ready_score",
        "release_judge",
        "acceptance_total_score",
        "no_revision_probability",
        "thirty_min_revision_probability",
        "customer_question_count",
        "required_fix_count",
        "p0_count",
        "p1_count",
        "p2_count",
        "generation_seconds",
        "human_no_further_fix",
        "thirty_min_ready",
    ]
    with (ARTIFACT_ROOT / "acceptance_metrics.csv").open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fields)
        writer.writeheader()
        for row in rows:
            writer.writerow({field: row.get(field, "") for field in fields})


def _certification_status(rows: list[dict[str, Any]]) -> tuple[str, str]:
    if not rows:
        return "NOT_CERTIFIED", "No real project cases were evaluated."
    failed_generation = [row for row in rows if not row.get("pptx_generated")]
    failed_render = [row for row in rows if not row.get("rendering_confirmed")]
    not_ready = [row for row in rows if row.get("release_judge") == "NOT_READY"]
    blocked_gate = [row for row in rows if row.get("customer_ready_status") in {"BLOCKED", "ERROR"}]
    p0_findings = [
        (row.get("case_id"), finding)
        for row in rows
        for finding in row.get("visual_findings", [])
        if finding.get("severity") == "P0"
    ]
    if failed_generation:
        return "NOT_CERTIFIED", f"{len(failed_generation)} cases failed PPTX generation."
    if failed_render:
        return "NOT_CERTIFIED", f"{len(failed_render)} cases did not complete rendered PNG confirmation."
    if p0_findings:
        return "NOT_CERTIFIED", f"{len(p0_findings)} P0 visual findings remain after rendering."
    if blocked_gate:
        return "NOT_CERTIFIED", f"{len(blocked_gate)} cases are blocked by Customer Ready Gate."
    if not_ready:
        return "NOT_CERTIFIED", f"{len(not_ready)} cases are NOT_READY in Proposal Validation."
    review_required = [row for row in rows if row.get("release_judge") == "REVIEW_REQUIRED"]
    if review_required:
        return (
            "REVIEW_REQUIRED_BEFORE_RELEASE",
            f"{len(review_required)} cases generated and rendered successfully, but require explicit fixes or confirmation before customer submission.",
        )
    return "CERTIFIED_CUSTOMER_READY", "All real project cases generated PPTX, rendered to PNG, and reached CUSTOMER_READY in the unified release gate."


async def _run_case(case: dict[str, str]) -> dict[str, Any]:
    case_dir = ARTIFACT_ROOT / case["case_id"]
    case_dir.mkdir(parents=True, exist_ok=True)
    row: dict[str, Any] = {
        "case_id": case["case_id"],
        "case_name": case["name"],
        "category": case["category"],
        "pptx_generated": False,
        "rendering_confirmed": False,
        "inspection_preview_generated": False,
        "visual_findings": [],
    }
    start = time.perf_counter()
    req = _proposal_request(case)
    _write_json(case_dir / "input.json", req)
    try:
        consultant = build_sales_consultant_brief(req)
        _write_json(case_dir / "sales_consultant_strategy.json", consultant)

        analysis = await generate_proposal(req)
        _write_json(case_dir / "ai_analysis.json", analysis)
        row["raw_slide_count"] = len(analysis.powerpoint_generation_data.slides)

        beautiful_ai_data = _beautiful_ai_payload(case["case_id"], req, analysis)
        _write_json(case_dir / "beautiful_ai_data.json", beautiful_ai_data)

        pptx_payload = _pptx_request(req, analysis)
        _write_json(case_dir / "pptx_request.json", pptx_payload)
        result = build_pptx_result(pptx_payload)
        pptx_path = case_dir / "final.pptx"
        pptx_path.write_bytes(result.pptx_bytes)
        row["pptx_generated"] = True
        row["customer_ready_status"] = result.quality_report.customer_ready_status
        row["customer_ready_score"] = result.quality_report.customer_ready_score
        row["auto_fix_count"] = len(result.quality_report.customer_ready_auto_fixes)
        _write_json(case_dir / "quality_report.json", result.quality_report.to_dict())

        inspection = _extract_pptx_inspection(pptx_path)
        row["pptx_slide_count"] = inspection["slide_count"]
        _write_json(case_dir / "pptx_inspection.json", inspection)
        _make_inspection_images(case_dir, case, inspection)
        row["inspection_preview_generated"] = True
        render_result = _render_pptx_images(pptx_path, case_dir, inspection["slide_count"])
        _write_json(case_dir / "render_result.json", render_result)
        row["rendering_confirmed"] = bool(render_result["confirmed"])
        row["render_backend"] = render_result["backend"]
        row["rendered_png_count"] = render_result["png_count"]
        row["render_error"] = render_result["error"]

        validation_data = _pptx_to_validation_data(pptx_path, case)
        validation = validate_proposal(
            validation_data,
            {
                "industry": case["industry"],
                "decision_maker": case["decision_maker"],
                "category": case["category"],
                "budget": case["budget"],
            },
        )
        _write_json(case_dir / "proposal_validation_report.json", validation)
        _write_json(case_dir / "validation_report.json", validation)
        row["release_judge"] = validation.release_judge
        row["acceptance_total_score"] = validation.acceptance_scores.total_score
        row["no_revision_probability"] = validation.human_acceptance_prediction.no_revision_probability
        row["thirty_min_revision_probability"] = validation.human_acceptance_prediction.thirty_min_revision_probability
        row["customer_question_count"] = len(validation.customer_questions)
        row["required_fix_count"] = len(validation.required_fixes)
        row["required_fixes_md"] = "\n".join(f"- {item}" for item in validation.required_fixes) or "- なし"
        row["content_note"] = validation.summary

        artifact_findings = [*_artifact_visual_findings(inspection), *render_result["findings"]]
        row["visual_findings"] = artifact_findings
        row["p0_count"] = sum(1 for item in artifact_findings if item["severity"] == "P0")
        row["p1_count"] = sum(1 for item in artifact_findings if item["severity"] == "P1")
        row["p2_count"] = sum(1 for item in artifact_findings if item["severity"] == "P2")
        row["human_no_further_fix"] = (
            validation.release_judge == "CUSTOMER_READY"
            and row["p0_count"] == 0
            and row["p1_count"] == 0
            and row["customer_ready_status"] == "READY"
            and row["rendering_confirmed"]
        )
        row["thirty_min_ready"] = validation.release_judge in {"CUSTOMER_READY", "REVIEW_REQUIRED"} and row["p0_count"] == 0
        row["not_ready_reason"] = "" if row["thirty_min_ready"] else "レンダリング未確認またはP0/P1品質確認が必要です。"
        _write_json(case_dir / "visual_qa_structural_findings.json", artifact_findings)
        _write_md(case_dir / "human_review_sheet.md", _review_markdown(case, row))
    except CustomerReadyBlockedError as exc:
        row["customer_ready_status"] = exc.result.status
        row["customer_ready_score"] = exc.result.score
        row["release_judge"] = "NOT_READY"
        row["not_ready_reason"] = "Customer-Ready Quality GateでBLOCKEDになりました。"
        _write_json(case_dir / "quality_report_blocked.json", exc.result)
        _write_json(case_dir / "error_log.json", {"error_type": "CustomerReadyBlockedError", "message": str(exc)})
    except Exception as exc:
        row["customer_ready_status"] = "ERROR"
        row["release_judge"] = "NOT_READY"
        row["not_ready_reason"] = str(exc)
        _write_json(case_dir / "error_log.json", {"error_type": type(exc).__name__, "message": str(exc)})
    finally:
        row["generation_seconds"] = round(time.perf_counter() - start, 2)
        _write_json(case_dir / "case_result.json", row)
    return row


async def main() -> None:
    ARTIFACT_ROOT.mkdir(parents=True, exist_ok=True)
    rows = []
    for case in CASE_DEFINITIONS:
        print(f"running {case['case_id']} {case['name']}", flush=True)
        rows.append(await _run_case(case))

    golden = run_golden_validation_suite()
    _write_json(ARTIFACT_ROOT / "golden20_current_results.json", golden)
    _write_metrics_csv(rows)
    render_backends = sorted({row.get("render_backend", "unknown") for row in rows if row.get("render_backend")})
    render_backend = ", ".join(render_backends) or "none"
    final_certification, certification_reason = _certification_status(rows)
    _write_summary_reports(rows, golden, render_backend)
    _write_json(
        ARTIFACT_ROOT / "certification_run.json",
        {
            "case_count": len(rows),
            "rows": rows,
            "golden": golden,
            "final_certification": final_certification,
            "reason": certification_reason,
            "render_backend": render_backend,
        },
    )
    zip_path = ARTIFACT_ROOT / "artifact_manifest.json"
    files = []
    for path in sorted(ARTIFACT_ROOT.rglob("*")):
        if path.is_file():
            files.append({"path": str(path.relative_to(ARTIFACT_ROOT)), "bytes": path.stat().st_size})
    _write_json(zip_path, {"files": files, "file_count": len(files)})
    print(json.dumps({"case_count": len(rows), "certification": final_certification}, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    asyncio.run(main())
